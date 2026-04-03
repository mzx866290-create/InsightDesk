"""
Deck generation, persistence, and export helpers.
"""

from __future__ import annotations

import io
import json
import os
import re
import sqlite3
import time
import uuid
from dataclasses import dataclass
from typing import Any, Literal

from langchain_core.documents import Document
from pydantic import BaseModel, Field

from agent_core import get_llm
from chat_store import connect_sqlite
from doc_pipeline import DocPipeline


DeckSourceMode = Literal["kb_plus_chat", "chat_only"]
DeckQualityState = Literal["supported", "weak_support", "manual"]


class DeckWarning(BaseModel):
    code: str
    message: str


class DeckMeta(BaseModel):
    title: str
    subtitle: str = ""
    language: str = "zh-CN"
    audience: str = "general"
    purpose: str = "briefing"
    author: str = "system"
    created_at: str
    session_id: str
    source_mode: DeckSourceMode
    generator_panel_id: str


class DeckGeneration(BaseModel):
    source: DeckSourceMode
    target_slide_count: int
    actual_slide_count: int = 0
    warnings: list[DeckWarning] = Field(default_factory=list)


class DeckBlock(BaseModel):
    id: str
    kind: str
    role: str
    content: dict[str, Any] = Field(default_factory=dict)
    editable: bool = True


class DeckEvidenceRef(BaseModel):
    id: str
    source_id: str
    source_title: str
    excerpt_id: str | None = None
    snippet: str = ""
    confidence: float = 0.0


class DeckSourceItem(BaseModel):
    id: str
    type: str
    title: str
    document_id: str | None = None
    uri: str | None = None
    metadata: dict[str, Any] = Field(default_factory=dict)


class DeckSlideStatus(BaseModel):
    locked: bool = False
    dirty: bool = False
    review_state: str = "draft"


class DeckSlide(BaseModel):
    id: str
    type: str
    title: str
    subtitle: str = ""
    layout: str
    intent: str = ""
    speaker_notes: str = ""
    blocks: list[DeckBlock] = Field(default_factory=list)
    evidence_refs: list[DeckEvidenceRef] = Field(default_factory=list)
    quality_state: DeckQualityState = "weak_support"
    status: DeckSlideStatus = Field(default_factory=DeckSlideStatus)


class DeckSpec(BaseModel):
    version: str = "1.0"
    deck_id: str
    status: str = "draft"
    meta: DeckMeta
    generation: DeckGeneration
    slides: list[DeckSlide]
    source_registry: list[DeckSourceItem] = Field(default_factory=list)


class OutlineSlidePlan(BaseModel):
    title: str
    objective: str
    section: str
    evidence_source_ids: list[str] = Field(default_factory=list)


class OutlinePlan(BaseModel):
    title: str
    subtitle: str
    core_message: str
    sections: list[str]
    content_slides: list[OutlineSlidePlan]


class DraftedContentSlide(BaseModel):
    title: str
    subtitle: str = ""
    key_points: list[str] = Field(default_factory=list)
    speaker_notes: str = ""
    evidence_excerpt_ids: list[str] = Field(default_factory=list)
    evidence_source_ids: list[str] = Field(default_factory=list)
    quality_state: DeckQualityState = "weak_support"


class DraftedSlideBundle(BaseModel):
    content_slides: list[DraftedContentSlide]


class SourceExcerpt(BaseModel):
    id: str
    source_id: str
    source_title: str
    snippet: str
    confidence: float = 0.85


@dataclass
class SourcePack:
    title_hint: str
    source_mode: DeckSourceMode
    qa_pairs: list[tuple[str, str]]
    chat_notes: list[str]
    excerpts: list[SourceExcerpt]
    source_registry: list[DeckSourceItem]
    warnings: list[DeckWarning]


_FALLBACK_SECTION_TITLES = ["主题概览", "关键拆解", "结论与建议"]
_FALLBACK_SLIDE_TITLES = [
    "主题概览与核心结论",
    "关键信息拆解",
    "适用场景与行动建议",
    "后续重点与风险提示",
    "补充观察",
]


def _now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%S%z")


def _clean_text(text: Any) -> str:
    return " ".join(str(text).strip().split())


def _truncate(text: str, limit: int) -> str:
    cleaned = _clean_text(text)
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1].rstrip() + "…"


def _stringify_llm_content(content: Any) -> str:
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, dict):
                text = item.get("text")
                if text:
                    parts.append(str(text))
            else:
                parts.append(str(item))
        return "\n".join(parts).strip()
    return str(content).strip()


def _extract_qa_pairs(messages: list[Any]) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    pending_question = ""
    for message in messages:
        role = getattr(message, "__class__", type(message)).__name__
        content = _clean_text(getattr(message, "content", ""))
        if not content:
            continue
        if role == "HumanMessage":
            pending_question = content
            continue
        if role == "AIMessage" and pending_question:
            pairs.append((pending_question, content))
            pending_question = ""
    return pairs


def _is_failed_answer(answer: str) -> bool:
    normalized = answer.strip().lower()
    if not normalized:
        return True
    failure_markers = (
        "agent stopped due to max iterations",
        "agent stopped due to iteration limit",
        "生成回答失败",
        "无法完成任务",
        "处理请求时发生异常",
        "模型工具调用次数超限",
        "internal_error",
    )
    return any(marker in normalized for marker in failure_markers)


def extract_successful_qa_pairs(messages: list[Any]) -> list[tuple[str, str]]:
    return [
        (question, answer)
        for question, answer in _extract_qa_pairs(messages)
        if not _is_failed_answer(answer)
    ]


def ensure_deckable_chat(messages: list[Any]) -> list[tuple[str, str]]:
    raw_pairs = _extract_qa_pairs(messages)
    qa_pairs = extract_successful_qa_pairs(messages)
    if qa_pairs:
        return qa_pairs
    if raw_pairs:
        raise ValueError(
            "该会话最近没有成功回答，当前只有失败结果，不能生成演示稿。"
        )
    raise ValueError("该会话没有可用于生成演示稿的成功问答内容。")


def _extract_json_payload(text: str) -> dict[str, Any]:
    candidate = text.strip()
    if candidate.startswith("```"):
        candidate = re.sub(r"^```(?:json)?\s*", "", candidate)
        candidate = re.sub(r"\s*```$", "", candidate)

    decoder = json.JSONDecoder()
    try:
        payload, _ = decoder.raw_decode(candidate)
        if not isinstance(payload, dict):
            raise ValueError("LLM JSON payload must be an object.")
        return payload
    except Exception:
        for start_char in ("{", "["):
            start = candidate.find(start_char)
            if start == -1:
                continue
            try:
                payload, _ = decoder.raw_decode(candidate[start:])
            except json.JSONDecodeError:
                continue
            if isinstance(payload, dict):
                return payload
        raise


async def _invoke_json(llm, prompt: str) -> dict[str, Any]:
    response = await llm.ainvoke(prompt)
    text = _stringify_llm_content(response.content)
    try:
        return _extract_json_payload(text)
    except Exception:
        repair_prompt = (
            "请把下面内容修正为合法 JSON。不要补充信息，不要输出解释，只输出 JSON。\n\n"
            f"{text}"
        )
        repaired = await llm.ainvoke(repair_prompt)
        repaired_text = _stringify_llm_content(repaired.content)
        return _extract_json_payload(repaired_text)


def _build_title_hint(session_id: str, messages: list[Any]) -> str:
    for message in messages:
        content = _clean_text(getattr(message, "content", ""))
        if (
            getattr(message, "__class__", type(message)).__name__ == "HumanMessage"
            and content
        ):
            return _truncate(content, 60)
    return f"Deck {session_id[:8]}"


def _first_non_empty(*values: Any) -> str:
    for value in values:
        cleaned = _clean_text(value)
        if cleaned:
            return cleaned
    return ""


def _default_sections() -> list[str]:
    return list(_FALLBACK_SECTION_TITLES)


def _fallback_slide_title(index: int) -> str:
    if index < len(_FALLBACK_SLIDE_TITLES):
        return _FALLBACK_SLIDE_TITLES[index]
    return f"核心主题 {index + 1}"


def _extract_answer_points(answer: str, limit: int = 4) -> list[str]:
    chunks = [
        _truncate(part, 72)
        for part in re.split(r"[。！？!?]\s+|\n+", answer)
        if _clean_text(part)
    ]
    deduped: list[str] = []
    seen: set[str] = set()
    for chunk in chunks:
        if chunk in seen:
            continue
        seen.add(chunk)
        deduped.append(chunk)
        if len(deduped) >= limit:
            break
    return deduped


def _fallback_outline(
    pack: SourcePack,
    content_slide_count: int,
) -> OutlinePlan:
    title = _truncate(pack.title_hint, 56)
    subtitle = (
        "基于知识库检索与成功对话整理"
        if pack.source_mode == "kb_plus_chat"
        else "基于成功对话整理"
    )
    core_message = _truncate(
        _first_non_empty(pack.qa_pairs[-1][1] if pack.qa_pairs else "", subtitle),
        120,
    )
    sections = _default_sections()
    content_slides: list[OutlineSlidePlan] = []

    for index in range(content_slide_count):
        question, answer = pack.qa_pairs[min(index, len(pack.qa_pairs) - 1)]
        evidence_source_ids: list[str] = []
        if pack.source_mode == "kb_plus_chat" and pack.excerpts:
            evidence_source_ids = [pack.excerpts[min(index, len(pack.excerpts) - 1)].source_id]

        content_slides.append(
            OutlineSlidePlan(
                title=_fallback_slide_title(index),
                objective=_truncate(_first_non_empty(answer, question), 100),
                section=sections[min(index, len(sections) - 1)],
                evidence_source_ids=evidence_source_ids,
            )
        )

    return OutlinePlan(
        title=title,
        subtitle=subtitle,
        core_message=core_message or "基于当前会话整理的核心结论",
        sections=sections,
        content_slides=content_slides,
    )


def _fallback_content_bundle(
    pack: SourcePack,
    outline: OutlinePlan,
) -> DraftedSlideBundle:
    content_slides: list[DraftedContentSlide] = []

    for index, plan in enumerate(outline.content_slides):
        question, answer = pack.qa_pairs[min(index, len(pack.qa_pairs) - 1)]
        points = _extract_answer_points(answer)
        if not points:
            points = [_truncate(_first_non_empty(answer, question, plan.objective), 72)]

        evidence_excerpt_ids: list[str] = []
        evidence_source_ids: list[str] = []
        quality_state: DeckQualityState = "manual"

        if pack.source_mode == "kb_plus_chat" and pack.excerpts:
            excerpt = pack.excerpts[min(index, len(pack.excerpts) - 1)]
            evidence_excerpt_ids = [excerpt.id]
            evidence_source_ids = [excerpt.source_id]
            quality_state = "supported"

        content_slides.append(
            DraftedContentSlide(
                title=plan.title or _fallback_slide_title(index),
                subtitle=_truncate(plan.objective, 72),
                key_points=points[:5],
                speaker_notes=_truncate(_first_non_empty(answer, question, plan.objective), 180),
                evidence_excerpt_ids=evidence_excerpt_ids,
                evidence_source_ids=evidence_source_ids,
                quality_state=quality_state,
            )
        )

    return DraftedSlideBundle(content_slides=content_slides)


def _dedupe_docs(docs: list[Document]) -> list[Document]:
    deduped: list[Document] = []
    seen: set[str] = set()
    for doc in docs:
        snippet = _truncate(doc.page_content, 180)
        key = f"{doc.metadata.get('source', 'unknown')}::{snippet}"
        if key in seen:
            continue
        seen.add(key)
        deduped.append(doc)
    return deduped


def _build_source_pack(
    session_id: str,
    messages: list[Any],
    knowledge_base_enabled: bool,
    vector_store_path: str | None,
    target_slide_count: int,
) -> SourcePack:
    qa_pairs = ensure_deckable_chat(messages)
    warnings: list[DeckWarning] = []
    title_hint = _build_title_hint(session_id, messages)
    chat_notes = [
        f"问题：{_truncate(question, 80)}\n回答摘要：{_truncate(answer, 180)}"
        for question, answer in qa_pairs[-4:]
    ]

    if not knowledge_base_enabled:
        warnings.append(
            DeckWarning(
                code="chat_only_mode",
                message="当前未启用知识库，本次演示稿将仅基于成功聊天答案生成，证据强度低于知识库模式。",
            )
        )
        return SourcePack(
            title_hint=title_hint,
            source_mode="chat_only",
            qa_pairs=qa_pairs,
            chat_notes=chat_notes,
            excerpts=[],
            source_registry=[],
            warnings=warnings,
        )

    default_store_path = os.getenv("VECTOR_STORE_PATH", "./vector_store")
    candidate_paths: list[str] = []
    if vector_store_path:
        candidate_paths.append(vector_store_path)
    if default_store_path not in candidate_paths:
        candidate_paths.append(default_store_path)

    pipeline = None
    resolved_store_path = None
    for path in candidate_paths:
        current_pipeline = DocPipeline(vector_store_path=path)
        if current_pipeline.load_store():
            pipeline = current_pipeline
            resolved_store_path = path
            break

    if pipeline is None:
        raise ValueError("当前知识库无法加载，暂时不能生成演示稿。")

    if not vector_store_path:
        warnings.append(
            DeckWarning(
                code="kb_fallback_default",
                message="当前角色未绑定专属知识库，已自动回退到默认知识库生成演示稿。",
            )
        )
    elif resolved_store_path != vector_store_path:
        warnings.append(
            DeckWarning(
                code="kb_fallback_default",
                message="当前角色绑定知识库加载失败，已自动回退到默认知识库生成演示稿。",
            )
        )

    first_question = qa_pairs[0][0]
    last_answer = qa_pairs[-1][1]
    queries = [
        title_hint,
        _truncate(first_question, 120),
        _truncate(last_answer, 200),
    ]

    collected: list[Document] = []
    for query in queries:
        if not query:
            continue
        collected.extend(pipeline.search_with_rerank(query, k=4, fetch_k=12))

    docs = _dedupe_docs(collected)[:8]
    if len(docs) < 3:
        raise ValueError(
            "当前知识库检索到的有效材料不足，至少需要 3 条资料片段后才能生成演示稿。"
        )

    source_id_map: dict[str, str] = {}
    source_registry: list[DeckSourceItem] = []
    excerpts: list[SourceExcerpt] = []

    for index, doc in enumerate(docs, start=1):
        source_title = str(doc.metadata.get("source", f"来源 {index}"))
        source_id = source_id_map.get(source_title)
        if not source_id:
            source_id = f"src_{len(source_id_map) + 1}"
            source_id_map[source_title] = source_id
            source_registry.append(
                DeckSourceItem(
                    id=source_id,
                    type="doc",
                    title=source_title,
                    document_id=str(doc.metadata.get("doc_id") or ""),
                    uri=str(doc.metadata.get("source") or ""),
                    metadata={
                        "page": doc.metadata.get("page"),
                        "chunk_index": doc.metadata.get("chunk_index"),
                    },
                )
            )
        excerpts.append(
            SourceExcerpt(
                id=f"ext_{index}",
                source_id=source_id,
                source_title=source_title,
                snippet=_truncate(doc.page_content, 260),
                confidence=0.88,
            )
        )

    if len(source_registry) < 2:
        raise ValueError(
            "当前知识库检索结果覆盖的来源文档不足，至少需要覆盖 2 个来源文档后才能生成演示稿。"
        )

    if len(excerpts) < target_slide_count - 2:
        warnings.append(
            DeckWarning(
                code="evidence_sparse",
                message="知识库证据量有限，本次会优先保证内容扎实，实际页数可能少于目标页数。",
            )
        )

    return SourcePack(
        title_hint=title_hint,
        source_mode="kb_plus_chat",
        qa_pairs=qa_pairs,
        chat_notes=chat_notes,
        excerpts=excerpts,
        source_registry=source_registry,
        warnings=warnings,
    )


def _decide_content_slide_count(pack: SourcePack, target_slide_count: int) -> int:
    if pack.source_mode == "kb_plus_chat":
        desired = max(2, min(target_slide_count - 3, 5))
        evidence_bound = max(2, min(5, len(pack.excerpts) // 2 + 1))
        return min(desired, evidence_bound)
    desired = max(2, min(target_slide_count - 2, 4))
    chat_bound = max(2, min(4, len(pack.qa_pairs) + 1))
    return min(desired, chat_bound)


def _serialize_source_pack(pack: SourcePack) -> str:
    if pack.source_mode == "kb_plus_chat":
        excerpts = "\n".join(
            f"- {excerpt.id} | {excerpt.source_id} | {excerpt.source_title}: {excerpt.snippet}"
            for excerpt in pack.excerpts
        )
        sources = "\n".join(
            f"- {source.id}: {source.title}" for source in pack.source_registry
        )
        return (
            f"成功问答摘要:\n{chr(10).join(pack.chat_notes)}\n\n"
            f"来源清单:\n{sources}\n\n"
            f"知识库片段:\n{excerpts}"
        )
    return "成功问答摘要:\n" + "\n".join(pack.chat_notes)


async def _generate_outline(
    llm,
    pack: SourcePack,
    target_slide_count: int,
    content_slide_count: int,
    system_prompt: str | None,
) -> OutlinePlan:
    evidence_rule = (
        "每个内容页必须从可用 source_id 中选择 1-2 个 evidence_source_ids。"
        if pack.source_mode == "kb_plus_chat"
        else "因为当前是 chat_only 模式，evidence_source_ids 统一返回空数组。"
    )
    prompt = f"""
你是一个专业的中文商业演示稿策划助手。请基于给定材料规划一个 PPT 大纲。

硬性要求：
1. 只输出 JSON，不要输出解释。
2. 不要使用 Q1、Q2、Part 这类问答标题。
3. 总体采用“主题页”结构，而不是复述聊天记录。
4. 内容页数量必须是 {content_slide_count} 页。
5. section 数量控制在 3-5 个。
6. {evidence_rule}
7. 标题必须像正式汇报标题，不要直接把用户原问题原样拿来当页标题。
8. 页数策略是宁少勿水，不要补空话。

返回 JSON Schema：
{{
  "title": "演示稿总标题",
  "subtitle": "副标题",
  "core_message": "一句话核心结论",
  "sections": ["章节1", "章节2", "章节3"],
  "content_slides": [
    {{
      "title": "主题页标题",
      "objective": "本页要说明什么",
      "section": "所属章节",
      "evidence_source_ids": ["src_1"]
    }}
  ]
}}

上下文补充：
- source_mode: {pack.source_mode}
- 目标总页数: {target_slide_count}
- 内容页数: {content_slide_count}
- 标题提示: {pack.title_hint}
- 角色上下文: {system_prompt or "无"}

材料：
{_serialize_source_pack(pack)}
"""
    payload = await _invoke_json(llm, prompt)
    try:
        return OutlinePlan.model_validate(payload)
    except Exception:
        return _fallback_outline(pack, content_slide_count)


async def _generate_content_slides(
    llm,
    pack: SourcePack,
    outline: OutlinePlan,
    system_prompt: str | None,
) -> DraftedSlideBundle:
    evidence_rule = (
        "每个内容页必须选择至少 1 个 evidence_excerpt_ids，并保证它们来自可用片段。quality_state 只能是 supported。"
        if pack.source_mode == "kb_plus_chat"
        else "当前是 chat_only 模式，evidence_excerpt_ids 和 evidence_source_ids 返回空数组，quality_state 统一使用 manual。"
    )
    prompt = f"""
你是一个专业的中文 PPT 内容撰写助手。请把下面的大纲扩展成内容页。

硬性要求：
1. 只输出 JSON，不要输出解释。
2. 只生成内容页，不生成封面、目录、附录。
3. 标题不能写成 Q1、Part 1 这种问答形式。
4. 每页 key_points 保持 3-5 条，每条一句话。
5. speaker_notes 用 1-2 句话提示讲述重点。
6. {evidence_rule}
7. 不要抄用户问题，不要注水。

返回 JSON Schema：
{{
  "content_slides": [
    {{
      "title": "主题页标题",
      "subtitle": "简短副标题",
      "key_points": ["要点1", "要点2", "要点3"],
      "speaker_notes": "讲述提示",
      "evidence_excerpt_ids": ["ext_1"],
      "evidence_source_ids": ["src_1"],
      "quality_state": "supported"
    }}
  ]
}}

角色上下文：{system_prompt or "无"}

大纲：
{outline.model_dump_json(indent=2, ensure_ascii=False)}

材料：
{_serialize_source_pack(pack)}
"""
    payload = await _invoke_json(llm, prompt)
    try:
        return DraftedSlideBundle.model_validate(payload)
    except Exception:
        return _fallback_content_bundle(pack, outline)


def _sanitize_slide_title(title: str, fallback: str) -> str:
    cleaned = _clean_text(title)
    cleaned = re.sub(r"^Q\d+[:：\-]\s*", "", cleaned, flags=re.I)
    cleaned = re.sub(r"^Part\s*\d+[:：\-]?\s*", "", cleaned, flags=re.I)
    cleaned = cleaned.strip(" -:：")
    return _truncate(cleaned or fallback, 48)


def _build_blocks(points: list[str], fallback: str) -> list[DeckBlock]:
    clean_points = [_truncate(point, 72) for point in points if _clean_text(point)]
    if clean_points:
        return [
            DeckBlock(
                id=f"block_{uuid.uuid4().hex[:8]}",
                kind="bullet_list",
                role="main_points",
                content={"items": clean_points[:5]},
            )
        ]
    return [
        DeckBlock(
            id=f"block_{uuid.uuid4().hex[:8]}",
            kind="paragraph",
            role="summary",
            content={"text": _truncate(fallback, 220)},
        )
    ]


def _pick_evidence_refs(
    slide: DraftedContentSlide,
    pack: SourcePack,
) -> list[DeckEvidenceRef]:
    excerpt_map = {excerpt.id: excerpt for excerpt in pack.excerpts}
    source_map = {source.id: source for source in pack.source_registry}
    refs: list[DeckEvidenceRef] = []

    for excerpt_id in slide.evidence_excerpt_ids:
        excerpt = excerpt_map.get(excerpt_id)
        if not excerpt:
            continue
        refs.append(
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        )

    if refs:
        return refs[:2]

    for source_id in slide.evidence_source_ids:
        if source_id not in source_map:
            continue
        excerpt = next((item for item in pack.excerpts if item.source_id == source_id), None)
        if excerpt is None:
            continue
        refs.append(
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        )
        if refs:
            return refs[:2]

    if pack.source_mode == "kb_plus_chat" and pack.excerpts:
        excerpt = pack.excerpts[0]
        return [
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        ]

    return []


def _resolve_quality_state(
    source_mode: DeckSourceMode,
    drafted_slide: DraftedContentSlide,
    evidence_refs: list[DeckEvidenceRef],
) -> DeckQualityState:
    if source_mode == "chat_only":
        return "manual"
    if evidence_refs:
        return "supported"
    if drafted_slide.quality_state in {"supported", "weak_support", "manual"}:
        return drafted_slide.quality_state
    return "weak_support"


def _build_appendix_slide(source_registry: list[DeckSourceItem]) -> DeckSlide:
    items = [source.title for source in source_registry]
    return DeckSlide(
        id="slide_appendix_sources",
        type="appendix_sources",
        title="附录：来源清单",
        subtitle="本次演示稿引用的知识库来源",
        layout="title-bullets",
        intent="appendix_sources",
        speaker_notes="附录页，供复核来源。",
        blocks=[
            DeckBlock(
                id="block_appendix_sources",
                kind="bullet_list",
                role="sources",
                content={"items": items},
            )
        ],
        evidence_refs=[],
        quality_state="supported",
        status=DeckSlideStatus(review_state="draft"),
    )


def _build_deck_from_generated(
    session_id: str,
    panel_config: Any,
    target_slide_count: int,
    pack: SourcePack,
    outline: OutlinePlan,
    drafted: DraftedSlideBundle,
) -> DeckSpec:
    slides: list[DeckSlide] = [
        DeckSlide(
            id="slide_cover",
            type="cover",
            title=_truncate(outline.title or pack.title_hint, 56),
            subtitle=_truncate(outline.subtitle or outline.core_message, 96),
            layout="hero-title",
            intent="cover",
            speaker_notes="封面页，用于建立整体主题。",
            blocks=[
                DeckBlock(
                    id="block_cover_message",
                    kind="paragraph",
                    role="core_message",
                    content={"text": _truncate(outline.core_message, 180)},
                )
            ],
            evidence_refs=[],
            quality_state="weak_support" if pack.source_mode == "kb_plus_chat" else "manual",
            status=DeckSlideStatus(review_state="draft"),
        ),
        DeckSlide(
            id="slide_outline",
            type="outline",
            title="汇报结构",
            subtitle="本次演示的核心章节",
            layout="title-bullets",
            intent="outline",
            speaker_notes="目录页，帮助听众理解结构。",
            blocks=[
                DeckBlock(
                    id="block_outline",
                    kind="bullet_list",
                    role="outline",
                    content={"items": [_truncate(section, 40) for section in outline.sections[:5]]},
                )
            ],
            evidence_refs=[],
            quality_state="weak_support" if pack.source_mode == "kb_plus_chat" else "manual",
            status=DeckSlideStatus(review_state="draft"),
        ),
    ]

    content_plans = outline.content_slides
    drafted_slides = drafted.content_slides
    content_count = min(len(content_plans), len(drafted_slides))

    for index in range(content_count):
        plan = content_plans[index]
        drafted_slide = drafted_slides[index]
        evidence_refs = _pick_evidence_refs(drafted_slide, pack)
        quality_state = _resolve_quality_state(pack.source_mode, drafted_slide, evidence_refs)
        title = _sanitize_slide_title(drafted_slide.title, plan.title)
        subtitle = _truncate(drafted_slide.subtitle or plan.objective, 72)
        blocks = _build_blocks(drafted_slide.key_points, plan.objective)
        slides.append(
            DeckSlide(
                id=f"slide_content_{index + 1}",
                type="content",
                title=title,
                subtitle=subtitle,
                layout="title-bullets",
                intent=plan.objective,
                speaker_notes=_truncate(drafted_slide.speaker_notes or plan.objective, 180),
                blocks=blocks,
                evidence_refs=evidence_refs,
                quality_state=quality_state,
                status=DeckSlideStatus(review_state="draft"),
            )
        )

    if pack.source_mode == "kb_plus_chat":
        slides.append(_build_appendix_slide(pack.source_registry))

    warnings = list(pack.warnings)
    if pack.source_mode == "chat_only":
        warnings.append(
            DeckWarning(
                code="manual_review_required",
                message="当前为仅聊天模式，导出前请重点人工复核结论和措辞。",
            )
        )

    deck = DeckSpec(
        deck_id=f"deck_{uuid.uuid4().hex}",
        meta=DeckMeta(
            title=_truncate(outline.title or pack.title_hint, 56),
            subtitle=_truncate(outline.subtitle, 96),
            created_at=_now_iso(),
            session_id=session_id,
            source_mode=pack.source_mode,
            generator_panel_id=getattr(panel_config, "panel_id", "panel"),
        ),
        generation=DeckGeneration(
            source=pack.source_mode,
            target_slide_count=target_slide_count,
            actual_slide_count=len(slides),
            warnings=warnings,
        ),
        slides=slides,
        source_registry=pack.source_registry,
    )
    return deck


async def build_deck(
    session_id: str,
    messages: list[Any],
    panel_config: Any,
    knowledge_base_enabled: bool,
    target_slide_count: int,
    vector_store_path: str | None = None,
    system_prompt: str | None = None,
) -> DeckSpec:
    qa_pairs = ensure_deckable_chat(messages)
    llm = get_llm(
        provider=getattr(panel_config, "provider", "local"),
        model_name=getattr(panel_config, "model", None),
        base_url=getattr(panel_config, "base_url", None),
        api_key=getattr(panel_config, "api_key", None) or None,
        temperature=float(getattr(panel_config, "temperature", 0.3) or 0.3),
    )
    pack = _build_source_pack(
        session_id=session_id,
        messages=messages,
        knowledge_base_enabled=knowledge_base_enabled,
        vector_store_path=vector_store_path,
        target_slide_count=target_slide_count,
    )
    if not qa_pairs:
        raise ValueError("该会话没有可用于生成演示稿的成功问答内容。")
    content_slide_count = _decide_content_slide_count(pack, target_slide_count)
    outline = await _generate_outline(
        llm=llm,
        pack=pack,
        target_slide_count=target_slide_count,
        content_slide_count=content_slide_count,
        system_prompt=system_prompt,
    )
    drafted = await _generate_content_slides(
        llm=llm,
        pack=pack,
        outline=outline,
        system_prompt=system_prompt,
    )
    return _build_deck_from_generated(
        session_id=session_id,
        panel_config=panel_config,
        target_slide_count=target_slide_count,
        pack=pack,
        outline=outline,
        drafted=drafted,
    )


def build_report_markdown(messages: list[Any], title: str) -> str:
    qa_pairs = ensure_deckable_chat(messages)
    lines = [
        "---",
        "theme: default",
        f"title: {title}",
        "class: text-center",
        "---",
        "",
        f"# {title}",
        "",
        "AI 对话报告",
        "",
    ]
    for index, (question, answer) in enumerate(qa_pairs, start=1):
        lines.append("---")
        lines.append("")
        lines.append(f"## 主题 {index}: {_truncate(question, 72)}")
        lines.append("")
        lines.append(answer)
        lines.append("")
    return "\n".join(lines)


def _slide_text(slide: DeckSlide) -> str:
    lines: list[str] = []
    for block in slide.blocks:
        if block.kind == "bullet_list":
            for item in block.content.get("items", []):
                value = _clean_text(item)
                if value:
                    lines.append(f"• {value}")
        else:
            text = _clean_text(block.content.get("text", ""))
            if text:
                lines.append(text)
    if slide.evidence_refs:
        lines.append("")
        lines.append(
            "Sources: "
            + ", ".join(ref.source_title for ref in slide.evidence_refs[:3])
        )
    return "\n".join(lines).strip()


def export_deck_to_pptx(deck: DeckSpec) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx 未安装，请在 requirements.txt 中添加 python-pptx 并重新安装依赖"
        ) from exc

    presentation = Presentation()
    for index, slide in enumerate(deck.slides):
        if index == 0:
            layout = presentation.slide_layouts[0]
            ppt_slide = presentation.slides.add_slide(layout)
            ppt_slide.shapes.title.text = slide.title or deck.meta.title
            if len(ppt_slide.placeholders) > 1:
                ppt_slide.placeholders[1].text = slide.subtitle or deck.meta.subtitle
            continue

        layout = presentation.slide_layouts[1]
        ppt_slide = presentation.slides.add_slide(layout)
        ppt_slide.shapes.title.text = slide.title
        text_frame = ppt_slide.placeholders[1].text_frame
        text_frame.word_wrap = True
        text_frame.text = _slide_text(slide)[:1600]
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


def build_export_filename(deck: DeckSpec, extension: str) -> str:
    safe_title = "".join(
        char for char in deck.meta.title if char.isalnum() or char in (" ", "-", "_")
    ).strip()[:40]
    return f"{safe_title or 'deck'}.{extension}"


class SQLiteDeckStore:
    def __init__(self, db_path: str = "./chat_history.db"):
        self.db_path = db_path
        self._init_db()

    def _init_db(self) -> None:
        with connect_sqlite(self.db_path) as conn:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS decks (
                    deck_id TEXT PRIMARY KEY,
                    session_id TEXT NOT NULL,
                    title TEXT NOT NULL,
                    spec_json TEXT NOT NULL,
                    created_at REAL NOT NULL,
                    updated_at REAL NOT NULL
                )
                """
            )
            conn.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_decks_session
                ON decks(session_id)
                """
            )
            conn.commit()

    def save(self, deck: DeckSpec) -> DeckSpec:
        now = time.time()
        payload = json.dumps(deck.model_dump(mode="json"), ensure_ascii=False)
        with connect_sqlite(self.db_path) as conn:
            conn.execute(
                """
                INSERT INTO decks (deck_id, session_id, title, spec_json, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?)
                ON CONFLICT(deck_id) DO UPDATE SET
                    session_id = excluded.session_id,
                    title = excluded.title,
                    spec_json = excluded.spec_json,
                    updated_at = excluded.updated_at
                """,
                (
                    deck.deck_id,
                    deck.meta.session_id,
                    deck.meta.title,
                    payload,
                    now,
                    now,
                ),
            )
            conn.commit()
        return deck

    def get(self, deck_id: str) -> DeckSpec:
        with connect_sqlite(self.db_path) as conn:
            cursor = conn.execute(
                "SELECT spec_json FROM decks WHERE deck_id = ?",
                (deck_id,),
            )
            row = cursor.fetchone()
        if not row:
            raise KeyError(deck_id)
        return DeckSpec.model_validate_json(row[0])
