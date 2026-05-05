import io
from pathlib import Path

from docx import Document as WordDocument
from fastapi.testclient import TestClient
from openpyxl import load_workbook
from pptx import Presentation

import backend.api_server as api_server
import backend.artifact_service as artifact_service
import backend.chat_store as chat_store
import backend.deck_service as deck_service


def _history_cls_for_db(db_path: Path):
    class TestSQLiteChatMessageHistory(chat_store.SQLiteChatMessageHistory):
        def __init__(self, session_id: str, db_path_arg: str = "./chat_history.db"):
            super().__init__(session_id=session_id, db_path=str(db_path))

    return TestSQLiteChatMessageHistory


def _deck(deck_id: str = "deck-artifact") -> deck_service.DeckSpec:
    return deck_service.DeckSpec(
        deck_id=deck_id,
        meta=deck_service.DeckMeta(
            title="Board Update",
            subtitle="Q2 snapshot",
            theme="default",
            created_at="2026-04-16T10:00:00+0800",
            session_id="session-artifact",
            source_mode="chat_only",
            generator_panel_id="panel-main",
            author="tester",
            audience="leaders",
            purpose="briefing",
        ),
        generation=deck_service.DeckGeneration(
            source="chat_only",
            target_slide_count=2,
            actual_slide_count=2,
        ),
        slides=[
            deck_service.DeckSlide(
                id="cover",
                type="cover",
                title="Board Update",
                subtitle="Q2 snapshot",
                layout="hero-title",
                blocks=[],
            ),
            deck_service.DeckSlide(
                id="content-1",
                type="content",
                title="Growth Signals",
                subtitle="What changed",
                layout="title-bullets",
                blocks=[],
            ),
        ],
        source_registry=[],
    )


def _panel_config_payload(panel_id: str = "panel-main") -> dict[str, object]:
    return {
        "panel_id": panel_id,
        "connection_type": "ollama",
        "provider": "ollama",
        "model": "qwen3.5-2B:latest",
        "base_url": "http://localhost:11434",
        "api_key": "",
        "temperature": 0.3,
        "agent_mode": "auto",
    }


def test_generate_report_persists_artifact_and_supports_exports(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)
    monkeypatch.setattr(
        api_server,
        "_artifact_store",
        artifact_service.SQLiteArtifactStore(db_path=str(db_path)),
    )

    history = history_cls("session-report-artifact")
    history.add_user_message("Board Update")
    history.add_ai_message("Revenue improved and churn declined.")

    client = TestClient(api_server.app)
    response = client.post(
        "/api/reports/generate",
        json={"session_id": "session-report-artifact"},
    )

    assert response.status_code == 200
    payload = response.json()
    artifact_id = payload["artifact_id"]
    assert artifact_id.startswith("artifact_")

    artifact_response = client.get(f"/api/artifacts/{artifact_id}")
    assert artifact_response.status_code == 200
    artifact_payload = artifact_response.json()
    assert artifact_payload["artifact_type"] == "report"
    assert artifact_payload["title"] == "Board Update"
    assert artifact_payload["content"]["markdown"] == payload["markdown"]
    assert artifact_payload["available_formats"] == ["md", "docx", "pptx"]

    list_response = client.get("/api/sessions/session-report-artifact/artifacts")
    assert list_response.status_code == 200
    listed_artifacts = list_response.json()["artifacts"]
    assert [item["artifact_id"] for item in listed_artifacts] == [artifact_id]

    global_list_response = client.get("/api/artifacts?artifact_type=report")
    assert global_list_response.status_code == 200
    global_artifacts = global_list_response.json()["artifacts"]
    assert [item["artifact_id"] for item in global_artifacts] == [artifact_id]

    export_md_response = client.get(f"/api/artifacts/{artifact_id}/export?format=md")
    assert export_md_response.status_code == 200
    assert export_md_response.text == payload["markdown"]

    export_pptx_response = client.get(f"/api/artifacts/{artifact_id}/export?format=pptx")
    assert export_pptx_response.status_code == 200
    presentation = Presentation(io.BytesIO(export_pptx_response.content))
    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Board Update"

    export_docx_response = client.get(f"/api/artifacts/{artifact_id}/export?format=docx")
    assert export_docx_response.status_code == 200
    assert export_docx_response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    document = WordDocument(io.BytesIO(export_docx_response.content))
    assert document.paragraphs[0].text == "Board Update"
    assert "Revenue improved and churn declined." in "\n".join(
        paragraph.text for paragraph in document.paragraphs
    )

    export_xlsx_response = client.get(f"/api/artifacts/{artifact_id}/export?format=xlsx")
    assert export_xlsx_response.status_code == 400
    assert "no table content" in export_xlsx_response.json()["detail"]

    update_response = client.patch(
        f"/api/artifacts/{artifact_id}",
        json={"title": "Updated Report", "markdown": "# Updated Report"},
    )
    assert update_response.status_code == 200
    updated_payload = update_response.json()
    assert updated_payload["title"] == "Updated Report"
    assert updated_payload["content"]["markdown"] == "# Updated Report"


def test_report_artifact_xlsx_available_only_when_markdown_table_exists(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)
    monkeypatch.setattr(
        api_server,
        "_artifact_store",
        artifact_service.SQLiteArtifactStore(db_path=str(db_path)),
    )

    history = history_cls("session-table-report-artifact")
    history.add_user_message("Summarize sales table")
    history.add_ai_message(
        "| Region | Revenue |\n| --- | ---: |\n| East | 120 |\n| West | 98 |"
    )

    client = TestClient(api_server.app)
    response = client.post(
        "/api/reports/generate",
        json={"session_id": "session-table-report-artifact"},
    )

    assert response.status_code == 200
    artifact_id = response.json()["artifact_id"]

    artifact_response = client.get(f"/api/artifacts/{artifact_id}")
    assert artifact_response.status_code == 200
    assert artifact_response.json()["available_formats"] == ["md", "docx", "xlsx", "pptx"]

    export_xlsx_response = client.get(f"/api/artifacts/{artifact_id}/export?format=xlsx")
    assert export_xlsx_response.status_code == 200
    workbook = load_workbook(io.BytesIO(export_xlsx_response.content))
    assert workbook.sheetnames[:2] == ["Report", "Table 1"]
    table_sheet = workbook["Table 1"]
    assert table_sheet["A1"].value == "Region"
    assert table_sheet["B2"].value == "120"


def test_build_research_archive_artifact_preserves_v2_evidence_chain():
    archive = artifact_service.build_research_archive_artifact(
        session_id="session-research-archive",
        title="AI slide market",
        task_id="task-research-1",
        agent_result={
            "output": "Research markdown.",
            "sources": [{"title": "Market report", "url": "https://example.com/report"}],
            "metadata": {"research_mode": "deep"},
            "artifacts": [
                {
                    "type": "research_report",
                    "version": "v2",
                    "query": "AI slide market",
                    "claim_evidence_chains": [
                        {
                            "claim_id": "claim-001",
                            "claim_text": "Market is growing.",
                            "status": "partial",
                        }
                    ],
                    "claim_verification_summary": {
                        "total_claims": 1,
                        "partial_claims": 1,
                    },
                }
            ],
        },
    )

    assert archive.artifact_type == "report"
    assert archive.linked_resource_type == "task"
    assert archive.linked_resource_id == "task-research-1"
    assert archive.content["research_archive"] is True
    assert archive.content["markdown"] == "Research markdown."
    assert archive.content["claim_evidence_chains"][0]["claim_id"] == "claim-001"
    assert archive.content["claim_verification_summary"]["total_claims"] == 1


def test_create_deck_persists_artifact_and_syncs_on_update(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)
    monkeypatch.setattr(
        api_server,
        "_artifact_store",
        artifact_service.SQLiteArtifactStore(db_path=str(db_path)),
    )
    monkeypatch.setattr(
        api_server,
        "_deck_store",
        deck_service.SQLiteDeckStore(db_path=str(db_path)),
    )

    history = history_cls("session-artifact")
    history.add_user_message("Trend scan")
    history.add_ai_message("Panel B research answer.", panel_id="panel-main")

    async def fake_build_deck(**kwargs):
        deck = _deck("deck-artifact-create")
        deck.meta.session_id = kwargs["session_id"]
        return deck

    monkeypatch.setattr(api_server, "build_deck", fake_build_deck)

    client = TestClient(api_server.app)
    response = client.post(
        "/api/decks",
        json={
            "session_id": "session-artifact",
            "panel_config": _panel_config_payload(),
            "knowledge_base_enabled": False,
            "target_slide_count": 6,
            "theme": "default",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["artifact_id"].startswith("artifact_")

    deck_list_response = client.get("/api/decks")
    assert deck_list_response.status_code == 200
    listed_decks = deck_list_response.json()["decks"]
    assert [item["deck_id"] for item in listed_decks] == ["deck-artifact-create"]

    artifact_list_response = client.get("/api/artifacts?artifact_type=deck")
    assert artifact_list_response.status_code == 200
    listed_deck_artifacts = artifact_list_response.json()["artifacts"]
    assert [item["artifact_id"] for item in listed_deck_artifacts] == [payload["artifact_id"]]

    artifact_response = client.get(f"/api/artifacts/{payload['artifact_id']}")
    assert artifact_response.status_code == 200
    artifact_payload = artifact_response.json()
    assert artifact_payload["artifact_type"] == "deck"
    assert artifact_payload["content"]["deck_id"] == "deck-artifact-create"
    assert artifact_payload["title"] == "Board Update"

    update_response = client.patch(
        "/api/decks/deck-artifact-create",
        json={"title": "Updated Board"},
    )
    assert update_response.status_code == 200

    synced_artifact_response = client.get(f"/api/artifacts/{payload['artifact_id']}")
    assert synced_artifact_response.status_code == 200
    synced_payload = synced_artifact_response.json()
    assert synced_payload["title"] == "Updated Board"

    export_response = client.get(f"/api/artifacts/{payload['artifact_id']}/export?format=pptx")
    assert export_response.status_code == 200
    exported = Presentation(io.BytesIO(export_response.content))
    assert len(exported.slides) == 2

    deck = api_server._deck_store.get("deck-artifact-create")
    deck.slides[1].blocks = [
        deck_service.DeckBlock(
            id="block-1",
            kind="paragraph",
            role="summary",
            content={"text": "Revenue rose", "evidence_ref_ids": ["ev-missing"]},
        )
    ]
    api_server._deck_store.save(deck)

    blocked_export_response = client.get(
        f"/api/artifacts/{payload['artifact_id']}/export?format=pptx"
    )
    assert blocked_export_response.status_code == 409
    blocked_detail = blocked_export_response.json()["detail"]
    assert blocked_detail["export_gate"]["blocked"] is True
    assert blocked_detail["blocking_slide_ids"] == ["content-1"]
    assert blocked_detail["citation_validation"]["missing_block_evidence_ref_ids"] == [
        "ev-missing"
    ]

    override_export_response = client.get(
        f"/api/artifacts/{payload['artifact_id']}/export",
        params={
            "format": "pptx",
            "allow_unsafe_export": "true",
            "override_reason": "manual legal review",
        },
    )
    assert override_export_response.status_code == 200
    overridden = Presentation(io.BytesIO(override_export_response.content))
    assert len(overridden.slides) == 2


def test_generate_artifact_endpoint_supports_report_and_deck(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)
    monkeypatch.setattr(
        api_server,
        "_artifact_store",
        artifact_service.SQLiteArtifactStore(db_path=str(db_path)),
    )
    monkeypatch.setattr(
        api_server,
        "_deck_store",
        deck_service.SQLiteDeckStore(db_path=str(db_path)),
    )

    history = history_cls("session-generate-artifact")
    history.add_user_message("Trend scan")
    history.add_ai_message("Scoped answer.", panel_id="panel-main")

    async def fake_build_deck(**kwargs):
        deck = _deck("deck-generate-artifact")
        deck.meta.session_id = kwargs["session_id"]
        return deck

    monkeypatch.setattr(api_server, "build_deck", fake_build_deck)

    client = TestClient(api_server.app)

    report_response = client.post(
        "/api/artifacts/generate",
        json={
            "artifact_type": "report",
            "session_id": "session-generate-artifact",
        },
    )
    assert report_response.status_code == 200
    assert report_response.json()["artifact_type"] == "report"

    deck_response = client.post(
        "/api/artifacts/generate",
        json={
            "artifact_type": "deck",
            "session_id": "session-generate-artifact",
            "panel_config": _panel_config_payload(),
            "knowledge_base_enabled": False,
            "target_slide_count": 6,
            "theme": "default",
        },
    )
    assert deck_response.status_code == 200
    assert deck_response.json()["artifact_type"] == "deck"
