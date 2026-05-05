import asyncio
import os
from pathlib import Path

import pytest
from langchain_core.documents import Document
from langchain_core.messages import AIMessage, HumanMessage
from pptx import Presentation
from pptx.util import Pt

import backend.helpers.document_helpers as document_helpers
from backend.stores.task_store import TaskStatus

api_document_helpers = document_helpers


class FakeUpload:
    def __init__(self, filename: str | None, payload: bytes = b"", error: Exception | None = None):
        self.filename = filename
        self._payload = payload
        self._error = error

    async def read(self) -> bytes:
        if self._error is not None:
            raise self._error
        return self._payload


def _install_tempfile_factory(monkeypatch, tmp_path: Path) -> None:
    counter = {"value": 0}

    def fake_mkstemp(
        *,
        suffix: str = "",
        dir: str | os.PathLike[str] | None = None,
    ) -> tuple[int, str]:
        counter["value"] += 1
        base_dir = Path(dir) if dir is not None else tmp_path
        base_dir.mkdir(parents=True, exist_ok=True)
        path = base_dir / f"upload-{counter['value']}{suffix}"
        fd = os.open(path, os.O_RDWR | os.O_CREAT | os.O_TRUNC)
        return fd, str(path)

    monkeypatch.setattr(document_helpers.tempfile, "mkstemp", fake_mkstemp)


def test_upload_file_suffix_handles_missing_and_nested_names():
    assert api_document_helpers.upload_file_suffix(None) == ""
    assert api_document_helpers.upload_file_suffix("") == ""
    assert api_document_helpers.upload_file_suffix(r"C:\docs\brief.txt") == ".txt"
    assert api_document_helpers.upload_file_suffix("/tmp/report.md") == ".md"


def test_cleanup_temp_paths_removes_existing_files(tmp_path):
    first = tmp_path / "first.tmp"
    second = tmp_path / "second.tmp"
    first.write_text("alpha", encoding="utf-8")
    second.write_text("beta", encoding="utf-8")

    api_document_helpers.cleanup_temp_paths([str(first), str(second)])

    assert not first.exists()
    assert not second.exists()


def test_stage_upload_files_writes_payloads(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    uploads = [
        FakeUpload("alpha.txt", b"alpha"),
        FakeUpload("brief.md", b"# brief"),
    ]

    temp_paths, file_names = asyncio.run(api_document_helpers.stage_upload_files(uploads))

    try:
        assert file_names == ["alpha.txt", "brief.md"]
        assert [Path(path).read_bytes() for path in temp_paths] == [b"alpha", b"# brief"]
    finally:
        api_document_helpers.cleanup_temp_paths(temp_paths)


def test_stage_upload_files_with_limits_uses_configured_staging_dir(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    staging_dir = tmp_path / "shared"
    uploads = [
        FakeUpload("alpha.txt", b"alpha"),
        FakeUpload("brief.md", b"# brief"),
    ]

    temp_paths, file_names = asyncio.run(
        api_document_helpers.stage_upload_files_with_limits(
            uploads,
            staging_dir=staging_dir,
        )
    )

    try:
        assert staging_dir.exists()
        assert file_names == ["alpha.txt", "brief.md"]
        assert all(
            Path(path).resolve().is_relative_to(staging_dir.resolve())
            for path in temp_paths
        )
        assert [Path(path).read_bytes() for path in temp_paths] == [b"alpha", b"# brief"]
    finally:
        api_document_helpers.cleanup_temp_paths(temp_paths)

    assert all(not Path(path).exists() for path in temp_paths)


def test_stage_upload_files_cleans_up_when_read_fails(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    uploads = [
        FakeUpload("alpha.txt", b"alpha"),
        FakeUpload("broken.txt", error=RuntimeError("read failed")),
    ]

    with pytest.raises(RuntimeError, match="read failed"):
        asyncio.run(api_document_helpers.stage_upload_files(uploads))

    assert list(tmp_path.iterdir()) == []


def test_stage_upload_files_rejects_unsupported_suffix(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    uploads = [FakeUpload("malware.exe", b"boom")]

    with pytest.raises(ValueError, match="不支持的文件类型"):
        asyncio.run(api_document_helpers.stage_upload_files(uploads))


def test_stage_upload_files_rejects_too_many_files(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    uploads = [FakeUpload(f"doc-{index}.txt", b"x") for index in range(3)]

    with pytest.raises(ValueError, match="单次最多上传 2 个文件"):
        asyncio.run(
            api_document_helpers.stage_upload_files_with_limits(
                uploads,
                max_file_count=2,
            )
        )


def test_stage_upload_files_rejects_oversized_payload(monkeypatch, tmp_path):
    _install_tempfile_factory(monkeypatch, tmp_path)
    uploads = [FakeUpload("large.txt", b"abcdef")]

    with pytest.raises(ValueError, match="文件过大"):
        asyncio.run(
            api_document_helpers.stage_upload_files_with_limits(
                uploads,
                max_file_bytes=4,
            )
        )


def test_build_upload_documents_task_record_sets_pending_status():
    record = api_document_helpers.build_upload_documents_task_record(
        temp_paths=["/tmp/a.txt"],
        file_names=["a.txt"],
        vector_store_path="vector/custom",
        task_id_factory=lambda: "task-123",
        current_time=lambda: 42.5,
    )

    assert record.task_id == "task-123"
    assert record.task_type == "upload_documents"
    assert record.status == TaskStatus.PENDING
    assert record.created_at == 42.5
    assert record.updated_at == 42.5
    assert record.params["temp_paths"] == ["/tmp/a.txt"]
    assert record.params["file_names"] == ["a.txt"]
    assert record.params["vector_store_path"] == "vector/custom"


def test_upload_documents_response_includes_task_metadata():
    record = api_document_helpers.build_upload_documents_task_record(
        temp_paths=["/tmp/a.txt"],
        file_names=["a.txt"],
        vector_store_path="vector/custom",
        task_id_factory=lambda: "task-456",
        current_time=lambda: 10.0,
    )

    payload = api_document_helpers.upload_documents_response(
        record,
        file_count=1,
        vector_store_path="vector/custom",
    )

    assert payload["ok"] is True
    assert payload["task_id"] == "task-456"
    assert payload["task_type"] == "upload_documents"
    assert payload["status"] == TaskStatus.PENDING
    assert "vector/custom" in payload["message"]


def test_build_chat_report_title_prefers_first_human_message():
    long_title = "A" * 70
    title = api_document_helpers.build_chat_report_title(
        [AIMessage(content="ignored"), HumanMessage(content=long_title)]
    )

    assert title == "A" * 50


def test_build_chat_report_title_uses_default_when_no_human_message():
    title = api_document_helpers.build_chat_report_title(
        [AIMessage(content="assistant only")],
        default_title="Fallback Title",
    )

    assert title == "Fallback Title"


def test_safe_report_filename_filters_invalid_characters():
    safe_name = api_document_helpers.safe_report_filename(
        "Board:/Q2?* Update<>",
        default_name="fallback",
    )

    assert safe_name == "BoardQ2 Update"
    assert api_document_helpers.safe_report_filename("///", default_name="fallback") == "fallback"


def test_populate_chat_report_presentation_creates_title_and_content_slides():
    presentation = Presentation()
    long_answer = "A" * 1305

    api_document_helpers.populate_chat_report_presentation(
        presentation,
        title="Board Update",
        qa_pairs=[("What changed?", long_answer)],
        body_font_size=Pt(12),
    )

    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Board Update"
    assert presentation.slides[0].placeholders[1].text
    assert "What changed?" in presentation.slides[1].shapes.title.text

    body_text = presentation.slides[1].placeholders[1].text
    assert body_text.startswith("A" * 1200)
    assert body_text != long_answer


def test_retrieval_test_payload_rejects_blank_queries():
    with pytest.raises(ValueError):
        api_document_helpers.retrieval_test_payload("   ", object())


def test_retrieval_test_payload_returns_error_when_store_not_loaded():
    class FakePipeline:
        def load_store(self):
            return False

    payload = api_document_helpers.retrieval_test_payload("alpha", FakePipeline())

    assert payload["results_count"] == 0
    assert payload["latency_ms"] == 0
    assert payload["top_scores"] == []
    assert payload["error"]


def test_retrieval_test_payload_returns_results_and_latency():
    class FakePipeline:
        def load_store(self):
            return True

        def search(self, query, k=5):
            assert query == "alpha"
            assert k == 5
            return [
                Document(page_content="Alpha findings and supporting detail", metadata={"source": "alpha.md"}),
                Document(page_content="Fallback source document", metadata={}),
            ]

    times = iter([100.0, 100.25])
    payload = api_document_helpers.retrieval_test_payload(
        "alpha",
        FakePipeline(),
        current_time=lambda: next(times),
    )

    assert payload["results_count"] == 2
    assert payload["latency_ms"] == 250.0
    assert payload["top_results"][0]["source"] == "alpha.md"
    assert payload["top_results"][0]["snippet"] == "Alpha findings and supporting detail"
    assert payload["top_results"][1]["source"]


def test_retrieval_test_payload_prefers_debug_retrieval_when_available():
    class FakePipeline:
        def load_store(self):
            return True

        def debug_retrieval(
            self,
            query,
            *,
            search_k,
            fetch_k,
            retrieval_mode,
            use_rerank,
        ):
            assert query == "alpha"
            assert search_k == 3
            assert fetch_k == 9
            assert retrieval_mode == "hybrid"
            assert use_rerank is True
            return {
                "results_count": 1,
                "search_mode": "hybrid_rerank",
                "retrieval_mode": "hybrid",
                "search_k": 3,
                "top_k": 3,
                "fetch_k": 9,
                "rewrite_query": "alpha",
                "rewrite_applied": False,
                "query_terms": ["alpha"],
                "coverage": {
                    "unique_sources": 1,
                    "source_ratio": 1.0,
                    "matched_terms": ["alpha"],
                    "matched_term_count": 1,
                },
                "top_results": [
                    {
                        "rank": 1,
                        "source": "alpha.md",
                        "snippet": "Alpha findings and supporting detail",
                        "score": 0.88,
                        "channel": "hybrid_rerank",
                    }
                ],
                "semantic_candidates": [],
                "keyword_candidates": [],
                "fused_candidates": [],
            }

    times = iter([10.0, 10.05])
    payload = api_document_helpers.retrieval_test_payload(
        "alpha",
        FakePipeline(),
        current_time=lambda: next(times),
        search_k=3,
        fetch_k=9,
        use_rerank=True,
        retrieval_mode="hybrid",
    )

    assert payload["results_count"] == 1
    assert payload["search_mode"] == "hybrid_rerank"
    assert payload["retrieval_mode"] == "hybrid"
    assert payload["coverage"]["matched_term_count"] == 1
    assert payload["top_results"][0]["source"] == "alpha.md"
    assert payload["latency_ms"] == 50.0


def test_retrieval_test_payload_returns_exception_payload():
    class FakePipeline:
        def load_store(self):
            return True

        def search(self, query, k=5):
            raise RuntimeError("search failed")

    payload = api_document_helpers.retrieval_test_payload("alpha", FakePipeline())

    assert payload["results_count"] == 0
    assert payload["latency_ms"] == 0
    assert payload["top_scores"] == []
    assert payload["error"] == "search failed"
