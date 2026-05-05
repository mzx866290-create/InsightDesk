import asyncio
import io
from types import SimpleNamespace

from pptx import Presentation
from langchain_core.messages import AIMessage, HumanMessage

import backend.deck_service as deck_service


def _panel_config():
    return SimpleNamespace(
        panel_id="panel-main",
        provider="local",
        model="test-model",
        base_url="http://localhost:11434",
        api_key="",
        temperature=0.3,
    )


def _messages():
    return [
        HumanMessage(content="请把这次复盘整理成演示稿"),
        AIMessage(content="先讲结论，再展开风险与行动建议。"),
    ]


def _deck(slide_title: str, *, theme: str = "default"):
    return deck_service.DeckSpec(
        deck_id="deck-test",
        meta=deck_service.DeckMeta(
            title="季度经营复盘",
            subtitle="聚焦风险与行动",
            theme=theme,
            created_at="2026-04-04T10:00:00+0800",
            session_id="session-1",
            source_mode="chat_only",
            generator_panel_id="panel-main",
            author="tester",
            audience="leadership",
            purpose="review",
        ),
        generation=deck_service.DeckGeneration(
            source="chat_only",
            target_slide_count=4,
            actual_slide_count=4,
        ),
        slides=[
            deck_service.DeckSlide(
                id="slide_cover",
                type="cover",
                title="季度经营复盘",
                subtitle="聚焦风险与行动",
                layout="hero-title",
                blocks=[
                    deck_service.DeckBlock(
                        id="cover-block",
                        kind="paragraph",
                        role="core_message",
                        content={"text": "先讲总览。"},
                    )
                ],
            ),
            deck_service.DeckSlide(
                id="slide_outline",
                type="outline",
                title="汇报结构",
                subtitle="先总览，再拆解",
                layout="title-bullets",
                blocks=[
                    deck_service.DeckBlock(
                        id="outline-block",
                        kind="bullet_list",
                        role="outline",
                        content={"items": ["整体表现", "核心风险", "行动建议"]},
                    )
                ],
            ),
            deck_service.DeckSlide(
                id="slide_content_1",
                type="content",
                title=slide_title,
                subtitle="旧副标题",
                layout="title-bullets",
                intent="risk-review",
                speaker_notes="旧讲稿",
                blocks=[
                    deck_service.DeckBlock(
                        id="content-block",
                        kind="bullet_list",
                        role="main_points",
                        content={"items": ["旧要点 1", "旧要点 2"]},
                    )
                ],
                quality_state="manual",
                status=deck_service.DeckSlideStatus(review_state="draft"),
            ),
            deck_service.DeckSlide(
                id="slide_content_2",
                type="content",
                title="第二页",
                subtitle="第二页副标题",
                layout="title-bullets",
                intent="actions",
                speaker_notes="第二页讲稿",
                blocks=[
                    deck_service.DeckBlock(
                        id="content-block-2",
                        kind="bullet_list",
                        role="main_points",
                        content={"items": ["行动 1", "行动 2"]},
                    )
                ],
                quality_state="manual",
                status=deck_service.DeckSlideStatus(review_state="draft"),
            ),
        ],
        source_registry=[],
    )


def _evidence_ref(
    ref_id: str,
    *,
    source_id: str = "src-1",
    source_title: str = "经营分析.md",
    snippet: str = "关键证据片段",
) -> deck_service.DeckEvidenceRef:
    return deck_service.DeckEvidenceRef(
        id=ref_id,
        source_id=source_id,
        source_title=source_title,
        snippet=snippet,
        confidence=0.9,
    )


def _source_item(
    source_id: str,
    title: str,
    *,
    uri: str | None = None,
) -> deck_service.DeckSourceItem:
    return deck_service.DeckSourceItem(
        id=source_id,
        type="doc",
        title=title,
        uri=uri,
    )


def _appendix_source_items(deck: deck_service.DeckSpec) -> list[str]:
    appendix = next(slide for slide in deck.slides if slide.type == "appendix_sources")
    items: list[str] = []
    for block in appendix.blocks:
        if block.kind == "bullet_list" and block.role == "sources":
            items.extend(str(item) for item in block.content.get("items", []))
    return items


def _replace_slide_and_refresh(
    deck: deck_service.DeckSpec,
    replacement: deck_service.DeckSlide,
) -> deck_service.DeckSpec:
    deck.slides = [
        replacement if slide.id == replacement.id else slide for slide in deck.slides
    ]
    return deck_service.refresh_deck_evidence_coverage(deck)


def _install_fake_build_deck(
    monkeypatch,
    regenerated_deck: deck_service.DeckSpec,
    *,
    expected_theme: str | None = None,
) -> None:
    async def fake_build_deck(**kwargs):
        if expected_theme is not None:
            assert kwargs["theme"] == expected_theme
        return regenerated_deck

    monkeypatch.setattr(deck_service, "build_deck", fake_build_deck)


def test_export_deck_to_pptx_uses_selected_theme_palette():
    deck = _deck("第一页", theme="midnight")

    pptx_bytes = deck_service.export_deck_to_pptx(deck)
    presentation = Presentation(io.BytesIO(pptx_bytes))

    assert str(presentation.slides[0].background.fill.fore_color.rgb) == "111827"


def test_normalize_deck_chart_blocks_coerces_chart_contract():
    deck = _deck("图表页")
    deck.slides[2].blocks.append(
        deck_service.DeckBlock(
            id="chart-1",
            kind="chart",
            role="dashboard_chart",
            content={
                "title": "  Revenue Trend  ",
                "chart_type": "bar",
                "labels": ["Q1", ""],
                "datasets": [
                    {"label": " Gross ", "data": ["1,200", "", "bad"]},
                    {"label": "", "data": ["5"]},
                ],
            },
        )
    )
    deck.slides[2].blocks.append(
        deck_service.DeckBlock(
            id="chart-invalid",
            kind="chart",
            role="dashboard_chart",
            content={"chart_type": "scatter", "datasets": []},
        )
    )

    report = deck_service.normalize_deck_chart_blocks(deck)
    normalized = deck.slides[2].blocks[-2].content
    invalid = deck.slides[2].blocks[-1].content

    assert report.normalized_block_ids == ["chart-1"]
    assert report.invalid_block_ids == ["chart-invalid"]
    assert normalized["normalization_status"] == "normalized"
    assert normalized["chart_type"] == "bar"
    assert normalized["labels"][0] == "Q1"
    assert normalized["labels"][1].endswith(" 2")
    assert normalized["datasets"] == [
        {"label": "Gross", "data": [1200.0, 0.0]},
        {"label": normalized["datasets"][1]["label"], "data": [5.0, 0.0]},
    ]
    assert normalized["datasets"][1]["label"].endswith(" 2")
    assert invalid["normalization_status"] == "invalid"
    assert invalid["normalization_issues"] == ["unsupported_chart_type"]


def test_regenerate_deck_slide_keeps_slot_and_updates_review_state(monkeypatch):
    current_deck = _deck("旧第一页", theme="sunrise")
    regenerated_deck = _deck("新第一页", theme="sunrise")
    regenerated_deck.slides[2].subtitle = "新的副标题"
    regenerated_deck.slides[2].speaker_notes = "新的讲稿"
    regenerated_deck.slides[2].blocks = [
        deck_service.DeckBlock(
            id="new-block",
            kind="bullet_list",
            role="main_points",
            content={"items": ["新要点 1", "新要点 2", "新要点 3"]},
        )
    ]
    regenerated_deck.slides[2].quality_state = "supported"

    _install_fake_build_deck(monkeypatch, regenerated_deck, expected_theme="sunrise")

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )

    assert slide.id == "slide_content_1"
    assert slide.title == "新第一页"
    assert slide.subtitle == "新的副标题"
    assert slide.quality_state == "manual"
    assert slide.status.review_state == "regenerated"
    assert slide.status.dirty is False


def test_regenerate_deck_slide_preserves_existing_evidence_when_new_refs_missing(
    monkeypatch,
):
    current_deck = _deck("旧第一页")
    old_ref = _evidence_ref("ev-old", snippet="旧页已有证据")
    current_deck.slides[2].evidence_refs = [old_ref]
    current_deck.slides[2].quality_state = "supported"
    regenerated_deck = _deck("新第一页")
    regenerated_deck.slides[2].evidence_refs = []
    regenerated_deck.slides[2].quality_state = "weak_support"
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )

    assert slide.title == "新第一页"
    assert [ref.id for ref in slide.evidence_refs] == ["ev-old"]
    assert slide.evidence_refs[0] is not old_ref
    assert slide.evidence_refs[0].snippet == "旧页已有证据"
    assert slide.quality_state == "supported"


def test_regenerate_deck_slide_replaces_new_evidence_and_refreshes_coverage(
    monkeypatch,
):
    current_deck = _deck("旧第一页")
    current_deck.slides[2].evidence_refs = [_evidence_ref("ev-old")]
    current_deck.slides[2].quality_state = "supported"
    regenerated_deck = _deck("新第一页")
    regenerated_deck.slides[2].evidence_refs = [
        _evidence_ref(
            "ev-new",
            source_id="src-2",
            source_title="客户反馈.md",
            snippet="新证据覆盖本页结论",
        )
    ]
    regenerated_deck.slides[2].quality_state = "weak_support"
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )
    _replace_slide_and_refresh(current_deck, slide)
    coverage = current_deck.generation.evidence_coverage

    assert [ref.id for ref in slide.evidence_refs] == ["ev-new"]
    assert slide.evidence_refs[0].source_title == "客户反馈.md"
    assert slide.quality_state == "supported"
    assert coverage.coverable_slide_count == 2
    assert coverage.slides_with_evidence == 1
    assert coverage.total_evidence_refs == 1
    assert coverage.coverage_ratio == 0.5
    assert coverage.unsupported_slide_ids == ["slide_content_2"]


def test_replace_deck_slide_refreshes_review_and_export_audit(monkeypatch):
    current_deck = _deck("旧第一页")
    current_deck.source_registry = [
        _source_item("src-1", "Q1经营分析.md"),
    ]
    regenerated_deck = _deck("新第一页")
    regenerated_deck.source_registry = [
        _source_item("src-1", "Q1经营分析.md"),
    ]
    regenerated_deck.slides[2].evidence_refs = [
        _evidence_ref("ev-1", source_id="src-1", source_title="Q1经营分析.md"),
    ]
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=True,
        )
    )

    from backend.api_deck_report_helpers import replace_deck_slide
    from backend.helpers.deck_report_helpers import build_deck_delivery_response

    replace_deck_slide(current_deck, slide)
    delivery_payload = build_deck_delivery_response(
        current_deck,
        focus_slide_id="slide_content_1",
    )
    slide_delivery = delivery_payload["slide_delivery"]

    assert current_deck.generation.evidence_review["citation_validation"]["status"] == "passed"
    assert current_deck.generation.citation_validation is not None
    assert current_deck.generation.citation_validation.can_export is True
    assert current_deck.citation_validation is not None
    assert current_deck.citation_validation.status == "passed"
    assert slide_delivery["slide_id"] == "slide_content_1"
    assert slide_delivery["review_state"] == "regenerated"
    assert slide_delivery["needs_review"] is False
    assert slide_delivery["blocks_export"] is False
    assert slide_delivery["deck_can_export"] is True


def test_regenerate_deck_slide_keeps_unsupported_coverage_when_no_refs_available(
    monkeypatch,
):
    current_deck = deck_service.refresh_deck_evidence_coverage(_deck("旧第一页"))
    regenerated_deck = _deck("新第一页")
    regenerated_deck.slides[2].evidence_refs = []
    regenerated_deck.slides[2].quality_state = "weak_support"
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )
    _replace_slide_and_refresh(current_deck, slide)
    coverage = current_deck.generation.evidence_coverage

    assert slide.evidence_refs == []
    assert slide.quality_state == "manual"
    assert coverage.slides_with_evidence == 0
    assert coverage.total_evidence_refs == 0
    assert coverage.coverage_ratio == 0.0
    assert coverage.unsupported_slide_ids == ["slide_content_1", "slide_content_2"]


def test_regenerate_deck_slide_adds_new_source_registry_and_appendix(monkeypatch):
    current_deck = _deck("Old content")
    regenerated_deck = _deck("New content")
    regenerated_deck.source_registry = [
        _source_item("src-new", "New Source", uri="https://example.com/new")
    ]
    regenerated_deck.slides[2].evidence_refs = [
        _evidence_ref("ev-new", source_id="src-new", source_title="New Source")
    ]
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )
    _replace_slide_and_refresh(current_deck, slide)

    assert [source.id for source in current_deck.source_registry] == ["src-new"]
    assert slide.evidence_refs[0].source_id == "src-new"
    assert _appendix_source_items(current_deck) == ["New Source"]
    appendix_coverage = next(
        item
        for item in current_deck.generation.evidence_coverage.slides
        if item.slide_type == "appendix_sources"
    )
    assert appendix_coverage.is_coverable is False


def test_regenerate_deck_slide_merges_registry_and_preserves_appendix_sources(
    monkeypatch,
):
    current_deck = _deck("Old content")
    current_deck.source_registry = [
        _source_item("src-old", "Existing Source", uri="https://example.com/old")
    ]
    current_deck.slides[2].evidence_refs = [
        _evidence_ref(
            "ev-old",
            source_id="src-old",
            source_title="Existing Source",
            snippet="Existing support.",
        )
    ]
    current_deck.slides[2].quality_state = "supported"
    current_deck.slides.append(
        deck_service._build_appendix_slide(current_deck.source_registry)
    )
    current_deck.generation.actual_slide_count = len(current_deck.slides)

    regenerated_deck = _deck("New content")
    regenerated_deck.source_registry = [
        _source_item("src-new", "New Source", uri="https://example.com/new")
    ]
    regenerated_deck.slides[2].evidence_refs = [
        _evidence_ref(
            "ev-new",
            source_id="src-new",
            source_title="New Source",
            snippet="New support.",
        )
    ]
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )
    _replace_slide_and_refresh(current_deck, slide)

    assert [source.id for source in current_deck.source_registry] == [
        "src-old",
        "src-new",
    ]
    assert _appendix_source_items(current_deck) == [
        "Existing Source",
        "New Source",
    ]
    assert slide.evidence_refs[0].source_id == "src-new"


def test_regenerate_deck_slide_remaps_conflicting_new_source_id(monkeypatch):
    current_deck = _deck("Old content")
    current_deck.source_registry = [
        _source_item("src-1", "Old Source", uri="https://example.com/old")
    ]
    current_deck.slides[2].evidence_refs = [
        _evidence_ref("ev-old", source_id="src-1", source_title="Old Source")
    ]

    regenerated_deck = _deck("New content")
    regenerated_deck.source_registry = [
        _source_item("src-1", "New Source", uri="https://example.com/new")
    ]
    regenerated_deck.slides[2].evidence_refs = [
        _evidence_ref("ev-new", source_id="src-1", source_title="New Source")
    ]
    _install_fake_build_deck(monkeypatch, regenerated_deck)

    slide = asyncio.run(
        deck_service.regenerate_deck_slide(
            deck=current_deck,
            slide_id="slide_content_1",
            messages=_messages(),
            panel_config=_panel_config(),
            knowledge_base_enabled=False,
        )
    )
    _replace_slide_and_refresh(current_deck, slide)

    assert slide.evidence_refs[0].source_id == "src-1_regen_1"
    assert [source.id for source in current_deck.source_registry] == [
        "src-1",
        "src-1_regen_1",
    ]
    assert _appendix_source_items(current_deck) == ["Old Source", "New Source"]
