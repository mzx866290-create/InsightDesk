import io
from types import SimpleNamespace

from pptx import Presentation
from backend.deck_service import (
    DeckBlock,
    DraftedContentSlide,
    DraftedSlideBundle,
    DeckEvidenceRef,
    DeckGeneration,
    DeckMeta,
    DeckSlide,
    DeckSlideStatus,
    DeckSourceItem,
    DeckSpec,
    OutlinePlan,
    OutlineSlidePlan,
    SourcePack,
    _build_deck_from_generated,
    export_deck_to_pptx,
)


def _shape_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if text:
            parts.append(text)
    return "\n".join(parts)


def test_export_deck_to_pptx_renders_structured_slides_and_notes():
    deck = DeckSpec(
        deck_id="deck_test_export",
        meta=DeckMeta(
            title="季度经营复盘",
            subtitle="围绕数据、风险和行动项的结构化汇报",
            created_at="2026-04-04T10:00:00+0800",
            session_id="session-1",
            source_mode="kb_plus_chat",
            generator_panel_id="panel-main",
            author="tester",
            audience="leadership",
            purpose="review",
        ),
        generation=DeckGeneration(
            source="kb_plus_chat",
            target_slide_count=4,
            actual_slide_count=4,
        ),
        slides=[
            DeckSlide(
                id="slide_cover",
                type="cover",
                title="季度经营复盘",
                subtitle="收入增长与交付风险并存",
                layout="hero-title",
                intent="cover",
                speaker_notes="先讲总体结论，再展开两个核心风险点。",
                blocks=[
                    DeckBlock(
                        id="cover_core",
                        kind="paragraph",
                        role="core_message",
                        content={"text": "Q1 收入达成较好，但项目交付和毛利率承压。"},
                    )
                ],
                quality_state="weak_support",
                status=DeckSlideStatus(review_state="draft"),
            ),
            DeckSlide(
                id="slide_outline",
                type="outline",
                title="汇报结构",
                subtitle="先总览，再拆重点",
                layout="title-bullets",
                intent="outline",
                speaker_notes="目录页只做导航，不展开细节。",
                blocks=[
                    DeckBlock(
                        id="outline_block",
                        kind="bullet_list",
                        role="outline",
                        content={"items": ["整体表现", "核心风险", "行动建议"]},
                    )
                ],
                quality_state="weak_support",
                status=DeckSlideStatus(review_state="draft"),
            ),
            DeckSlide(
                id="slide_content_1",
                type="content",
                title="核心风险与改进方向",
                subtitle="交付效率、资源配置、毛利率",
                layout="title-bullets",
                intent="risk_review",
                speaker_notes="强调这页要先给结论，再落到责任人与时间点。",
                blocks=[
                    DeckBlock(
                        id="content_points",
                        kind="bullet_list",
                        role="main_points",
                        content={
                            "items": [
                                "重点项目延期 2 周，导致回款节奏后移。",
                                "人力投放集中在高复杂度需求，毛利率被稀释。",
                                "建议以项目群方式统一排期，并提高交付复用率。",
                            ]
                        },
                    )
                ],
                evidence_refs=[
                    DeckEvidenceRef(
                        id="ev-1",
                        source_id="src-1",
                        source_title="Q1经营分析.md",
                        snippet="重点项目延期 2 周，直接影响 4 月回款节奏。",
                        confidence=0.93,
                    ),
                    DeckEvidenceRef(
                        id="ev-2",
                        source_id="src-2",
                        source_title="项目周报-交付风险.xlsx",
                        snippet="交付团队反馈复杂需求占用核心人力，复用率不足。",
                        confidence=0.88,
                    )
                ],
                quality_state="supported",
                status=DeckSlideStatus(review_state="draft"),
            ),
            DeckSlide(
                id="slide_appendix",
                type="appendix_sources",
                title="附录：来源清单",
                subtitle="供复核的文档来源",
                layout="title-bullets",
                intent="appendix_sources",
                speaker_notes="附录页不口头展开，需要时再回看。",
                blocks=[
                    DeckBlock(
                        id="appendix_sources",
                        kind="bullet_list",
                        role="sources",
                        content={
                            "items": [
                                "Q1经营分析.md",
                                "项目周报-交付风险.xlsx",
                                "客户反馈汇总.docx",
                            ]
                        },
                    )
                ],
                quality_state="supported",
                status=DeckSlideStatus(review_state="draft"),
            ),
        ],
        source_registry=[
            DeckSourceItem(id="src-1", type="doc", title="Q1经营分析.md"),
            DeckSourceItem(id="src-2", type="sheet", title="项目周报-交付风险.xlsx"),
        ],
    )

    pptx_bytes = export_deck_to_pptx(deck)
    presentation = Presentation(io.BytesIO(pptx_bytes))

    assert len(presentation.slides) == 4
    assert presentation.slide_width == 12192000
    assert presentation.slide_height == 6858000

    cover_text = _shape_text(presentation.slides[0])
    assert "季度经营复盘" in cover_text
    assert "Q1 收入达成较好，但项目交付和毛利率承压。" in cover_text

    outline_text = _shape_text(presentation.slides[1])
    assert "证据来源" not in outline_text
    assert "[1]" not in outline_text

    content_text = _shape_text(presentation.slides[2])
    assert "核心风险与改进方向" in content_text
    assert "关键摘要" in content_text
    assert "• 重点项目延期 2 周，导致回款节奏后移。 [1][2]" in content_text
    assert "证据来源" in content_text
    assert "[1] Q1经营分析.md (93%)" in content_text
    assert "[2] 项目周报-交付风险.xlsx (88%)" in content_text
    assert "证据充分" in content_text

    appendix_text = _shape_text(presentation.slides[3])
    assert "附录：来源清单" in appendix_text
    assert "共 3 个来源" in appendix_text
    assert "项目周报-交付风险.xlsx" in appendix_text

    notes_text = presentation.slides[2].notes_slide.notes_text_frame.text
    assert "强调这页要先给结论，再落到责任人与时间点。" in notes_text
    assert "证据来源" in notes_text
    assert "Q1经营分析.md" in notes_text
