import io
import os
from pathlib import Path

from fastapi.testclient import TestClient
from langchain_core.documents import Document
from pptx import Presentation

import backend.api_server as api_server
import backend.chat_store as chat_store
import backend.deck_service as deck_service


def _history_cls_for_db(db_path: Path):
    class TestSQLiteChatMessageHistory(chat_store.SQLiteChatMessageHistory):
        def __init__(self, session_id: str, db_path_arg: str = "./chat_history.db"):
            super().__init__(session_id=session_id, db_path=str(db_path))

    return TestSQLiteChatMessageHistory


def _deck(deck_id: str = "deck-api") -> deck_service.DeckSpec:
    return deck_service.DeckSpec(
        deck_id=deck_id,
        meta=deck_service.DeckMeta(
            title="Board Update",
            subtitle="Q2 snapshot",
            theme="default",
            created_at="2026-04-12T10:00:00+0800",
            session_id="session-report",
            source_mode="chat_only",
            generator_panel_id="panel-main",
            author="tester",
            audience="leaders",
            purpose="briefing",
        ),
        generation=deck_service.DeckGeneration(
            source="chat_only",
            target_slide_count=2,
            actual_slide_count=2,
        ),
        slides=[
            deck_service.DeckSlide(
                id="cover",
                type="cover",
                title="Board Update",
                subtitle="Q2 snapshot",
                layout="hero-title",
                blocks=[],
            ),
            deck_service.DeckSlide(
                id="content-1",
                type="content",
                title="Growth Signals",
                subtitle="What changed",
                layout="title-bullets",
                blocks=[
                    deck_service.DeckBlock(
                        id="content-block",
                        kind="bullet_list",
                        role="main_points",
                        content={"items": ["Revenue +18%", "Churn -4%"]},
                    )
                ],
            ),
        ],
        source_registry=[],
    )


def _panel_config_payload(panel_id: str = "panel-main") -> dict[str, object]:
    return {
        "panel_id": panel_id,
        "connection_type": "ollama",
        "provider": "ollama",
        "model": "qwen3.5-2B:latest",
        "base_url": "http://localhost:11434",
        "api_key": "",
        "temperature": 0.3,
        "agent_mode": "auto",
    }


def test_upload_documents_endpoint_stages_files_and_persists_task(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    task_store = api_server.SQLiteTaskStore(db_path=str(db_path))
    scheduled: list[object] = []
    task_id = ""

    monkeypatch.setattr(api_server, "_task_store", task_store)
    monkeypatch.setattr(api_server, "_tasks", {})
    monkeypatch.setattr(
        api_server,
        "_effective_vector_store_path",
        lambda path=None: f"resolved::{path or 'current'}",
    )
    monkeypatch.setattr(
        api_server.asyncio,
        "create_task",
        lambda coro: scheduled.append(coro) or object(),
    )

    client = TestClient(api_server.app)

    try:
        response = client.post(
            "/api/documents/upload",
            data={"vector_store_path": "kb_custom"},
            files=[
                ("files", ("alpha.txt", b"alpha", "text/plain")),
                ("files", ("brief.md", b"# brief", "text/markdown")),
            ],
        )

        assert response.status_code == 200
        payload = response.json()
        task_id = payload["task_id"]
        assert payload["task_type"] == "upload_documents"
        assert payload["status"] == api_server.TaskStatus.PENDING
        assert "resolved::kb_custom" in payload["message"]

        record = api_server._tasks[task_id]
        assert record.params["file_names"] == ["alpha.txt", "brief.md"]
        assert record.params["vector_store_path"] == "resolved::kb_custom"
        assert all(Path(path).exists() for path in record.params["temp_paths"])

        stored = task_store.get(task_id)
        assert stored is not None
        assert stored.task_type == "upload_documents"
        assert stored.params["file_names"] == ["alpha.txt", "brief.md"]
    finally:
        for coro in scheduled:
            coro.close()
        record = api_server._tasks.get(task_id)
        if record is not None:
            for temp_path in record.params.get("temp_paths", []):
                try:
                    os.remove(temp_path)
                except OSError:
                    pass


def test_upload_documents_endpoint_cleans_temp_files_when_task_setup_fails(
    monkeypatch,
    tmp_path,
):
    staged_paths: list[str] = []

    async def fake_stage_upload_files(files, **kwargs):
        paths = [tmp_path / "upload-1.txt", tmp_path / "upload-2.md"]
        for index, path in enumerate(paths, start=1):
            path.write_bytes(f"payload-{index}".encode("utf-8"))
        staged_paths.extend(str(path) for path in paths)
        return staged_paths, ["upload-1.txt", "upload-2.md"]

    monkeypatch.setattr(api_server, "stage_upload_files", fake_stage_upload_files)
    monkeypatch.setattr(
        api_server,
        "build_upload_documents_task_record",
        lambda **kwargs: (_ for _ in ()).throw(RuntimeError("record failed")),
    )
    monkeypatch.setattr(api_server, "_effective_vector_store_path", lambda path=None: "resolved::kb")

    client = TestClient(api_server.app)
    response = client.post(
        "/api/documents/upload",
        files=[("files", ("alpha.txt", b"alpha", "text/plain"))],
    )

    assert response.status_code == 500
    assert response.json()["detail"] == "record failed"
    assert staged_paths
    assert all(not Path(path).exists() for path in staged_paths)


def test_upload_documents_endpoint_rejects_unsupported_extension():
    client = TestClient(api_server.app)

    response = client.post(
        "/api/documents/upload",
        files=[("files", ("payload.exe", b"boom", "application/octet-stream"))],
    )

    assert response.status_code == 400
    assert "不支持的文件类型" in response.json()["detail"]


def test_download_report_endpoint_returns_pptx(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-report")
    history.add_user_message("Board Update")
    history.add_ai_message("Revenue improved and churn declined.")

    client = TestClient(api_server.app)
    response = client.get("/api/reports/download/session-report")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert 'attachment; filename="Board Update.pptx"' == response.headers["content-disposition"]

    presentation = Presentation(io.BytesIO(response.content))
    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Board Update"
    assert "Revenue improved and churn declined." in presentation.slides[1].placeholders[1].text


def test_generate_report_endpoint_returns_markdown(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-generate-report")
    history.add_user_message("Board Update")
    history.add_ai_message("Revenue improved and churn declined.")

    client = TestClient(api_server.app)
    response = client.post(
        "/api/reports/generate",
        json={"session_id": "session-generate-report"},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["title"] == "Board Update"
    assert "Board Update" in payload["markdown"]
    assert "Revenue improved and churn declined." in payload["markdown"]


def test_generate_report_endpoint_supports_answer_group_scope(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-scoped-report")
    history.add_user_message("Trend scan", answer_group_id="group-1")
    history.add_ai_message(
        "Panel A answer should be ignored.",
        panel_id="panel-a",
        answer_group_id="group-1",
    )
    history.add_ai_message(
        "Panel B research answer should be used.",
        panel_id="panel-b",
        answer_group_id="group-1",
        sources=[
            {
                "title": "Example Source",
                "url": "https://example.com/live",
                "snippet": "Fresh update from the web",
            }
        ],
        task_type="web_research",
        model_id="web_research",
    )

    client = TestClient(api_server.app)
    response = client.post(
        "/api/reports/generate",
        json={
            "session_id": "session-scoped-report",
            "answer_group_id": "group-1",
            "panel_id": "panel-b",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["title"] == "Trend scan"
    assert "Panel B research answer should be used." in payload["markdown"]
    assert "参考来源" in payload["markdown"]
    assert "https://example.com/live" in payload["markdown"]
    assert "Panel A answer should be ignored." not in payload["markdown"]


def test_download_report_endpoint_supports_answer_group_scope(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-scoped-download")
    history.add_user_message("Trend scan", answer_group_id="group-2")
    history.add_ai_message(
        "Old comparison answer.",
        panel_id="panel-a",
        answer_group_id="group-2",
    )
    history.add_ai_message(
        "Scoped research answer.",
        panel_id="panel-b",
        answer_group_id="group-2",
        sources=[{"title": "Live feed", "snippet": "Realtime update"}],
        task_type="web_research",
        model_id="web_research",
    )

    client = TestClient(api_server.app)
    response = client.get(
        "/api/reports/download/session-scoped-download",
        params={
            "answer_group_id": "group-2",
            "panel_id": "panel-b",
        },
    )

    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.content))
    assert len(presentation.slides) == 2
    slide_text = presentation.slides[1].placeholders[1].text
    assert "Scoped research answer." in slide_text
    assert "Old comparison answer." not in slide_text


def test_create_deck_endpoint_supports_answer_group_scope(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-scoped-deck")
    history.add_user_message("Trend scan", answer_group_id="group-3")
    history.add_ai_message(
        "Panel A answer should be ignored.",
        panel_id="panel-a",
        answer_group_id="group-3",
    )
    history.add_ai_message(
        "Panel B research answer should be used.",
        panel_id="panel-b",
        answer_group_id="group-3",
        sources=[{"title": "Live feed", "snippet": "Realtime update"}],
        task_type="web_research",
        model_id="web_research",
    )

    captured: dict[str, object] = {}

    async def fake_build_deck(**kwargs):
        captured.update(kwargs)
        deck = _deck("deck-scoped-create")
        deck.meta.session_id = kwargs["session_id"]
        deck.meta.source_answer_group_id = kwargs["source_answer_group_id"]
        deck.meta.source_panel_id = kwargs["source_panel_id"]
        return deck

    monkeypatch.setattr(api_server, "build_deck", fake_build_deck)

    client = TestClient(api_server.app)
    response = client.post(
        "/api/decks",
        json={
            "session_id": "session-scoped-deck",
            "panel_config": _panel_config_payload("panel-b"),
            "knowledge_base_enabled": False,
            "target_slide_count": 6,
            "theme": "default",
            "answer_group_id": "group-3",
            "panel_id": "panel-b",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["deck_id"] == "deck-scoped-create"
    assert payload["meta"]["source_answer_group_id"] == "group-3"
    assert payload["meta"]["source_panel_id"] == "panel-b"

    messages = captured["messages"]
    assert len(messages) == 2
    assert messages[0].content == "Trend scan"
    assert "Panel B research answer should be used." in messages[1].content
    assert "Panel A answer should be ignored." not in messages[1].content
    assert captured["source_answer_group_id"] == "group-3"
    assert captured["source_panel_id"] == "panel-b"


def test_regenerate_deck_endpoint_reuses_saved_scope(monkeypatch, tmp_path):
    db_path = tmp_path / "chat_history.db"
    history_cls = _history_cls_for_db(db_path)
    monkeypatch.setattr(chat_store, "SQLiteChatMessageHistory", history_cls)

    history = history_cls("session-scoped-regenerate")
    history.add_user_message("Trend scan", answer_group_id="group-4")
    history.add_ai_message(
        "Panel A answer should be ignored.",
        panel_id="panel-a",
        answer_group_id="group-4",
    )
    history.add_ai_message(
        "Scoped research answer.",
        panel_id="panel-b",
        answer_group_id="group-4",
        sources=[{"title": "Live feed", "snippet": "Realtime update"}],
        task_type="web_research",
        model_id="web_research",
    )

    deck = _deck("deck-scoped-regenerate")
    deck.meta.session_id = "session-scoped-regenerate"
    deck.meta.source_answer_group_id = "group-4"
    deck.meta.source_panel_id = "panel-b"
    api_server._deck_store.save(deck)

    captured: dict[str, object] = {}

    async def fake_regenerate_deck_slide(*, messages, **kwargs):
        captured["messages"] = messages
        captured["kwargs"] = kwargs
        return deck.slides[1]

    monkeypatch.setattr(
        api_server,
        "build_regenerate_deck_kwargs",
        lambda *args, **kwargs: {
            "panel_config": args[1].panel_config,
            "knowledge_base_enabled": False,
        },
    )
    monkeypatch.setattr(api_server, "regenerate_deck_slide", fake_regenerate_deck_slide)

    client = TestClient(api_server.app)
    response = client.post(
        "/api/decks/deck-scoped-regenerate/slides/content-1/regenerate",
        json={"panel_config": _panel_config_payload("panel-b")},
    )

    assert response.status_code == 200
    messages = captured["messages"]
    assert len(messages) == 2
    assert messages[0].content == "Trend scan"
    assert "Scoped research answer." in messages[1].content
    assert "Panel A answer should be ignored." not in messages[1].content


def test_update_deck_endpoint_updates_title_theme_and_cover():
    deck = _deck("deck-update")
    api_server._deck_store.save(deck)

    client = TestClient(api_server.app)
    response = client.patch(
        "/api/decks/deck-update",
        json={
            "title": "Updated Board",
            "theme": "midnight",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["meta"]["title"] == "Updated Board"
    assert payload["meta"]["theme"] == "midnight"
    assert payload["slides"][0]["title"] == "Updated Board"


def test_export_deck_endpoint_returns_pptx_bytes():
    deck = _deck("deck-export")
    api_server._deck_store.save(deck)

    client = TestClient(api_server.app)
    response = client.get("/api/decks/deck-export/export")

    assert response.status_code == 200
    assert response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert 'attachment; filename="Board Update.pptx"' == response.headers["content-disposition"]

    presentation = Presentation(io.BytesIO(response.content))
    assert len(presentation.slides) == 2


def test_test_retrieval_endpoint_returns_results(monkeypatch):
    import backend.doc_pipeline as doc_pipeline

    events: list[tuple[str, object]] = []

    class FakePipeline:
        def __init__(self, vector_store_path=None):
            self.vector_store_path = vector_store_path
            events.append(("init", vector_store_path))

        def load_store(self):
            events.append(("load_store", self.vector_store_path))
            return True

        def search(self, query, k=5):
            events.append(("search", query, k))
            return [
                Document(
                    page_content="Alpha findings with detail",
                    metadata={"source": "alpha.md"},
                )
            ]

    monkeypatch.setattr(doc_pipeline, "DocPipeline", FakePipeline)
    monkeypatch.setattr(
        api_server,
        "_effective_vector_store_path",
        lambda path=None: f"resolved::{path or 'current'}",
    )

    client = TestClient(api_server.app)
    response = client.post(
        "/api/knowledge-base/test-retrieval",
        json={"query": "alpha", "vector_store_path": "kb_custom"},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["results_count"] == 1
    assert payload["top_results"][0]["source"] == "alpha.md"
    assert payload["top_results"][0]["snippet"] == "Alpha findings with detail"
    assert payload["latency_ms"] >= 0
    assert events == [
        ("init", "resolved::kb_custom"),
        ("load_store", "resolved::kb_custom"),
        ("search", "alpha", 5),
    ]


def test_test_retrieval_endpoint_supports_hybrid_mode(monkeypatch):
    import backend.doc_pipeline as doc_pipeline

    events: list[tuple[str, object]] = []

    class FakePipeline:
        def __init__(self, vector_store_path=None):
            self.vector_store_path = vector_store_path
            events.append(("init", vector_store_path))

        def load_store(self):
            events.append(("load_store", self.vector_store_path))
            return True

        def debug_retrieval(
            self,
            query,
            *,
            search_k,
            fetch_k,
            retrieval_mode,
            use_rerank,
        ):
            events.append(
                ("debug_retrieval", query, search_k, fetch_k, retrieval_mode, use_rerank)
            )
            return {
                "results_count": 1,
                "search_mode": "hybrid_rerank",
                "retrieval_mode": "hybrid",
                "search_k": search_k,
                "top_k": search_k,
                "fetch_k": fetch_k,
                "rewrite_query": query,
                "rewrite_applied": False,
                "query_terms": ["alpha"],
                "coverage": {
                    "unique_sources": 1,
                    "source_ratio": 1.0,
                    "matched_terms": ["alpha"],
                    "matched_term_count": 1,
                },
                "top_results": [
                    {
                        "rank": 1,
                        "source": "alpha.md",
                        "snippet": "Alpha findings with detail",
                        "score": 0.91,
                        "channel": "hybrid_rerank",
                    }
                ],
                "semantic_candidates": [],
                "keyword_candidates": [],
                "fused_candidates": [],
            }

    monkeypatch.setattr(doc_pipeline, "DocPipeline", FakePipeline)
    monkeypatch.setattr(api_server, "_effective_vector_store_path", lambda path=None: "resolved::kb")

    client = TestClient(api_server.app)
    response = client.post(
        "/api/knowledge-base/test-retrieval",
        json={
            "query": "alpha",
            "retrieval_mode": "hybrid",
            "search_k": 3,
            "fetch_k": 7,
            "use_rerank": True,
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["search_mode"] == "hybrid_rerank"
    assert payload["retrieval_mode"] == "hybrid"
    assert payload["coverage"]["matched_term_count"] == 1
    assert payload["top_results"][0]["source"] == "alpha.md"
    assert events == [
        ("init", "resolved::kb"),
        ("load_store", "resolved::kb"),
        ("debug_retrieval", "alpha", 3, 7, "hybrid", True),
    ]


def test_test_retrieval_endpoint_rejects_blank_query(monkeypatch):
    import backend.doc_pipeline as doc_pipeline

    class FakePipeline:
        def __init__(self, vector_store_path=None):
            self.vector_store_path = vector_store_path

    monkeypatch.setattr(doc_pipeline, "DocPipeline", FakePipeline)
    monkeypatch.setattr(api_server, "_effective_vector_store_path", lambda path=None: "resolved::kb")

    client = TestClient(api_server.app)
    response = client.post(
        "/api/knowledge-base/test-retrieval",
        json={"query": "   ", "vector_store_path": "kb_custom"},
    )

    assert response.status_code == 400
    assert response.json()["detail"]
