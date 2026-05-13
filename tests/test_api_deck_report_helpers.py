import io
from types import SimpleNamespace

import pytest
from langchain_core.messages import AIMessage, HumanMessage
from pptx import Presentation
from pptx.util import Pt

import backend.helpers.deck_report_helpers as api_deck_report_helpers
import backend.deck_service as deck_service


def _deck() -> deck_service.DeckSpec:
    return deck_service.DeckSpec(
        deck_id="deck-helper",
        meta=deck_service.DeckMeta(
            title="Board Update",
            subtitle="Q2 snapshot",
            theme="default",
            created_at="2026-04-12T10:00:00+0800",
            session_id="session-1",
            source_mode="chat_only",
            generator_panel_id="panel-main",
            author="tester",
            audience="leaders",
            purpose="briefing",
        ),
        generation=deck_service.DeckGeneration(
            source="chat_only",
            target_slide_count=2,
            actual_slide_count=2,
        ),
        slides=[
            deck_service.DeckSlide(
                id="cover",
                type="cover",
                title="Board Update",
                subtitle="Q2 snapshot",
                layout="hero-title",
                blocks=[],
            ),
            deck_service.DeckSlide(
                id="content-1",
                type="content",
                title="Growth",
                subtitle="What changed",
                layout="title-bullets",
                blocks=[],
            ),
        ],
        source_registry=[],
    )


def test_create_share_link_payload_builds_token_and_url():
    payload = api_deck_report_helpers.create_share_link_payload(
        "deck",
        "deck-1",
        SimpleNamespace(base_url="ignored"),
        secret="secret",
        encode_share_token=lambda resource_type, resource_id, secret: (
            f"{resource_type}:{resource_id}:{secret}"
        ),
        build_share_url=lambda request, token: f"http://test/{token}",
    )

    assert payload == {
        "resource_type": "deck",
        "resource_id": "deck-1",
        "share_token": "deck:deck-1:secret",
        "share_url": "http://test/deck:deck-1:secret",
    }


def test_apply_report_template_metadata_updates_existing_frontmatter_idempotently():
    markdown = (
        "---\n"
        "template: legacy_template\n"
        "template_options_json: '{\"legacy\": true}'\n"
        "theme: default\n"
        "---\n"
        "\n"
        "# Report\n"
    )

    updated = api_deck_report_helpers.apply_report_template_metadata(
        markdown,
        template_id="executive_report",
        template_options={
            "scope": "answer_group",
            "include_citations": True,
            "nested": {"ignored": True},
        },
    )
    reapplied = api_deck_report_helpers.apply_report_template_metadata(
        updated,
        template_id="executive_report",
        template_options={
            "scope": "answer_group",
            "include_citations": True,
            "nested": {"ignored": True},
        },
    )

    assert updated == reapplied
    assert updated.count("template: executive_report") == 1
    assert updated.count("template_options_json:") == 1
    assert "legacy_template" not in updated
    assert '"legacy": true' not in updated
    assert updated.splitlines()[:5] == [
        "---",
        "template: executive_report",
        'template_options_json: \'{"include_citations": true, "scope": "answer_group"}\'',
        "theme: default",
        "---",
    ]


def test_build_create_deck_kwargs_resolves_prompt_runtime_and_theme():
    request = SimpleNamespace(
        session_id="session-1",
        panel_config="panel-config",
        knowledge_base_enabled=True,
        target_slide_count=6,
        theme="sunrise",
        answer_group_id=" group-1 ",
        panel_id=" panel-main ",
    )

    kwargs = api_deck_report_helpers.build_create_deck_kwargs(
        request,
        resolve_active_prompt_runtime=lambda enabled: (
            "prompt-content" if enabled else None,
            "kb-path" if enabled else None,
            {},
        ),
        normalize_deck_theme=lambda theme: theme.upper(),
    )

    assert kwargs == {
        "session_id": "session-1",
        "panel_config": "panel-config",
        "knowledge_base_enabled": True,
        "target_slide_count": 6,
        "vector_store_path": "kb-path",
        "system_prompt": "prompt-content",
        "theme": "SUNRISE",
        "source_answer_group_id": "group-1",
        "source_panel_id": "panel-main",
    }


def test_apply_deck_update_updates_cover_title_theme_and_slide_count():
    deck = _deck()
    replacement_slides = [
        deck_service.DeckSlide(
            id="cover",
            type="cover",
            title="Old title",
            subtitle="New subtitle",
            layout="hero-title",
            blocks=[],
        )
    ]
    request = SimpleNamespace(
        title="Updated Board",
        theme="midnight",
        slides=replacement_slides,
    )

    api_deck_report_helpers.apply_deck_update(
        deck,
        request,
        normalize_deck_theme=lambda theme: theme.upper(),
    )

    assert deck.meta.title == "Updated Board"
    assert deck.slides[0].title == "Updated Board"
    assert deck.meta.theme == "MIDNIGHT"
    assert deck.generation.actual_slide_count == 1


def test_build_regenerate_deck_kwargs_defaults_kb_mode_from_deck_source():
    deck = _deck()
    deck.meta.source_mode = "kb_plus_chat"
    request = SimpleNamespace(
        panel_config="panel-config",
        knowledge_base_enabled=None,
    )

    kwargs = api_deck_report_helpers.build_regenerate_deck_kwargs(
        deck,
        request,
        normalize_model_config=lambda config: {"normalized": config},
        resolve_active_prompt_runtime=lambda enabled: (
            "prompt-content" if enabled else None,
            "kb-path" if enabled else None,
            {},
        ),
    )

    assert kwargs == {
        "panel_config": {"normalized": "panel-config"},
        "knowledge_base_enabled": True,
        "vector_store_path": "kb-path",
        "system_prompt": "prompt-content",
    }


def test_replace_deck_slide_swaps_only_matching_slide():
    deck = _deck()
    regenerated = deck_service.DeckSlide(
        id="content-1",
        type="content",
        title="Regenerated",
        subtitle="Fresh",
        layout="title-bullets",
        blocks=[],
    )

    api_deck_report_helpers.replace_deck_slide(deck, regenerated)

    assert deck.slides[0].title == "Board Update"
    assert deck.slides[1].title == "Regenerated"


def test_export_deck_payload_delegates_to_exporter_and_filename_builder():
    deck = _deck()

    payload = api_deck_report_helpers.export_deck_payload(
        deck,
        export_deck_to_pptx=lambda deck_obj: f"pptx:{deck_obj.deck_id}".encode("utf-8"),
        build_export_filename=lambda deck_obj, extension: f"{deck_obj.deck_id}.{extension}",
    )

    assert payload == {
        "content": b"pptx:deck-helper",
        "filename": "deck-helper.pptx",
        "evidence_coverage": {
            "total_slides": 2,
            "coverable_slide_count": 1,
            "slides_with_evidence": 0,
            "total_evidence_refs": 0,
            "coverage_ratio": 0.0,
            "unsupported_slide_ids": ["content-1"],
            "slides": [
                {
                    "slide_id": "cover",
                    "slide_type": "cover",
                    "evidence_ref_count": 0,
                    "has_evidence": False,
                    "is_coverable": False,
                    "quality_state": "weak_support",
                },
                {
                    "slide_id": "content-1",
                    "slide_type": "content",
                    "evidence_ref_count": 0,
                    "has_evidence": False,
                    "is_coverable": True,
                    "quality_state": "weak_support",
                },
            ],
        },
        "evidence_review": {
            "status": "needs_review",
            "coverage_ratio": 0.0,
            "coverable_slide_count": 1,
            "slides_with_evidence": 0,
            "unsupported_slide_ids": ["content-1"],
            "needs_review_slide_ids": ["content-1"],
            "action_count": 2,
            "action_items": [
                {
                    "code": "add_missing_slide_evidence",
                    "severity": "warning",
                    "message": "Add evidence references to unsupported slides.",
                    "slide_ids": ["content-1"],
                },
                {
                    "code": "attach_deck_sources",
                    "severity": "warning",
                    "message": "Attach at least one source before final export.",
                    "slide_ids": ["content-1"],
                },
            ],
            "slides": [
                {
                    "slide_id": "cover",
                    "title": "Board Update",
                    "slide_type": "cover",
                    "is_coverable": False,
                    "has_evidence": False,
                    "evidence_ref_count": 0,
                    "quality_state": "weak_support",
                    "needs_review": False,
                    "source_ids": [],
                    "source_titles": [],
                },
                {
                    "slide_id": "content-1",
                    "title": "Growth",
                    "slide_type": "content",
                    "is_coverable": True,
                    "has_evidence": False,
                    "evidence_ref_count": 0,
                    "quality_state": "weak_support",
                    "needs_review": True,
                    "source_ids": [],
                    "source_titles": [],
                },
            ],
            "citation_validation": {
                "status": "passed",
                "can_export": True,
                "issue_count": 0,
                "missing_source_ids": [],
                "missing_block_evidence_ref_ids": [],
                "issues": [],
            },
        },
        "citation_validation": {
            "status": "passed",
            "can_export": True,
            "issue_count": 0,
            "missing_source_ids": [],
            "missing_block_evidence_ref_ids": [],
            "issues": [],
        },
        "export_gate": {
            "blocked": False,
            "overridden": False,
            "can_export": True,
            "reason": "",
            "message": "",
        },
    }


def test_build_deck_evidence_review_payload_reports_supported_sources():
    deck = _deck()
    deck.source_registry = [
        deck_service.DeckSourceItem(
            id="src-1",
            type="doc",
            title="Q2 Board Memo",
        )
    ]
    deck.slides[1].evidence_refs = [
        deck_service.DeckEvidenceRef(
            id="ev-1",
            source_id="src-1",
            source_title="Q2 Board Memo",
            snippet="Revenue rose in Q2.",
            confidence=0.92,
        )
    ]
    deck.slides[1].quality_state = "supported"

    payload = api_deck_report_helpers.build_deck_evidence_review_payload(deck)

    assert payload["status"] == "supported"
    assert payload["coverage_ratio"] == 1.0
    assert payload["unsupported_slide_ids"] == []
    assert payload["action_items"] == []
    content_review = payload["slides"][1]
    assert content_review["slide_id"] == "content-1"
    assert content_review["source_ids"] == ["src-1"]
    assert content_review["source_titles"] == ["Q2 Board Memo"]
    assert content_review["needs_review"] is False
    assert payload["citation_validation"]["status"] == "passed"


def test_deck_citation_validation_reports_missing_sources_and_block_refs():
    deck = _deck()
    deck.slides[1].evidence_refs = [
        deck_service.DeckEvidenceRef(
            id="ev-1",
            source_id="missing-src",
            source_title="Missing Source",
            snippet="Revenue rose in Q2.",
            confidence=0.92,
        )
    ]
    deck.slides[1].blocks = [
        deck_service.DeckBlock(
            id="block-1",
            kind="bullet_list",
            role="main_points",
            content={
                "items": ["Revenue rose"],
                "evidence_ref_ids": ["ev-1", "ev-missing"],
            },
        )
    ]

    validation = deck_service.validate_deck_citation_consistency(deck)

    assert validation.status == "failed"
    assert validation.can_export is False
    assert validation.missing_source_ids == ["missing-src"]
    assert validation.missing_block_evidence_ref_ids == ["ev-missing"]
    assert [issue.code for issue in validation.issues] == [
        "missing_source_registry_entry",
        "missing_slide_evidence_ref",
    ]


def test_update_deck_block_refs_refreshes_citation_gate_payload():
    deck = _deck()
    deck.source_registry = [
        deck_service.DeckSourceItem(id="src-1", type="doc", title="Q2 Board Memo")
    ]
    deck.slides[1].evidence_refs = [
        deck_service.DeckEvidenceRef(
            id="ev-1",
            source_id="src-1",
            source_title="Q2 Board Memo",
            snippet="Revenue rose in Q2.",
            confidence=0.92,
        )
    ]
    deck.slides[1].blocks = [
        deck_service.DeckBlock(
            id="block-1",
            kind="paragraph",
            role="summary",
            content={"text": "Revenue rose", "evidence_ref_ids": ["ev-missing"]},
        )
    ]

    before = deck_service.validate_deck_citation_consistency(deck)
    result = api_deck_report_helpers.update_deck_block_refs(
        deck,
        "content-1",
        "block-1",
        {
            "evidence_ref_ids": ["ev-1", "ev-1", ""],
            "source_ref_ids": ["src-1"],
        },
    )

    assert before.status == "failed"
    assert deck.slides[1].blocks[0].content["evidence_ref_ids"] == ["ev-1"]
    assert deck.slides[1].blocks[0].content["evidence_source_ids"] == ["src-1"]
    assert deck.slides[1].status.dirty is True
    assert result["citation_validation"]["status"] == "passed"
    assert result["export_gate"] == {
        "blocked": False,
        "overridden": False,
        "can_export": True,
        "reason": "",
        "message": "",
    }
    assert result["evidence_review"]["citation_validation"]["status"] == "passed"


def test_export_deck_payload_blocks_failed_citation_validation_by_default():
    deck = _deck()
    deck.slides[1].blocks = [
        deck_service.DeckBlock(
            id="block-1",
            kind="paragraph",
            role="summary",
            content={"text": "Revenue rose", "evidence_ref_ids": ["ev-missing"]},
        )
    ]

    with pytest.raises(api_deck_report_helpers.DeckExportGateError) as exc_info:
        api_deck_report_helpers.export_deck_payload(
            deck,
            export_deck_to_pptx=lambda deck_obj: f"pptx:{deck_obj.deck_id}".encode("utf-8"),
            build_export_filename=lambda deck_obj, extension: f"{deck_obj.deck_id}.{extension}",
        )

    payload = exc_info.value.payload
    validation = payload["citation_validation"]
    assert validation["status"] == "failed"
    assert validation["can_export"] is False
    assert validation["missing_block_evidence_ref_ids"] == ["ev-missing"]
    assert payload["evidence_review"]["citation_validation"] == validation
    assert payload["export_gate"] == {
        "blocked": True,
        "overridden": False,
        "can_export": False,
        "reason": "",
        "message": "Deck export is blocked because citation validation failed.",
    }


def test_export_deck_payload_allows_explicit_unsafe_override_with_metadata():
    deck = _deck()
    deck.slides[1].blocks = [
        deck_service.DeckBlock(
            id="block-1",
            kind="paragraph",
            role="summary",
            content={"text": "Revenue rose", "evidence_ref_ids": ["ev-missing"]},
        )
    ]

    payload = api_deck_report_helpers.export_deck_payload(
        deck,
        export_deck_to_pptx=lambda deck_obj: f"pptx:{deck_obj.deck_id}".encode("utf-8"),
        build_export_filename=lambda deck_obj, extension: f"{deck_obj.deck_id}.{extension}",
        allow_unsafe_export=True,
        override_reason="manual legal review",
    )

    assert payload["content"] == b"pptx:deck-helper"
    assert payload["citation_validation"]["can_export"] is False
    assert payload["export_gate"] == {
        "blocked": False,
        "overridden": True,
        "can_export": False,
        "reason": "manual legal review",
        "message": "",
    }


def test_report_markdown_payload_uses_title_builder_and_renderer():
    messages = [HumanMessage(content="Board Update"), AIMessage(content="Revenue up.")]

    payload = api_deck_report_helpers.report_markdown_payload(
        messages,
        ensure_deckable_chat=lambda msgs: [("Board Update", "Revenue up.")],
        build_chat_report_title=lambda msgs: "Board Update",
        build_report_markdown=lambda msgs, title: f"# {title}\n\nRevenue up.",
    )

    assert payload == {
        "markdown": "# Board Update\n\nRevenue up.",
        "title": "Board Update",
    }


def test_report_download_payload_builds_presentation_and_filename():
    messages = [HumanMessage(content="Board Update"), AIMessage(content="Revenue up.")]

    payload = api_deck_report_helpers.report_download_payload(
        messages,
        ensure_deckable_chat=lambda msgs: [("Board Update", "Revenue up.")],
        build_chat_report_title=lambda msgs: "Board Update",
        presentation_factory=Presentation,
        body_font_size=Pt(12),
        populate_chat_report_presentation=lambda prs, **kwargs: (
            prs.slides.add_slide(prs.slide_layouts[0]),
            setattr(prs.slides[0].shapes.title, "text", kwargs["title"]),
        ),
        safe_report_filename=lambda title: title.replace(" ", "_"),
    )

    assert payload["filename"] == "Board_Update.pptx"
    presentation = payload["presentation"]
    buffer = io.BytesIO()
    presentation.save(buffer)
    assert buffer.getvalue()


def test_build_scoped_report_messages_selects_requested_panel_and_appends_sources():
    records = [
        {
            "type": "human",
            "content": "Board Update",
            "answer_group_id": "group-1",
            "timestamp": 1,
        },
        {
            "type": "ai",
            "content": "Other panel answer.",
            "panel_id": "panel-alt",
            "answer_group_id": "group-1",
            "timestamp": 2,
            "sources": [],
            "task_type": "",
            "model_id": "assistant",
        },
        {
            "type": "ai",
            "content": "Research answer.",
            "panel_id": "panel-main",
            "answer_group_id": "group-1",
            "timestamp": 3,
            "sources": [
                {
                    "title": "Example Source",
                    "url": "https://example.com/report",
                    "snippet": "Fresh market signal",
                }
            ],
            "task_type": "web_research",
            "model_id": "web_research",
        },
    ]

    messages = api_deck_report_helpers.build_scoped_report_messages(
        records,
        answer_group_id="group-1",
        panel_id="panel-main",
        human_message_factory=lambda content: HumanMessage(content=content),
        ai_message_factory=lambda content: AIMessage(content=content),
    )

    assert [message.__class__.__name__ for message in messages] == [
        "HumanMessage",
        "AIMessage",
    ]
    assert messages[0].content == "Board Update"
    assert "Research answer." in messages[1].content
    assert "参考来源" in messages[1].content
    assert "https://example.com/report" in messages[1].content
    assert "Other panel answer." not in messages[1].content


def test_build_scoped_report_messages_prefers_web_research_when_panel_is_omitted():
    records = [
        {
            "type": "human",
            "content": "Trend scan",
            "answer_group_id": "group-2",
            "timestamp": 1,
        },
        {
            "type": "ai",
            "content": "Generic answer.",
            "panel_id": "panel-a",
            "answer_group_id": "group-2",
            "timestamp": 5,
            "sources": [],
            "task_type": "",
            "model_id": "assistant",
        },
        {
            "type": "ai",
            "content": "Research-first answer.",
            "panel_id": "panel-b",
            "answer_group_id": "group-2",
            "timestamp": 2,
            "sources": [{"title": "Live feed", "snippet": "Realtime update"}],
            "task_type": "web_research",
            "model_id": "web_research",
        },
    ]

    messages = api_deck_report_helpers.build_scoped_report_messages(
        records,
        answer_group_id="group-2",
        human_message_factory=lambda content: HumanMessage(content=content),
        ai_message_factory=lambda content: AIMessage(content=content),
    )

    assert messages[0].content == "Trend scan"
    assert "Research-first answer." in messages[1].content
    assert "Generic answer." not in messages[1].content


def test_build_scoped_report_messages_preserves_structured_evidence_metadata():
    records = [
        {
            "type": "human",
            "content": "Build a deck",
            "answer_group_id": "group-evidence",
            "timestamp": 1,
        },
        {
            "type": "ai",
            "content": "Evidence-backed answer.",
            "panel_id": "panel-main",
            "answer_group_id": "group-evidence",
            "timestamp": 2,
            "sources": [
                {
                    "title": "Primary source",
                    "url": "https://example.com/source",
                    "snippet": "Important cited fact.",
                }
            ],
            "claim_evidence_chains": [
                {
                    "claim_id": "claim-1",
                    "claim_text": "Important cited fact.",
                    "evidence_strength": "high",
                    "sources": [{"title": "Primary source"}],
                }
            ],
            "task_type": "web_research",
            "model_id": "web_research",
        },
    ]

    messages = api_deck_report_helpers.build_scoped_report_messages(
        records,
        answer_group_id="group-evidence",
        panel_id="panel-main",
        human_message_factory=lambda content: HumanMessage(content=content),
        ai_message_factory=lambda content: AIMessage(content=content),
    )

    assert messages[1].additional_kwargs["sources"][0]["title"] == "Primary source"
    assert (
        messages[1].additional_kwargs["claim_evidence_chains"][0]["claim_text"]
        == "Important cited fact."
    )
