import io
from types import SimpleNamespace

from pptx import Presentation

from deck_service import (
    DeckBlock,
    DeckGeneration,
    DeckMeta,
    DeckSlide,
    DeckSlideStatus,
    DeckSpec,
    DraftedContentSlide,
    DraftedSlideBundle,
    OutlinePlan,
    OutlineSlidePlan,
    SourcePack,
    _build_deck_from_generated,
    export_deck_to_pptx,
)


def _shape_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if text:
            parts.append(text)
    return "\n".join(parts)


def test_export_deck_to_pptx_renders_chart_blocks_and_manual_confirmation():
    deck = DeckSpec(
        deck_id="deck_chart_export",
        meta=DeckMeta(
            title="经营看板",
            subtitle="包含图表的导出验证",
            created_at="2026-04-13T10:00:00+0800",
            session_id="session-chart",
            source_mode="chat_only",
            generator_panel_id="panel-main",
            author="tester",
            audience="leadership",
            purpose="review",
        ),
        generation=DeckGeneration(
            source="chat_only",
            target_slide_count=3,
            actual_slide_count=3,
        ),
        slides=[
            DeckSlide(
                id="slide_cover",
                type="cover",
                title="经营看板",
                subtitle="包含图表的导出验证",
                layout="hero-title",
                intent="cover",
                speaker_notes="封面页",
                blocks=[
                    DeckBlock(
                        id="cover_core",
                        kind="paragraph",
                        role="core_message",
                        content={"text": "这是一个包含 dashboard 图表的 deck。"},
                    )
                ],
                quality_state="manual",
                status=DeckSlideStatus(review_state="draft"),
            ),
            DeckSlide(
                id="slide_content_1",
                type="content",
                title="收入趋势",
                subtitle="按季度观察增长变化",
                layout="title-bullets",
                intent="dashboard_summary",
                speaker_notes="这页虽然来源是聊天模式，但已人工确认。",
                blocks=[
                    DeckBlock(
                        id="content_points",
                        kind="bullet_list",
                        role="main_points",
                        content={"items": ["收入逐季增长", "Q4 增幅最明显"]},
                    ),
                    DeckBlock(
                        id="content_chart",
                        kind="chart",
                        role="dashboard_chart",
                        content={
                            "title": "季度收入趋势",
                            "description": "图表来自 dashboard-card",
                            "chart_type": "bar",
                            "labels": ["Q1", "Q2", "Q3", "Q4"],
                            "datasets": [
                                {
                                    "label": "收入",
                                    "data": [120.0, 148.0, 176.0, 210.0],
                                }
                            ],
                        },
                        editable=False,
                    ),
                ],
                quality_state="manual",
                status=DeckSlideStatus(review_state="confirmed"),
            ),
            DeckSlide(
                id="slide_appendix",
                type="appendix_sources",
                title="附录",
                subtitle="",
                layout="title-bullets",
                intent="appendix",
                speaker_notes="",
                blocks=[],
                quality_state="supported",
                status=DeckSlideStatus(review_state="draft"),
            ),
        ],
        source_registry=[],
    )

    pptx_bytes = export_deck_to_pptx(deck)
    presentation = Presentation(io.BytesIO(pptx_bytes))
    content_slide = presentation.slides[1]

    assert any(getattr(shape, "has_chart", False) for shape in content_slide.shapes)
    assert "已人工确认" in _shape_text(content_slide)
    assert "季度收入趋势" in _shape_text(content_slide)
    assert "已人工确认" in content_slide.notes_slide.notes_text_frame.text


def test_build_deck_from_generated_extracts_dashboard_chart_blocks_from_answer():
    pack = SourcePack(
        title_hint="经营看板",
        source_mode="chat_only",
        qa_pairs=[
            (
                "请总结经营情况",
                "\n".join(
                    [
                        "整体经营稳步增长，建议重点关注收入趋势。",
                        ":::dashboard-card",
                        '{"title":"经营总览","charts":[{"title":"季度收入趋势","type":"bar","description":"来自 dashboard-card","chart_data":{"type":"bar","labels":["Q1","Q2","Q3"],"datasets":[{"label":"收入","data":[120,150,180]}]}}]}',
                        ":::",
                    ]
                ),
            )
        ],
        chat_notes=[],
        excerpts=[],
        source_registry=[],
        warnings=[],
    )
    outline = OutlinePlan(
        title="经营看板",
        subtitle="按季度复盘",
        core_message="收入持续增长",
        sections=["总览"],
        content_slides=[
            OutlineSlidePlan(
                title="收入趋势",
                objective="说明季度收入变化",
                section="总览",
                evidence_source_ids=[],
            )
        ],
    )
    drafted = DraftedSlideBundle(
        content_slides=[
            DraftedContentSlide(
                title="收入趋势",
                subtitle="按季度观察增长变化",
                key_points=["收入逐季增长", "Q3 增长明显"],
                speaker_notes="先讲趋势，再讲原因。",
                evidence_excerpt_ids=[],
                evidence_source_ids=[],
                quality_state="manual",
            )
        ]
    )

    deck = _build_deck_from_generated(
        session_id="session-chart",
        panel_config=SimpleNamespace(panel_id="panel-main"),
        target_slide_count=3,
        pack=pack,
        outline=outline,
        drafted=drafted,
    )

    content_slide = next(slide for slide in deck.slides if slide.type == "content")
    chart_blocks = [block for block in content_slide.blocks if block.kind == "chart"]

    assert len(chart_blocks) == 1
    assert chart_blocks[0].content["chart_type"] == "bar"
    assert chart_blocks[0].content["labels"] == ["Q1", "Q2", "Q3"]
