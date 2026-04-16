"""
FastAPI 后端 API 服务
提供 REST + SSE 端点，包装现有 agent_core / chat_store / doc_pipeline 模块
"""

import asyncio
from dataclasses import dataclass
import hashlib
import ipaddress
import json
import logging
import os
from pathlib import Path
import re
import secrets
import shutil
import sys
import threading
import time
from urllib.parse import quote
import uuid
from typing import Any, AsyncGenerator, Literal, Optional

import httpx
from artifact_service import (
    SQLiteArtifactStore,
    artifact_export_formats,
    build_deck_artifact,
    build_report_artifact,
    sync_deck_artifact,
)
from agent_mcp_helpers import (
    default_mcp_server_names,
    list_mcp_server_catalog,
    normalize_mcp_server_names,
)
from api_session_helpers import (
    build_answer_group_review_payload as _build_answer_group_review_payload,
    build_session_messages_payload as _build_session_messages_payload,
    collect_session_attachments as _collect_session_attachments,
    find_session_attachment as _find_session_attachment,
    render_shared_deck_html as _render_shared_deck_html,
    render_shared_session_html as _render_shared_session_html,
)
from api_attachment_route_helpers import (
    prepare_attachment_promotion,
    session_attachments_payload,
)
from api_chat_input_helpers import (
    build_message_with_files as _build_message_with_files_impl,
    build_user_input as _build_user_input_impl,
    chat_file_suffix as _chat_file_suffix_impl,
    clip_attachment_preview_text as _clip_attachment_preview_text_impl,
    decode_data_url as _decode_data_url_impl,
    model_supports_images as _model_supports_images_impl,
    stringify_user_input as _stringify_user_input_impl,
    user_input_has_images as _user_input_has_images_impl,
    validate_chat_payload as _validate_chat_payload_impl,
)
from api_chat_file_helpers import (
    ChatFileConfig,
    prepare_chat_files as _prepare_chat_files_impl,
)
from api_agent_stream_helpers import (
    dashboard_prompt_excerpt,
    fail_dashboard_task,
    finalize_dashboard_task,
    resolve_non_stream_agent_result,
    stream_agent_item,
    task_created_event,
)
from api_chat_route_helpers import (
    build_parallel_agent_streams,
    build_single_agent_stream,
    prepare_chat_route_runtime,
    sse_streaming_response,
)
from api_workspace_session_helpers import (
    create_session_record,
    reorder_sessions_payload,
    session_update_requested,
    workspaces_payload,
)
from api_document_helpers import (
    build_chat_report_title,
    build_upload_documents_task_record,
    cleanup_temp_paths,
    populate_chat_report_presentation,
    retrieval_test_payload,
    safe_report_filename,
    stage_upload_files,
    upload_documents_response,
)
from api_deck_report_helpers import (
    apply_deck_update,
    build_create_deck_kwargs,
    build_regenerate_deck_kwargs,
    create_share_link_payload,
    export_deck_payload,
    replace_deck_slide,
    resolve_report_messages,
    report_download_payload,
    report_markdown_payload,
)
from api_chat_stream_helpers import (
    answer_chunks,
    build_agent_config_payload,
    done_event as _done_event,
    panel_event as _panel_event,
    stream_parallel_sse,
    stream_single_sse,
)
from api_session_memory_helpers import (
    build_phase_summary_content,
    build_phase_summary_llm_prompt,
    covered_turns_from_summary,
    latest_auto_summary,
    normalize_llm_text_content,
    summarize_window_meta,
    summary_llm_enabled,
    summary_llm_timeout_seconds,
    summary_turns,
)
from api_session_memory_route_helpers import (
    delete_session_memory_payload,
    pin_session_memory_payload,
    session_memory_payload,
    session_memory_updates,
    summarize_session_memory_payload,
    update_session_memory_payload,
)
from api_kb_helpers import (
    filter_kb_chunks,
    kb_collect_chunks as _kb_collect_chunks,
    kb_docstore_dict as _kb_docstore_dict,
    kb_rebuild_from_documents as _kb_rebuild_from_documents,
    kb_safe_metadata as _kb_safe_metadata,
)
from api_kb_delete_helpers import delete_kb_directory
from api_kb_management_helpers import (
    kb_health_payload,
    knowledge_bases_payload,
)
from api_kb_chunk_route_helpers import (
    delete_kb_chunk_payload,
    list_kb_chunks_payload,
    update_kb_chunk_payload,
)
from api_share_helpers import (
    build_share_url as _build_share_url,
    decode_share_token as _decode_share_token,
    encode_share_token as _encode_share_token,
    SQLiteShareLinkStore,
)
from api_shared_resource_helpers import open_shared_resource_payload
from api_task_helpers import (
    contains_dashboard_card as _contains_dashboard_card,
    create_inline_task_record,
    prune_task_records,
    set_inline_task_state,
    should_start_dashboard_task as _should_start_dashboard_task,
    summarize_dashboard_task_error as _summarize_dashboard_task_error,
    summarize_dashboard_task_result as _summarize_dashboard_task_result,
)
from api_task_runtime_helpers import (
    attach_current_kb_status,
    enqueue_task,
    list_tasks_payload,
    task_record_payload,
)
from api_task_execution_helpers import (
    persist_web_research_task_placeholder,
    persist_web_research_task_result,
    run_analyze_knowledge_base_task,
    run_generate_deck_task,
    run_generate_report_task,
    run_placeholder_task,
    run_promote_attachment_to_kb_task,
    run_upload_documents_task,
    run_web_research_task,
)
from api_task_store import SQLiteTaskStore, TaskRecord, TaskStatus
from deck_service import (
    DeckSlide,
    SQLiteDeckStore,
    build_deck,
    build_report_markdown,
    build_export_filename,
    ensure_deckable_chat,
    export_deck_to_pptx,
    normalize_deck_theme,
    regenerate_deck_slide,
)
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(name)s] %(levelname)s %(message)s",
    stream=sys.stderr,
)
logger = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).resolve().parent
UPLOAD_CHUNK_SIZE = 1024 * 1024
TASK_HISTORY_LIMIT = 200
TASK_HISTORY_TTL_SECONDS = int(os.getenv("TASK_HISTORY_TTL_SECONDS", str(6 * 60 * 60)))
KB_METADATA_TTL_SECONDS = 30
CHAT_FILE_CONTEXT_START_MARKER = "[[CHAT_FILE_CONTEXT_START]]"
CHAT_FILE_CONTEXT_END_MARKER = "[[CHAT_FILE_CONTEXT_END]]"
CHAT_FILE_MAX_COUNT = int(os.getenv("CHAT_FILE_MAX_COUNT", "6"))
CHAT_FILE_MAX_BYTES = int(os.getenv("CHAT_FILE_MAX_BYTES", str(10 * 1024 * 1024)))
CHAT_FILE_MAX_CHARS_PER_FILE = int(
    os.getenv("CHAT_FILE_MAX_CHARS_PER_FILE", "8000")
)
CHAT_FILE_MAX_TOTAL_CHARS = int(os.getenv("CHAT_FILE_MAX_TOTAL_CHARS", "24000"))
CHAT_ATTACHMENT_PREVIEW_CHARS = int(
    os.getenv("CHAT_ATTACHMENT_PREVIEW_CHARS", "4000")
)
SHARE_LINK_SECRET = os.getenv("SHARE_LINK_SECRET", "local-share-secret")
SHARE_LINK_TTL_SECONDS = int(os.getenv("SHARE_LINK_TTL_SECONDS", str(7 * 24 * 60 * 60)))
DEFAULT_SHARE_LINK_SECRET = "local-share-secret"
MIN_SHARE_LINK_SECRET_LENGTH = 16
SESSION_MEMORY_AUTO_SUMMARY_MIN_TURNS = int(
    os.getenv("SESSION_MEMORY_AUTO_SUMMARY_MIN_TURNS", "12")
)
SESSION_MEMORY_AUTO_SUMMARY_MIN_NEW_TURNS = int(
    os.getenv("SESSION_MEMORY_AUTO_SUMMARY_MIN_NEW_TURNS", "6")
)
SESSION_MEMORY_AUTO_SUMMARY_WINDOW_SIZE = int(
    os.getenv("SESSION_MEMORY_AUTO_SUMMARY_WINDOW_SIZE", "14")
)
SESSION_MEMORY_AUTO_SUMMARY_MAX_CONTENT_CHARS = int(
    os.getenv("SESSION_MEMORY_AUTO_SUMMARY_MAX_CONTENT_CHARS", "900")
)
SUPPORTED_CHAT_FILE_EXTENSIONS = {
    ".pdf",
    ".doc",
    ".docx",
    ".txt",
    ".md",
    ".csv",
    ".xls",
    ".xlsx",
}
DOCUMENT_UPLOAD_MAX_COUNT = int(os.getenv("DOCUMENT_UPLOAD_MAX_COUNT", "20"))
DOCUMENT_UPLOAD_MAX_FILE_BYTES = int(
    os.getenv("DOCUMENT_UPLOAD_MAX_FILE_BYTES", str(25 * 1024 * 1024))
)
DOCUMENT_UPLOAD_MAX_TOTAL_BYTES = int(
    os.getenv("DOCUMENT_UPLOAD_MAX_TOTAL_BYTES", str(100 * 1024 * 1024))
)


@dataclass
class KBMetadataCacheEntry:
    expires_at: float
    value: dict[str, Any]


def _env_flag(name: str, default: bool = False) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _env_int(name: str, default: int) -> int:
    raw = os.getenv(name)
    if raw is None:
        return default
    try:
        return int(raw.strip())
    except (TypeError, ValueError):
        return default


def _cors_settings() -> tuple[list[str], bool]:
    raw = os.getenv("CORS_ALLOW_ORIGINS", "").strip()
    if not raw:
        return (
            [
                "http://127.0.0.1:5173",
                "http://localhost:5173",
                "http://127.0.0.1:8000",
                "http://localhost:8000",
            ],
            True,
        )

    origins = [item.strip() for item in raw.split(",") if item.strip()]
    if "*" in origins:
        return ["*"], False
    return origins, True


def _is_loopback_host(host: Optional[str]) -> bool:
    if not host:
        return False
    if host in {"localhost", "testclient"}:
        return True
    try:
        return ipaddress.ip_address(host).is_loopback
    except ValueError:
        return False


def _hash_secret(secret: str) -> str:
    if not secret:
        return "no-key"
    return hashlib.sha256(secret.encode("utf-8")).hexdigest()[:12]


def _token_fingerprint(token: str) -> str:
    normalized = str(token or "").strip()
    if not normalized:
        return "empty"
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:12]


def _token_preview(token: str) -> str:
    normalized = str(token or "").strip()
    if len(normalized) <= 10:
        return normalized
    return f"{normalized[:6]}...{normalized[-4:]}"


def _request_client_ip(request: Request) -> str:
    client = getattr(request, "client", None)
    host = getattr(client, "host", "") if client is not None else ""
    return str(host or "").strip()


def _request_user_agent(request: Request) -> str:
    return str(request.headers.get("user-agent") or "").strip()


def _request_is_local(request: Request) -> bool:
    client = getattr(request, "client", None)
    host = getattr(client, "host", "") if client is not None else ""
    return _is_loopback_host(host)


def _current_admin_api_token() -> str:
    return str(os.getenv("ADMIN_API_TOKEN", "") or "").strip()


def _extract_admin_token(request: Request) -> str:
    header_token = str(request.headers.get("X-Admin-Token") or "").strip()
    if header_token:
        return header_token
    authorization = str(request.headers.get("Authorization") or "").strip()
    if authorization.lower().startswith("bearer "):
        return authorization[7:].strip()
    return ""


def _admin_auth_mode(request: Request) -> str:
    if _request_is_local(request):
        return "local"
    if str(request.headers.get("X-Admin-Token") or "").strip():
        return "header"
    authorization = str(request.headers.get("Authorization") or "").strip()
    if authorization.lower().startswith("bearer "):
        return "bearer"
    return "missing"


def _sanitize_log_value(value: Any, *, max_length: int = 256) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    text = re.sub(r"[\r\n\t]+", " ", text)
    text = re.sub(r"\s{2,}", " ", text)
    if len(text) > max_length:
        return f"{text[: max_length - 3]}..."
    return text


def _request_user_id(request: Request) -> str:
    return _sanitize_log_value(request.headers.get("X-User-Id"), max_length=128)


def _request_user_role(request: Request) -> str:
    return _sanitize_log_value(request.headers.get("X-User-Role"), max_length=64)


def _sanitize_request_path(path: str) -> str:
    normalized = str(path or "").strip() or "/"
    if normalized.startswith("/api/share-links/"):
        return "/api/share-links/<token>"
    if normalized.startswith("/shared/"):
        return "/shared/<token>"
    return normalized


def _current_share_link_secret() -> str:
    return str(os.getenv("SHARE_LINK_SECRET", SHARE_LINK_SECRET) or "").strip()


def _share_link_secret_is_weak() -> bool:
    secret = _current_share_link_secret()
    return (
        not secret
        or secret == DEFAULT_SHARE_LINK_SECRET
        or len(secret) < MIN_SHARE_LINK_SECRET_LENGTH
    )


def _require_remote_admin(request: Request) -> None:
    if not ALLOW_REMOTE_CLIENTS or _request_is_local(request):
        return

    configured_token = _current_admin_api_token()
    if not configured_token:
        _audit_security_event(
            "remote_admin_guard",
            request,
            result="blocked",
            details="reason=admin_token_not_configured",
        )
        raise HTTPException(
            status_code=503,
            detail="远程管理接口已禁用，请先配置 ADMIN_API_TOKEN",
        )

    provided_token = _extract_admin_token(request)
    if not provided_token or not secrets.compare_digest(provided_token, configured_token):
        reason = "missing_admin_token" if not provided_token else "invalid_admin_token"
        _audit_security_event(
            "remote_admin_guard",
            request,
            result="rejected",
            details=f"reason={reason}",
        )
        raise HTTPException(status_code=403, detail="缺少有效的管理令牌")


def _require_remote_share_secret(request: Request) -> None:
    if not ALLOW_REMOTE_CLIENTS or _request_is_local(request):
        return
    if _share_link_secret_is_weak():
        _audit_security_event(
            "remote_share_guard",
            request,
            result="blocked",
            details="reason=weak_share_link_secret",
        )
        raise HTTPException(
            status_code=503,
            detail="远程分享已禁用，请配置强 SHARE_LINK_SECRET 后再启用",
        )


def _content_hash(value: Any) -> str:
    if value in (None, ""):
        return ""
    if isinstance(value, str):
        payload = value
    else:
        payload = json.dumps(
            value,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        )
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _audit_security_event(
    action: str,
    request: Request,
    *,
    result: str = "ok",
    details: str = "",
) -> None:
    request_id = getattr(request.state, "request_id", "")
    logger.info(
        (
            "security_event action=%s result=%s request_id=%s ip=%s local=%s "
            "auth=%s user_id=%s user_role=%s details=%s"
        ),
        action,
        result,
        request_id,
        _request_client_ip(request),
        _request_is_local(request),
        _admin_auth_mode(request),
        _request_user_id(request),
        _request_user_role(request),
        _sanitize_log_value(details, max_length=512),
    )


def _share_link_audit_payload(record: Any, *, now: Optional[float] = None) -> dict[str, Any]:
    current_time = time.time() if now is None else float(now)
    revoked_at = getattr(record, "revoked_at", None)
    expires_at = float(getattr(record, "expires_at", 0) or 0)
    share_token = str(getattr(record, "share_token", "") or "").strip()
    return {
        "resource_type": str(getattr(record, "resource_type", "") or "").strip(),
        "resource_id": str(getattr(record, "resource_id", "") or "").strip(),
        "created_at": float(getattr(record, "created_at", 0) or 0),
        "expires_at": expires_at,
        "revoked_at": float(revoked_at) if revoked_at is not None else None,
        "is_active": revoked_at is None and expires_at > current_time,
        "created_by_ip": str(getattr(record, "created_by_ip", "") or "").strip(),
        "created_user_agent": str(getattr(record, "created_user_agent", "") or "").strip(),
        "access_count": int(getattr(record, "access_count", 0) or 0),
        "last_accessed_at": (
            float(getattr(record, "last_accessed_at", 0))
            if getattr(record, "last_accessed_at", None) is not None
            else None
        ),
        "last_accessed_ip": str(getattr(record, "last_accessed_ip", "") or "").strip(),
        "last_accessed_user_agent": str(
            getattr(record, "last_accessed_user_agent", "") or ""
        ).strip(),
        "share_token_preview": _token_preview(share_token),
        "share_token_fingerprint": _token_fingerprint(share_token),
    }


def _security_status_payload() -> dict[str, Any]:
    cors_origins, cors_allow_credentials = _cors_settings()
    admin_token_configured = bool(_current_admin_api_token())
    share_link_secret_healthy = not _share_link_secret_is_weak()
    return {
        "allow_remote_clients": ALLOW_REMOTE_CLIENTS,
        "local_only_mode": not ALLOW_REMOTE_CLIENTS,
        "admin_token_configured": admin_token_configured,
        "remote_admin_ready": (not ALLOW_REMOTE_CLIENTS) or admin_token_configured,
        "share_link_secret_healthy": share_link_secret_healthy,
        "remote_share_ready": (not ALLOW_REMOTE_CLIENTS) or share_link_secret_healthy,
        "share_link_ttl_seconds": int(SHARE_LINK_TTL_SECONDS),
        "share_link_ttl_hours": round(float(SHARE_LINK_TTL_SECONDS) / 3600.0, 2),
        "cors_allow_credentials": cors_allow_credentials,
        "cors_allowed_origins": cors_origins,
        "request_id_header": "X-Request-ID",
        "process_time_header": "X-Process-Time-Ms",
        "chat_file_limits": {
            "max_count": int(CHAT_FILE_MAX_COUNT),
            "max_bytes": int(CHAT_FILE_MAX_BYTES),
            "max_chars_per_file": int(CHAT_FILE_MAX_CHARS_PER_FILE),
            "max_total_chars": int(CHAT_FILE_MAX_TOTAL_CHARS),
            "preview_chars": int(CHAT_ATTACHMENT_PREVIEW_CHARS),
        },
        "document_upload_limits": {
            "max_count": int(DOCUMENT_UPLOAD_MAX_COUNT),
            "max_file_bytes": int(DOCUMENT_UPLOAD_MAX_FILE_BYTES),
            "max_total_bytes": int(DOCUMENT_UPLOAD_MAX_TOTAL_BYTES),
        },
    }


def _resolve_project_subdir(candidate: str) -> Path:
    raw_path = Path(candidate).expanduser()
    if not raw_path.is_absolute():
        raw_path = PROJECT_ROOT / raw_path
    resolved = raw_path.resolve()
    try:
        resolved.relative_to(PROJECT_ROOT)
    except ValueError as exc:
        raise HTTPException(
            status_code=403, detail="不允许访问项目目录之外的路径"
        ) from exc
    if resolved == PROJECT_ROOT:
        raise HTTPException(status_code=403, detail="不允许直接操作项目根目录")
    return resolved


def _faiss_safe_store_path(path: str | Path) -> str:
    target = Path(path)
    if not target.is_absolute():
        target = (PROJECT_ROOT / target).resolve()
    else:
        target = target.resolve()

    try:
        return str(target.relative_to(PROJECT_ROOT))
    except ValueError:
        return str(target)


def _resolve_deletable_knowledge_base(candidate: str) -> Path:
    target_path = _resolve_project_subdir(candidate)

    if not target_path.exists():
        raise HTTPException(status_code=404, detail="知识库路径不存在")
    if not target_path.is_dir():
        raise HTTPException(status_code=400, detail="知识库路径必须是目录")
    if not (target_path / "index.faiss").is_file():
        raise HTTPException(
            status_code=400,
            detail="只能删除包含 index.faiss 的目录",
        )

    return target_path


def _active_vector_store_id() -> str | None:
    _, vector_store_path, _ = _resolve_active_prompt_runtime(True)
    return vector_store_path


def _resolve_active_prompt_runtime(
    knowledge_base_enabled: bool,
) -> tuple[Optional[str], Optional[str], dict[str, Any]]:
    from chat_store import get_active_system_prompt

    active_prompt = get_active_system_prompt() or {}

    system_prompt_content_raw = active_prompt.get("content")
    system_prompt_content = (
        str(system_prompt_content_raw).strip()
        if isinstance(system_prompt_content_raw, str)
        else ""
    ) or None

    vector_store_path = str(active_prompt.get("vector_store_id") or "").strip() or None
    if not knowledge_base_enabled:
        vector_store_path = None

    dashboard_template_raw = active_prompt.get("dashboard_template", {})
    dashboard_template = (
        dict(dashboard_template_raw)
        if isinstance(dashboard_template_raw, dict)
        else {}
    )

    return system_prompt_content, vector_store_path, dashboard_template


def _dashboard_feature_enabled(dashboard_template: Optional[dict[str, Any]]) -> bool:
    if not isinstance(dashboard_template, dict):
        return True
    return dashboard_template.get("enabled") is not False


def _effective_vector_store_path(candidate: Optional[str] = None) -> str:
    raw = str(candidate or "").strip()
    if not raw:
        raw = _active_vector_store_id() or os.getenv("VECTOR_STORE_PATH", "./vector_store")
    return _faiss_safe_store_path(_resolve_project_subdir(raw))


def _require_workspace_session(
    session_id: str,
    workspace_id: Optional[str] = None,
) -> dict[str, Any]:
    from chat_store import DEFAULT_WORKSPACE_ID, get_session, get_workspace

    normalized_session_id = str(session_id or "").strip()
    session = get_session(normalized_session_id)
    if session is None:
        raise HTTPException(status_code=404, detail="会话不存在")

    normalized_workspace_id = str(workspace_id or "").strip()
    if not normalized_workspace_id:
        return session

    if get_workspace(normalized_workspace_id) is None:
        raise HTTPException(status_code=400, detail="工作区不存在")

    session_workspace_id = str(session.get("workspace_id") or DEFAULT_WORKSPACE_ID)
    if session_workspace_id != normalized_workspace_id:
        raise HTTPException(status_code=404, detail="当前工作区中不存在该会话")

    return session


def _build_download_content_disposition(filename: str) -> str:
    raw_filename = str(filename or "").strip() or "download"
    ascii_only = all(ord(char) < 128 for char in raw_filename)
    if ascii_only:
        escaped = raw_filename.replace("\\", "\\\\").replace('"', r"\"")
        return f'attachment; filename="{escaped}"'

    ascii_fallback = re.sub(r'[^A-Za-z0-9._-]+', "_", raw_filename).strip("._")
    if not ascii_fallback:
        ascii_fallback = "download"
    encoded = quote(raw_filename, safe="")
    return f"attachment; filename={ascii_fallback}; filename*=UTF-8''{encoded}"


ALLOW_REMOTE_CLIENTS = _env_flag("ALLOW_REMOTE_CLIENTS", False)
_cors_origins, _cors_allow_credentials = _cors_settings()

app = FastAPI(title="InsightDesk API", version="2.0.0")
_deck_store = SQLiteDeckStore()
_artifact_store = SQLiteArtifactStore()
_share_link_store = SQLiteShareLinkStore()

app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins,
    allow_credentials=_cors_allow_credentials,
    allow_methods=["*"],
    allow_headers=["*"],
)


def _artifact_payload(artifact: Any) -> dict[str, Any]:
    payload = artifact.model_dump(mode="json")
    payload["available_formats"] = artifact_export_formats(artifact)
    return payload


def _create_report_artifact(
    *,
    session_id: str,
    messages: list[Any],
    answer_group_id: str = "",
    panel_id: str = "",
) -> tuple[Any, str, str]:
    qa_pairs = ensure_deckable_chat(messages)
    title = build_chat_report_title(messages)
    markdown = build_report_markdown(messages, title)
    artifact = build_report_artifact(
        session_id=session_id,
        title=title,
        markdown=markdown,
        qa_pairs=qa_pairs,
        answer_group_id=answer_group_id,
        panel_id=panel_id,
    )
    _artifact_store.save(artifact)
    return artifact, title, markdown


def _create_deck_artifact(deck: Any) -> Any:
    artifact = build_deck_artifact(deck)
    _artifact_store.save(artifact)
    return artifact


def _sync_deck_artifacts(deck: Any) -> None:
    deck_id = str(getattr(deck, "deck_id", "") or "").strip()
    if not deck_id:
        return
    for artifact in _artifact_store.list_by_linked_resource("deck", deck_id):
        sync_deck_artifact(artifact, deck)
        _artifact_store.save(artifact)


@app.middleware("http")
async def restrict_remote_clients(request: Request, call_next):
    request_id = request.headers.get("X-Request-ID", "").strip() or str(uuid.uuid4())
    request.state.request_id = request_id
    started_at = time.perf_counter()
    if ALLOW_REMOTE_CLIENTS:
        response = await call_next(request)

    else:
        client_host = request.client.host if request.client else None
        if _is_loopback_host(client_host):
            response = await call_next(request)

        else:
            logger.warning(
                "Blocked non-local request request_id=%s host=%s path=%s",
                request_id,
                client_host,
                _sanitize_request_path(request.url.path),
            )
            response = JSONResponse(
                status_code=403,
                content={
            "code": "LOCAL_ONLY",
            "message": "当前服务默认只允许本机访问",
            "suggestion": "请使用 localhost/127.0.0.1 访问，或显式设置 ALLOW_REMOTE_CLIENTS=true 后再对外暴露",
                },
            )

    process_time_ms = (time.perf_counter() - started_at) * 1000.0
    response.headers["X-Request-ID"] = request_id
    response.headers["X-Process-Time-Ms"] = f"{process_time_ms:.2f}"
    response.headers.setdefault("X-Content-Type-Options", "nosniff")
    response.headers.setdefault("X-Frame-Options", "DENY")
    response.headers.setdefault("Referrer-Policy", "same-origin")
    response.headers.setdefault(
        "Permissions-Policy",
        "camera=(), microphone=(), geolocation=()",
    )
    logger.info(
        "request_id=%s method=%s path=%s status=%s latency_ms=%.2f",
        request_id,
        request.method,
        _sanitize_request_path(request.url.path),
        response.status_code,
        process_time_ms,
    )
    return response


@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    err = _classify_error(exc)
    request_id = getattr(request.state, "request_id", "")
    logger.exception("Unhandled exception request_id=%s on %s", request_id, request.url.path)
    return JSONResponse(
        status_code=500,
        content={
            "code": err["code"],
            "message": err["message"],
            "suggestion": err["suggestion"],
            "request_id": request_id,
        },
    )


def _classify_error(e: Exception) -> dict:
    """将异常映射为用户友好的错误信息 + 错误码 + 恢复建议"""
    msg = str(e)
    if isinstance(e, ConnectionError) or any(
        kw in msg
        for kw in ("Connection refused", "ConnectError", "connect ECONNREFUSED")
    ):
        return {
            "code": "MODEL_UNAVAILABLE",
            "message": "模型服务暂时不可用",
            "suggestion": "请检查 Ollama 是否已启动，或稍后重试",
        }
    if any(
        kw in msg
        for kw in ("API Key", "api_key", "401", "Unauthorized", "authentication")
    ):
        return {
            "code": "AUTH_FAILED",
            "message": "API 认证失败",
            "suggestion": "请在设置中检查 API Key 是否正确",
        }
    if any(kw in msg for kw in ("timeout", "Timeout", "TimeoutError", "timed out")):
        return {
            "code": "TIMEOUT",
            "message": "请求超时，模型响应过慢",
            "suggestion": "请稍后重试，或尝试切换为响应更快的模型",
        }
    if any(kw in msg for kw in ("rate limit", "429", "quota")):
        return {
            "code": "RATE_LIMIT",
            "message": "API 调用频率已达上限",
            "suggestion": "请稍等片刻后重试",
        }
    if any(
        kw in msg.lower()
        for kw in (
            "does not support image",
            "doesn't support image",
            "image input",
            "vision",
            "multimodal",
            "input image",
            "unsupported image",
            "invalid image",
            "content type image_url",
            "image_url",
        )
    ):
        return {
            "code": "MODEL_NO_VISION",
            "message": "当前模型未接受图片输入，可能不支持视觉能力或当前接入方式不兼容。",
            "suggestion": "系统已实际尝试发送图片。请检查模型本身是否支持读图，以及当前 API/Base URL 是否支持多模态输入。",
        }
    if any(kw in msg for kw in ("Invalid model", "model not found", "404")):
        return {
            "code": "MODEL_NOT_FOUND",
            "message": "指定的模型不存在",
            "suggestion": "请在面板顶部的模型选择器中确认模型名称是否正确",
        }
    if any(kw in msg for kw in ("max iterations", "Agent stopped", "iteration limit")):
        return {
            "code": "MAX_ITERATIONS",
            "message": "模型工具调用次数超限，无法完成任务",
            "suggestion": "请尝试简化问题、减少工具依赖，或切换至其他模型后重试",
        }
    return {
        "code": "INTERNAL_ERROR",
        "message": "处理请求时发生异常",
        "suggestion": "请尝试清除上下文或新建对话后重试",
    }


def _is_max_iterations_output(text: str) -> bool:
    """检测 AgentExecutor 返回的 max iterations 错误字符串"""
    lower = text.lower()
    return (
        "agent stopped due to max iterations" in lower
        or "agent stopped due to iteration limit" in lower
        or (lower.startswith("agent stopped") and "iteration" in lower)
    )


async def _fallback_generate(
    mc: "ModelConfig", user_message: str, tool_outputs: str
) -> str:
    """当 Agent 达到迭代上限时，用单次 LLM 调用基于已有工具结果生成回答"""
    from agent_core import get_llm

    try:
        mc = _normalize_model_config(mc)
        llm = get_llm(
            provider=mc.provider,
            model_name=mc.model,
            base_url=mc.base_url,
            api_key=mc.api_key if mc.api_key else None,
            temperature=mc.temperature,
        )
        fallback_prompt = f"""你是一个企业知识库助手。以下是工具已收集到的信息，请基于这些信息直接回答用户问题，用中文作答。

用户问题：{user_message}

工具收集到的信息：
{tool_outputs[:4000]}

请基于以上信息给出完整、清晰的回答："""
        response = await llm.ainvoke(fallback_prompt)
        return response.content.strip()
    except Exception as exc:
        logger.warning("Fallback generate failed: %s", exc)
        return (
            "抱歉，模型处理时间过长，未能生成完整回答。请尝试简化问题或切换模型后重试。"
        )


# ─────────────────────────────────────────────
# Pydantic 请求/响应模型
# ─────────────────────────────────────────────


class ModelConfig(BaseModel):
    panel_id: str
    connection_type: Optional[str] = None
    provider: str = "ollama"
    model: str = "qwen3.5-2B:latest"
    base_url: str = "http://localhost:11434"
    api_key: str = ""
    temperature: float = 0.3
    agent_mode: str = "auto"


def _model_config_payload(mc: ModelConfig) -> dict[str, Any]:
    if hasattr(mc, "model_dump"):
        return mc.model_dump()
    return mc.dict()


def _base_model_payload(model: Any) -> dict[str, Any]:
    if hasattr(model, "model_dump"):
        return model.model_dump()
    if hasattr(model, "dict"):
        return model.dict()
    return dict(model)


def _normalize_model_config(mc: ModelConfig) -> ModelConfig:
    from agent_core import (
        default_base_url_for_connection_type,
        default_model_for_connection_type,
        normalize_connection_type,
    )

    data = _model_config_payload(mc)
    connection_type = normalize_connection_type(
        data.get("connection_type") or data.get("provider"),
        data.get("base_url"),
    )
    data["connection_type"] = connection_type
    data["provider"] = connection_type
    data["base_url"] = (
        str(data.get("base_url") or default_base_url_for_connection_type(connection_type)).strip()
    )
    data["model"] = (
        str(data.get("model") or default_model_for_connection_type(connection_type)).strip()
    )
    data["api_key"] = str(data.get("api_key") or "").strip()
    return ModelConfig(**data)


class ImageInput(BaseModel):
    name: str
    media_type: str
    data_url: str


class FileInput(BaseModel):
    name: str
    media_type: str
    data_url: str
    size_bytes: int = 0
    extracted_text: str = ""


class ChatRequest(BaseModel):
    session_id: str
    message: str = ""
    images: list[ImageInput] = Field(default_factory=list)
    files: list[FileInput] = Field(default_factory=list)
    models: list[ModelConfig]
    web_search_enabled: bool = False
    knowledge_base_enabled: bool = True
    enabled_mcp_servers: list[str] = Field(default_factory=list)
    answer_group_id: Optional[str] = None


class SingleChatRequest(BaseModel):
    session_id: str
    message: str = ""
    images: list[ImageInput] = Field(default_factory=list)
    files: list[FileInput] = Field(default_factory=list)
    panel_config: ModelConfig
    web_search_enabled: bool = False
    knowledge_base_enabled: bool = True
    enabled_mcp_servers: list[str] = Field(default_factory=list)
    answer_group_id: Optional[str] = None
    persist_user_history: bool = True
    persist_ai_history: bool = True
    replace_ai_history: bool = False
    exclude_ai_answer_group_id: Optional[str] = None


class CreateSessionRequest(BaseModel):
    title: str = ""
    workspace_id: Optional[str] = None


class UpdateSessionRequest(BaseModel):
    title: Optional[str] = None
    is_archived: Optional[bool] = None
    is_favorite: Optional[bool] = None
    is_pinned: Optional[bool] = None
    tags: Optional[list[str]] = None
    workspace_id: Optional[str] = None


class ReorderSessionsRequest(BaseModel):
    session_ids: list[str]
    workspace_id: Optional[str] = None


class CreateBookmarkRequest(BaseModel):
    session_id: str
    role: str
    message_id: Optional[int] = None
    panel_id: str = ""
    answer_group_id: str = ""
    content: str = ""
    model_id: str = ""
    session_title: str = ""


class ShareLinkResponse(BaseModel):
    resource_type: str
    resource_id: str
    share_token: str
    share_url: str
    expires_at: float


class RevokeShareLinkResponse(BaseModel):
    ok: bool


class ShareLinkAuditRecord(BaseModel):
    resource_type: str
    resource_id: str
    created_at: float
    expires_at: float
    revoked_at: Optional[float] = None
    is_active: bool
    created_by_ip: str = ""
    created_user_agent: str = ""
    access_count: int = 0
    last_accessed_at: Optional[float] = None
    last_accessed_ip: str = ""
    last_accessed_user_agent: str = ""
    share_token_preview: str
    share_token_fingerprint: str


class ShareLinkAuditListResponse(BaseModel):
    share_links: list[ShareLinkAuditRecord]
    total: int
    active_count: int


class SecurityStatusResponse(BaseModel):
    allow_remote_clients: bool
    local_only_mode: bool
    admin_token_configured: bool
    remote_admin_ready: bool
    share_link_secret_healthy: bool
    remote_share_ready: bool
    share_link_ttl_seconds: int
    share_link_ttl_hours: float
    cors_allow_credentials: bool
    cors_allowed_origins: list[str]
    request_id_header: str
    process_time_header: str
    chat_file_limits: dict[str, int]
    document_upload_limits: dict[str, int]


class WorkspaceToolConfigRequest(BaseModel):
    web_search_enabled: bool = False
    knowledge_base_enabled: bool = True
    mcp_servers_enabled: list[str] = Field(default_factory=default_mcp_server_names)


class WorkspaceOutputPresetRequest(BaseModel):
    deck_theme: Literal["default", "midnight", "sunrise"] = "default"
    target_slide_count: int = Field(default=8, ge=4, le=10)


class WorkspacePresetRequest(BaseModel):
    default_panels: list[ModelConfig] = Field(default_factory=list)
    tool_config: WorkspaceToolConfigRequest = Field(
        default_factory=WorkspaceToolConfigRequest
    )
    output_preset: WorkspaceOutputPresetRequest = Field(
        default_factory=WorkspaceOutputPresetRequest
    )


class CreateWorkspaceRequest(BaseModel):
    name: str
    description: str = ""
    color: str = "blue"
    activate: bool = True
    preset: Optional[WorkspacePresetRequest] = None


class UpdateWorkspaceRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    color: Optional[str] = None
    preset: Optional[WorkspacePresetRequest] = None


class SetMessageFeedbackRequest(BaseModel):
    value: int
    message_id: Optional[int] = None
    panel_id: str = ""
    answer_group_id: str = ""


class TruncateSessionMessagesRequest(BaseModel):
    answer_group_id: str
    content: str = ""
    images: list[ImageInput] = Field(default_factory=list)
    files: list[FileInput] = Field(default_factory=list)


class ImportedFileInput(BaseModel):
    name: str
    media_type: str
    data_url: str = ""
    size_bytes: int = 0
    extracted_text: str = ""


class ImportSessionMessageRequest(BaseModel):
    role: Literal["user", "assistant"]
    content: str = ""
    images: list[ImageInput] = Field(default_factory=list)
    files: list[ImportedFileInput] = Field(default_factory=list)
    sources: list[dict[str, Any]] = Field(default_factory=list)
    model_id: str = ""
    panel_id: str = ""
    answer_group_id: str = ""
    workflow_nodes: list[dict[str, Any]] = Field(default_factory=list)
    task_id: str = ""
    task_type: str = ""


class ImportSessionMessagesRequest(BaseModel):
    panels: list[ModelConfig] = Field(default_factory=list)
    messages: list[ImportSessionMessageRequest] = Field(default_factory=list)


class SetRetrievalFeedbackRequest(BaseModel):
    panel_id: str
    answer_group_id: str
    source: dict[str, Any]
    value: int


class PinSessionMemoryRequest(BaseModel):
    content: str
    kind: str = Field(default="fact")


class UpdateSessionMemoryRequest(BaseModel):
    content: Optional[str] = None
    kind: Optional[str] = None


class SaveConfigRequest(BaseModel):
    tavily_api_key: Optional[str] = None
    embedding_model: Optional[str] = None


class CreatePromptRequest(BaseModel):
    name: str
    content: str
    vector_store_id: Optional[str] = None
    dashboard_template: dict[str, Any] = Field(default_factory=dict)


class UpdatePromptRequest(BaseModel):
    name: str
    content: str
    vector_store_id: Optional[str] = None
    dashboard_template: dict[str, Any] = Field(default_factory=dict)


class GenerateReportRequest(BaseModel):
    session_id: str
    answer_group_id: Optional[str] = None
    panel_id: Optional[str] = None


def _resolve_report_messages(
    history: Any,
    *,
    answer_group_id: Optional[str] = None,
    panel_id: Optional[str] = None,
) -> list[Any]:
    from langchain_core.messages import AIMessage, HumanMessage

    return resolve_report_messages(
        history,
        answer_group_id=str(answer_group_id or "").strip(),
        panel_id=str(panel_id or "").strip(),
        human_message_factory=lambda content: HumanMessage(content=content),
        ai_message_factory=lambda content: AIMessage(content=content),
    )


class CreateDeckRequest(BaseModel):
    session_id: str
    panel_config: ModelConfig
    knowledge_base_enabled: bool = True
    target_slide_count: int = Field(default=8, ge=4, le=10)
    theme: str = "default"
    answer_group_id: Optional[str] = None
    panel_id: Optional[str] = None


class GenerateArtifactRequest(BaseModel):
    artifact_type: Literal["report", "deck"]
    session_id: str
    answer_group_id: Optional[str] = None
    panel_id: Optional[str] = None
    panel_config: Optional[ModelConfig] = None
    knowledge_base_enabled: bool = True
    target_slide_count: int = Field(default=8, ge=4, le=10)
    theme: str = "default"


class UpdateArtifactRequest(BaseModel):
    title: Optional[str] = None
    markdown: Optional[str] = None


class UpdateDeckRequest(BaseModel):
    title: Optional[str] = None
    theme: Optional[str] = None
    slides: Optional[list[DeckSlide]] = None


class RegenerateDeckSlideRequest(BaseModel):
    panel_config: ModelConfig
    knowledge_base_enabled: Optional[bool] = None


class TestRetrievalRequest(BaseModel):
    query: str
    vector_store_path: Optional[str] = None
    search_k: Optional[int] = None
    fetch_k: Optional[int] = None
    use_rerank: Optional[bool] = None
    retrieval_mode: Optional[str] = None


class UpdateKBChunkRequest(BaseModel):
    content: Optional[str] = None
    source: Optional[str] = None


class CreateTaskRequest(BaseModel):
    task_type: str  # e.g. "analyze_knowledge_base", "generate_report"
    params: dict = {}
    session_id: Optional[str] = None


# ─────────────────────────────────────────────
# 异步任务状态机
# ─────────────────────────────────────────────


def _update_task_progress(record: TaskRecord, progress: int) -> None:
    """线程安全之外的轻量进度更新，供同步处理阶段回调使用。"""
    if record.task_id in _suppressed_task_ids:
        return
    record.progress = max(0, min(100, progress))
    record.updated_at = time.time()
    _persist_task_record(record)


async def _drop_suppressed_task(record: TaskRecord) -> bool:
    if record.task_id not in _suppressed_task_ids:
        return False
    async with _tasks_lock:
        _tasks.pop(record.task_id, None)
        _prune_task_records_locked()
    return True


async def _create_inline_task_record(
    task_type: str,
    params: dict[str, Any],
    *,
    session_id: Optional[str] = None,
    progress: int = 10,
) -> TaskRecord:
    return await create_inline_task_record(
        _tasks,
        _tasks_lock,
        task_type=task_type,
        params=params,
        session_id=session_id,
        progress=progress,
        prune_in_memory=_prune_task_records_locked,
        persist_record=_persist_task_record,
        prune_persisted=_prune_persisted_tasks,
    )


async def _set_inline_task_state(
    record: TaskRecord,
    *,
    status: TaskStatus,
    progress: Optional[int] = None,
    result: Optional[str] = None,
    error: Optional[str] = None,
) -> TaskRecord:
    if record.task_id in _suppressed_task_ids:
        async with _tasks_lock:
            _tasks.pop(record.task_id, None)
            _prune_task_records_locked()
        record.status = status
        record.updated_at = time.time()
        if progress is not None:
            record.progress = max(0, min(100, progress))
        if result is not None:
            record.result = result
        if error is not None:
            record.error = error
        return record
    return await set_inline_task_state(
        _tasks,
        _tasks_lock,
        record=record,
        status=status,
        progress=progress,
        result=result,
        error=error,
        prune_in_memory=_prune_task_records_locked,
        persist_record=_persist_task_record,
        prune_persisted=_prune_persisted_tasks,
    )




_tasks: dict[str, TaskRecord] = {}
_tasks_lock = asyncio.Lock()
_suppressed_task_ids: set[str] = set()
_task_store: SQLiteTaskStore | None = None
_task_store_init_lock = threading.Lock()


def _get_task_store() -> SQLiteTaskStore:
    global _task_store
    if _task_store is None:
        with _task_store_init_lock:
            if _task_store is None:
                _task_store = SQLiteTaskStore(
                    history_limit=TASK_HISTORY_LIMIT,
                    ttl_seconds=TASK_HISTORY_TTL_SECONDS,
                )
    return _task_store


def _persist_task_record(record: TaskRecord) -> None:
    if record.task_id in _suppressed_task_ids:
        logger.info("Skip persisting suppressed task record: %s", record.task_id)
        return
    try:
        _get_task_store().save(record)
    except Exception:
        logger.exception("Failed to persist task record: %s", record.task_id)


def _prune_persisted_tasks() -> None:
    try:
        _get_task_store().prune()
    except Exception:
        logger.exception("Failed to prune persisted task records")


def _prune_task_records_locked(now: float | None = None) -> None:
    """Prune expired and excess terminal tasks while holding _tasks_lock."""
    prune_task_records(
        _tasks,
        history_limit=TASK_HISTORY_LIMIT,
        ttl_seconds=TASK_HISTORY_TTL_SECONDS,
        now=now,
        logger=logger,
    )


async def _run_task(record: TaskRecord) -> None:
    """后台执行任务并更新状态"""
    async with _tasks_lock:
        _prune_task_records_locked()
        record.status = TaskStatus.RUNNING
        record.updated_at = time.time()
        record.progress = 10
    if await _drop_suppressed_task(record):
        return
    _persist_task_record(record)

    try:
        task_type = record.task_type

        async def _set_progress(progress: int) -> None:
            if await _drop_suppressed_task(record):
                return
            async with _tasks_lock:
                record.progress = progress
                record.updated_at = time.time()
            _persist_task_record(record)

        if task_type == "analyze_knowledge_base":
            await run_analyze_knowledge_base_task(
                record,
                set_progress=_set_progress,
                effective_vector_store_path=_effective_vector_store_path,
            )

        elif task_type == "generate_report":
            await run_generate_report_task(
                record,
                set_progress=_set_progress,
                resolve_report_messages=_resolve_report_messages,
                ensure_deckable_chat=ensure_deckable_chat,
                build_chat_report_title=build_chat_report_title,
                build_report_markdown=build_report_markdown,
                build_report_artifact=build_report_artifact,
                save_artifact=_artifact_store.save,
            )

        elif task_type == "generate_deck":
            await run_generate_deck_task(
                record,
                set_progress=_set_progress,
                resolve_report_messages=_resolve_report_messages,
                normalize_model_config=_normalize_model_config,
                resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
                build_deck=build_deck,
                save_deck=_deck_store.save,
                build_deck_artifact=build_deck_artifact,
                save_artifact=_artifact_store.save,
            )

        elif task_type == "upload_documents":
            await run_upload_documents_task(
                record,
                set_progress=_set_progress,
                update_progress=_update_task_progress,
                effective_vector_store_path=_effective_vector_store_path,
                clear_agent_cache=_clear_agent_cache,
                logger=logger,
            )

        elif task_type == "promote_attachment_to_kb":
            await run_promote_attachment_to_kb_task(
                record,
                set_progress=_set_progress,
                update_progress=_update_task_progress,
                effective_vector_store_path=_effective_vector_store_path,
                chat_file_suffix=_chat_file_suffix,
                decode_data_url=_decode_data_url,
                clear_agent_cache=_clear_agent_cache,
                logger=logger,
            )

        elif task_type == "web_research":
            from agent_core import get_llm

            await run_web_research_task(
                record,
                set_progress=_set_progress,
                normalize_model_config=_normalize_model_config,
                create_llm=get_llm,
            )

        else:
            await run_placeholder_task(
                record,
                set_progress=_set_progress,
            )

        if await _drop_suppressed_task(record):
            return
        async with _tasks_lock:
            record.status = TaskStatus.COMPLETED
            record.progress = 100
            record.updated_at = time.time()
            _prune_task_records_locked(record.updated_at)
        _persist_task_record(record)
        _prune_persisted_tasks()

    except Exception as exc:
        logger.exception(
            "task_id=%s task_type=%s 执行失败", record.task_id, record.task_type
        )
        if record.task_type == "upload_documents":
            for temp_path in record.params.get("temp_paths", []):
                try:
                    os.remove(str(temp_path))
                except OSError:
                    pass
        if await _drop_suppressed_task(record):
            return
        async with _tasks_lock:
            record.status = TaskStatus.FAILED
            record.error = str(exc)
            record.updated_at = time.time()
            _prune_task_records_locked(record.updated_at)
        _persist_task_record(record)
        if record.task_type == "web_research":
            persist_web_research_task_result(
                record,
                content=f"联网研究任务失败：{record.error}",
                sources=[],
            )
        _prune_persisted_tasks()


# ─────────────────────────────────────────────
# 内部工具函数
# ─────────────────────────────────────────────

# 缓存已构建的 agent 实例，key = (provider, model, base_url, api_key, temperature, agent_mode)
_agent_cache: dict[str, Any] = {}
_agent_cache_lock = asyncio.Lock()


async def _clear_agent_cache() -> None:
    async with _agent_cache_lock:
        cleared_count = len(_agent_cache)
        _agent_cache.clear()
    if cleared_count:
        logger.info("Cleared %d cached agent(s)", cleared_count)


def _validate_chat_payload(
    message: str,
    images: list[ImageInput],
    files: list[FileInput],
) -> None:
    _validate_chat_payload_impl(message, images, files)


def _chat_file_suffix(name: str) -> str:
    return _chat_file_suffix_impl(name)


def _decode_data_url(data_url: str, file_name: str) -> bytes:
    return _decode_data_url_impl(data_url, file_name)


def _clip_attachment_preview_text(text: str, limit: int = CHAT_ATTACHMENT_PREVIEW_CHARS) -> str:
    return _clip_attachment_preview_text_impl(text, limit)


def _prepare_chat_files(files: list[FileInput]) -> tuple[list[dict[str, Any]], str]:
    return _prepare_chat_files_impl(
        files,
        config=ChatFileConfig(
            context_end_marker=CHAT_FILE_CONTEXT_END_MARKER,
            context_start_marker=CHAT_FILE_CONTEXT_START_MARKER,
            max_bytes=CHAT_FILE_MAX_BYTES,
            max_chars_per_file=CHAT_FILE_MAX_CHARS_PER_FILE,
            max_count=CHAT_FILE_MAX_COUNT,
            max_total_chars=CHAT_FILE_MAX_TOTAL_CHARS,
            preview_chars=CHAT_ATTACHMENT_PREVIEW_CHARS,
            supported_extensions=frozenset(SUPPORTED_CHAT_FILE_EXTENSIONS),
        ),
        logger=logger,
    )


def _build_message_with_files(message: str, attachment_context: str = "") -> str:
    return _build_message_with_files_impl(message, attachment_context)


def _build_user_input(
    message: str,
    images: list[ImageInput],
    attachment_context: str = "",
) -> Any:
    return _build_user_input_impl(message, images, attachment_context)


def _user_input_has_images(user_input: Any) -> bool:
    return _user_input_has_images_impl(user_input)


def _model_supports_images(provider: str, model_name: str) -> bool:
    return _model_supports_images_impl(provider, model_name)


def _stringify_user_input(user_input: Any) -> str:
    return _stringify_user_input_impl(user_input)


def _request_field_set(model: BaseModel) -> set[str]:
    fields = getattr(model, "model_fields_set", None)
    if fields is None:
        fields = getattr(model, "__fields_set__", set())
    return set(fields or set())


def _clip_text(text: Any, limit: int) -> str:
    normalized = re.sub(r"\s+", " ", str(text or "")).strip()
    if limit <= 0 or len(normalized) <= limit:
        return normalized
    return normalized[: max(1, limit - 3)].rstrip() + "..."


def _summary_llm_enabled() -> bool:
    return summary_llm_enabled()


def _summary_llm_timeout_seconds() -> float:
    return summary_llm_timeout_seconds(12.0)


def _normalize_llm_text_content(content: Any) -> str:
    return normalize_llm_text_content(content)


def _resolve_summary_model_config(
    session_id: str,
    preferred_model_config: Optional[dict[str, Any]] = None,
) -> Optional[ModelConfig]:
    from chat_store import get_session_panels

    if preferred_model_config:
        try:
            return _normalize_model_config(ModelConfig(**preferred_model_config))
        except Exception:
            logger.warning("Invalid preferred model config for summary session_id=%s", session_id)

    panels = get_session_panels(session_id)
    if not panels:
        return None

    selected = next((item for item in panels if item.get("is_primary")), panels[0])
    model_config = dict(selected.get("model_config") or {})
    if not str(model_config.get("panel_id") or "").strip():
        model_config["panel_id"] = str(selected.get("panel_id") or "summary")

    try:
        return _normalize_model_config(ModelConfig(**model_config))
    except Exception:
        logger.warning("Cannot resolve summary model config from session panels session_id=%s", session_id)
        return None


def _build_phase_summary_llm_prompt(
    turns: list[dict[str, Any]],
    *,
    total_turns: int,
) -> str:
    return build_phase_summary_llm_prompt(turns, total_turns=total_turns)


async def _try_llm_phase_summary_content(
    session_id: str,
    turns: list[dict[str, Any]],
    *,
    total_turns: int,
    preferred_model_config: Optional[dict[str, Any]] = None,
) -> Optional[str]:
    if not _summary_llm_enabled():
        return None

    model_config = _resolve_summary_model_config(
        session_id,
        preferred_model_config=preferred_model_config,
    )
    if model_config is None:
        return None

    from agent_core import get_llm

    try:
        llm = get_llm(
            provider=model_config.connection_type or model_config.provider,
            model_name=model_config.model,
            base_url=model_config.base_url,
            api_key=model_config.api_key if model_config.api_key else None,
            temperature=min(max(model_config.temperature, 0.0), 0.4),
        )
        prompt = _build_phase_summary_llm_prompt(turns, total_turns=total_turns)
        timeout_seconds = _summary_llm_timeout_seconds()
        response = await asyncio.wait_for(llm.ainvoke(prompt), timeout=timeout_seconds)
        content = _normalize_llm_text_content(getattr(response, "content", response))
        content = _clip_text(content, max(120, SESSION_MEMORY_AUTO_SUMMARY_MAX_CONTENT_CHARS))
        if not content:
            return None
        return content
    except Exception:
        logger.warning(
            "LLM summary generation failed, fallback to rule summary session_id=%s",
            session_id,
            exc_info=True,
        )
        return None


def _summary_turns(message_records: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return summary_turns(message_records, clip_text=_clip_text)


def _build_phase_summary_content(turns: list[dict[str, Any]], *, total_turns: int) -> str:
    return build_phase_summary_content(
        turns,
        total_turns=total_turns,
        clip_text=_clip_text,
        max_chars=SESSION_MEMORY_AUTO_SUMMARY_MAX_CONTENT_CHARS,
    )


async def _generate_session_phase_summary_memory(
    session_id: str,
    *,
    trigger: str,
    force: bool = False,
    preferred_model_config: Optional[dict[str, Any]] = None,
) -> Optional[dict[str, Any]]:
    from chat_store import (
        SQLiteChatMessageHistory,
        list_session_memory,
        pin_session_memory,
    )

    history = SQLiteChatMessageHistory(session_id=session_id)
    turns = _summary_turns(history.get_all_message_records())
    total_turns = len(turns)
    min_turns = max(2, SESSION_MEMORY_AUTO_SUMMARY_MIN_TURNS)

    if not force and total_turns < min_turns:
        raise ValueError(
            f"Need at least {min_turns} conversation turns before generating a phase summary memory."
        )

    summaries = list_session_memory(
        session_id,
        kind="summary",
        newest_first=True,
        db_path=history.db_path,
    )
    latest_auto = latest_auto_summary(summaries)
    covered_turns = covered_turns_from_summary(latest_auto)

    new_turns = max(0, total_turns - covered_turns)
    min_new_turns = max(1, SESSION_MEMORY_AUTO_SUMMARY_MIN_NEW_TURNS)
    if latest_auto and not force and new_turns < min_new_turns:
        return {
            "created": False,
            "memory": latest_auto,
            "reason": "up_to_date",
            "stats": {
                "total_turns": total_turns,
                "new_turns": new_turns,
                "required_new_turns": min_new_turns,
            },
        }

    window_size = max(2, SESSION_MEMORY_AUTO_SUMMARY_WINDOW_SIZE)
    window_turns = turns[-window_size:]
    content = _build_phase_summary_content(window_turns, total_turns=total_turns)
    generator = "rules"
    llm_content = await _try_llm_phase_summary_content(
        session_id,
        window_turns,
        total_turns=total_turns,
        preferred_model_config=preferred_model_config,
    )
    if llm_content:
        content = llm_content
        generator = "llm"
    meta = summarize_window_meta(
        total_turns=total_turns,
        trigger=trigger,
        generator=generator,
        window_turns=window_turns,
    )
    result = pin_session_memory(
        session_id,
        content=content,
        kind="summary",
        meta=meta,
        db_path=history.db_path,
    )
    if not result:
        return None
    return {
        **result,
        "reason": "created" if result.get("created") else "deduped",
        "stats": {
            "total_turns": total_turns,
            "new_turns": new_turns,
            "required_new_turns": min_new_turns,
        },
    }


async def _auto_generate_phase_summary_memory(
    session_id: str,
    *,
    trigger: str,
    preferred_model_config: Optional[dict[str, Any]] = None,
) -> None:
    try:
        result = await _generate_session_phase_summary_memory(
            session_id,
            trigger=trigger,
            force=False,
            preferred_model_config=preferred_model_config,
        )
        if result and result.get("created"):
            logger.info("Auto summary memory created session_id=%s", session_id)
    except ValueError:
        # Session not long enough yet.
        return
    except Exception:
        logger.exception("Auto summary memory generation failed session_id=%s", session_id)


def _get_attachment_promotion_task(
    attachment_id: str,
    vector_store_path: str,
) -> TaskRecord | None:
    attachment_id = str(attachment_id or "").strip()
    vector_store_path = str(vector_store_path or "").strip()
    if not attachment_id or not vector_store_path:
        return None

    try:
        return _get_task_store().get_attachment_promotion_task(
            attachment_id,
            vector_store_path,
        )
    except Exception:
        logger.exception(
            "Failed to resolve attachment promotion state: %s -> %s",
            attachment_id,
            vector_store_path,
        )
        return None


async def _get_or_build_agent(
    mc: ModelConfig,
    system_prompt: Optional[str] = None,
    web_search_enabled: bool = True,
    knowledge_base_enabled: bool = True,
    vector_store_path: Optional[str] = None,
    dashboard_template: Optional[dict[str, Any]] = None,
    enabled_mcp_servers: Optional[list[str]] = None,
):
    """获取或构建 Agent（带缓存）"""
    from agent_core import build_agent

    mc = _normalize_model_config(mc)
    api_key_hash = _hash_secret(mc.api_key)
    cache_key = _content_hash(
        {
            "connection_type": mc.connection_type or mc.provider,
            "provider": mc.provider,
            "model": mc.model,
            "base_url": mc.base_url,
            "api_key_hash": api_key_hash,
            "temperature": mc.temperature,
            "agent_mode": mc.agent_mode,
            "web_search_enabled": web_search_enabled,
            "knowledge_base_enabled": knowledge_base_enabled,
            "vector_store_path": str(_resolve_project_subdir(vector_store_path))
            if vector_store_path
            else "",
            "system_prompt": system_prompt or "",
            "dashboard_template": dashboard_template or {},
            "enabled_mcp_servers": normalize_mcp_server_names(enabled_mcp_servers),
        }
    )
    async with _agent_cache_lock:
        if cache_key not in _agent_cache:
            logger.info("Building new agent: %s", cache_key[:12])
            _agent_cache[cache_key] = await build_agent(
                provider=mc.provider,
                model_name=mc.model,
                base_url=mc.base_url,
                api_key=mc.api_key if mc.api_key else None,
                temperature=mc.temperature,
                agent_mode=mc.agent_mode,
                system_prompt=system_prompt,
                web_search_enabled=web_search_enabled,
                knowledge_base_enabled=knowledge_base_enabled,
                vector_store_path=vector_store_path,
                dashboard_template=dashboard_template,
                enabled_mcp_servers=normalize_mcp_server_names(enabled_mcp_servers),
            )
        return _agent_cache[cache_key]


async def _invoke_agent_stream(
    panel_id: str,
    mc: ModelConfig,
    message: Any,
    session_id: str,
    web_search_enabled: bool,
    knowledge_base_enabled: bool,
    system_prompt: Optional[str] = None,
    vector_store_path: Optional[str] = None,
    dashboard_template: Optional[dict[str, Any]] = None,
    enabled_mcp_servers: Optional[list[str]] = None,
    persist_history: bool = True,
    persist_user_history: bool = True,
    persist_ai_history: bool = True,
    replace_ai_history: bool = False,
    exclude_ai_answer_group_id: str = "",
    answer_group_id: str = "",
    raw_user_message: str = "",
    raw_images: Optional[list[dict[str, Any]]] = None,
    raw_files: Optional[list[dict[str, Any]]] = None,
    auto_summary_trigger: bool = False,
) -> AsyncGenerator[str, None]:
    """调用单个 agent，将结果封装为 SSE 事件流"""
    try:
        mc = _normalize_model_config(mc)
        agent = await _get_or_build_agent(
            mc,
            system_prompt=system_prompt,
            web_search_enabled=web_search_enabled,
            knowledge_base_enabled=knowledge_base_enabled,
            vector_store_path=vector_store_path,
            dashboard_template=dashboard_template,
            enabled_mcp_servers=normalize_mcp_server_names(enabled_mcp_servers),
        )
        answer_parts: list[str] = []
        dashboard_task_record: TaskRecord | None = None

        dashboard_requested = _should_start_dashboard_task(
            raw_user_message or message,
            knowledge_base_enabled=knowledge_base_enabled,
            logger=logger,
        ) and _dashboard_feature_enabled(dashboard_template)
        if dashboard_requested:
            prompt_excerpt = dashboard_prompt_excerpt(raw_user_message)
            dashboard_task_record = await _create_inline_task_record(
                "generate_dashboard",
                {
                    "panel_id": panel_id,
                    "prompt_excerpt": prompt_excerpt,
                },
                session_id=session_id,
                progress=20,
            )
            yield task_created_event(panel_id, dashboard_task_record)

        config_payload = build_agent_config_payload(
            session_id=session_id,
            persist_history=persist_history,
            persist_user_history=persist_user_history,
            persist_ai_history=persist_ai_history,
            replace_ai_history=replace_ai_history,
            exclude_ai_answer_group_id=exclude_ai_answer_group_id,
            panel_id=panel_id,
            model_id=mc.model,
            answer_group_id=answer_group_id,
            raw_user_message=raw_user_message,
            raw_images=raw_images or [],
            raw_files=raw_files or [],
            task_id=dashboard_task_record.task_id if dashboard_task_record else "",
            task_type=dashboard_task_record.task_type if dashboard_task_record else "",
        )

        # 尝试流式输出（LangGraph wrapper 支持 astream_answer）
        if hasattr(agent, "astream_answer"):
            async for item in agent.astream_answer(
                message,
                config=config_payload,
            ):
                event, chunk = stream_agent_item(panel_id, item)
                if chunk:
                    answer_parts.append(chunk)
                yield event
        else:
            # 非流式：一次性返回
            result = await agent.ainvoke(
                {"input": message},
                config=config_payload,
            )
            if _is_max_iterations_output(result.get("output", str(result))):
                logger.warning(
                    "panel_id=%s max iterations reached, attempting fallback", panel_id
                )
            outcome = await resolve_non_stream_agent_result(
                panel_id,
                result,
                mc=mc,
                message=message,
                is_max_iterations_output=_is_max_iterations_output,
                stringify_user_input=_stringify_user_input,
                fallback_generate=_fallback_generate,
            )
            for event in outcome.events:
                yield event
            if outcome.should_stop:
                if dashboard_task_record is not None and outcome.dashboard_error:
                    await fail_dashboard_task(
                        dashboard_task_record,
                        error=outcome.dashboard_error,
                        set_inline_task_state=_set_inline_task_state,
                    )
                return

            answer = outcome.answer
            sources = outcome.sources

            answer_parts.clear()
            answer_parts.append(answer)

            # Emit sources before answer chunks
            if sources:
                yield _panel_event(panel_id, "sources", sources=sources)

            # 将答案分块模拟流式效果（每20字符一块）
            for chunk in answer_chunks(answer, chunk_size=20):
                yield _panel_event(panel_id, "chunk", content=chunk)
                await asyncio.sleep(0.01)

        # 完成信号
        if dashboard_task_record is not None:
            final_answer = "".join(answer_parts)
            await finalize_dashboard_task(
                dashboard_task_record,
                final_answer,
                contains_dashboard_card=_contains_dashboard_card,
                summarize_dashboard_task_result=_summarize_dashboard_task_result,
                summarize_dashboard_task_error=_summarize_dashboard_task_error,
                set_inline_task_state=_set_inline_task_state,
            )

        if auto_summary_trigger and persist_ai_history:
            asyncio.create_task(
                _auto_generate_phase_summary_memory(
                    session_id,
                    trigger=f"chat_stream:{panel_id}",
                    preferred_model_config=_model_config_payload(mc),
                )
            )

        yield _done_event(panel_id)

    except Exception as e:
        logger.exception("Agent invocation failed panel_id=%s", panel_id)
        if 'dashboard_task_record' in locals() and dashboard_task_record is not None:
            await fail_dashboard_task(
                dashboard_task_record,
                error=str(e),
                set_inline_task_state=_set_inline_task_state,
            )
        err = _classify_error(e)
        yield _panel_event(
            panel_id,
            "error",
            content=err["message"],
            error_code=err["code"],
            suggestion=err["suggestion"],
        )


# ─────────────────────────────────────────────
# API 端点
# ─────────────────────────────────────────────


@app.get("/api/health")
async def health():
    return {"status": "ok", "timestamp": time.time()}


# ── 聊天 ──────────────────────────────────────


@app.post("/api/chat/parallel")
async def chat_parallel(request: ChatRequest, http_request: Request):
    """
    多模型并行对话 SSE 端点。
    每个面板的事件包含 panel_id 字段，前端据此路由到对应面板。
    """
    runtime = prepare_chat_route_runtime(
        request,
        resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
        validate_chat_payload=_validate_chat_payload,
        prepare_chat_files=_prepare_chat_files,
        build_user_input=_build_user_input,
        base_model_payload=_base_model_payload,
    )

    if not request.models:
        raise HTTPException(status_code=400, detail="至少需要选择一个模型")

    normalized_models = [_normalize_model_config(mc) for mc in request.models]

    from chat_store import replace_session_panels

    replace_session_panels(
        request.session_id,
        [_model_config_payload(mc) for mc in normalized_models],
    )

    async def event_generator() -> AsyncGenerator[str, None]:
        generators = build_parallel_agent_streams(
            normalized_models,
            runtime=runtime,
            request=request,
            invoke_agent_stream=_invoke_agent_stream,
        )
        async for item in stream_parallel_sse(
            generators,
            is_disconnected=http_request.is_disconnected,
            logger=logger,
        ):
            yield item

    return sse_streaming_response(event_generator())


@app.post("/api/chat/single")
async def chat_single(request: SingleChatRequest, http_request: Request):
    """单模型对话 SSE 端点"""
    runtime = prepare_chat_route_runtime(
        request,
        resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
        validate_chat_payload=_validate_chat_payload,
        prepare_chat_files=_prepare_chat_files,
        build_user_input=_build_user_input,
        base_model_payload=_base_model_payload,
    )
    normalized_panel_config = _normalize_model_config(request.panel_config)

    from chat_store import upsert_session_panel

    upsert_session_panel(
        request.session_id,
        _model_config_payload(normalized_panel_config),
    )

    async def event_generator():
        async for chunk in stream_single_sse(
            build_single_agent_stream(
                normalized_panel_config,
                runtime=runtime,
                request=request,
                invoke_agent_stream=_invoke_agent_stream,
            ),
            is_disconnected=http_request.is_disconnected,
        ):
            yield chunk

    return sse_streaming_response(event_generator())


# ── 会话管理 ──────────────────────────────────


@app.get("/api/connectors/mcp")
async def list_mcp_connectors():
    return {
        "connectors": list_mcp_server_catalog(),
        "default_enabled": default_mcp_server_names(),
    }


@app.get("/api/workspaces")
async def get_workspaces():
    from chat_store import list_workspaces

    workspaces = list_workspaces()
    return workspaces_payload(workspaces)


@app.post("/api/workspaces")
async def create_workspace_endpoint(request: CreateWorkspaceRequest):
    from chat_store import create_workspace

    try:
        preset = request.preset
        workspace = create_workspace(
            request.name,
            description=request.description,
            color=request.color,
            default_panels=[
                _base_model_payload(_normalize_model_config(panel))
                for panel in (preset.default_panels if preset else [])
            ],
            tool_config=(
                _base_model_payload(preset.tool_config)
                if preset is not None
                else None
            ),
            output_preset=(
                _base_model_payload(preset.output_preset)
                if preset is not None
                else None
            ),
            activate=request.activate,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    return {"ok": True, "workspace": workspace}


@app.patch("/api/workspaces/{workspace_id}")
async def update_workspace_endpoint(workspace_id: str, request: UpdateWorkspaceRequest):
    from chat_store import update_workspace

    field_set = _request_field_set(request)
    if not field_set:
        raise HTTPException(status_code=400, detail="至少需要提供一个工作区字段")

    try:
        preset = request.preset
        workspace = update_workspace(
            workspace_id,
            name=request.name if "name" in field_set else None,
            description=request.description if "description" in field_set else None,
            color=request.color if "color" in field_set else None,
            default_panels=(
                [
                    _base_model_payload(_normalize_model_config(panel))
                    for panel in preset.default_panels
                ]
                if "preset" in field_set and preset is not None
                else None
            ),
            tool_config=(
                _base_model_payload(preset.tool_config)
                if "preset" in field_set and preset is not None
                else None
            ),
            output_preset=(
                _base_model_payload(preset.output_preset)
                if "preset" in field_set and preset is not None
                else None
            ),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not workspace:
        raise HTTPException(status_code=404, detail="未找到工作区")

    return {"ok": True, "workspace": workspace}


@app.post("/api/workspaces/{workspace_id}/activate")
async def activate_workspace_endpoint(workspace_id: str):
    from chat_store import activate_workspace

    workspace = activate_workspace(workspace_id)
    if not workspace:
        raise HTTPException(status_code=404, detail="未找到工作区")

    return {"ok": True, "workspace": workspace}


@app.delete("/api/workspaces/{workspace_id}")
async def delete_workspace_endpoint(
    workspace_id: str,
    target_workspace_id: Optional[str] = None,
):
    from chat_store import delete_workspace

    try:
        result = delete_workspace(
            workspace_id,
            target_workspace_id=target_workspace_id,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not result:
        raise HTTPException(status_code=404, detail="未找到工作区")

    return {"ok": True, **result}


@app.get("/api/sessions")
async def get_sessions(
    query: str = "",
    archived: Optional[bool] = None,
    favorite: Optional[bool] = None,
    tag: str = "",
    workspace_id: Optional[str] = None,
):
    from chat_store import get_all_sessions

    sessions = get_all_sessions(
        query=query,
        archived=archived,
        favorite=favorite,
        tag=tag,
        workspace_id=workspace_id,
    )
    return {"sessions": sessions}


@app.post("/api/sessions")
async def create_session(request: CreateSessionRequest):
    from chat_store import (
        SQLiteChatMessageHistory,
        connect_sqlite,
        get_session,
        get_workspace,
        update_session_meta,
    )

    try:
        return create_session_record(
            request,
            history_factory=SQLiteChatMessageHistory,
            connect_sqlite=connect_sqlite,
            get_session=get_session,
            get_workspace=get_workspace,
            update_session_meta=update_session_meta,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.patch("/api/sessions/{session_id}")
async def update_session_endpoint(session_id: str, request: UpdateSessionRequest):
    from chat_store import update_session_meta

    if not session_update_requested(request):
        raise HTTPException(status_code=400, detail="至少需要提供一个可更新字段")

    try:
        session = update_session_meta(
            session_id,
            title=request.title,
            is_archived=request.is_archived,
            is_favorite=request.is_favorite,
            is_pinned=request.is_pinned,
            tags=request.tags,
            workspace_id=request.workspace_id,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not session:
        raise HTTPException(status_code=404, detail="会话不存在")

    return {"ok": True, "session": session}


@app.post("/api/sessions/reorder")
async def reorder_sessions_endpoint(request: ReorderSessionsRequest):
    from chat_store import get_all_sessions, reorder_sessions

    try:
        result = reorder_sessions(
            request.session_ids,
            workspace_id=request.workspace_id,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    sessions = get_all_sessions(workspace_id=request.workspace_id)
    return reorder_sessions_payload(result, sessions)


@app.delete("/api/sessions/{session_id}")
async def delete_session_endpoint(session_id: str):
    from chat_store import delete_session

    deck_ids = _deck_store.list_ids_by_session(session_id)
    async with _tasks_lock:
        stale_task_ids = [
            task_id
            for task_id, record in _tasks.items()
            if str(record.session_id or "").strip() == session_id
        ]
        _suppressed_task_ids.update(stale_task_ids)
        for task_id in stale_task_ids:
            _tasks.pop(task_id, None)
        _prune_task_records_locked()

    _share_link_store.delete_for_resource("session", session_id)
    for deck_id in deck_ids:
        _share_link_store.delete_for_resource("deck", deck_id)
    _get_task_store().delete_for_session(session_id)
    _artifact_store.delete_by_session(session_id)
    _deck_store.delete_by_session(session_id)
    delete_session(session_id)
    return {"ok": True}


@app.get("/api/bookmarks")
async def list_bookmarks_endpoint(session_id: Optional[str] = None):
    from chat_store import list_bookmarks

    return {"bookmarks": list_bookmarks(session_id=session_id)}


@app.post("/api/bookmarks")
async def create_bookmark_endpoint(request: CreateBookmarkRequest):
    from chat_store import create_or_update_bookmark

    try:
        bookmark = create_or_update_bookmark(
            request.session_id,
            role=request.role,
            message_id=request.message_id,
            panel_id=request.panel_id,
            answer_group_id=request.answer_group_id,
            content=request.content,
            model_id=request.model_id,
            session_title=request.session_title,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not bookmark:
        raise HTTPException(status_code=404, detail="未找到消息")

    return {"ok": True, "bookmark": bookmark}


@app.delete("/api/bookmarks/{bookmark_id}")
async def delete_bookmark_endpoint(bookmark_id: str):
    from chat_store import delete_bookmark

    ok = delete_bookmark(bookmark_id)
    if not ok:
        raise HTTPException(status_code=404, detail="未找到书签")

    return {"ok": True}


@app.post("/api/sessions/{session_id}/share", response_model=ShareLinkResponse)
async def create_session_share_link(session_id: str, request: Request):
    from chat_store import get_session

    _require_remote_share_secret(request)
    share_secret = _current_share_link_secret()
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="未找到会话")

    payload = create_share_link_payload(
        "session",
        session_id,
        request,
        secret=share_secret,
        encode_share_token=_encode_share_token,
        build_share_url=_build_share_url,
    )
    record = _share_link_store.upsert(
        share_token=payload["share_token"],
        resource_type="session",
        resource_id=session_id,
        expires_at=time.time() + SHARE_LINK_TTL_SECONDS,
        created_by_ip=_request_client_ip(request),
        created_user_agent=_request_user_agent(request),
    )
    _audit_security_event(
        "create_session_share_link",
        request,
        details=f"session_id={session_id}",
    )
    return ShareLinkResponse(**payload, expires_at=record.expires_at)


@app.get("/api/sessions/{session_id}/messages")
async def get_session_messages(session_id: str):
    payload = _build_session_messages_payload(session_id)
    return {
        "messages": payload["messages"],
        "context_limit": payload["context_limit"],
        "total_messages": payload["total_messages"],
        "panels": payload["panels"],
        "panel_messages": payload["panel_messages"],
    }


@app.post("/api/sessions/{session_id}/messages/import")
async def import_session_messages_endpoint(
    session_id: str,
    request: ImportSessionMessagesRequest,
):
    from chat_store import SQLiteChatMessageHistory, replace_session_panels

    history = SQLiteChatMessageHistory(session_id=session_id)
    if history.get_all_message_records():
        raise HTTPException(
            status_code=400,
            detail="Session already has messages; import only supports empty sessions.",
        )

    normalized_panels: list[dict[str, Any]] = []
    panel_ids: set[str] = set()
    for panel in request.panels:
        normalized_panel = _base_model_payload(_normalize_model_config(panel))
        panel_id = str(normalized_panel.get("panel_id") or "").strip()
        if not panel_id:
            raise HTTPException(status_code=400, detail="Imported panels require panel_id.")
        if panel_id in panel_ids:
            raise HTTPException(status_code=400, detail="Imported panels contain duplicate panel_id.")
        panel_ids.add(panel_id)
        normalized_panels.append(normalized_panel)

    if normalized_panels:
        replace_session_panels(
            session_id,
            normalized_panels,
            db_path=history.db_path,
        )

    for message in request.messages:
        images = [_base_model_payload(image) for image in message.images]
        files = [_base_model_payload(file) for file in message.files]
        answer_group_id = str(message.answer_group_id or "").strip()

        if message.role == "user":
            history.add_user_message(
                message.content,
                answer_group_id=answer_group_id,
                images=images,
                files=files,
            )
            continue

        panel_id = str(message.panel_id or "").strip()
        if not panel_id:
            raise HTTPException(
                status_code=400,
                detail="Imported assistant messages require panel_id.",
            )
        if panel_ids and panel_id not in panel_ids:
            raise HTTPException(
                status_code=400,
                detail="Imported assistant message references an unknown panel_id.",
            )

        history.add_ai_message(
            message.content,
            model_id=str(message.model_id or "").strip(),
            panel_id=panel_id,
            answer_group_id=answer_group_id,
            images=images,
            files=files,
            sources=[dict(item) for item in message.sources if isinstance(item, dict)],
            workflow_nodes=[
                dict(item) for item in message.workflow_nodes if isinstance(item, dict)
            ],
            task_id=str(message.task_id or "").strip(),
            task_type=str(message.task_type or "").strip(),
        )

    payload = _build_session_messages_payload(session_id)
    return {
        "messages": payload["messages"],
        "context_limit": payload["context_limit"],
        "total_messages": payload["total_messages"],
        "panels": payload["panels"],
        "panel_messages": payload["panel_messages"],
    }


@app.post("/api/sessions/{session_id}/messages/feedback")
async def set_message_feedback_endpoint(
    session_id: str,
    request: SetMessageFeedbackRequest,
):
    from chat_store import set_message_feedback

    try:
        result = set_message_feedback(
            session_id,
            feedback_value=request.value,
            message_id=request.message_id,
            panel_id=request.panel_id,
            answer_group_id=request.answer_group_id,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not result:
        raise HTTPException(status_code=404, detail="未找到消息")

    return {"ok": True, "feedback": result}


@app.post("/api/sessions/{session_id}/messages/truncate")
async def truncate_session_messages_endpoint(
    session_id: str,
    request: TruncateSessionMessagesRequest,
):
    from chat_store import truncate_session_from_answer_group

    _validate_chat_payload(request.content, request.images, request.files)
    try:
        result = truncate_session_from_answer_group(
            session_id,
            answer_group_id=request.answer_group_id,
            content=request.content,
            images=[_base_model_payload(image) for image in request.images],
            files=[_base_model_payload(file) for file in request.files],
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not result:
        raise HTTPException(status_code=404, detail="未找到回答分组")

    return {"ok": True, "result": result}


@app.post("/api/sessions/{session_id}/retrieval-feedback")
async def set_retrieval_feedback_endpoint(
    session_id: str,
    request: SetRetrievalFeedbackRequest,
):
    from chat_store import set_retrieval_feedback

    try:
        result = set_retrieval_feedback(
            session_id,
            panel_id=request.panel_id,
            answer_group_id=request.answer_group_id,
            source=request.source,
            feedback_value=request.value,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    return {"ok": True, "feedback": result}


@app.get("/api/sessions/{session_id}/retrieval-feedback")
async def list_retrieval_feedback_endpoint(
    session_id: str,
    panel_id: str,
    answer_group_id: str,
):
    from chat_store import list_retrieval_feedback

    try:
        feedback = list_retrieval_feedback(
            session_id,
            panel_id=panel_id,
            answer_group_id=answer_group_id,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    return {
        "session_id": session_id,
        "panel_id": panel_id,
        "answer_group_id": answer_group_id,
        "feedback": feedback,
    }


@app.get("/api/sessions/{session_id}/memory")
async def get_session_memory(session_id: str, kind: str = ""):
    from chat_store import get_session, list_session_memory

    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="会话不存在")

    try:
        memories = list_session_memory(
            session_id,
            kind=kind or None,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    try:
        return session_memory_payload(
            session_id=session_id,
            session=session,
            memories=memories,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="会话记忆不存在") from exc


@app.post("/api/sessions/{session_id}/memory/pin")
async def pin_session_memory_endpoint(session_id: str, request: PinSessionMemoryRequest):
    from chat_store import pin_session_memory

    try:
        result = pin_session_memory(
            session_id,
            content=request.content,
            kind=request.kind,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if not result:
        raise HTTPException(status_code=404, detail="会话不存在")

    return pin_session_memory_payload(result)


@app.patch("/api/sessions/{session_id}/memory/{memory_id}")
async def update_session_memory_endpoint(
    session_id: str,
    memory_id: str,
    request: UpdateSessionMemoryRequest,
):
    from chat_store import update_session_memory

    try:
        updates = session_memory_updates(
            request,
            field_set=_request_field_set(request),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    try:
        memory = update_session_memory(
            session_id,
            memory_id,
            **updates,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    try:
        return update_session_memory_payload(memory)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到会话记忆") from exc


@app.post("/api/sessions/{session_id}/memory/summarize")
async def summarize_session_memory_endpoint(session_id: str, force: bool = False):
    from chat_store import get_session

    session = get_session(session_id)

    try:
        result = await _generate_session_phase_summary_memory(
            session_id,
            trigger="manual_api",
            force=force,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    try:
        return summarize_session_memory_payload(
            session=session,
            result=result,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到会话") from exc


@app.delete("/api/sessions/{session_id}/memory/{memory_id}")
async def delete_session_memory_endpoint(session_id: str, memory_id: str):
    from chat_store import delete_session_memory

    deleted = delete_session_memory(session_id, memory_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="会话记忆不存在")

    return delete_session_memory_payload(deleted)


@app.get("/api/sessions/{session_id}/attachments")
async def get_session_attachments(
    session_id: str,
    vector_store_path: Optional[str] = None,
    workspace_id: Optional[str] = None,
):
    from chat_store import SQLiteChatMessageHistory

    _require_workspace_session(session_id, workspace_id)
    history = SQLiteChatMessageHistory(session_id=session_id)
    return session_attachments_payload(
        session_id=session_id,
        message_records=history.get_all_message_records(),
        preview_char_limit=CHAT_ATTACHMENT_PREVIEW_CHARS,
        vector_store_path=_effective_vector_store_path(vector_store_path),
        collect_attachments=_collect_session_attachments,
        attach_current_kb_status=attach_current_kb_status,
        lookup_task=_get_attachment_promotion_task,
    )


@app.post("/api/sessions/{session_id}/attachments/{attachment_id}/promote")
async def promote_session_attachment_to_kb(
    session_id: str,
    attachment_id: str,
    vector_store_path: Optional[str] = None,
    workspace_id: Optional[str] = None,
):
    _require_workspace_session(session_id, workspace_id)
    attachment = _find_session_attachment(
        session_id,
        attachment_id,
        preview_char_limit=CHAT_ATTACHMENT_PREVIEW_CHARS,
    )
    target_vector_store_path = _effective_vector_store_path(vector_store_path)
    try:
        promotion_request = prepare_attachment_promotion(
            session_id=session_id,
            attachment_id=attachment_id,
            attachment=attachment,
            target_vector_store_path=target_vector_store_path,
            workspace_id=workspace_id,
            existing_task=_get_attachment_promotion_task(
                attachment_id,
                target_vector_store_path,
            ),
            task_record_payload=task_record_payload,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc.args[0])) from exc
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    if promotion_request["dedupe_payload"] is not None:
        return promotion_request["dedupe_payload"]

    return await enqueue_task(
        _tasks,
        _tasks_lock,
        **promotion_request["enqueue_kwargs"],
        prune_in_memory=_prune_task_records_locked,
        persist_record=_persist_task_record,
        prune_persisted=_prune_persisted_tasks,
        run_task=_run_task,
        spawn_background_task=asyncio.create_task,
        logger=logger,
    )


@app.post("/api/sessions/{session_id}/answer-groups/{answer_group_id}/promote")
async def promote_answer_group(session_id: str, answer_group_id: str, panel_id: str):
    from chat_store import SQLiteChatMessageHistory, promote_panel_answer

    history = SQLiteChatMessageHistory(session_id=session_id)
    promoted = promote_panel_answer(
        session_id,
        answer_group_id,
        panel_id,
        db_path=history.db_path,
    )
    if not promoted:
        raise HTTPException(status_code=404, detail="未找到回答分组或面板消息")
    return {"ok": True, **promoted}


@app.get("/api/sessions/{session_id}/answer-groups/{answer_group_id}/review")
async def review_answer_group(session_id: str, answer_group_id: str):
    try:
        return _build_answer_group_review_payload(session_id, answer_group_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该回答分组的面板答案") from exc


@app.post("/api/sessions/{session_id}/answer-groups/{answer_group_id}/promote/recommended")
async def promote_recommended_answer_group(session_id: str, answer_group_id: str):
    from chat_store import SQLiteChatMessageHistory, promote_panel_answer

    try:
        review = _build_answer_group_review_payload(session_id, answer_group_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该回答分组的面板答案") from exc

    panel_id = str(review.get("recommended_panel_id") or "").strip()
    history = SQLiteChatMessageHistory(session_id=session_id)
    promoted = promote_panel_answer(
        session_id,
        answer_group_id,
        panel_id,
        db_path=history.db_path,
    )
    if not promoted:
        raise HTTPException(status_code=404, detail="未找到该回答分组的面板答案")
    return {"ok": True, "review": review, **promoted}


@app.delete("/api/sessions/{session_id}/messages")
async def clear_session_messages(session_id: str):
    from agent_core import clear_session_history

    clear_session_history(session_id)
    return {"ok": True}


# ── 文档管理 ──────────────────────────────────


@app.post("/api/documents/upload")
async def upload_documents(
    request: Request,
    files: list[UploadFile] = File(...),
    vector_store_path: Optional[str] = Form(default=None),
):
    _require_remote_admin(request)
    temp_paths: list[str] = []
    try:
        effective_vector_store_path = _effective_vector_store_path(vector_store_path)
        temp_paths, file_names = await stage_upload_files(files)
        record = build_upload_documents_task_record(
            temp_paths=temp_paths,
            file_names=file_names,
            vector_store_path=effective_vector_store_path,
        )
        async with _tasks_lock:
            _tasks[record.task_id] = record
            _prune_task_records_locked(record.created_at)
        _persist_task_record(record)
        _prune_persisted_tasks()

        asyncio.create_task(_run_task(record))
        logger.info("task_id=%s task_type=upload_documents 已创建", record.task_id)
        _audit_security_event(
            "upload_documents",
            request,
            details=(
                f"file_count={len(file_names)} "
                f"vector_store_path={effective_vector_store_path}"
            ),
        )
        return upload_documents_response(
            record,
            file_count=len(file_names),
            vector_store_path=effective_vector_store_path,
        )
    except ValueError as e:
        if temp_paths:
            cleanup_temp_paths(temp_paths)
        _audit_security_event(
            "upload_documents",
            request,
            result="rejected",
            details=str(e),
        )
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        if temp_paths:
            cleanup_temp_paths(temp_paths)
        logger.exception("Document upload failed")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/documents/stats")
async def get_document_stats(request: Request, path: Optional[str] = None):
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    pipeline = DocPipeline(vector_store_path=_effective_vector_store_path(path))
    try:
        pipeline.load_store()
        stats = pipeline.get_stats()
        stats.setdefault("store_path", pipeline.vector_store_path)
        _audit_security_event(
            "get_document_stats",
            request,
            details=f"path={pipeline.vector_store_path}",
        )
        return stats
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── 模型发现 ──────────────────────────────────


@app.get("/api/models/ollama")
async def get_ollama_models(base_url: str = "http://localhost:11434"):
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            resp = await client.get(f"{base_url}/api/tags")
        resp.raise_for_status()
        models = resp.json().get("models", [])
        return {"models": [m["name"] for m in models]}
    except httpx.HTTPError as e:
        logger.warning("Cannot reach Ollama: %s", e)
        return {"models": [], "error": str(e)}


# ── 配置管理 ──────────────────────────────────


@app.get("/api/security/status", response_model=SecurityStatusResponse)
async def get_security_status(request: Request):
    _require_remote_admin(request)
    payload = _security_status_payload()
    _audit_security_event(
        "get_security_status",
        request,
        details=(
            f"remote_clients={payload['allow_remote_clients']} "
            f"admin_ready={payload['remote_admin_ready']} "
            f"share_ready={payload['remote_share_ready']}"
        ),
    )
    return payload


@app.get("/api/config")
async def get_config(request: Request):
    _require_remote_admin(request)
    payload = {
        "tavily_api_key_set": bool(os.environ.get("TAVILY_API_KEY")),
        "llm_provider": os.environ.get("LLM_PROVIDER", "ollama"),
        "ollama_model": os.environ.get("OLLAMA_MODEL", "qwen3.5-2B:latest"),
        "ollama_base_url": os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434"),
        "openrouter_model": os.environ.get("OPENROUTER_MODEL", ""),
        "openrouter_base_url": os.environ.get(
            "OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"
        ),
    }
    _audit_security_event("get_config", request)
    return payload


@app.post("/api/config")
async def save_config(request: Request, payload: SaveConfigRequest):
    _require_remote_admin(request)
    if payload.tavily_api_key is not None:
        os.environ["TAVILY_API_KEY"] = payload.tavily_api_key
        # 清除 agent cache 以便下次重建时使用新 key
    _audit_security_event(
        "save_config",
        request,
        details=f"tavily_api_key_set={payload.tavily_api_key is not None}",
    )
    return {"ok": True}


@app.post("/api/agents/reset")
async def reset_agents(request: Request):
    """清除 agent 缓存，下次请求时重新构建"""
    _require_remote_admin(request)
    await _clear_agent_cache()
    _audit_security_event("reset_agents", request)
    return {"ok": True, "message": "智能体缓存已清除"}


# ── System Prompts ─────────────────────────────


@app.get("/api/prompts")
async def list_prompts(request: Request):
    from chat_store import get_all_system_prompts

    _require_remote_admin(request)
    payload = {"prompts": get_all_system_prompts()}
    _audit_security_event(
        "list_prompts",
        request,
        details=f"prompt_count={len(payload['prompts'])}",
    )
    return payload


@app.post("/api/prompts")
async def create_prompt(request: Request, payload: CreatePromptRequest):
    from chat_store import create_system_prompt

    _require_remote_admin(request)
    prompt = create_system_prompt(
        payload.name,
        payload.content,
        vector_store_id=payload.vector_store_id or "",
        dashboard_template=payload.dashboard_template or {},
    )
    return prompt


@app.put("/api/prompts/{prompt_id}")
async def update_prompt(prompt_id: str, request: Request, payload: UpdatePromptRequest):
    from chat_store import update_system_prompt

    _require_remote_admin(request)
    prompt = update_system_prompt(
        prompt_id,
        payload.name,
        payload.content,
        vector_store_id=payload.vector_store_id or "",
        dashboard_template=payload.dashboard_template or {},
    )
    if not prompt:
        raise HTTPException(status_code=404, detail="未找到提示词")
    await _clear_agent_cache()
    return prompt


@app.delete("/api/prompts/{prompt_id}")
async def delete_prompt(prompt_id: str, request: Request):
    from chat_store import delete_system_prompt

    _require_remote_admin(request)
    ok = delete_system_prompt(prompt_id)
    if not ok:
        raise HTTPException(
            status_code=404, detail="未找到提示词，或该提示词是最后一个默认项"
        )
    # Rebuild agents since active prompt may have changed
    await _clear_agent_cache()
    return {"ok": True}


@app.post("/api/prompts/{prompt_id}/activate")
async def activate_prompt(prompt_id: str, request: Request):
    from chat_store import activate_system_prompt, get_all_system_prompts
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    ok = activate_system_prompt(prompt_id)
    if not ok:
        raise HTTPException(status_code=404, detail="未找到提示词")
    # Clear agent cache so next request rebuilds with new prompt
    await _clear_agent_cache()

    # Load knowledge base bound to this role (if any)
    prompts = get_all_system_prompts()
    target = next((p for p in prompts if p["id"] == prompt_id), None)
    kb_status = "none"
    if target and target.get("vector_store_id"):
        try:
            pipeline = DocPipeline(vector_store_path=target["vector_store_id"])
            loaded = pipeline.load_store()
            kb_status = "loaded" if loaded else "error"
        except Exception as e:
            logger.warning("Failed to load KB for prompt %s: %s", prompt_id, e)
            kb_status = "error"
    return {"ok": True, "kb_status": kb_status}


# ── 异步任务管理 ────────────────────────────────


@app.post("/api/tasks")
async def create_task(request: CreateTaskRequest):
    """创建异步任务，立即返回 task_id，任务在后台执行"""
    payload = await enqueue_task(
        _tasks,
        _tasks_lock,
        task_type=request.task_type,
        params=request.params,
        session_id=request.session_id,
        prune_in_memory=_prune_task_records_locked,
        persist_record=_persist_task_record,
        prune_persisted=_prune_persisted_tasks,
        run_task=_run_task,
        spawn_background_task=asyncio.create_task,
        logger=logger,
        on_record_created=(
            persist_web_research_task_placeholder
            if request.task_type == "web_research"
            else None
        ),
    )
    return payload


@app.get("/api/tasks/{task_id}")
async def get_task(task_id: str):
    """轮询任务状态"""
    async with _tasks_lock:
        _prune_task_records_locked()
        record = _tasks.get(task_id)
    if record is None:
        _prune_persisted_tasks()
        record = _get_task_store().get(task_id)

    if record is None:
        raise HTTPException(status_code=404, detail="未找到任务")

    return task_record_payload(record)


@app.get("/api/tasks")
async def list_tasks(limit: int = 20):
    """列出最近的任务（按创建时间倒序）"""
    async with _tasks_lock:
        _prune_task_records_locked()
        in_memory_tasks = list(_tasks.values())

    _prune_persisted_tasks()
    return list_tasks_payload(
        in_memory_tasks=in_memory_tasks,
        persisted_tasks=_get_task_store().list_recent(limit=max(limit, TASK_HISTORY_LIMIT)),
        limit=limit,
    )


# ── 会话重置 ──────────────────────────────────


@app.post("/api/sessions/{session_id}/reset")
async def reset_session(session_id: str):
    """彻底清除会话上下文和 Agent 缓存，确保下次是全新起点"""
    from agent_core import clear_session_history
    from chat_store import clear_session_memory

    clear_session_history(session_id)
    clear_session_memory(session_id)
    await _clear_agent_cache()
    return {"ok": True, "message": "会话已完全重置"}


@app.post("/api/decks")
async def create_deck(request: CreateDeckRequest):
    """Create a structured deck draft from session history."""
    from chat_store import SQLiteChatMessageHistory

    history = SQLiteChatMessageHistory(session_id=request.session_id)
    try:
        messages = _resolve_report_messages(
            history,
            answer_group_id=request.answer_group_id,
            panel_id=request.panel_id,
        )
    except KeyError as exc:
        raise HTTPException(
            status_code=404, detail="Requested deck scope was not found."
        ) from exc
    if not messages:
        raise HTTPException(status_code=400, detail="该会话没有对话记录")

    try:
        deck = await build_deck(
            messages=messages,
            **build_create_deck_kwargs(
                request,
                resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
                normalize_deck_theme=normalize_deck_theme,
            ),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    _deck_store.save(deck)
    artifact = _create_deck_artifact(deck)
    payload = deck.model_dump(mode="json")
    payload["artifact_id"] = artifact.artifact_id
    return payload


@app.get("/api/decks/{deck_id}")
async def get_deck(deck_id: str):
    """Return a previously generated deck draft."""
    try:
        deck = _deck_store.get(deck_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc
    return deck.model_dump(mode="json")


@app.patch("/api/decks/{deck_id}")
async def update_deck(deck_id: str, request: UpdateDeckRequest):
    """Persist edited deck content from the in-app editor."""
    try:
        deck = _deck_store.get(deck_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc

    if request.slides is not None:
        if not request.slides:
            raise HTTPException(status_code=400, detail="演示稿至少需要保留一页")
    apply_deck_update(
        deck,
        request,
        normalize_deck_theme=normalize_deck_theme,
    )

    _deck_store.save(deck)
    _sync_deck_artifacts(deck)
    return deck.model_dump(mode="json")


@app.post("/api/decks/{deck_id}/slides/{slide_id}/regenerate")
async def regenerate_saved_deck_slide(
    deck_id: str,
    slide_id: str,
    request: RegenerateDeckSlideRequest,
):
    """Regenerate a single slide in an existing deck draft."""
    from chat_store import SQLiteChatMessageHistory

    try:
        deck = _deck_store.get(deck_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc

    history = SQLiteChatMessageHistory(session_id=deck.meta.session_id)
    try:
        messages = _resolve_report_messages(
            history,
            answer_group_id=getattr(deck.meta, "source_answer_group_id", None),
            panel_id=getattr(deck.meta, "source_panel_id", None),
        )
    except KeyError as exc:
        raise HTTPException(
            status_code=404, detail="Requested deck scope was not found."
        ) from exc
    if not messages:
        raise HTTPException(status_code=400, detail="该会话没有对话记录")

    regenerate_kwargs = build_regenerate_deck_kwargs(
        deck,
        request,
        normalize_model_config=_normalize_model_config,
        resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
    )

    try:
        regenerated_slide = await regenerate_deck_slide(
            deck=deck,
            slide_id=slide_id,
            messages=messages,
            **regenerate_kwargs,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到页面") from exc
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    replace_deck_slide(deck, regenerated_slide)
    _deck_store.save(deck)
    _sync_deck_artifacts(deck)
    return deck.model_dump(mode="json")


@app.get("/api/decks/{deck_id}/export")
async def export_deck(deck_id: str, format: str = "pptx"):
    """Export a saved deck draft."""
    if format != "pptx":
        raise HTTPException(status_code=400, detail="当前仅支持导出 PPTX")

    try:
        deck = _deck_store.get(deck_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc

    try:
        export_payload = export_deck_payload(
            deck,
            export_deck_to_pptx=export_deck_to_pptx,
            build_export_filename=build_export_filename,
        )
    except RuntimeError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    return Response(
        content=export_payload["content"],
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": _build_download_content_disposition(
                export_payload["filename"]
            )
        },
    )


@app.post("/api/decks/{deck_id}/share", response_model=ShareLinkResponse)
async def create_deck_share_link(deck_id: str, request: Request):
    _require_remote_share_secret(request)
    share_secret = _current_share_link_secret()
    try:
        _deck_store.get(deck_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到演示稿") from exc

    payload = create_share_link_payload(
        "deck",
        deck_id,
        request,
        secret=share_secret,
        encode_share_token=_encode_share_token,
        build_share_url=_build_share_url,
    )
    record = _share_link_store.upsert(
        share_token=payload["share_token"],
        resource_type="deck",
        resource_id=deck_id,
        expires_at=time.time() + SHARE_LINK_TTL_SECONDS,
        created_by_ip=_request_client_ip(request),
        created_user_agent=_request_user_agent(request),
    )
    _audit_security_event(
        "create_deck_share_link",
        request,
        details=f"deck_id={deck_id}",
    )
    return ShareLinkResponse(**payload, expires_at=record.expires_at)


@app.post("/api/reports/generate")
async def generate_report(request: GenerateReportRequest):
    """根据会话历史生成 Slidev 格式 Markdown 报告"""
    from chat_store import SQLiteChatMessageHistory

    history = SQLiteChatMessageHistory(session_id=request.session_id)
    try:
        msgs = _resolve_report_messages(
            history,
            answer_group_id=request.answer_group_id,
            panel_id=request.panel_id,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到对应的报告消息范围。") from exc
    if not msgs:
        raise HTTPException(status_code=400, detail="该会话没有对话记录")

    # Use first user message as title
    title = "AI 对话报告"
    title = build_chat_report_title(msgs)

    try:
        artifact, title, markdown = _create_report_artifact(
            session_id=request.session_id,
            messages=msgs,
            answer_group_id=str(request.answer_group_id or "").strip(),
            panel_id=str(request.panel_id or "").strip(),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    return {
        "markdown": markdown,
        "title": title,
        "artifact_id": artifact.artifact_id,
    }


@app.get("/api/reports/download/{session_id}")
async def download_report_pptx(
    session_id: str,
    answer_group_id: Optional[str] = None,
    panel_id: Optional[str] = None,
):
    """将会话历史转换为 PPTX 文件并下载"""
    from chat_store import SQLiteChatMessageHistory
    import io

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx 未安装，请在 requirements.txt 中添加 python-pptx 并重新安装依赖",
        )

    history = SQLiteChatMessageHistory(session_id=session_id)
    try:
        msgs = _resolve_report_messages(
            history,
            answer_group_id=answer_group_id,
            panel_id=panel_id,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="未找到对应的报告消息范围。") from exc
    try:
        qa_pairs = ensure_deckable_chat(msgs)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    if not msgs:
        raise HTTPException(status_code=400, detail="该会话没有对话记录")

    report_payload = report_download_payload(
        msgs,
        ensure_deckable_chat=ensure_deckable_chat,
        build_chat_report_title=build_chat_report_title,
        presentation_factory=Presentation,
        body_font_size=Pt(12),
        populate_chat_report_presentation=populate_chat_report_presentation,
        safe_report_filename=safe_report_filename,
    )
    prs = report_payload["presentation"]

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": _build_download_content_disposition(
                report_payload["filename"]
            )
        },
    )


# ── 知识库管理 ─────────────────────────────────


@app.get("/api/sessions/{session_id}/artifacts")
async def list_session_artifacts(
    session_id: str,
    artifact_type: Optional[str] = None,
):
    artifacts = _artifact_store.list_by_session(
        session_id,
        artifact_type=str(artifact_type or "").strip(),
    )
    return {"artifacts": [_artifact_payload(artifact) for artifact in artifacts]}


@app.get("/api/artifacts/{artifact_id}")
async def get_artifact(artifact_id: str):
    try:
        artifact = _artifact_store.get(artifact_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
    return _artifact_payload(artifact)


@app.patch("/api/artifacts/{artifact_id}")
async def update_artifact(artifact_id: str, request: UpdateArtifactRequest):
    try:
        artifact = _artifact_store.get(artifact_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="Artifact was not found.") from exc

    next_title = str(request.title or "").strip()

    if artifact.artifact_type == "report":
        if next_title:
            artifact.title = next_title
        if request.markdown is not None:
            artifact.content["markdown"] = str(request.markdown or "").strip()
        _artifact_store.save(artifact)
        return _artifact_payload(artifact)

    if artifact.artifact_type == "deck":
        if request.markdown is not None:
            raise HTTPException(
                status_code=400,
                detail="Deck artifact does not support markdown patching.",
            )
        deck_id = str(
            artifact.linked_resource_id or artifact.content.get("deck_id") or ""
        ).strip()
        if not deck_id:
            raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
        try:
            deck = _deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        if next_title:
            deck.meta.title = next_title
            if deck.slides and deck.slides[0].type == "cover":
                deck.slides[0].title = next_title
            _deck_store.save(deck)
        _sync_deck_artifacts(deck)
        return _artifact_payload(_artifact_store.get(artifact_id))

    raise HTTPException(status_code=400, detail="Unsupported artifact type.")


@app.get("/api/artifacts/{artifact_id}/export")
async def export_artifact(artifact_id: str, format: str = ""):
    import io

    try:
        artifact = _artifact_store.get(artifact_id)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="Artifact was not found.") from exc

    if artifact.artifact_type == "report":
        export_format = str(format or "md").strip().lower() or "md"
        if export_format == "md":
            filename = f"{safe_report_filename(artifact.title)}.md"
            return Response(
                content=str(artifact.content.get("markdown") or ""),
                media_type="text/markdown; charset=utf-8",
                headers={
                    "Content-Disposition": _build_download_content_disposition(filename)
                },
            )
        if export_format != "pptx":
            raise HTTPException(status_code=400, detail="Report artifact only supports md / pptx.")

        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError as exc:
            raise HTTPException(
                status_code=500,
                detail="python-pptx is not installed.",
            ) from exc

        raw_pairs = (
            artifact.content.get("qa_pairs")
            if isinstance(artifact.content.get("qa_pairs"), list)
            else []
        )
        qa_pairs = [
            (
                str(item.get("question") or "").strip(),
                str(item.get("answer") or "").strip(),
            )
            for item in raw_pairs
            if isinstance(item, dict)
            and (
                str(item.get("question") or "").strip()
                or str(item.get("answer") or "").strip()
            )
        ]
        if not qa_pairs:
            raise HTTPException(status_code=400, detail="Report artifact has no exportable content.")

        presentation = Presentation()
        populate_chat_report_presentation(
            presentation,
            title=artifact.title,
            qa_pairs=qa_pairs,
            body_font_size=Pt(12),
        )
        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        filename = f"{safe_report_filename(artifact.title)}.pptx"
        return Response(
            content=buffer.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": _build_download_content_disposition(filename)
            },
        )

    if artifact.artifact_type == "deck":
        if str(format or "pptx").strip().lower() != "pptx":
            raise HTTPException(status_code=400, detail="Deck artifact only supports pptx.")
        deck_id = str(
            artifact.linked_resource_id or artifact.content.get("deck_id") or ""
        ).strip()
        if not deck_id:
            raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
        try:
            deck = _deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        export_payload = export_deck_payload(
            deck,
            export_deck_to_pptx=export_deck_to_pptx,
            build_export_filename=build_export_filename,
        )
        return Response(
            content=export_payload["content"],
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": _build_download_content_disposition(
                    export_payload["filename"]
                )
            },
        )

    raise HTTPException(status_code=400, detail="Unsupported artifact type.")


@app.post("/api/artifacts/generate")
async def generate_artifact(request: GenerateArtifactRequest):
    from chat_store import SQLiteChatMessageHistory

    history = SQLiteChatMessageHistory(session_id=request.session_id)
    try:
        messages = _resolve_report_messages(
            history,
            answer_group_id=request.answer_group_id,
            panel_id=request.panel_id,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail="Requested artifact scope was not found.") from exc
    if not messages:
        raise HTTPException(status_code=400, detail="No usable messages were found for artifact generation.")

    if request.artifact_type == "report":
        try:
            artifact, _, _ = _create_report_artifact(
                session_id=request.session_id,
                messages=messages,
                answer_group_id=str(request.answer_group_id or "").strip(),
                panel_id=str(request.panel_id or "").strip(),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        return _artifact_payload(artifact)

    if request.artifact_type == "deck":
        if request.panel_config is None:
            raise HTTPException(status_code=400, detail="Deck artifact requires panel_config.")
        try:
            deck = await build_deck(
                messages=messages,
                **build_create_deck_kwargs(
                    request,
                    resolve_active_prompt_runtime=_resolve_active_prompt_runtime,
                    normalize_deck_theme=normalize_deck_theme,
                ),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        _deck_store.save(deck)
        artifact = _create_deck_artifact(deck)
        return _artifact_payload(artifact)

    raise HTTPException(status_code=400, detail="Unsupported artifact type.")


@app.get("/api/knowledge-base/chunks")
async def list_knowledge_base_chunks(
    request: Request,
    path: Optional[str] = None,
    query: str = "",
    source: str = "",
    offset: int = 0,
    limit: int = 20,
):
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    store_path = _effective_vector_store_path(path)
    abs_store = str(_resolve_project_subdir(store_path))
    payload = list_kb_chunks_payload(
        store_path=store_path,
        abs_store=abs_store,
        index_exists=(Path(abs_store) / "index.faiss").exists(),
        query=query,
        source=source,
        offset=offset,
        limit=limit,
        pipeline_factory=lambda vector_store_path: DocPipeline(
            vector_store_path=vector_store_path
        ),
        collect_chunks=_kb_collect_chunks,
        filter_chunks=filter_kb_chunks,
    )
    _audit_security_event(
        "list_knowledge_base_chunks",
        request,
        details=(
            f"path={store_path} query={query or '<empty>'} source={source or '<empty>'} "
            f"offset={offset} limit={limit}"
        ),
    )
    return payload


@app.patch("/api/knowledge-base/chunks/{chunk_id}")
async def update_knowledge_base_chunk(
    chunk_id: str,
    request: Request,
    payload: UpdateKBChunkRequest,
    path: Optional[str] = None,
):
    from doc_pipeline import DocPipeline
    from langchain_core.documents import Document

    _require_remote_admin(request)
    store_path = _effective_vector_store_path(path)
    result = update_kb_chunk_payload(
        chunk_id=chunk_id,
        request=payload,
        field_set=_request_field_set(payload),
        pipeline_factory=lambda: DocPipeline(vector_store_path=store_path),
        docstore_dict=_kb_docstore_dict,
        safe_metadata=_kb_safe_metadata,
        rebuild_from_documents=_kb_rebuild_from_documents,
        doc_factory=lambda page_content, metadata: Document(
            page_content=page_content,
            metadata=metadata,
        ),
        current_time=time.time,
    )
    _audit_security_event(
        "update_knowledge_base_chunk",
        request,
        details=f"chunk_id={chunk_id} path={store_path}",
    )
    return result


@app.delete("/api/knowledge-base/chunks/{chunk_id}")
async def delete_knowledge_base_chunk(
    chunk_id: str,
    request: Request,
    path: Optional[str] = None,
):
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    store_path = _effective_vector_store_path(path)
    result = delete_kb_chunk_payload(
        chunk_id=chunk_id,
        pipeline_factory=lambda: DocPipeline(vector_store_path=store_path),
        docstore_dict=_kb_docstore_dict,
        rebuild_from_documents=_kb_rebuild_from_documents,
    )
    _audit_security_event(
        "delete_knowledge_base_chunk",
        request,
        details=f"chunk_id={chunk_id} path={store_path}",
    )
    return result


@app.get("/api/knowledge-bases")
async def list_knowledge_bases(request: Request):
    """列出可用的知识库目录"""
    import glob as glob_module
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    base_dir = os.path.dirname(os.path.abspath(__file__))
    try:
        current_effective_path = _effective_vector_store_path()
    except HTTPException:
        current_effective_path = None
    payload = knowledge_bases_payload(
        base_dir=base_dir,
        active_vector_store_id=_active_vector_store_id(),
        current_effective_path=current_effective_path,
        env_vector_store_path=os.getenv("VECTOR_STORE_PATH", "./vector_store"),
        sibling_paths=list(glob_module.glob(os.path.join(base_dir, "vector_store*"))),
        resolve_project_subdir=_resolve_project_subdir,
        faiss_safe_store_path=_faiss_safe_store_path,
        pipeline_factory=lambda vector_store_path: DocPipeline(
            vector_store_path=vector_store_path
        ),
    )
    _audit_security_event(
        "list_knowledge_bases",
        request,
        details=f"knowledge_base_count={len(payload.get('knowledge_bases', []))}",
    )
    return payload


@app.get("/api/knowledge-base/health")
async def get_kb_health(request: Request, path: Optional[str] = None):
    """知识库健康检查 - 返回详细状态信息"""
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    target_path = _resolve_project_subdir(_effective_vector_store_path(path))
    payload = kb_health_payload(
        target_path,
        embedding_model=os.getenv("EMBEDDING_MODEL", "BAAI/bge-large-zh-v1.5"),
        faiss_safe_store_path=_faiss_safe_store_path,
        pipeline_factory=lambda vector_store_path: DocPipeline(
            vector_store_path=vector_store_path
        ),
        logger=logger,
    )
    _audit_security_event(
        "get_knowledge_base_health",
        request,
        details=f"path={target_path}",
    )
    return payload


@app.post("/api/knowledge-base/test-retrieval")
async def test_retrieval(request: Request, payload: TestRetrievalRequest):
    """执行测试检索，用于诊断知识库"""
    from doc_pipeline import DocPipeline

    _require_remote_admin(request)
    vector_store_path = _effective_vector_store_path(payload.vector_store_path)
    pipeline = DocPipeline(vector_store_path=vector_store_path)
    try:
        result = retrieval_test_payload(
            payload.query,
            pipeline,
            current_time=time.time,
            search_k=payload.search_k or 5,
            fetch_k=payload.fetch_k or 10,
            use_rerank=payload.use_rerank or False,
            retrieval_mode=payload.retrieval_mode or "semantic",
        )
        _audit_security_event(
            "test_knowledge_base_retrieval",
            request,
            details=(
                f"path={vector_store_path} query_hash={_content_hash(payload.query)} "
                f"results_count={result.get('results_count', 0)}"
            ),
        )
        return result
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception as e:
        return {"results_count": 0, "top_scores": [], "latency_ms": 0, "error": str(e)}


@app.delete("/api/knowledge-base")
async def delete_knowledge_base(request: Request, path: Optional[str] = None):
    """删除当前默认知识库"""
    _require_remote_admin(request)
    target_path = _resolve_deletable_knowledge_base(_effective_vector_store_path(path))
    try:
        result = await delete_kb_directory(
            target_path,
            remove_tree=shutil.rmtree,
            clear_agent_cache=_clear_agent_cache,
            success_message="知识库已删除：{path}",
            on_success=lambda abs_path: logger.info("Knowledge base deleted: %s", abs_path),
            on_failure=lambda: logger.exception("Failed to delete knowledge base"),
        )
        _audit_security_event(
            "delete_knowledge_base",
            request,
            details=f"path={target_path}",
        )
        return result
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.delete("/api/knowledge-base/by-path")
async def delete_knowledge_base_by_path(request: Request, path: str):
    """按路径删除指定知识库"""

    _require_remote_admin(request)
    target_path = _resolve_deletable_knowledge_base(path)
    try:
        result = await delete_kb_directory(
            target_path,
            remove_tree=shutil.rmtree,
            clear_agent_cache=_clear_agent_cache,
            success_message="已删除：{path}",
        )
        _audit_security_event(
            "delete_knowledge_base_by_path",
            request,
            details=f"path={target_path}",
        )
        return result
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


# ── 静态文件（生产模式）─────────────────────────

@app.get("/api/share-links", response_model=ShareLinkAuditListResponse)
async def list_share_links(
    request: Request,
    resource_type: str = "",
    active_only: bool = False,
    limit: int = 100,
    offset: int = 0,
):
    _require_remote_admin(request)
    records = _share_link_store.list_links(
        resource_type=resource_type,
        active_only=active_only,
        limit=limit,
        offset=offset,
    )
    payload_records = [_share_link_audit_payload(record) for record in records]
    payload = {
        "share_links": payload_records,
        "total": len(payload_records),
        "active_count": sum(1 for item in payload_records if item["is_active"]),
    }
    _audit_security_event(
        "list_share_links",
        request,
        details=(
            f"resource_type={resource_type or '<all>'} "
            f"active_only={active_only} total={payload['total']}"
        ),
    )
    return payload


@app.delete("/api/share-links/{share_token}", response_model=RevokeShareLinkResponse)
async def revoke_share_link(share_token: str, request: Request):
    _require_remote_admin(request)
    if not _share_link_store.revoke(share_token):
        raise HTTPException(status_code=404, detail="未找到分享链接")
    _audit_security_event(
        "revoke_share_link",
        request,
        details=f"share_token_fp={_token_fingerprint(share_token)}",
    )
    return RevokeShareLinkResponse(ok=True)


@app.get("/shared/{share_token}")
async def open_shared_resource(share_token: str, request: Request):
    _require_remote_share_secret(request)
    share_secret = _current_share_link_secret()
    try:
        link_record = _share_link_store.get_active(share_token)
        if link_record is None:
            raise ValueError("分享链接不存在、已过期或已撤销")
        decoded_type, decoded_id = _decode_share_token(share_token, share_secret)
        if link_record.resource_type != decoded_type or link_record.resource_id != decoded_id:
            raise ValueError("分享链接无效")
        shared_payload = open_shared_resource_payload(
            share_token,
            request,
            secret=share_secret,
            decode_share_token=_decode_share_token,
            build_share_url=_build_share_url,
            build_session_messages_payload=_build_session_messages_payload,
            render_shared_session_html=_render_shared_session_html,
            get_deck=_deck_store.get,
            render_shared_deck_html=_render_shared_deck_html,
        )
        _share_link_store.record_access(
            share_token,
            accessed_ip=_request_client_ip(request),
            accessed_user_agent=_request_user_agent(request),
        )
        _audit_security_event(
            "open_shared_resource",
            request,
            details=(
                f"resource_type={decoded_type} "
                f"share_token_fp={_token_fingerprint(share_token)}"
            ),
        )
    except ValueError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except KeyError as exc:
        detail = str(exc.args[0]) if exc.args else "Not found"
        raise HTTPException(status_code=404, detail=detail) from exc

    return Response(
        content=shared_payload["content"],
        media_type=shared_payload["media_type"],
    )


_frontend_dist = os.path.join(os.path.dirname(__file__), "frontend", "dist")
if os.path.isdir(_frontend_dist):
    from fastapi.responses import FileResponse

    app.mount(
        "/assets",
        StaticFiles(directory=os.path.join(_frontend_dist, "assets")),
        name="assets",
    )

    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str):
        index = os.path.join(_frontend_dist, "index.html")
        return FileResponse(index)


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(
        app,
        host=os.getenv("SERVER_HOST", "127.0.0.1"),
        port=int(os.getenv("SERVER_PORT", "8000")),
        reload=True,
    )
