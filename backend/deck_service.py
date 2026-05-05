"""
Deck generation, persistence, and export helpers.
"""

from __future__ import annotations

import io
import json
import os
import re
import time
import uuid
from dataclasses import dataclass
from types import SimpleNamespace
from typing import Any, Literal

from langchain_core.documents import Document
from pydantic import BaseModel, Field

from backend.agent_core import get_llm
from backend.chat_store import connect_sqlite
from backend.core.storage_runtime import app_database_path
from backend.doc_pipeline import DocPipeline


DeckSourceMode = Literal["kb_plus_chat", "chat_only"]
DeckQualityState = Literal["supported", "weak_support", "manual"]
DeckThemeName = Literal["default", "midnight", "sunrise"]
DeckChartType = Literal["bar", "line", "pie"]

_DASHBOARD_CARD_BLOCK_RE = re.compile(
    r":::dashboard-card\s*\n([\s\S]*?)\n:::",
    flags=re.IGNORECASE,
)

DECK_THEME_PALETTES: dict[str, dict[str, str]] = {
    "default": {
        "bg": "F6F8FC",
        "surface": "FFFFFF",
        "surface_alt": "EEF4FF",
        "border": "D8E0EE",
        "title": "162033",
        "body": "2A3547",
        "muted": "60708A",
        "accent": "2563EB",
        "accent_soft": "DBEAFE",
        "success": "1F9D68",
        "warning": "D97706",
        "danger": "C2410C",
    },
    "midnight": {
        "bg": "0F172A",
        "surface": "111827",
        "surface_alt": "1E293B",
        "border": "334155",
        "title": "F8FAFC",
        "body": "E2E8F0",
        "muted": "94A3B8",
        "accent": "38BDF8",
        "accent_soft": "082F49",
        "success": "34D399",
        "warning": "FBBF24",
        "danger": "F87171",
    },
    "sunrise": {
        "bg": "FFF7ED",
        "surface": "FFFBF5",
        "surface_alt": "FDE7D6",
        "border": "F4C7A1",
        "title": "7C2D12",
        "body": "9A3412",
        "muted": "C2410C",
        "accent": "EA580C",
        "accent_soft": "FED7AA",
        "success": "2F855A",
        "warning": "D97706",
        "danger": "C2410C",
    },
}


class DeckWarning(BaseModel):
    code: str
    message: str


class DeckSlideEvidenceCoverage(BaseModel):
    slide_id: str
    slide_type: str
    evidence_ref_count: int = 0
    has_evidence: bool = False
    is_coverable: bool = True
    quality_state: DeckQualityState = "weak_support"


class DeckEvidenceCoverage(BaseModel):
    total_slides: int = 0
    coverable_slide_count: int = 0
    slides_with_evidence: int = 0
    total_evidence_refs: int = 0
    coverage_ratio: float = 0.0
    unsupported_slide_ids: list[str] = Field(default_factory=list)
    slides: list[DeckSlideEvidenceCoverage] = Field(default_factory=list)


class DeckCitationValidationIssue(BaseModel):
    code: str
    message: str
    slide_id: str = ""
    block_id: str = ""
    evidence_ref_id: str = ""
    source_id: str = ""


class DeckCitationValidation(BaseModel):
    status: Literal["passed", "failed"] = "passed"
    can_export: bool = True
    issue_count: int = 0
    missing_source_ids: list[str] = Field(default_factory=list)
    missing_block_evidence_ref_ids: list[str] = Field(default_factory=list)
    issues: list[DeckCitationValidationIssue] = Field(default_factory=list)


class DeckMeta(BaseModel):
    title: str
    subtitle: str = ""
    language: str = "zh-CN"
    audience: str = "general"
    purpose: str = "briefing"
    author: str = "system"
    theme: DeckThemeName = "default"
    created_at: str
    session_id: str
    source_mode: DeckSourceMode
    generator_panel_id: str
    source_answer_group_id: str = ""
    source_panel_id: str = ""


class DeckGeneration(BaseModel):
    source: DeckSourceMode
    target_slide_count: int
    actual_slide_count: int = 0
    warnings: list[DeckWarning] = Field(default_factory=list)
    evidence_coverage: DeckEvidenceCoverage = Field(default_factory=DeckEvidenceCoverage)
    evidence_review: dict[str, Any] = Field(default_factory=dict)
    citation_validation: DeckCitationValidation | None = None


class DeckBlock(BaseModel):
    id: str
    kind: str
    role: str
    content: dict[str, Any] = Field(default_factory=dict)
    editable: bool = True


class DeckChartNormalizationReport(BaseModel):
    normalized_block_count: int = 0
    invalid_block_count: int = 0
    normalized_block_ids: list[str] = Field(default_factory=list)
    invalid_block_ids: list[str] = Field(default_factory=list)


class DeckEvidenceRef(BaseModel):
    id: str
    source_id: str
    source_title: str
    excerpt_id: str | None = None
    snippet: str = ""
    confidence: float = 0.0


class DeckSourceItem(BaseModel):
    id: str
    type: str
    title: str
    document_id: str | None = None
    uri: str | None = None
    metadata: dict[str, Any] = Field(default_factory=dict)


class DeckSlideStatus(BaseModel):
    locked: bool = False
    dirty: bool = False
    review_state: str = "draft"


class DeckSlide(BaseModel):
    id: str
    type: str
    title: str
    subtitle: str = ""
    layout: str
    intent: str = ""
    speaker_notes: str = ""
    blocks: list[DeckBlock] = Field(default_factory=list)
    evidence_refs: list[DeckEvidenceRef] = Field(default_factory=list)
    quality_state: DeckQualityState = "weak_support"
    status: DeckSlideStatus = Field(default_factory=DeckSlideStatus)


class DeckSpec(BaseModel):
    version: str = "1.0"
    deck_id: str
    status: str = "draft"
    meta: DeckMeta
    generation: DeckGeneration
    slides: list[DeckSlide]
    source_registry: list[DeckSourceItem] = Field(default_factory=list)
    citation_validation: DeckCitationValidation | None = None

    def refresh_evidence_coverage(self) -> "DeckSpec":
        self.generation.evidence_coverage = build_deck_evidence_coverage(self)
        return self


_EVIDENCE_COVERAGE_EXCLUDED_SLIDE_TYPES = {"cover", "outline", "appendix_sources"}


def _is_evidence_coverable_slide(slide: DeckSlide, evidence_ref_count: int) -> bool:
    if evidence_ref_count > 0:
        return True
    return slide.type not in _EVIDENCE_COVERAGE_EXCLUDED_SLIDE_TYPES


def build_deck_evidence_coverage(deck: DeckSpec) -> DeckEvidenceCoverage:
    slide_stats: list[DeckSlideEvidenceCoverage] = []
    coverable_slide_count = 0
    slides_with_evidence = 0
    total_evidence_refs = 0
    unsupported_slide_ids: list[str] = []

    for slide in deck.slides:
        evidence_ref_count = len(slide.evidence_refs or [])
        has_evidence = evidence_ref_count > 0
        is_coverable = _is_evidence_coverable_slide(slide, evidence_ref_count)

        total_evidence_refs += evidence_ref_count
        if is_coverable:
            coverable_slide_count += 1
            if has_evidence:
                slides_with_evidence += 1
            else:
                unsupported_slide_ids.append(slide.id)

        slide_stats.append(
            DeckSlideEvidenceCoverage(
                slide_id=slide.id,
                slide_type=slide.type,
                evidence_ref_count=evidence_ref_count,
                has_evidence=has_evidence,
                is_coverable=is_coverable,
                quality_state=slide.quality_state,
            )
        )

    coverage_ratio = (
        round(slides_with_evidence / coverable_slide_count, 4)
        if coverable_slide_count
        else 0.0
    )
    return DeckEvidenceCoverage(
        total_slides=len(deck.slides),
        coverable_slide_count=coverable_slide_count,
        slides_with_evidence=slides_with_evidence,
        total_evidence_refs=total_evidence_refs,
        coverage_ratio=coverage_ratio,
        unsupported_slide_ids=unsupported_slide_ids,
        slides=slide_stats,
    )


def refresh_deck_evidence_coverage(deck: DeckSpec) -> DeckSpec:
    return deck.refresh_evidence_coverage()


def _coerce_evidence_id_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, (str, int, float)):
        items = [value]
    elif isinstance(value, list | tuple | set):
        items = list(value)
    else:
        return []
    return [str(item).strip() for item in items if str(item or "").strip()]


def _block_bound_evidence_ref_ids(block: DeckBlock) -> list[str]:
    content = block.content if isinstance(block.content, dict) else {}
    ids = _coerce_evidence_id_list(content.get("evidence_ref_ids"))
    if ids:
        return ids
    return _coerce_evidence_id_list(content.get("evidence_refs"))


def validate_deck_citation_consistency(deck: DeckSpec) -> DeckCitationValidation:
    """Validate citation references before review/export payloads are consumed."""

    source_ids = {
        source.id.strip()
        for source in deck.source_registry
        if isinstance(source.id, str) and source.id.strip()
    }
    issues: list[DeckCitationValidationIssue] = []
    missing_source_ids: list[str] = []
    missing_block_evidence_ref_ids: list[str] = []

    for slide in deck.slides:
        slide_evidence_ids = {
            ref.id.strip()
            for ref in slide.evidence_refs
            if isinstance(ref.id, str) and ref.id.strip()
        }

        for ref in slide.evidence_refs:
            source_id = ref.source_id.strip() if isinstance(ref.source_id, str) else ""
            if source_id and source_id not in source_ids:
                missing_source_ids.append(source_id)
                issues.append(
                    DeckCitationValidationIssue(
                        code="missing_source_registry_entry",
                        message="Evidence reference points to a source_id missing from source_registry.",
                        slide_id=slide.id,
                        evidence_ref_id=ref.id,
                        source_id=source_id,
                    )
                )

        for block in slide.blocks:
            for evidence_ref_id in _block_bound_evidence_ref_ids(block):
                if evidence_ref_id in slide_evidence_ids:
                    continue
                missing_block_evidence_ref_ids.append(evidence_ref_id)
                issues.append(
                    DeckCitationValidationIssue(
                        code="missing_slide_evidence_ref",
                        message="Block evidence_ref_id is not present in the parent slide evidence_refs.",
                        slide_id=slide.id,
                        block_id=block.id,
                        evidence_ref_id=evidence_ref_id,
                    )
                )

    issue_count = len(issues)
    return DeckCitationValidation(
        status="failed" if issue_count else "passed",
        can_export=issue_count == 0,
        issue_count=issue_count,
        missing_source_ids=list(dict.fromkeys(missing_source_ids)),
        missing_block_evidence_ref_ids=list(
            dict.fromkeys(missing_block_evidence_ref_ids)
        ),
        issues=issues,
    )


class OutlineSlidePlan(BaseModel):
    title: str
    objective: str
    section: str
    evidence_source_ids: list[str] = Field(default_factory=list)


class OutlinePlan(BaseModel):
    title: str
    subtitle: str
    core_message: str
    sections: list[str]
    content_slides: list[OutlineSlidePlan]


class DraftedContentSlide(BaseModel):
    title: str
    subtitle: str = ""
    key_points: list[str] = Field(default_factory=list)
    speaker_notes: str = ""
    evidence_excerpt_ids: list[str] = Field(default_factory=list)
    evidence_source_ids: list[str] = Field(default_factory=list)
    quality_state: DeckQualityState = "weak_support"


class DraftedSlideBundle(BaseModel):
    content_slides: list[DraftedContentSlide]


class SourceExcerpt(BaseModel):
    id: str
    source_id: str
    source_title: str
    snippet: str
    confidence: float = 0.85


@dataclass
class SourcePack:
    title_hint: str
    source_mode: DeckSourceMode
    qa_pairs: list[tuple[str, str]]
    chat_notes: list[str]
    excerpts: list[SourceExcerpt]
    source_registry: list[DeckSourceItem]
    warnings: list[DeckWarning]


_FALLBACK_SECTION_TITLES = ["Topic Overview", "Key Findings", "Recommendations"]
_FALLBACK_SLIDE_TITLES = [
    "Topic overview and core conclusion",
    "鍏抽敭淇℃伅鎷嗚В",
    "Use cases and action recommendations",
    "Follow-up priorities and risk notes",
    "琛ュ厖瑙傚療",
]


def _now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%S%z")


def _clean_text(text: Any) -> str:
    return " ".join(str(text).strip().split())


def _metadata_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def normalize_deck_theme(theme: Any) -> DeckThemeName:
    raw = _clean_text(theme).lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "": "default",
        "default": "default",
        "classic": "default",
        "midnight": "midnight",
        "night": "midnight",
        "dark": "midnight",
        "sunrise": "sunrise",
        "warm": "sunrise",
    }
    return aliases.get(raw, "default")


def _truncate(text: str, limit: int) -> str:
    cleaned = _clean_text(text)
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1].rstrip() + "..."


def _truncate_multiline(text: Any, limit: int) -> str:
    normalized = str(text).replace("\r\n", "\n").replace("\r", "\n").strip()
    if len(normalized) <= limit:
        return normalized
    return normalized[: limit - 1].rstrip() + "..."


def _stringify_llm_content(content: Any) -> str:
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, dict):
                text = item.get("text")
                if text:
                    parts.append(str(text))
            else:
                parts.append(str(item))
        return "\n".join(parts).strip()
    return str(content).strip()


def _normalize_message_text(content: Any) -> str:
    return (
        _stringify_llm_content(content)
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )


def _strip_dashboard_card_blocks(text: Any) -> str:
    normalized = _normalize_message_text(text)
    if not normalized:
        return ""
    stripped = _DASHBOARD_CARD_BLOCK_RE.sub("", normalized)
    stripped = re.sub(r"\n{3,}", "\n\n", stripped)
    return stripped.strip()


def _extract_dashboard_card_payloads(text: Any) -> list[dict[str, Any]]:
    payloads: list[dict[str, Any]] = []
    normalized = _normalize_message_text(text)
    if not normalized:
        return payloads

    for match in _DASHBOARD_CARD_BLOCK_RE.finditer(normalized):
        raw_payload = match.group(1).strip()
        if not raw_payload:
            continue
        try:
            payload = json.loads(raw_payload)
        except json.JSONDecodeError:
            continue
        if isinstance(payload, dict):
            payloads.append(payload)
    return payloads


def _coerce_chart_number(value: Any) -> float | None:
    if isinstance(value, bool):
        return float(value)
    if isinstance(value, (int, float)):
        return float(value)
    cleaned = _clean_text(value).replace(",", "")
    if not cleaned:
        return None
    try:
        return float(cleaned)
    except ValueError:
        return None


def _normalize_chart_labels(values: Any) -> list[str]:
    if not isinstance(values, list):
        return []
    labels: list[str] = []
    for index, value in enumerate(values, start=1):
        cleaned = _clean_text(value)
        labels.append(cleaned or f"绫诲埆 {index}")
    return labels


def _normalize_chart_datasets(values: Any, label_count: int) -> list[dict[str, Any]]:
    if not isinstance(values, list):
        return []

    datasets: list[dict[str, Any]] = []
    for dataset_index, value in enumerate(values, start=1):
        if not isinstance(value, dict):
            continue

        raw_points = value.get("data")
        if not isinstance(raw_points, list):
            continue

        clean_points = [_coerce_chart_number(item) for item in raw_points]
        if not any(item is not None for item in clean_points):
            continue

        target_size = label_count or len(clean_points)
        normalized_points = [
            _coerce_chart_number(item) or 0.0 for item in raw_points[:target_size]
        ]
        if len(normalized_points) < target_size:
            normalized_points.extend([0.0] * (target_size - len(normalized_points)))

        label = _clean_text(value.get("label")) or f"绯诲垪 {dataset_index}"
        datasets.append(
            {
                "label": label,
                "data": normalized_points,
            }
        )
    return datasets


def _chart_summary_from_answer(answer: Any) -> str:
    summaries: list[str] = []
    for payload in _extract_dashboard_card_payloads(answer):
        charts = payload.get("charts")
        if not isinstance(charts, list):
            continue
        for chart in charts:
            if not isinstance(chart, dict):
                continue
            title = _clean_text(chart.get("title"))
            chart_type = _clean_text(chart.get("type")).lower()
            if chart_type in {"bar", "line", "pie"}:
                summaries.append(title or f"{chart_type} 鍥捐〃")
    if not summaries:
        return ""
    return "鍖呭惈鍥捐〃: " + " / ".join(summaries[:3])


def _answer_plaintext(answer: Any) -> str:
    stripped = _strip_dashboard_card_blocks(answer)
    if stripped:
        return stripped
    return _chart_summary_from_answer(answer)


def _build_chart_block(
    chart: dict[str, Any],
    *,
    fallback_title: str,
    block_index: int,
) -> DeckBlock | None:
    chart_data = chart.get("chart_data")
    if not isinstance(chart_data, dict):
        return None

    raw_type = _clean_text(chart_data.get("type") or chart.get("type")).lower()
    if raw_type not in {"bar", "line", "pie"}:
        return None

    labels = _normalize_chart_labels(chart_data.get("labels"))
    datasets = _normalize_chart_datasets(chart_data.get("datasets"), len(labels))
    if not datasets:
        return None

    if not labels:
        labels = [f"绫诲埆 {index + 1}" for index in range(len(datasets[0]["data"]))]
        datasets = _normalize_chart_datasets(chart_data.get("datasets"), len(labels))
        if not datasets:
            return None

    return DeckBlock(
        id=f"block_chart_{block_index}_{uuid.uuid4().hex[:6]}",
        kind="chart",
        role="dashboard_chart",
        content={
            "title": _truncate(_clean_text(chart.get("title")) or fallback_title, 48),
            "description": _truncate(_clean_text(chart.get("description")), 120),
            "chart_type": raw_type,
            "labels": labels[:12],
            "datasets": [
                {
                    "label": _truncate(item["label"], 32),
                    "data": item["data"][:12],
                }
                for item in datasets[:4]
            ],
        },
        editable=False,
    )


def normalize_deck_chart_block(block: DeckBlock) -> tuple[DeckBlock, str]:
    """Normalize chart block payloads into the export/frontend contract."""

    if block.kind != "chart":
        return block, "skipped"

    content = block.content if isinstance(block.content, dict) else {}
    raw_type = _clean_text(content.get("chart_type") or content.get("type")).lower()
    if raw_type not in {"bar", "line", "pie"}:
        return (
            block.model_copy(
                update={
                    "content": {
                        **content,
                        "normalization_status": "invalid",
                        "normalization_issues": ["unsupported_chart_type"],
                    }
                },
                deep=True,
            ),
            "invalid",
        )

    labels = _normalize_chart_labels(content.get("labels"))
    datasets = _normalize_chart_datasets(content.get("datasets"), len(labels))
    if not datasets:
        return (
            block.model_copy(
                update={
                    "content": {
                        **content,
                        "chart_type": raw_type,
                        "normalization_status": "invalid",
                        "normalization_issues": ["missing_or_invalid_datasets"],
                    }
                },
                deep=True,
            ),
            "invalid",
        )

    if not labels:
        labels = [f"类别 {index + 1}" for index in range(len(datasets[0]["data"]))]
        datasets = _normalize_chart_datasets(content.get("datasets"), len(labels))
        if not datasets:
            return (
                block.model_copy(
                    update={
                        "content": {
                            **content,
                            "chart_type": raw_type,
                            "normalization_status": "invalid",
                            "normalization_issues": ["missing_or_invalid_datasets"],
                        }
                    },
                    deep=True,
                ),
                "invalid",
            )

    normalized_content = {
        **content,
        "title": _truncate(_clean_text(content.get("title")) or "数据图表", 48),
        "description": _truncate(_clean_text(content.get("description")), 120),
        "chart_type": raw_type,
        "labels": labels[:12],
        "datasets": [
            {
                "label": _truncate(item["label"], 32),
                "data": item["data"][:12],
            }
            for item in datasets[:4]
        ],
        "normalization_status": "normalized",
        "normalization_version": "chart-block-v1",
    }
    normalized_content.pop("normalization_issues", None)
    return block.model_copy(update={"content": normalized_content}, deep=True), "normalized"


def normalize_deck_chart_blocks(deck: DeckSpec) -> DeckChartNormalizationReport:
    report = DeckChartNormalizationReport()
    for slide in deck.slides:
        next_blocks: list[DeckBlock] = []
        for block in slide.blocks:
            normalized_block, status = normalize_deck_chart_block(block)
            next_blocks.append(normalized_block)
            if status == "normalized":
                report.normalized_block_count += 1
                report.normalized_block_ids.append(block.id)
            elif status == "invalid":
                report.invalid_block_count += 1
                report.invalid_block_ids.append(block.id)
        slide.blocks = next_blocks
    return report


def _extract_dashboard_chart_blocks(answer: Any, limit: int = 1) -> list[DeckBlock]:
    blocks: list[DeckBlock] = []
    for payload in _extract_dashboard_card_payloads(answer):
        charts = payload.get("charts")
        if not isinstance(charts, list):
            continue
        fallback_title = _clean_text(payload.get("title")) or "鏁版嵁鍥捐〃"
        for chart in charts:
            if not isinstance(chart, dict):
                continue
            block = _build_chart_block(
                chart,
                fallback_title=fallback_title,
                block_index=len(blocks) + 1,
            )
            if block is None:
                continue
            blocks.append(block)
            if len(blocks) >= limit:
                return blocks
    return blocks


def _extract_qa_pairs(messages: list[Any]) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    pending_question = ""
    for message in messages:
        role = getattr(message, "__class__", type(message)).__name__
        content = _normalize_message_text(getattr(message, "content", ""))
        if not content:
            continue
        if role == "HumanMessage":
            pending_question = _clean_text(content)
            continue
        if role == "AIMessage" and pending_question:
            pairs.append((pending_question, content))
            pending_question = ""
    return pairs


def _is_failed_answer(answer: str) -> bool:
    normalized = answer.strip().lower()
    if not normalized:
        return True
    failure_markers = (
        "agent stopped due to max iterations",
        "agent stopped due to iteration limit",
        "鐢熸垚鍥炵瓟澶辫触",
        "鏃犳硶瀹屾垚浠诲姟",
        "request processing error",
        "妯″瀷宸ュ叿璋冪敤娆℃暟瓒呴檺",
        "internal_error",
    )
    return any(marker in normalized for marker in failure_markers)


def extract_successful_qa_pairs(messages: list[Any]) -> list[tuple[str, str]]:
    return [
        (question, answer)
        for question, answer in _extract_qa_pairs(messages)
        if not _is_failed_answer(answer)
    ]


def ensure_deckable_chat(messages: list[Any]) -> list[tuple[str, str]]:
    raw_pairs = _extract_qa_pairs(messages)
    qa_pairs = extract_successful_qa_pairs(messages)
    if qa_pairs:
        return qa_pairs
    if raw_pairs:
        raise ValueError("The latest chat answers are failed results and cannot be converted into a deck.")
    raise ValueError("This session has no successful Q&A content for deck generation.")


def _extract_json_payload(text: str) -> dict[str, Any]:
    candidate = text.strip()
    if candidate.startswith("```"):
        candidate = re.sub(r"^```(?:json)?\s*", "", candidate)
        candidate = re.sub(r"\s*```$", "", candidate)

    decoder = json.JSONDecoder()
    try:
        payload, _ = decoder.raw_decode(candidate)
        if not isinstance(payload, dict):
            raise ValueError("LLM JSON payload must be an object.")
        return payload
    except Exception:
        for start_char in ("{", "["):
            start = candidate.find(start_char)
            if start == -1:
                continue
            try:
                payload, _ = decoder.raw_decode(candidate[start:])
            except json.JSONDecodeError:
                continue
            if isinstance(payload, dict):
                return payload
        raise


def _is_authentication_failure(exc: Exception) -> bool:
    message = str(exc).lower()
    markers = (
        "authenticationerror",
        "401",
        "invalid api key",
        "invalid token",
        "unauthorized",
        "invalid credential",
        "invalid_api_key",
    )
    return any(marker in message for marker in markers)


def _local_ollama_panel_config(panel_config: Any) -> Any:
    fallback_model = os.getenv("OLLAMA_MODEL", "qwen3.5:4b").strip() or "qwen3.5:4b"
    fallback_base_url = (
        os.getenv("OLLAMA_BASE_URL", "http://localhost:11434").strip()
        or "http://localhost:11434"
    )
    return SimpleNamespace(
        panel_id=getattr(panel_config, "panel_id", "panel"),
        provider="ollama",
        connection_type="ollama",
        model=fallback_model,
        base_url=fallback_base_url,
        api_key="",
        temperature=float(getattr(panel_config, "temperature", 0.3) or 0.3),
    )


async def _invoke_json(llm, prompt: str) -> dict[str, Any]:
    response = await llm.ainvoke(prompt)
    text = _stringify_llm_content(response.content)
    try:
        return _extract_json_payload(text)
    except Exception:
        repair_prompt = (
            "璇锋妸涓嬮潰鍐呭淇涓哄悎娉?JSON銆備笉瑕佽ˉ鍏呬俊鎭紝涓嶈杈撳嚭瑙ｉ噴锛屽彧杈撳嚭 JSON銆俓n\n"
            f"{text}"
        )
        repaired = await llm.ainvoke(repair_prompt)
        repaired_text = _stringify_llm_content(repaired.content)
        return _extract_json_payload(repaired_text)


def _build_title_hint(session_id: str, messages: list[Any]) -> str:
    for message in messages:
        content = _clean_text(getattr(message, "content", ""))
        if (
            getattr(message, "__class__", type(message)).__name__ == "HumanMessage"
            and content
        ):
            return _truncate(content, 60)
    return f"Deck {session_id[:8]}"


def _first_non_empty(*values: Any) -> str:
    for value in values:
        cleaned = _clean_text(value)
        if cleaned:
            return cleaned
    return ""


def _default_sections() -> list[str]:
    return list(_FALLBACK_SECTION_TITLES)


def _fallback_slide_title(index: int) -> str:
    if index < len(_FALLBACK_SLIDE_TITLES):
        return _FALLBACK_SLIDE_TITLES[index]
    return f"鏍稿績涓婚 {index + 1}"


def _extract_answer_points(answer: str, limit: int = 4) -> list[str]:
    clean_answer = _answer_plaintext(answer)
    answer = clean_answer
    chunks = [
        _truncate(part, 72)
        for part in re.split(r"[銆傦紒锛??]\s+|\n+", answer)
        if _clean_text(part)
    ]
    deduped: list[str] = []
    seen: set[str] = set()
    for chunk in chunks:
        if chunk in seen:
            continue
        seen.add(chunk)
        deduped.append(chunk)
        if len(deduped) >= limit:
            break
    return deduped


def _fallback_outline(
    pack: SourcePack,
    content_slide_count: int,
) -> OutlinePlan:
    title = _truncate(pack.title_hint, 56)
    subtitle = (
        "鍩轰簬鐭ヨ瘑搴撴绱笌鎴愬姛瀵硅瘽鏁寸悊"
        if pack.source_mode == "kb_plus_chat"
        else "鍩轰簬鎴愬姛瀵硅瘽鏁寸悊"
    )
    core_message = _truncate(
        _first_non_empty(
            _answer_plaintext(pack.qa_pairs[-1][1]) if pack.qa_pairs else "",
            subtitle,
        ),
        120,
    )
    sections = _default_sections()
    content_slides: list[OutlineSlidePlan] = []

    for index in range(content_slide_count):
        question, answer = pack.qa_pairs[min(index, len(pack.qa_pairs) - 1)]
        plain_answer = _answer_plaintext(answer)
        evidence_source_ids: list[str] = []
        if pack.source_mode == "kb_plus_chat" and pack.excerpts:
            evidence_source_ids = [
                pack.excerpts[min(index, len(pack.excerpts) - 1)].source_id
            ]

        content_slides.append(
            OutlineSlidePlan(
                title=_fallback_slide_title(index),
                objective=_truncate(_first_non_empty(plain_answer, question), 100),
                section=sections[min(index, len(sections) - 1)],
                evidence_source_ids=evidence_source_ids,
            )
        )

    return OutlinePlan(
        title=title,
        subtitle=subtitle,
        core_message=core_message or "Core conclusion based on the current conversation.",
        sections=sections,
        content_slides=content_slides,
    )


def _fallback_content_bundle(
    pack: SourcePack,
    outline: OutlinePlan,
) -> DraftedSlideBundle:
    content_slides: list[DraftedContentSlide] = []

    for index, plan in enumerate(outline.content_slides):
        question, answer = pack.qa_pairs[min(index, len(pack.qa_pairs) - 1)]
        plain_answer = _answer_plaintext(answer)
        points = _extract_answer_points(answer)
        if not points:
            points = [
                _truncate(_first_non_empty(plain_answer, question, plan.objective), 72)
            ]

        evidence_excerpt_ids: list[str] = []
        evidence_source_ids: list[str] = []
        quality_state: DeckQualityState = "manual"

        if pack.excerpts:
            excerpt = pack.excerpts[min(index, len(pack.excerpts) - 1)]
            evidence_excerpt_ids = [excerpt.id]
            evidence_source_ids = [excerpt.source_id]
            quality_state = "supported"

        content_slides.append(
            DraftedContentSlide(
                title=plan.title or _fallback_slide_title(index),
                subtitle=_truncate(plan.objective, 72),
                key_points=points[:5],
                speaker_notes=_truncate(
                    _first_non_empty(plain_answer, question, plan.objective),
                    180,
                ),
                evidence_excerpt_ids=evidence_excerpt_ids,
                evidence_source_ids=evidence_source_ids,
                quality_state=quality_state,
            )
        )

    return DraftedSlideBundle(content_slides=content_slides)


def _dedupe_docs(docs: list[Document]) -> list[Document]:
    deduped: list[Document] = []
    seen: set[str] = set()
    for doc in docs:
        snippet = _truncate(doc.page_content, 180)
        key = f"{doc.metadata.get('source', 'unknown')}::{snippet}"
        if key in seen:
            continue
        seen.add(key)
        deduped.append(doc)
    return deduped


def _build_chat_only_source_pack(
    title_hint: str,
    qa_pairs: list[tuple[str, str]],
    chat_notes: list[str],
    warnings: list[DeckWarning],
    chat_only_message: str,
    excerpts: list[SourceExcerpt] | None = None,
    source_registry: list[DeckSourceItem] | None = None,
) -> SourcePack:
    return SourcePack(
        title_hint=title_hint,
        source_mode="chat_only",
        qa_pairs=qa_pairs,
        chat_notes=chat_notes,
        excerpts=excerpts or [],
        source_registry=source_registry or [],
        warnings=[
            *warnings,
            DeckWarning(
                code="chat_only_mode",
                message=chat_only_message,
            ),
        ],
    )


def _message_metadata(message: Any) -> dict[str, Any]:
    metadata: dict[str, Any] = {}
    for attr in ("additional_kwargs", "response_metadata"):
        value = getattr(message, attr, None)
        if isinstance(value, dict):
            metadata.update(value)
    return metadata


def _source_identity(source: dict[str, Any], fallback: str) -> str:
    return _first_non_empty(
        source.get("url"),
        source.get("uri"),
        source.get("source"),
        source.get("title"),
        fallback,
    )


def _source_title(source: dict[str, Any], fallback: str) -> str:
    return _truncate(
        _first_non_empty(
            source.get("title"),
            source.get("name"),
            source.get("domain"),
            source.get("url"),
            source.get("uri"),
            fallback,
        ),
        96,
    )


def _upsert_report_source(
    source: dict[str, Any],
    *,
    source_id_map: dict[str, str],
    source_registry: list[DeckSourceItem],
    fallback_title: str,
) -> tuple[str, str]:
    identity = _source_identity(source, fallback_title)
    source_id = source_id_map.get(identity)
    if source_id:
        existing = next(item for item in source_registry if item.id == source_id)
        return source_id, existing.title

    source_id = f"src_report_{len(source_id_map) + 1}"
    source_id_map[identity] = source_id
    title = _source_title(source, fallback_title)
    uri = _first_non_empty(source.get("url"), source.get("uri"), source.get("source"))
    source_registry.append(
        DeckSourceItem(
            id=source_id,
            type=_first_non_empty(source.get("type"), "web"),
            title=title,
            document_id=_first_non_empty(source.get("document_id"), source.get("doc_id"))
            or None,
            uri=uri or None,
            metadata={
                key: value
                for key, value in source.items()
                if key not in {"title", "name", "url", "uri", "source", "snippet"}
            },
        )
    )
    return source_id, title


def _extract_report_evidence_from_messages(
    messages: list[Any],
) -> tuple[list[SourceExcerpt], list[DeckSourceItem]]:
    source_registry: list[DeckSourceItem] = []
    excerpts: list[SourceExcerpt] = []
    source_id_map: dict[str, str] = {}

    for message in messages:
        if getattr(message, "__class__", type(message)).__name__ != "AIMessage":
            continue

        metadata = _message_metadata(message)
        raw_sources = metadata.get("sources")
        if isinstance(raw_sources, list):
            for source_index, raw_source in enumerate(raw_sources, start=1):
                source = _metadata_dict(raw_source)
                if not source:
                    continue
                source_id, source_title = _upsert_report_source(
                    source,
                    source_id_map=source_id_map,
                    source_registry=source_registry,
                    fallback_title=f"Report Source {source_index}",
                )
                snippet = _truncate(
                    _first_non_empty(
                        source.get("snippet"),
                        source.get("excerpt"),
                        source.get("summary"),
                        source.get("selection_reason"),
                    ),
                    260,
                )
                if snippet:
                    excerpts.append(
                        SourceExcerpt(
                            id=f"ext_report_{len(excerpts) + 1}",
                            source_id=source_id,
                            source_title=source_title,
                            snippet=snippet,
                            confidence=0.82,
                        )
                    )

        raw_chains = metadata.get("claim_evidence_chains")
        if not isinstance(raw_chains, list):
            continue
        for chain_index, raw_chain in enumerate(raw_chains, start=1):
            chain = _metadata_dict(raw_chain)
            chain_sources = chain.get("sources")
            if not isinstance(chain_sources, list):
                continue
            claim_text = _first_non_empty(chain.get("claim_text"), chain.get("claim"))
            verification_note = _first_non_empty(chain.get("verification_note"))
            strength = _clean_text(chain.get("evidence_strength")).lower()
            confidence = 0.9 if strength == "high" else 0.78 if strength == "medium" else 0.65
            for chain_source_index, raw_source in enumerate(chain_sources, start=1):
                source = _metadata_dict(raw_source)
                if not source:
                    continue
                source_id, source_title = _upsert_report_source(
                    source,
                    source_id_map=source_id_map,
                    source_registry=source_registry,
                    fallback_title=f"Claim Source {chain_index}.{chain_source_index}",
                )
                snippet = _truncate(_first_non_empty(claim_text, verification_note), 260)
                if snippet:
                    excerpts.append(
                        SourceExcerpt(
                            id=f"ext_report_{len(excerpts) + 1}",
                            source_id=source_id,
                            source_title=source_title,
                            snippet=snippet,
                            confidence=confidence,
                        )
                    )

    return excerpts[:12], source_registry[:12]


def _build_source_pack(
    session_id: str,
    messages: list[Any],
    knowledge_base_enabled: bool,
    vector_store_path: str | None,
    target_slide_count: int,
) -> SourcePack:
    qa_pairs = ensure_deckable_chat(messages)
    warnings: list[DeckWarning] = []
    title_hint = _build_title_hint(session_id, messages)
    chat_notes = [
        f"Question: {_truncate(question, 80)}\nAnswer summary: {_truncate(answer, 180)}"
        for question, answer in qa_pairs[-4:]
    ]

    chat_notes = [
        f"Question: {_truncate(question, 80)}\nAnswer summary: {_truncate(_answer_plaintext(answer), 180)}"
        for question, answer in qa_pairs[-4:]
    ]

    report_excerpts, report_source_registry = _extract_report_evidence_from_messages(
        messages
    )

    if not knowledge_base_enabled:
        return _build_chat_only_source_pack(
            title_hint=title_hint,
            qa_pairs=qa_pairs,
            chat_notes=chat_notes,
            warnings=warnings,
            chat_only_message="Knowledge base is disabled; the deck will be generated from successful chat answers only.",
            excerpts=report_excerpts,
            source_registry=report_source_registry,
        )

    default_store_path = os.getenv("VECTOR_STORE_PATH", "./vector_store")
    candidate_paths: list[str] = []
    if vector_store_path:
        candidate_paths.append(vector_store_path)
    if default_store_path not in candidate_paths:
        candidate_paths.append(default_store_path)

    pipeline = None
    resolved_store_path = None
    for path in candidate_paths:
        current_pipeline = DocPipeline(vector_store_path=path)
        if current_pipeline.load_store():
            pipeline = current_pipeline
            resolved_store_path = path
            break

    if pipeline is None:
        warnings.append(
            DeckWarning(
                code="kb_unavailable_fallback",
                message="Knowledge base is unavailable, so the deck will be generated from answer content only.",
            )
        )
        return _build_chat_only_source_pack(
            title_hint=title_hint,
            qa_pairs=qa_pairs,
            chat_notes=chat_notes,
            warnings=warnings,
            chat_only_message="The deck will be generated from successful chat answers only; evidence strength is lower than KB mode.",
            excerpts=report_excerpts,
            source_registry=report_source_registry,
        )

    if not vector_store_path:
        warnings.append(
            DeckWarning(
                code="kb_fallback_default",
                message="No role-specific knowledge base is bound; falling back to the default knowledge base.",
            )
        )
    elif resolved_store_path != vector_store_path:
        warnings.append(
            DeckWarning(
                code="kb_fallback_default",
                message="The bound knowledge base failed to load; falling back to the default knowledge base.",
            )
        )

    first_question = qa_pairs[0][0]
    last_answer = qa_pairs[-1][1]
    last_answer = _answer_plaintext(last_answer)
    queries = [
        title_hint,
        _truncate(first_question, 120),
        _truncate(last_answer, 200),
    ]

    collected: list[Document] = []
    for query in queries:
        if not query:
            continue
        collected.extend(pipeline.search_with_rerank(query, k=4, fetch_k=12))

    docs = _dedupe_docs(collected)[:8]
    if len(docs) < 3:
        warnings.append(
            DeckWarning(
                code="kb_insufficient_material_fallback",
                message="Knowledge base retrieval returned insufficient material; falling back to answer content only.",
            )
        )
        return _build_chat_only_source_pack(
            title_hint=title_hint,
            qa_pairs=qa_pairs,
            chat_notes=chat_notes,
            warnings=warnings,
            chat_only_message="The deck will be generated from successful chat answers only; evidence strength is lower than KB mode.",
            excerpts=report_excerpts,
            source_registry=report_source_registry,
        )

    source_id_map: dict[str, str] = {}
    source_registry: list[DeckSourceItem] = []
    excerpts: list[SourceExcerpt] = []

    for index, doc in enumerate(docs, start=1):
        source_title = str(doc.metadata.get("source", f"鏉ユ簮 {index}"))
        source_id = source_id_map.get(source_title)
        if not source_id:
            source_id = f"src_{len(source_id_map) + 1}"
            source_id_map[source_title] = source_id
            source_registry.append(
                DeckSourceItem(
                    id=source_id,
                    type="doc",
                    title=source_title,
                    document_id=str(doc.metadata.get("doc_id") or ""),
                    uri=str(doc.metadata.get("source") or ""),
                    metadata={
                        "page": doc.metadata.get("page"),
                        "chunk_index": doc.metadata.get("chunk_index"),
                    },
                )
            )
        excerpts.append(
            SourceExcerpt(
                id=f"ext_{index}",
                source_id=source_id,
                source_title=source_title,
                snippet=_truncate(doc.page_content, 260),
                confidence=0.88,
            )
        )

    if len(source_registry) < 2:
        warnings.append(
            DeckWarning(
                code="kb_insufficient_source_coverage",
                message="Knowledge base retrieval has insufficient source coverage; falling back to answer content only.",
            )
        )
        return _build_chat_only_source_pack(
            title_hint=title_hint,
            qa_pairs=qa_pairs,
            chat_notes=chat_notes,
            warnings=warnings,
            chat_only_message="The deck will be generated from successful chat answers only; evidence strength is lower than KB mode.",
            excerpts=report_excerpts,
            source_registry=report_source_registry,
        )

    if len(excerpts) < target_slide_count - 2:
        warnings.append(
            DeckWarning(
                code="evidence_sparse",
                message="Knowledge base evidence is sparse; slide count may be lower than requested.",
            )
        )

    return SourcePack(
        title_hint=title_hint,
        source_mode="kb_plus_chat",
        qa_pairs=qa_pairs,
        chat_notes=chat_notes,
        excerpts=excerpts,
        source_registry=source_registry,
        warnings=warnings,
    )


def _decide_content_slide_count(pack: SourcePack, target_slide_count: int) -> int:
    recent_qa_pairs = "\n\n---\n\n".join(
        (
            f"Question: {_truncate(question, 120)}\n"
            f"回答原文: \n{_truncate_multiline(_answer_plaintext(answer), 5000 if pack.source_mode == 'chat_only' else 2400)}"
        )
        for question, answer in pack.qa_pairs[-2:]
    )

    if pack.source_mode == "kb_plus_chat":
        desired = max(2, min(target_slide_count - 3, 5))
        evidence_bound = max(2, min(5, len(pack.excerpts) // 2 + 1))
        return min(desired, evidence_bound)
    desired = max(2, min(target_slide_count - 2, 4))
    chat_bound = max(2, min(4, len(pack.qa_pairs) + 1))
    return min(desired, chat_bound)


def _serialize_source_pack(pack: SourcePack) -> str:
    recent_qa_pairs = "\n\n---\n\n".join(
        (
            f"Question: {_truncate(question, 120)}\n"
            f"回答原文: \n{_truncate_multiline(answer, 5000 if pack.source_mode == 'chat_only' else 2400)}"
        )
        for question, answer in pack.qa_pairs[-2:]
    )
    if pack.source_mode == "kb_plus_chat":
        excerpts = "\n".join(
            f"- {excerpt.id} | {excerpt.source_id} | {excerpt.source_title}: {excerpt.snippet}"
            for excerpt in pack.excerpts
        )
        sources = "\n".join(
            f"- {source.id}: {source.title}" for source in pack.source_registry
        )
        return (
            f"成功问答摘要:\n{chr(10).join(pack.chat_notes)}\n\n"
            f"最近成功问答原文\n{recent_qa_pairs}\n\n"
            f"来源清单:\n{sources}\n\n"
            f"知识库片段\n{excerpts}"
        )
    if pack.excerpts:
        excerpts = "\n".join(
            f"- {excerpt.id} | {excerpt.source_id} | {excerpt.source_title}: {excerpt.snippet}"
            for excerpt in pack.excerpts
        )
        sources = "\n".join(
            f"- {source.id}: {source.title}" for source in pack.source_registry
        )
        return (
            "最近成功问答原文\n"
            f"{recent_qa_pairs}\n\n"
            f"报告来源清单:\n{sources}\n\n"
            f"报告证据片段:\n{excerpts}"
        )
    return "最近成功问答原文\n" + recent_qa_pairs


async def _generate_outline(
    llm,
    pack: SourcePack,
    target_slide_count: int,
    content_slide_count: int,
    system_prompt: str | None,
) -> OutlinePlan:
    evidence_rule = (
        "Each content slide must choose 1-2 available source IDs in evidence_source_ids."
        if pack.source_mode == "kb_plus_chat" or pack.excerpts
        else "Current mode has no structured evidence; return empty evidence_source_ids."
    )
    structure_rule = (
        "Use the Q&A summary and knowledge base evidence to build a formal briefing structure."
        if pack.source_mode == "kb_plus_chat"
        else "Reuse the original answer structure and do not invent unsupported conclusions."
    )
    prompt = f"""
浣犳槸涓€涓笓涓氱殑涓枃鍟嗕笟婕旂ず绋跨瓥鍒掑姪鎵嬨€傝鍩轰簬缁欏畾鏉愭枡瑙勫垝涓€涓?PPT 澶х翰銆?

纭€ц姹傦細
1. 鍙緭鍑?JSON锛屼笉瑕佽緭鍑鸿В閲娿€?
2. 涓嶈浣跨敤 Q1銆丵2銆丳art 杩欑被闂瓟鏍囬銆?
3. 鎬讳綋閲囩敤鈥滀富棰橀〉鈥濈粨鏋勶紝鑰屼笉鏄杩拌亰澶╄褰曘€?
4. 鍐呭椤垫暟閲忓繀椤绘槸 {content_slide_count} 椤点€?
5. section 鏁伴噺鎺у埗鍦?3-5 涓€?
6. {evidence_rule}
7. {structure_rule}
8. 鏍囬蹇呴』鍍忔寮忔眹鎶ユ爣棰橈紝涓嶈鐩存帴鎶婄敤鎴峰師闂鍘熸牱鎷挎潵褰撻〉鏍囬銆?
9. 椤垫暟绛栫暐鏄畞灏戝嬁姘达紝涓嶈琛ョ┖璇濄€?

杩斿洖 JSON Schema锛?
{{
  "title": "婕旂ず绋挎€绘爣棰?,
  "subtitle": "鍓爣棰?,
  "core_message": "涓€鍙ヨ瘽鏍稿績缁撹",
  "sections": ["绔犺妭1", "绔犺妭2", "绔犺妭3"],
  "content_slides": [
    {{
      "title": "涓婚椤垫爣棰?,
      "objective": "鏈〉瑕佽鏄庝粈涔?,
      "section": "鎵€灞炵珷鑺?,
      "evidence_source_ids": ["src_1"]
    }}
  ]
}}

涓婁笅鏂囪ˉ鍏咃細
- source_mode: {pack.source_mode}
- 鐩爣鎬婚〉鏁? {target_slide_count}
- 鍐呭椤垫暟: {content_slide_count}
- 鏍囬鎻愮ず: {pack.title_hint}
- Role context: {system_prompt or "None"}

鏉愭枡锛?
{_serialize_source_pack(pack)}
"""
    payload = await _invoke_json(llm, prompt)
    try:
        return OutlinePlan.model_validate(payload)
    except Exception:
        return _fallback_outline(pack, content_slide_count)


async def _generate_content_slides(
    llm,
    pack: SourcePack,
    outline: OutlinePlan,
    system_prompt: str | None,
) -> DraftedSlideBundle:
    evidence_rule = (
        "Each content slide must choose at least one available evidence_excerpt_id and use quality_state=supported."
        if pack.source_mode == "kb_plus_chat" or pack.excerpts
        else "Current mode has no structured evidence; return empty evidence fields and use quality_state=manual."
    )
    fidelity_rule = (
        "Combine evidence and answer content into concise briefing-ready slide language."
        if pack.source_mode == "kb_plus_chat"
        else "Reuse the original answer headings, order, and key points when possible."
    )
    prompt = f"""
浣犳槸涓€涓笓涓氱殑涓枃 PPT 鍐呭鎾板啓鍔╂墜銆傝鎶婁笅闈㈢殑澶х翰鎵╁睍鎴愬唴瀹归〉銆?

纭€ц姹傦細
1. 鍙緭鍑?JSON锛屼笉瑕佽緭鍑鸿В閲娿€?
2. 鍙敓鎴愬唴瀹归〉锛屼笉鐢熸垚灏侀潰銆佺洰褰曘€侀檮褰曘€?
3. 鏍囬涓嶈兘鍐欐垚 Q1銆丳art 1 杩欑闂瓟褰㈠紡銆?
4. 姣忛〉 key_points 淇濇寔 3-5 鏉★紝姣忔潯涓€鍙ヨ瘽銆?
5. speaker_notes 鐢?1-2 鍙ヨ瘽鎻愮ず璁茶堪閲嶇偣銆?
6. {evidence_rule}
7. {fidelity_rule}
8. 涓嶈鎶勭敤鎴烽棶棰橈紝涓嶈娉ㄦ按銆?

杩斿洖 JSON Schema锛?
{{
  "content_slides": [
    {{
      "title": "涓婚椤垫爣棰?,
      "subtitle": "绠€鐭壇鏍囬",
      "key_points": ["瑕佺偣1", "瑕佺偣2", "瑕佺偣3"],
      "speaker_notes": "璁茶堪鎻愮ず",
      "evidence_excerpt_ids": ["ext_1"],
      "evidence_source_ids": ["src_1"],
      "quality_state": "supported"
    }}
  ]
}}

Role context: {system_prompt or "None"}

澶х翰锛?
{outline.model_dump_json(indent=2, ensure_ascii=False)}

鏉愭枡锛?
{_serialize_source_pack(pack)}
"""
    payload = await _invoke_json(llm, prompt)
    try:
        return DraftedSlideBundle.model_validate(payload)
    except Exception:
        return _fallback_content_bundle(pack, outline)


def _sanitize_slide_title(title: str, fallback: str) -> str:
    cleaned = _clean_text(title)
    cleaned = re.sub(r"^Q\d+[:锛歕-]\s*", "", cleaned, flags=re.I)
    cleaned = re.sub(r"^Part\s*\d+[:锛歕-]?\s*", "", cleaned, flags=re.I)
    cleaned = cleaned.strip(" -:")
    return _truncate(cleaned or fallback, 48)


def _block_evidence_binding_content(
    content: dict[str, Any],
    evidence_refs: list[DeckEvidenceRef] | None = None,
) -> dict[str, Any]:
    refs = list(evidence_refs or [])
    if not refs:
        return content

    bound_content = dict(content)
    bound_content.setdefault("evidence_ref_ids", [ref.id for ref in refs if ref.id])
    bound_content.setdefault(
        "evidence_source_ids",
        list(dict.fromkeys(ref.source_id for ref in refs if ref.source_id)),
    )
    excerpt_ids = [ref.excerpt_id for ref in refs if ref.excerpt_id]
    if excerpt_ids:
        bound_content.setdefault("evidence_excerpt_ids", excerpt_ids)
    return bound_content


def _build_blocks(
    points: list[str],
    fallback: str,
    evidence_refs: list[DeckEvidenceRef] | None = None,
) -> list[DeckBlock]:
    clean_points = [_truncate(point, 72) for point in points if _clean_text(point)]
    if clean_points:
        return [
            DeckBlock(
                id=f"block_{uuid.uuid4().hex[:8]}",
                kind="bullet_list",
                role="main_points",
                content=_block_evidence_binding_content(
                    {"items": clean_points[:5]},
                    evidence_refs,
                ),
            )
        ]
    return [
        DeckBlock(
            id=f"block_{uuid.uuid4().hex[:8]}",
            kind="paragraph",
            role="summary",
            content=_block_evidence_binding_content(
                {"text": _truncate(fallback, 220)},
                evidence_refs,
            ),
        )
    ]


def _pick_evidence_refs(
    slide: DraftedContentSlide,
    pack: SourcePack,
) -> list[DeckEvidenceRef]:
    excerpt_map = {excerpt.id: excerpt for excerpt in pack.excerpts}
    source_map = {source.id: source for source in pack.source_registry}
    refs: list[DeckEvidenceRef] = []

    for excerpt_id in slide.evidence_excerpt_ids:
        excerpt = excerpt_map.get(excerpt_id)
        if not excerpt:
            continue
        refs.append(
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        )

    if refs:
        return refs[:2]

    for source_id in slide.evidence_source_ids:
        if source_id not in source_map:
            continue
        excerpt = next(
            (item for item in pack.excerpts if item.source_id == source_id), None
        )
        if excerpt is None:
            continue
        refs.append(
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        )
        if refs:
            return refs[:2]

    if pack.excerpts:
        excerpt = pack.excerpts[0]
        return [
            DeckEvidenceRef(
                id=f"ev_{uuid.uuid4().hex[:8]}",
                source_id=excerpt.source_id,
                source_title=excerpt.source_title,
                excerpt_id=excerpt.id,
                snippet=excerpt.snippet,
                confidence=excerpt.confidence,
            )
        ]

    return []


def _resolve_quality_state(
    source_mode: DeckSourceMode,
    drafted_slide: DraftedContentSlide,
    evidence_refs: list[DeckEvidenceRef],
) -> DeckQualityState:
    if evidence_refs:
        return "supported"
    if source_mode == "chat_only":
        return "manual"
    if drafted_slide.quality_state in {"supported", "weak_support", "manual"}:
        return drafted_slide.quality_state
    return "weak_support"


def _appendix_source_items(source_registry: list[DeckSourceItem]) -> list[str]:
    items: list[str] = []
    seen: set[str] = set()
    for source in source_registry:
        title = _clean_text(source.title)
        if not title or title in seen:
            continue
        items.append(title)
        seen.add(title)
    return items


def _existing_appendix_source_items(slide: DeckSlide) -> list[str]:
    items: list[str] = []
    for block in slide.blocks:
        if block.kind != "bullet_list" or block.role != "sources":
            continue
        raw_items = block.content.get("items")
        if not isinstance(raw_items, list):
            continue
        for item in raw_items:
            cleaned = _clean_text(item)
            if cleaned:
                items.append(cleaned)
    return items


def _merge_appendix_source_items(
    source_registry: list[DeckSourceItem],
    existing_items: list[str] | None = None,
) -> list[str]:
    items = _appendix_source_items(source_registry)
    seen = set(items)
    for item in existing_items or []:
        cleaned = _clean_text(item)
        if not cleaned or cleaned in seen:
            continue
        items.append(cleaned)
        seen.add(cleaned)
    return items


def _build_appendix_slide(
    source_registry: list[DeckSourceItem],
    existing_items: list[str] | None = None,
) -> DeckSlide:
    items = _merge_appendix_source_items(source_registry, existing_items)
    return DeckSlide(
        id="slide_appendix_sources",
        type="appendix_sources",
        title="Appendix: Sources",
        subtitle="Sources cited by this deck",
        layout="title-bullets",
        intent="appendix_sources",
        speaker_notes="Appendix slide for source review.",
        blocks=[
            DeckBlock(
                id="block_appendix_sources",
                kind="bullet_list",
                role="sources",
                content={"items": items},
            )
        ],
        evidence_refs=[],
        quality_state="supported",
        status=DeckSlideStatus(review_state="draft"),
    )


def _source_identity_key(source: DeckSourceItem) -> str:
    uri = _clean_text(source.uri).lower()
    if uri:
        return f"uri:{uri}"
    document_id = _clean_text(source.document_id).lower()
    if document_id:
        return f"document:{document_id}"
    title = _clean_text(source.title).lower()
    if title:
        return f"title:{title}"
    return ""


def _source_items_match(left: DeckSourceItem, right: DeckSourceItem) -> bool:
    left_uri = _clean_text(left.uri).lower()
    right_uri = _clean_text(right.uri).lower()
    if left_uri and right_uri:
        return left_uri == right_uri

    left_document_id = _clean_text(left.document_id).lower()
    right_document_id = _clean_text(right.document_id).lower()
    if left_document_id and right_document_id:
        return left_document_id == right_document_id

    left_title = _clean_text(left.title).lower()
    right_title = _clean_text(right.title).lower()
    if left_title and right_title:
        return left_title == right_title

    return left.id == right.id


def _source_item_from_evidence_ref(ref: DeckEvidenceRef) -> DeckSourceItem:
    title = _truncate(_first_non_empty(ref.source_title, ref.source_id), 96)
    return DeckSourceItem(
        id=_first_non_empty(ref.source_id, f"src_ref_{uuid.uuid4().hex[:8]}"),
        type="evidence",
        title=title or "Regenerated evidence source",
        metadata={
            "from_evidence_ref": True,
            "excerpt_id": ref.excerpt_id,
        },
    )


def _unique_source_id(base_source_id: str, used_source_ids: set[str]) -> str:
    base = _first_non_empty(base_source_id, "src_regenerated")
    if base not in used_source_ids:
        return base

    suffix = 1
    while f"{base}_regen_{suffix}" in used_source_ids:
        suffix += 1
    return f"{base}_regen_{suffix}"


def _align_evidence_ref_source_registry(
    source_registry: list[DeckSourceItem],
    ref: DeckEvidenceRef,
    regenerated_sources_by_id: dict[str, DeckSourceItem],
) -> None:
    used_source_ids = {source.id for source in source_registry}
    existing_sources_by_id = {source.id: source for source in source_registry}
    existing_sources_by_identity = {
        identity: source
        for source in source_registry
        if (identity := _source_identity_key(source))
    }

    candidate = regenerated_sources_by_id.get(ref.source_id)
    if candidate is None:
        candidate = _source_item_from_evidence_ref(ref)

    existing = existing_sources_by_id.get(ref.source_id)
    if existing is not None and _source_items_match(existing, candidate):
        ref.source_title = existing.title or ref.source_title
        return

    identity = _source_identity_key(candidate)
    identity_match = existing_sources_by_identity.get(identity) if identity else None
    if identity_match is not None:
        ref.source_id = identity_match.id
        ref.source_title = identity_match.title or ref.source_title
        return

    next_source_id = _unique_source_id(candidate.id or ref.source_id, used_source_ids)
    next_source = candidate.model_copy(deep=True, update={"id": next_source_id})
    source_registry.append(next_source)
    ref.source_id = next_source.id
    ref.source_title = next_source.title or ref.source_title


def _sync_appendix_sources_slide(deck: DeckSpec) -> None:
    if not deck.source_registry:
        return

    appendix_slide = next(
        (slide for slide in deck.slides if slide.type == "appendix_sources"),
        None,
    )
    if appendix_slide is None:
        deck.slides.append(_build_appendix_slide(deck.source_registry))
        return

    existing_items = _existing_appendix_source_items(appendix_slide)
    refreshed_appendix = _build_appendix_slide(
        deck.source_registry,
        existing_items=existing_items,
    )
    appendix_slide.subtitle = appendix_slide.subtitle or refreshed_appendix.subtitle
    appendix_slide.layout = appendix_slide.layout or refreshed_appendix.layout
    appendix_slide.intent = "appendix_sources"
    appendix_slide.speaker_notes = (
        appendix_slide.speaker_notes or refreshed_appendix.speaker_notes
    )
    appendix_slide.blocks = refreshed_appendix.blocks
    appendix_slide.evidence_refs = []
    appendix_slide.quality_state = "supported"


def _sync_deck_sources_after_regeneration(
    deck: DeckSpec,
    regenerated_deck: DeckSpec,
    replacement: DeckSlide,
) -> None:
    source_registry = [
        source.model_copy(deep=True) for source in (deck.source_registry or [])
    ]
    regenerated_sources_by_id = {
        source.id: source for source in (regenerated_deck.source_registry or [])
    }

    for ref in replacement.evidence_refs or []:
        _align_evidence_ref_source_registry(
            source_registry,
            ref,
            regenerated_sources_by_id,
        )

    deck.source_registry = source_registry
    _sync_appendix_sources_slide(deck)
    deck.generation.actual_slide_count = len(deck.slides)
    refresh_deck_evidence_coverage(deck)


def _build_deck_from_generated(
    session_id: str,
    panel_config: Any,
    target_slide_count: int,
    pack: SourcePack,
    outline: OutlinePlan,
    drafted: DraftedSlideBundle,
    theme: DeckThemeName = "default",
    source_answer_group_id: str = "",
    source_panel_id: str = "",
) -> DeckSpec:
    slides: list[DeckSlide] = [
        DeckSlide(
            id="slide_cover",
            type="cover",
            title=_truncate(outline.title or pack.title_hint, 56),
            subtitle=_truncate(outline.subtitle or outline.core_message, 96),
            layout="hero-title",
            intent="cover",
            speaker_notes="Cover slide for establishing the overall topic.",
            blocks=[
                DeckBlock(
                    id="block_cover_message",
                    kind="paragraph",
                    role="core_message",
                    content={"text": _truncate(outline.core_message, 180)},
                )
            ],
            evidence_refs=[],
            quality_state="weak_support"
            if pack.source_mode == "kb_plus_chat"
            else "manual",
            status=DeckSlideStatus(review_state="draft"),
        ),
        DeckSlide(
            id="slide_outline",
            type="outline",
            title="姹囨姤缁撴瀯",
            subtitle="Core sections of this deck",
            layout="title-bullets",
            intent="outline",
            speaker_notes="Outline slide for navigation.",
            blocks=[
                DeckBlock(
                    id="block_outline",
                    kind="bullet_list",
                    role="outline",
                    content={
                        "items": [
                            _truncate(section, 40) for section in outline.sections[:5]
                        ]
                    },
                )
            ],
            evidence_refs=[],
            quality_state="weak_support"
            if pack.source_mode == "kb_plus_chat"
            else "manual",
            status=DeckSlideStatus(review_state="draft"),
        ),
    ]

    content_plans = outline.content_slides
    drafted_slides = drafted.content_slides
    content_count = min(len(content_plans), len(drafted_slides))

    for index in range(content_count):
        plan = content_plans[index]
        drafted_slide = drafted_slides[index]
        raw_answer = pack.qa_pairs[min(index, len(pack.qa_pairs) - 1)][1]
        evidence_refs = _pick_evidence_refs(drafted_slide, pack)
        quality_state = _resolve_quality_state(
            pack.source_mode, drafted_slide, evidence_refs
        )
        title = _sanitize_slide_title(drafted_slide.title, plan.title)
        subtitle = _truncate(drafted_slide.subtitle or plan.objective, 72)
        blocks = _build_blocks(drafted_slide.key_points, plan.objective, evidence_refs)
        blocks.extend(_extract_dashboard_chart_blocks(raw_answer, limit=1))
        slides.append(
            DeckSlide(
                id=f"slide_content_{index + 1}",
                type="content",
                title=title,
                subtitle=subtitle,
                layout="title-bullets",
                intent=plan.objective,
                speaker_notes=_truncate(
                    drafted_slide.speaker_notes or plan.objective, 180
                ),
                blocks=blocks,
                evidence_refs=evidence_refs,
                quality_state=quality_state,
                status=DeckSlideStatus(review_state="draft"),
            )
        )

    if pack.source_registry:
        slides.append(_build_appendix_slide(pack.source_registry))

    warnings = list(pack.warnings)
    if pack.source_mode == "chat_only":
        warnings.append(
            DeckWarning(
                code="manual_review_required",
                message="Chat-only mode requires manual review before export.",
            )
        )

    deck = DeckSpec(
        deck_id=f"deck_{uuid.uuid4().hex}",
        meta=DeckMeta(
            title=_truncate(outline.title or pack.title_hint, 56),
            subtitle=_truncate(outline.subtitle, 96),
            theme=normalize_deck_theme(theme),
            created_at=_now_iso(),
            session_id=session_id,
            source_mode=pack.source_mode,
            generator_panel_id=getattr(panel_config, "panel_id", "panel"),
            source_answer_group_id=str(source_answer_group_id or "").strip(),
            source_panel_id=str(source_panel_id or "").strip(),
        ),
        generation=DeckGeneration(
            source=pack.source_mode,
            target_slide_count=target_slide_count,
            actual_slide_count=len(slides),
            warnings=warnings,
        ),
        slides=slides,
        source_registry=pack.source_registry,
    )
    return refresh_deck_evidence_coverage(deck)


async def build_deck(
    session_id: str,
    messages: list[Any],
    panel_config: Any,
    knowledge_base_enabled: bool,
    target_slide_count: int,
    vector_store_path: str | None = None,
    system_prompt: str | None = None,
    theme: DeckThemeName = "default",
    source_answer_group_id: str = "",
    source_panel_id: str = "",
) -> DeckSpec:
    qa_pairs = ensure_deckable_chat(messages)
    pack = _build_source_pack(
        session_id=session_id,
        messages=messages,
        knowledge_base_enabled=knowledge_base_enabled,
        vector_store_path=vector_store_path,
        target_slide_count=target_slide_count,
    )
    if not qa_pairs:
        raise ValueError("This session has no successful Q&A content for deck generation.")
    content_slide_count = _decide_content_slide_count(pack, target_slide_count)

    async def run_generation(
        active_panel_config: Any,
    ) -> tuple[OutlinePlan, DraftedSlideBundle]:
        llm = get_llm(
            provider=getattr(active_panel_config, "provider", "local"),
            model_name=getattr(active_panel_config, "model", None),
            base_url=getattr(active_panel_config, "base_url", None),
            api_key=getattr(active_panel_config, "api_key", None) or None,
            temperature=float(getattr(active_panel_config, "temperature", 0.3) or 0.3),
        )
        outline = await _generate_outline(
            llm=llm,
            pack=pack,
            target_slide_count=target_slide_count,
            content_slide_count=content_slide_count,
            system_prompt=system_prompt,
        )
        drafted = await _generate_content_slides(
            llm=llm,
            pack=pack,
            outline=outline,
            system_prompt=system_prompt,
        )
        return outline, drafted

    try:
        outline, drafted = await run_generation(panel_config)
    except Exception as exc:
        if not _is_authentication_failure(exc):
            raise
        outline, drafted = await run_generation(
            _local_ollama_panel_config(panel_config)
        )
    return _build_deck_from_generated(
        session_id=session_id,
        panel_config=panel_config,
        target_slide_count=target_slide_count,
        pack=pack,
        outline=outline,
        drafted=drafted,
        theme=normalize_deck_theme(theme),
        source_answer_group_id=source_answer_group_id,
        source_panel_id=source_panel_id,
    )


def _select_regenerated_slide(
    current_deck: DeckSpec,
    regenerated_deck: DeckSpec,
    target_slide_id: str,
) -> tuple[int, DeckSlide]:
    current_index = next(
        (
            index
            for index, slide in enumerate(current_deck.slides)
            if slide.id == target_slide_id
        ),
        -1,
    )
    if current_index < 0:
        raise KeyError(target_slide_id)

    current_slide = current_deck.slides[current_index]
    if current_slide.type == "cover":
        return current_index, regenerated_deck.slides[0]

    if current_slide.type == "outline":
        outline_slide = next(
            (slide for slide in regenerated_deck.slides if slide.type == "outline"),
            regenerated_deck.slides[min(1, len(regenerated_deck.slides) - 1)],
        )
        return current_index, outline_slide

    if current_slide.type.startswith("appendix"):
        appendix_slide = next(
            (
                slide
                for slide in regenerated_deck.slides
                if slide.type.startswith("appendix")
            ),
            regenerated_deck.slides[-1],
        )
        return current_index, appendix_slide

    current_content_index = sum(
        1 for slide in current_deck.slides[:current_index] if slide.type == "content"
    )
    regenerated_content = [
        slide for slide in regenerated_deck.slides if slide.type == "content"
    ]
    if not regenerated_content:
        raise ValueError("Regenerated deck did not contain any content slides.")

    target_content_index = min(current_content_index, len(regenerated_content) - 1)
    return current_index, regenerated_content[target_content_index]


def _reconcile_regenerated_slide_evidence(
    current_slide: DeckSlide,
    replacement: DeckSlide,
) -> DeckSlide:
    if replacement.evidence_refs:
        replacement.quality_state = "supported"
        return replacement

    # Single-slide regeneration may rewrite content without being able to cite
    # sources again; keep the previous support state so coverage stays stable.
    replacement.evidence_refs = [
        ref.model_copy(deep=True) for ref in (current_slide.evidence_refs or [])
    ]
    replacement.quality_state = current_slide.quality_state
    return replacement


async def regenerate_deck_slide(
    deck: DeckSpec,
    slide_id: str,
    messages: list[Any],
    panel_config: Any,
    knowledge_base_enabled: bool,
    vector_store_path: str | None = None,
    system_prompt: str | None = None,
) -> DeckSlide:
    regenerated_deck = await build_deck(
        session_id=deck.meta.session_id,
        messages=messages,
        panel_config=panel_config,
        knowledge_base_enabled=knowledge_base_enabled,
        target_slide_count=max(deck.generation.target_slide_count, len(deck.slides)),
        vector_store_path=vector_store_path,
        system_prompt=system_prompt,
        theme=deck.meta.theme,
        source_answer_group_id=deck.meta.source_answer_group_id,
        source_panel_id=deck.meta.source_panel_id,
    )
    current_index, replacement = _select_regenerated_slide(
        deck, regenerated_deck, slide_id
    )
    current_slide = deck.slides[current_index]
    replacement = replacement.model_copy(deep=True)
    replacement.id = current_slide.id
    replacement = _reconcile_regenerated_slide_evidence(current_slide, replacement)
    _sync_deck_sources_after_regeneration(deck, regenerated_deck, replacement)
    replacement.status = DeckSlideStatus(
        locked=current_slide.status.locked,
        dirty=False,
        review_state="regenerated",
    )
    return replacement


def build_report_markdown(messages: list[Any], title: str) -> str:
    qa_pairs = ensure_deckable_chat(messages)
    lines = [
        "---",
        "theme: default",
        f"title: {title}",
        "class: text-center",
        "---",
        "",
        f"# {title}",
        "",
        "AI 瀵硅瘽鎶ュ憡",
        "",
    ]
    for index, (question, answer) in enumerate(qa_pairs, start=1):
        lines.append("---")
        lines.append("")
        lines.append(f"## 涓婚 {index}: {_truncate(question, 72)}")
        lines.append("")
        lines.append(answer)
        lines.append("")
    return "\n".join(lines)


def _slide_text(slide: DeckSlide) -> str:
    lines: list[str] = []
    for block in slide.blocks:
        if block.kind == "bullet_list":
            for item in block.content.get("items", []):
                value = _clean_text(item)
                if value:
                    lines.append(f"• {value}")
        elif block.kind == "chart":
            title = _clean_text(block.content.get("title")) or "鍥捐〃"
            chart_type = _clean_text(block.content.get("chart_type")).lower()
            lines.append(f"[鍥捐〃] {title} ({chart_type or 'chart'})")
        else:
            text = _clean_text(block.content.get("text", ""))
            if text:
                lines.append(text)
    chart_specs = _extract_chart_specs(slide)
    if chart_specs:
        lines.append("")
        lines.append("鍥捐〃")
        for chart in chart_specs:
            chart_line = f"- {chart['title']} [{chart['chart_type']}]"
            if chart["description"]:
                chart_line += f": {chart['description']}"
            lines.append(chart_line)

    if slide.evidence_refs:
        lines.append("")
        lines.append(
            "Sources: " + ", ".join(ref.source_title for ref in slide.evidence_refs[:3])
        )
    return "\n".join(lines).strip()


def _quality_state_color_key(quality_state: DeckQualityState) -> str:
    if quality_state == "supported":
        return "success"
    if quality_state == "manual":
        return "danger"
    return "warning"


def _is_slide_manually_confirmed(slide: DeckSlide) -> bool:
    return (
        slide.quality_state != "supported" and slide.status.review_state == "confirmed"
    )


def _slide_quality_label(slide: DeckSlide) -> str:
    if _is_slide_manually_confirmed(slide):
        return "Manually confirmed"
    return _quality_state_label(slide.quality_state)


def _slide_quality_color_key(slide: DeckSlide) -> str:
    if _is_slide_manually_confirmed(slide):
        return "success"
    return _quality_state_color_key(slide.quality_state)


def _extract_chart_specs(slide: DeckSlide) -> list[dict[str, Any]]:
    specs: list[dict[str, Any]] = []
    for block in slide.blocks:
        if block.kind != "chart":
            continue
        chart_type = _clean_text(
            block.content.get("chart_type") or block.content.get("type")
        ).lower()
        if chart_type not in {"bar", "line", "pie"}:
            continue

        labels = _normalize_chart_labels(block.content.get("labels"))
        datasets = _normalize_chart_datasets(block.content.get("datasets"), len(labels))
        if not datasets:
            continue

        if not labels:
            labels = [f"绫诲埆 {index + 1}" for index in range(len(datasets[0]["data"]))]
            datasets = _normalize_chart_datasets(
                block.content.get("datasets"), len(labels)
            )
            if not datasets:
                continue

        specs.append(
            {
                "title": _clean_text(block.content.get("title")) or "鏁版嵁鍥捐〃",
                "description": _clean_text(block.content.get("description")),
                "chart_type": chart_type,
                "labels": labels[:12],
                "datasets": datasets[:4],
            }
        )
    return specs


def _extract_slide_sections(slide: DeckSlide) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    for block in slide.blocks:
        if block.kind == "bullet_list":
            items = [
                _clean_text(item)
                for item in block.content.get("items", [])
                if _clean_text(item)
            ]
            if items:
                sections.append(
                    {
                        "kind": "bullet_list",
                        "role": block.role,
                        "items": items,
                    }
                )
        else:
            text = _clean_text(block.content.get("text", ""))
            if text:
                sections.append(
                    {
                        "kind": "paragraph",
                        "role": block.role,
                        "text": text,
                    }
                )
    return sections


def _quality_state_label(quality_state: DeckQualityState) -> str:
    if quality_state == "supported":
        return "证据充分"
    if quality_state == "manual":
        return "需人工确认"
    return "证据偏弱"


def _split_evenly(items: list[str], column_count: int = 2) -> list[list[str]]:
    if column_count <= 1:
        return [items]
    if not items:
        return [[] for _ in range(column_count)]

    per_column = (len(items) + column_count - 1) // column_count
    return [
        items[index * per_column : (index + 1) * per_column]
        for index in range(column_count)
    ]


def _slide_citation_marker(slide: DeckSlide, *, limit: int = 3) -> str:
    # Keep PPTX body citations compact and aligned with the evidence side panel.
    ref_count = len(slide.evidence_refs or [])
    if ref_count <= 0:
        return ""

    marker_count = min(ref_count, max(1, limit))
    marker = "".join(f"[{index}]" for index in range(1, marker_count + 1))
    return f"{marker}+" if ref_count > marker_count else marker


def _append_citation_marker(
    text: str,
    marker: str,
    *,
    limit: int | None = None,
) -> str:
    cleaned = _clean_text(text)
    if not cleaned:
        return ""
    if not marker or cleaned.endswith(marker):
        return _truncate(cleaned, limit) if limit else cleaned
    if limit:
        suffix = f" {marker}"
        cleaned = _truncate(cleaned, max(1, limit - len(suffix)))
        return f"{cleaned}{suffix}"
    return f"{cleaned} {marker}"


def _compose_export_notes(slide: DeckSlide) -> str:
    lines: list[str] = [slide.title.strip() or "Untitled Slide"]

    if slide.subtitle.strip():
        lines.append(slide.subtitle.strip())

    lines.append(f"质量状态: {_quality_state_label(slide.quality_state)}")

    lines[-1] = f"质量状态: {_slide_quality_label(slide)}"
    if slide.speaker_notes.strip():
        lines.extend(["", "璁茶堪澶囨敞", slide.speaker_notes.strip()])
        lines.extend(["", "演讲备注", slide.speaker_notes.strip()])
    sections = _extract_slide_sections(slide)
    if sections:
        lines.append("")
        lines.append("页面内容")
        for section in sections:
            role = _clean_text(section.get("role", ""))
            if role and role not in {"main_points", "summary", "outline", "sources"}:
                lines.append(f"[{role}]")
            if section["kind"] == "bullet_list":
                lines.extend(f"- {item}" for item in section["items"])
            else:
                lines.append(section["text"])

    if slide.evidence_refs:
        lines.append("")
        lines.append("证据来源")
        for ref in slide.evidence_refs:
            confidence = (
                f" ({round(ref.confidence * 100)}%)"
                if ref.confidence and ref.confidence > 0
                else ""
            )
            snippet = _truncate(_clean_text(ref.snippet), 220)
            if snippet:
                lines.append(f"- {ref.source_title}{confidence}: {snippet}")
            else:
                lines.append(f"- {ref.source_title}{confidence}")

    return "\n".join(line for line in lines if line is not None).strip()


def export_deck_to_pptx(deck: DeckSpec) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is not installed. Install requirements.txt and retry."
        ) from exc

    theme = DECK_THEME_PALETTES[normalize_deck_theme(deck.meta.theme)]

    def rgb(color_key: str) -> RGBColor:
        return RGBColor.from_string(theme.get(color_key, color_key))

    def set_slide_background(ppt_slide, color_key: str = "bg") -> None:
        fill = ppt_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color_key)

    def add_textbox(
        ppt_slide,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill_color: str | None = None,
        line_color: str | None = None,
        margins: tuple[float, float, float, float] = (0.08, 0.05, 0.08, 0.05),
        vertical_anchor=MSO_ANCHOR.TOP,
    ):
        from pptx.util import Inches

        shape = ppt_slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.vertical_anchor = vertical_anchor
        text_frame.margin_left = Inches(margins[0])
        text_frame.margin_top = Inches(margins[1])
        text_frame.margin_right = Inches(margins[2])
        text_frame.margin_bottom = Inches(margins[3])

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill_color)
        else:
            shape.fill.background()

        if line_color:
            shape.line.color.rgb = rgb(line_color)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()

        return shape, text_frame

    def set_text_frame_content(text_frame, paragraphs: list[dict[str, Any]]) -> None:
        text_frame.clear()
        if not paragraphs:
            return

        for index, spec in enumerate(paragraphs):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = spec.get("text", "")
            paragraph.alignment = spec.get("align", PP_ALIGN.LEFT)
            paragraph.space_after = Pt(spec.get("space_after", 6))
            paragraph.space_before = Pt(spec.get("space_before", 0))
            paragraph.line_spacing = spec.get("line_spacing", 1.15)
            font = paragraph.font
            font.name = spec.get("font_name", "Microsoft YaHei")
            font.size = Pt(spec.get("font_size", 18))
            font.bold = spec.get("bold", False)
            font.italic = spec.get("italic", False)
            font.color.rgb = rgb(spec.get("color", "body"))

    def add_badge(
        ppt_slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill_color: str,
        text_color: str = "surface",
    ) -> None:
        from pptx.util import Inches

        shape = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        shape.line.fill.background()
        text_frame = shape.text_frame
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        set_text_frame_content(
            text_frame,
            [
                {
                    "text": text,
                    "font_size": 10,
                    "bold": True,
                    "color": text_color,
                    "align": PP_ALIGN.CENTER,
                    "space_after": 0,
                }
            ],
        )

    def quality_color(slide: DeckSlide) -> str:
        return _slide_quality_color_key(slide)

    def add_chart_panel(
        ppt_slide,
        chart: dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
    ) -> None:
        from pptx.util import Inches

        panel = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb("surface_alt")
        panel.line.color.rgb = rgb("border")
        panel.line.width = Pt(1)

        title_height = 0.36
        description_height = 0.0
        description = _clean_text(chart.get("description"))
        if description:
            description_height = 0.36

        _, title_box = add_textbox(
            ppt_slide,
            x + 0.18,
            y + 0.12,
            w - 0.36,
            title_height,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            title_box,
            [
                {
                    "text": chart["title"],
                    "font_size": 11.5,
                    "bold": True,
                    "color": "title",
                    "space_after": 0,
                }
            ],
        )

        if description:
            _, desc_box = add_textbox(
                ppt_slide,
                x + 0.18,
                y + 0.46,
                w - 0.36,
                description_height,
                margins=(0.0, 0.0, 0.0, 0.0),
            )
            set_text_frame_content(
                desc_box,
                [
                    {
                        "text": _truncate(description, 110),
                        "font_size": 9.5,
                        "color": "muted",
                        "space_after": 0,
                    }
                ],
            )

        chart_box_y = y + 0.56 + description_height
        chart_box_h = max(1.2, h - (chart_box_y - y) - 0.18)
        chart_data = CategoryChartData()
        chart_data.categories = chart["labels"]
        datasets = chart["datasets"]
        if chart["chart_type"] == "pie":
            datasets = datasets[:1]
        for dataset in datasets:
            chart_data.add_series(dataset["label"], tuple(dataset["data"]))

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        graphic_frame = ppt_slide.shapes.add_chart(
            chart_type_map[chart["chart_type"]],
            Inches(x + 0.16),
            Inches(chart_box_y),
            Inches(w - 0.32),
            Inches(chart_box_h),
            chart_data,
        )
        ppt_chart = graphic_frame.chart
        ppt_chart.has_title = False
        ppt_chart.has_legend = len(datasets) > 1 or chart["chart_type"] == "pie"
        if ppt_chart.has_legend:
            ppt_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            ppt_chart.legend.font.size = Pt(9)

        palette = ["accent", "success", "warning", "danger"]
        try:
            if chart["chart_type"] == "pie":
                series = ppt_chart.series[0]
                for point_index, point in enumerate(series.points):
                    color_key = palette[point_index % len(palette)]
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = rgb(color_key)
                    point.format.line.color.rgb = rgb("surface")
            else:
                for series_index, series in enumerate(ppt_chart.series):
                    color_key = palette[series_index % len(palette)]
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = rgb(color_key)
                    series.format.line.color.rgb = rgb(color_key)
        except Exception:
            pass

    def add_footer(
        ppt_slide,
        slide_index: int,
        total_slides: int,
        slide: DeckSlide,
    ) -> None:
        from pptx.util import Inches

        divider = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.75),
            Inches(6.78),
            Inches(11.85),
            Inches(0.02),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = rgb("border")
        divider.line.fill.background()

        _, left_footer = add_textbox(
            ppt_slide,
            0.78,
            6.84,
            7.6,
            0.28,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            left_footer,
            [
                {
                    "text": f"{deck.meta.source_mode} | {slide.type} | {slide.layout}",
                    "font_size": 9.5,
                    "color": "muted",
                    "space_after": 0,
                }
            ],
        )

        _, right_footer = add_textbox(
            ppt_slide,
            10.55,
            6.82,
            1.75,
            0.3,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            right_footer,
            [
                {
                    "text": f"{slide_index + 1}/{total_slides}",
                    "font_size": 10,
                    "bold": True,
                    "color": "muted",
                    "align": PP_ALIGN.RIGHT,
                    "space_after": 0,
                }
            ],
        )

    def add_notes(ppt_slide, slide: DeckSlide) -> None:
        notes_text = _compose_export_notes(slide)
        if notes_text:
            ppt_slide.notes_slide.notes_text_frame.text = notes_text

    def format_date_label(raw_value: str) -> str:
        raw = _clean_text(raw_value)
        if not raw:
            return "鏈煡鏃堕棿"
        if "T" in raw:
            return raw.split("T", 1)[0]
        return raw[:10]

    def collect_bullet_items(slide: DeckSlide) -> list[str]:
        items: list[str] = []
        for section in _extract_slide_sections(slide):
            if section["kind"] == "bullet_list":
                items.extend(section["items"])
        return items

    def collect_paragraph_texts(slide: DeckSlide) -> list[str]:
        items: list[str] = []
        for section in _extract_slide_sections(slide):
            if section["kind"] == "paragraph":
                items.append(section["text"])
        return items

    def render_cover_slide(ppt_slide, slide_index: int, slide: DeckSlide) -> None:
        from pptx.util import Inches

        set_slide_background(ppt_slide, "surface")
        top_bar = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.22),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = rgb("accent")
        top_bar.line.fill.background()

        accent_panel = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(10.85),
            Inches(0.45),
            Inches(1.85),
            Inches(5.55),
        )
        accent_panel.fill.solid()
        accent_panel.fill.fore_color.rgb = rgb("surface_alt")
        accent_panel.line.fill.background()

        core_message = ""
        for section in _extract_slide_sections(slide):
            if section["kind"] == "paragraph":
                core_message = section["text"]
                break
            if section["kind"] == "bullet_list":
                core_message = "; ".join(section["items"][:2])
                break

        _, eyebrow = add_textbox(
            ppt_slide,
            0.9,
            0.55,
            3.4,
            0.35,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            eyebrow,
            [
                {
                    "text": "AI Knowledge Base Deck",
                    "font_size": 11,
                    "bold": True,
                    "color": "accent",
                    "space_after": 0,
                }
            ],
        )

        _, title_box = add_textbox(
            ppt_slide,
            0.9,
            1.0,
            10.9,
            1.4,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            title_box,
            [
                {
                    "text": slide.title or deck.meta.title,
                    "font_size": 26,
                    "bold": True,
                    "color": "title",
                    "line_spacing": 1.05,
                    "space_after": 0,
                }
            ],
        )

        _, subtitle_box = add_textbox(
            ppt_slide,
            0.92,
            2.35,
            9.1,
            0.75,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            subtitle_box,
            [
                {
                    "text": slide.subtitle
                    or deck.meta.subtitle
                    or "AI 鑷姩鐢熸垚鐨勭粨鏋勫寲姹囨姤鑽夌",
                    "font_size": 16,
                    "color": "muted",
                    "line_spacing": 1.15,
                    "space_after": 0,
                }
            ],
        )

        if core_message:
            _, message_box = add_textbox(
                ppt_slide,
                0.9,
                3.25,
                6.7,
                1.45,
                fill_color="surface_alt",
                line_color="accent_soft",
                margins=(0.18, 0.12, 0.18, 0.12),
            )
            set_text_frame_content(
                message_box,
                [
                    {
                        "text": "鏍稿績缁撹",
                        "font_size": 11,
                        "bold": True,
                        "color": "accent",
                        "space_after": 6,
                    },
                    {
                        "text": core_message,
                        "font_size": 15,
                        "color": "body",
                        "line_spacing": 1.18,
                        "space_after": 0,
                    },
                ],
            )

        _, meta_box = add_textbox(
            ppt_slide,
            8.4,
            3.28,
            3.9,
            2.05,
            fill_color="bg",
            line_color="border",
            margins=(0.16, 0.12, 0.16, 0.12),
        )
        set_text_frame_content(
            meta_box,
            [
                {
                    "text": "瀵煎嚭鎽樿",
                    "font_size": 11,
                    "bold": True,
                    "color": "title",
                    "space_after": 8,
                },
                {
                    "text": f"鏉ユ簮妯″紡: {deck.meta.source_mode}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 4,
                },
                {
                    "text": f"鍙椾紬: {deck.meta.audience}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 4,
                },
                {
                    "text": f"椤垫暟: {deck.generation.actual_slide_count}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 4,
                },
                {
                    "text": f"鐢熸垚闈㈡澘: {deck.meta.generator_panel_id}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 4,
                },
                {
                    "text": f"鏃ユ湡: {format_date_label(deck.meta.created_at)}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 0,
                },
            ],
        )

        add_badge(
            ppt_slide,
            _slide_quality_label(slide),
            10.55,
            0.62,
            1.72,
            0.38,
            fill_color=quality_color(slide),
        )
        add_footer(ppt_slide, slide_index, len(deck.slides), slide)
        add_notes(ppt_slide, slide)

    def render_outline_slide(ppt_slide, slide_index: int, slide: DeckSlide) -> None:
        set_slide_background(ppt_slide)
        add_badge(
            ppt_slide,
            _slide_quality_label(slide),
            10.55,
            0.55,
            1.72,
            0.38,
            fill_color=quality_color(slide),
        )

        _, title_box = add_textbox(
            ppt_slide,
            0.8,
            0.65,
            8.9,
            0.65,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            title_box,
            [
                {
                    "text": slide.title,
                    "font_size": 24,
                    "bold": True,
                    "color": "title",
                    "space_after": 0,
                }
            ],
        )

        if slide.subtitle:
            _, subtitle_box = add_textbox(
                ppt_slide,
                0.82,
                1.18,
                8.8,
                0.45,
                margins=(0.0, 0.0, 0.0, 0.0),
            )
            set_text_frame_content(
                subtitle_box,
                [
                    {
                        "text": slide.subtitle,
                        "font_size": 13,
                        "color": "muted",
                        "space_after": 0,
                    }
                ],
            )

        outline_items: list[str] = []
        for section in _extract_slide_sections(slide):
            if section["kind"] == "bullet_list":
                outline_items.extend(section["items"])

        visible_items = outline_items[:6]
        card_width = 3.35
        card_height = 1.08
        gap_x = 0.25
        gap_y = 0.22
        start_x = 0.82
        start_y = 1.86

        if not visible_items:
            _, body_box = add_textbox(
                ppt_slide,
                0.82,
                1.8,
                7.1,
                4.55,
                fill_color="surface",
                line_color="border",
                margins=(0.18, 0.14, 0.18, 0.12),
            )
            set_text_frame_content(
                body_box,
                [
                    {
                        "text": "No outline content was generated.",
                        "font_size": 16,
                        "color": "muted",
                        "space_after": 0,
                    }
                ],
            )
        else:
            for item_index, item in enumerate(visible_items, start=1):
                row = (item_index - 1) // 2
                col = (item_index - 1) % 2
                x = start_x + col * (card_width + gap_x)
                y = start_y + row * (card_height + gap_y)
                fill_color = "surface_alt" if item_index == 1 else "surface"
                line_color = "accent_soft" if item_index == 1 else "border"
                _, card_box = add_textbox(
                    ppt_slide,
                    x,
                    y,
                    card_width,
                    card_height,
                    fill_color=fill_color,
                    line_color=line_color,
                    margins=(0.18, 0.14, 0.18, 0.12),
                )
                set_text_frame_content(
                    card_box,
                    [
                        {
                            "text": f"Part {item_index}",
                            "font_size": 10.5,
                            "bold": True,
                            "color": "accent",
                            "space_after": 8,
                        },
                        {
                            "text": item,
                            "font_size": 15.5,
                            "bold": item_index == 1,
                            "color": "body",
                            "line_spacing": 1.12,
                            "space_after": 0,
                        },
                    ],
                )

        _, info_box = add_textbox(
            ppt_slide,
            8.3,
            1.8,
            4.05,
            4.55,
            fill_color="surface_alt",
            line_color="border",
            margins=(0.16, 0.14, 0.16, 0.12),
        )
        warnings_text = (
            "; ".join(w.message for w in deck.generation.warnings[:2])
            if deck.generation.warnings
            else "DeckSpec structured export is being used."
        )
        set_text_frame_content(
            info_box,
            [
                {
                    "text": "瀵煎嚭淇℃伅",
                    "font_size": 12,
                    "bold": True,
                    "color": "title",
                    "space_after": 10,
                },
                {
                    "text": f"鍙椾紬: {deck.meta.audience}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 6,
                },
                {
                    "text": f"鐩殑: {deck.meta.purpose}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 6,
                },
                {
                    "text": f"鏉ユ簮妯″紡: {deck.meta.source_mode}",
                    "font_size": 11,
                    "color": "body",
                    "space_after": 6,
                },
                {
                    "text": f"椋庨櫓鎻愮ず: {warnings_text}",
                    "font_size": 10.5,
                    "color": "muted",
                    "line_spacing": 1.2,
                    "space_after": 0,
                },
            ],
        )

        add_footer(ppt_slide, slide_index, len(deck.slides), slide)
        add_notes(ppt_slide, slide)

    def render_content_slide(ppt_slide, slide_index: int, slide: DeckSlide) -> None:
        set_slide_background(ppt_slide, "surface")
        add_badge(
            ppt_slide,
            _slide_quality_label(slide),
            10.55,
            0.48,
            1.72,
            0.38,
            fill_color=quality_color(slide),
        )

        from pptx.util import Inches

        accent_rule = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.8),
            Inches(0.48),
            Inches(0.09),
            Inches(0.82),
        )
        accent_rule.fill.solid()
        accent_rule.fill.fore_color.rgb = rgb("accent")
        accent_rule.line.fill.background()

        _, title_box = add_textbox(
            ppt_slide,
            1.0,
            0.56,
            9.1,
            0.65,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            title_box,
            [
                {
                    "text": slide.title,
                    "font_size": 23,
                    "bold": True,
                    "color": "title",
                    "space_after": 0,
                }
            ],
        )

        if slide.subtitle:
            _, subtitle_box = add_textbox(
                ppt_slide,
                1.02,
                1.08,
                8.7,
                0.42,
                margins=(0.0, 0.0, 0.0, 0.0),
            )
            set_text_frame_content(
                subtitle_box,
                [
                    {
                        "text": slide.subtitle,
                        "font_size": 12.5,
                        "color": "muted",
                        "space_after": 0,
                    }
                ],
            )

        sections = _extract_slide_sections(slide)
        chart_specs = _extract_chart_specs(slide)
        primary_chart = chart_specs[0] if chart_specs else None
        citation_marker = _slide_citation_marker(slide)
        bullet_items = collect_bullet_items(slide)
        paragraph_texts = collect_paragraph_texts(slide)
        summary_text = (
            paragraph_texts[0]
            if paragraph_texts
            else _clean_text(slide.subtitle or slide.intent)
        )

        body_x = 0.8
        body_y = 1.65
        body_width = 7.85 if slide.evidence_refs else 11.45
        body_height = 4.95

        if summary_text:
            _, summary_box = add_textbox(
                ppt_slide,
                body_x,
                1.62,
                body_width,
                0.78,
                fill_color="surface_alt",
                line_color="accent_soft",
                margins=(0.18, 0.12, 0.18, 0.1),
            )
            set_text_frame_content(
                summary_box,
                [
                    {
                        "text": "关键摘要",
                        "font_size": 10.5,
                        "bold": True,
                        "color": "accent",
                        "space_after": 5,
                    },
                    {
                        "text": _append_citation_marker(
                            summary_text,
                            citation_marker,
                            limit=120,
                        ),
                        "font_size": 13.5,
                        "color": "body",
                        "line_spacing": 1.15,
                        "space_after": 0,
                    },
                ],
            )
            body_y = 2.55
            body_height = 4.05

        chart_panel: tuple[float, float, float, float] | None = None
        if primary_chart and slide.evidence_refs:
            body_height = 1.55
            chart_panel = (0.8, 4.15, 7.85, 2.25)
        elif primary_chart:
            body_width = 5.15
            chart_panel = (6.2, body_y, 6.05, body_height)

        bullet_count = len(bullet_items)
        max_bullet_length = max((len(item) for item in bullet_items), default=0)
        use_cards = (
            not slide.evidence_refs
            and primary_chart is None
            and 1 < bullet_count <= 4
            and max_bullet_length <= 34
        )
        use_two_columns = (
            not slide.evidence_refs and primary_chart is None and bullet_count >= 5
        )

        if use_cards:
            card_width = 5.55
            card_height = 1.24 if bullet_count <= 2 else 1.08
            gap_x = 0.35
            gap_y = 0.22
            for item_index, item in enumerate(bullet_items[:4], start=1):
                row = (item_index - 1) // 2
                col = (item_index - 1) % 2
                x = body_x + col * (card_width + gap_x)
                y = body_y + row * (card_height + gap_y)
                _, card_box = add_textbox(
                    ppt_slide,
                    x,
                    y,
                    card_width,
                    card_height,
                    fill_color="bg" if item_index % 2 == 1 else "surface_alt",
                    line_color="border",
                    margins=(0.16, 0.12, 0.16, 0.1),
                )
                set_text_frame_content(
                    card_box,
                    [
                        {
                            "text": f"瑕佺偣 {item_index}",
                            "font_size": 10.5,
                            "bold": True,
                            "color": "accent",
                            "space_after": 7,
                        },
                        {
                            "text": item,
                            "font_size": 16,
                            "color": "body",
                            "line_spacing": 1.15,
                            "space_after": 0,
                        },
                    ],
                )
        elif use_two_columns:
            split_columns = _split_evenly(bullet_items, column_count=2)
            column_width = 5.55
            gap_x = 0.35
            for column_index, items in enumerate(split_columns):
                x = body_x + column_index * (column_width + gap_x)
                _, column_box = add_textbox(
                    ppt_slide,
                    x,
                    body_y,
                    column_width,
                    body_height,
                    fill_color="bg",
                    line_color="border",
                    margins=(0.18, 0.14, 0.18, 0.12),
                )
                column_paragraphs: list[dict[str, Any]] = [
                    {
                        "text": f"鍏抽敭瑕佺偣 {column_index + 1}",
                        "font_size": 10.5,
                        "bold": True,
                        "color": "accent",
                        "space_after": 8,
                    }
                ]
                for item in items:
                    column_paragraphs.append(
                        {
                            "text": f"• {item}",
                            "font_size": 14.5,
                            "color": "body",
                            "line_spacing": 1.16,
                            "space_after": 8,
                        }
                    )
                set_text_frame_content(column_box, column_paragraphs)
        else:
            _, body_box = add_textbox(
                ppt_slide,
                body_x,
                body_y,
                body_width,
                body_height,
                fill_color="bg",
                line_color="border",
                margins=(0.18, 0.14, 0.18, 0.12),
            )

            body_font_size = 17 if bullet_count <= 4 else 15
            body_paragraphs: list[dict[str, Any]] = []
            for section in sections:
                role = _clean_text(section.get("role", ""))
                if role and role not in {"main_points", "summary"}:
                    body_paragraphs.append(
                        {
                            "text": role.replace("_", " ").title(),
                            "font_size": 10.5,
                            "bold": True,
                            "color": "accent",
                            "space_after": 6,
                        }
                    )
                if section["kind"] == "bullet_list":
                    for item in section["items"]:
                        display_item = _append_citation_marker(item, citation_marker)
                        body_paragraphs.append(
                            {
                                "text": f"• {display_item}",
                                "font_size": body_font_size,
                                "color": "body",
                                "line_spacing": 1.18,
                                "space_after": 10,
                            }
                        )
                else:
                    display_text = _append_citation_marker(
                        section["text"], citation_marker
                    )
                    body_paragraphs.append(
                        {
                            "text": display_text,
                            "font_size": 15,
                            "color": "body",
                            "line_spacing": 1.2,
                            "space_after": 10,
                        }
                    )
            set_text_frame_content(body_box, body_paragraphs)

        if primary_chart and chart_panel:
            add_chart_panel(
                ppt_slide,
                primary_chart,
                chart_panel[0],
                chart_panel[1],
                chart_panel[2],
                chart_panel[3],
            )

        if slide.evidence_refs:
            _, evidence_box = add_textbox(
                ppt_slide,
                8.95,
                1.62,
                3.55,
                4.98,
                fill_color="surface_alt",
                line_color="border",
                margins=(0.16, 0.14, 0.16, 0.12),
            )
            evidence_paragraphs: list[dict[str, Any]] = [
                {
                    "text": "证据来源",
                    "font_size": 12,
                    "bold": True,
                    "color": "title",
                    "space_after": 8,
                }
            ]
            for ref_index, ref in enumerate(slide.evidence_refs[:3], start=1):
                confidence = (
                    f" ({round(ref.confidence * 100)}%)"
                    if ref.confidence and ref.confidence > 0
                    else ""
                )
                evidence_paragraphs.append(
                    {
                        "text": f"[{ref_index}] {ref.source_title}{confidence}",
                        "font_size": 10.5,
                        "bold": True,
                        "color": "body",
                        "space_after": 4,
                    }
                )
                evidence_paragraphs.append(
                    {
                        "text": _truncate(_clean_text(ref.snippet), 150) or "No excerpt.",
                        "font_size": 9.5,
                        "color": "muted",
                        "line_spacing": 1.2,
                        "space_after": 10,
                    }
                )
            evidence_paragraphs.append(
                {
                    "text": f"状态: {_slide_quality_label(slide)}",
                    "font_size": 10,
                    "bold": True,
                    "color": quality_color(slide),
                    "space_after": 0,
                }
            )
            set_text_frame_content(evidence_box, evidence_paragraphs)

        add_footer(ppt_slide, slide_index, len(deck.slides), slide)
        add_notes(ppt_slide, slide)

    def render_appendix_slide(ppt_slide, slide_index: int, slide: DeckSlide) -> None:
        set_slide_background(ppt_slide)
        add_badge(
            ppt_slide,
            _slide_quality_label(slide),
            10.55,
            0.55,
            1.72,
            0.38,
            fill_color=quality_color(slide),
        )

        _, title_box = add_textbox(
            ppt_slide,
            0.8,
            0.65,
            10.0,
            0.62,
            margins=(0.0, 0.0, 0.0, 0.0),
        )
        set_text_frame_content(
            title_box,
            [
                {
                    "text": slide.title,
                    "font_size": 22,
                    "bold": True,
                    "color": "title",
                    "space_after": 0,
                }
            ],
        )

        if slide.subtitle:
            _, subtitle_box = add_textbox(
                ppt_slide,
                0.82,
                1.14,
                9.4,
                0.4,
                margins=(0.0, 0.0, 0.0, 0.0),
            )
            set_text_frame_content(
                subtitle_box,
                [
                    {
                        "text": slide.subtitle,
                        "font_size": 12.5,
                        "color": "muted",
                        "space_after": 0,
                    }
                ],
            )

        appendix_items: list[str] = []
        for section in _extract_slide_sections(slide):
            if section["kind"] == "bullet_list":
                appendix_items.extend(section["items"])
            else:
                appendix_items.append(section["text"])

        total_items = len(appendix_items)
        _, stat_box = add_textbox(
            ppt_slide,
            9.58,
            0.72,
            2.88,
            0.82,
            fill_color="surface_alt",
            line_color="border",
            margins=(0.14, 0.1, 0.14, 0.08),
        )
        set_text_frame_content(
            stat_box,
            [
                {
                    "text": "鏉ユ簮缁熻",
                    "font_size": 10.5,
                    "bold": True,
                    "color": "accent",
                    "space_after": 4,
                    "align": PP_ALIGN.CENTER,
                },
                {
                    "text": f"共 {total_items} 个来源",
                    "font_size": 13,
                    "bold": True,
                    "color": "title",
                    "align": PP_ALIGN.CENTER,
                    "space_after": 0,
                },
            ],
        )

        if total_items <= 12:
            card_width = 5.55
            card_height = 0.68 if total_items <= 8 else 0.56
            gap_x = 0.3
            gap_y = 0.16
            for item_index, item in enumerate(appendix_items, start=1):
                row = (item_index - 1) // 2
                col = (item_index - 1) % 2
                x = 0.8 + col * (card_width + gap_x)
                y = 1.78 + row * (card_height + gap_y)
                _, item_box = add_textbox(
                    ppt_slide,
                    x,
                    y,
                    card_width,
                    card_height,
                    fill_color="surface" if item_index % 2 else "bg",
                    line_color="border",
                    margins=(0.14, 0.08, 0.14, 0.06),
                )
                set_text_frame_content(
                    item_box,
                    [
                        {
                            "text": f"{item_index:02d}",
                            "font_size": 9.5,
                            "bold": True,
                            "color": "accent",
                            "space_after": 3,
                        },
                        {
                            "text": item,
                            "font_size": 11.5 if total_items <= 8 else 10.5,
                            "color": "body",
                            "space_after": 0,
                            "line_spacing": 1.05,
                        },
                    ],
                )
        else:
            columns = _split_evenly(appendix_items, column_count=2)
            font_size = 11
            for column_index, items in enumerate(columns):
                x = 0.8 + column_index * 6.1
                _, column_box = add_textbox(
                    ppt_slide,
                    x,
                    1.78,
                    5.45,
                    4.9,
                    fill_color="surface",
                    line_color="border",
                    margins=(0.16, 0.12, 0.16, 0.1),
                )
                column_paragraphs = [
                    {
                        "text": f"鏉ユ簮 {column_index + 1}",
                        "font_size": 11,
                        "bold": True,
                        "color": "accent",
                        "space_after": 8,
                    }
                ]
                column_paragraphs.extend(
                    {
                        "text": f"• {item}",
                        "font_size": font_size,
                        "color": "body",
                        "line_spacing": 1.12,
                        "space_after": 6,
                    }
                    for item in items
                )
                set_text_frame_content(column_box, column_paragraphs)

        add_footer(ppt_slide, slide_index, len(deck.slides), slide)
        add_notes(ppt_slide, slide)

    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    presentation.core_properties.title = deck.meta.title
    presentation.core_properties.author = deck.meta.author or "system"
    presentation.core_properties.subject = deck.meta.subtitle or deck.meta.purpose
    presentation.core_properties.keywords = "deck,pptx,insightdesk"

    blank_layout = presentation.slide_layouts[6]
    for index, slide in enumerate(deck.slides):
        ppt_slide = presentation.slides.add_slide(blank_layout)
        if slide.type == "cover" or index == 0:
            render_cover_slide(ppt_slide, index, slide)
        elif slide.type == "outline":
            render_outline_slide(ppt_slide, index, slide)
        elif slide.type.startswith("appendix"):
            render_appendix_slide(ppt_slide, index, slide)
        else:
            render_content_slide(ppt_slide, index, slide)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


def build_export_filename(deck: DeckSpec, extension: str) -> str:
    safe_title = "".join(
        char for char in deck.meta.title if char.isalnum() or char in (" ", "-", "_")
    ).strip()[:40]
    return f"{safe_title or 'deck'}.{extension}"


class SQLiteDeckStore:
    def __init__(self, db_path: str | None = None):
        self.db_path = str(db_path or app_database_path()).strip()
        self._init_db()

    def _init_db(self) -> None:
        with connect_sqlite(self.db_path) as conn:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS decks (
                    deck_id TEXT PRIMARY KEY,
                    session_id TEXT NOT NULL,
                    title TEXT NOT NULL,
                    spec_json TEXT NOT NULL,
                    created_at REAL NOT NULL,
                    updated_at REAL NOT NULL
                )
                """
            )
            conn.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_decks_session
                ON decks(session_id)
                """
            )
            conn.commit()

    def save(self, deck: DeckSpec) -> DeckSpec:
        self._init_db()
        refresh_deck_evidence_coverage(deck)
        now = time.time()
        payload = json.dumps(deck.model_dump(mode="json"), ensure_ascii=False)
        with connect_sqlite(self.db_path) as conn:
            conn.execute(
                """
                INSERT INTO decks (deck_id, session_id, title, spec_json, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?)
                ON CONFLICT(deck_id) DO UPDATE SET
                    session_id = excluded.session_id,
                    title = excluded.title,
                    spec_json = excluded.spec_json,
                    updated_at = excluded.updated_at
                """,
                (
                    deck.deck_id,
                    deck.meta.session_id,
                    deck.meta.title,
                    payload,
                    now,
                    now,
                ),
            )
            conn.commit()
        return deck

    def get(self, deck_id: str) -> DeckSpec:
        self._init_db()
        with connect_sqlite(self.db_path) as conn:
            cursor = conn.execute(
                "SELECT spec_json FROM decks WHERE deck_id = ?",
                (deck_id,),
            )
            row = cursor.fetchone()
        if not row:
            raise KeyError(deck_id)
        return refresh_deck_evidence_coverage(DeckSpec.model_validate_json(row[0]))

    def list_recent(self, *, limit: int = 100) -> list[DeckSpec]:
        self._init_db()
        safe_limit = max(1, min(500, int(limit or 100)))
        with connect_sqlite(self.db_path) as conn:
            rows = conn.execute(
                "SELECT deck_id FROM decks ORDER BY updated_at DESC, created_at DESC LIMIT ?",
                (safe_limit,),
            ).fetchall()
        return [self.get(str(row[0] or "")) for row in rows if row and row[0]]

    def list_ids_by_session(self, session_id: str) -> list[str]:
        normalized_session_id = str(session_id or "").strip()
        if not normalized_session_id:
            return []

        self._init_db()
        with connect_sqlite(self.db_path) as conn:
            rows = conn.execute(
                "SELECT deck_id FROM decks WHERE session_id = ? ORDER BY updated_at DESC, created_at DESC",
                (normalized_session_id,),
            ).fetchall()
        return [str(row[0] or "") for row in rows if row and row[0]]

    def delete_by_session(self, session_id: str) -> int:
        normalized_session_id = str(session_id or "").strip()
        if not normalized_session_id:
            return 0

        self._init_db()
        with connect_sqlite(self.db_path) as conn:
            cursor = conn.execute(
                "DELETE FROM decks WHERE session_id = ?",
                (normalized_session_id,),
            )
            conn.commit()
            return int(cursor.rowcount or 0)
