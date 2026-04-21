from __future__ import annotations

import io
import logging
import time
from typing import Any, Awaitable, Callable, Optional

from fastapi import APIRouter, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import Response


def build_content_router(
    *,
    artifact_store: Any,
    deck_store: Any,
    share_link_store: Any,
    tasks: dict[str, Any],
    tasks_lock: Any,
    suppressed_task_ids: set[str],
    prune_task_records_locked: Callable[..., None],
    persist_task_record: Callable[..., None],
    prune_persisted_tasks: Callable[[], None],
    get_task_store: Callable[[], Any],
    run_task: Callable[..., Awaitable[None]],
    enqueue_task: Callable[..., Any],
    task_record_payload: Callable[..., dict[str, Any]],
    list_tasks_payload: Callable[..., dict[str, Any]],
    task_history_limit: int,
    artifact_payload: Callable[[Any], dict[str, Any]],
    artifact_export_formats: Callable[[Any], list[str]],
    build_deck_artifact: Callable[..., Any],
    build_report_artifact: Callable[..., Any],
    sync_deck_artifact: Callable[..., None],
    require_remote_viewer: Callable[[Request], dict[str, Any]],
    require_remote_editor: Callable[[Request], dict[str, Any]],
    require_remote_admin: Callable[[Request], dict[str, Any]],
    require_remote_share_secret: Callable[[Request], None],
    current_share_link_secret: Callable[[], str],
    audit_security_event: Callable[..., Any],
    token_fingerprint: Callable[[str], str],
    encode_share_token: Callable[..., str],
    decode_share_token: Callable[..., tuple[str, str]],
    build_share_url: Callable[..., str],
    create_share_link_payload_fn: Callable[..., dict[str, Any]],
    share_link_ttl_seconds: int,
    request_client_ip: Callable[[Request], str],
    request_user_agent: Callable[[Request], str],
    share_link_audit_payload: Callable[[Any], dict[str, Any]],
    share_link_response_model: type,
    revoke_share_link_response_model: type,
    share_link_audit_list_response_model: type,
    open_shared_resource_payload: Callable[..., dict[str, Any]],
    build_session_messages_payload: Callable[..., dict[str, Any]],
    render_shared_session_html: Callable[..., Any],
    render_shared_deck_html: Callable[..., Any],
    build_download_content_disposition: Callable[[str], str],
    build_chat_report_title: Callable[..., str],
    build_report_markdown: Callable[..., str],
    ensure_deckable_chat: Callable[..., Any],
    populate_chat_report_presentation: Callable[..., None],
    safe_report_filename: Callable[[str], str],
    stage_upload_files: Callable[..., Awaitable[Any]],
    build_upload_documents_task_record: Callable[..., Any],
    cleanup_temp_paths: Callable[..., None],
    upload_documents_response: Callable[..., dict[str, Any]],
    effective_vector_store_path: Callable[[Optional[str]], str],
    resolve_report_messages: Callable[..., list[Any]],
    resolve_active_prompt_runtime: Callable[..., Any],
    normalize_model_config: Callable[..., Any],
    build_deck: Callable[..., Awaitable[Any]],
    build_create_deck_kwargs: Callable[..., dict[str, Any]],
    build_regenerate_deck_kwargs: Callable[..., dict[str, Any]],
    apply_deck_update: Callable[..., None],
    replace_deck_slide: Callable[..., None],
    export_deck_payload: Callable[..., dict[str, Any]],
    export_deck_to_pptx: Callable[..., Any],
    build_export_filename: Callable[..., str],
    normalize_deck_theme: Callable[..., str],
    regenerate_deck_slide: Callable[..., Awaitable[Any]],
    sync_deck_artifacts: Callable[..., None],
    create_deck_artifact: Callable[..., Any],
    report_download_payload: Callable[..., dict[str, Any]],
    resolve_report_messages_fn: Callable[..., list[Any]],
    create_report_artifact: Callable[..., tuple[Any, str, str]],
    persist_web_research_task_placeholder: Callable[..., None],
    document_upload_max_count: int,
    document_upload_max_file_bytes: int,
    document_upload_max_total_bytes: int,
    create_task_request_model: type,
    create_deck_request_model: type,
    update_deck_request_model: type,
    regenerate_deck_slide_request_model: type,
    generate_report_request_model: type,
    update_artifact_request_model: type,
    generate_artifact_request_model: type,
    logger: logging.Logger,
) -> APIRouter:
    import asyncio

    router = APIRouter()

    # ── 文档管理 ──────────────────────────────────

    @router.post("/api/documents/upload")
    async def upload_documents(
        request: Request,
        files: list[UploadFile] = File(...),
        vector_store_path: Optional[str] = Form(default=None),
    ):
        require_remote_editor(request)
        temp_paths: list[str] = []
        try:
            evsp = effective_vector_store_path(vector_store_path)
            temp_paths, file_names = await stage_upload_files(
                files,
                max_file_count=document_upload_max_count,
                max_file_bytes=document_upload_max_file_bytes,
                max_total_bytes=document_upload_max_total_bytes,
            )
            record = build_upload_documents_task_record(
                temp_paths=temp_paths,
                file_names=file_names,
                vector_store_path=evsp,
            )
            async with tasks_lock:
                tasks[record.task_id] = record
                prune_task_records_locked(record.created_at)
            persist_task_record(record)
            prune_persisted_tasks()
            asyncio.create_task(run_task(record))
            logger.info("task_id=%s task_type=upload_documents 已创建", record.task_id)
            audit_security_event(
                "upload_documents", request,
                details=f"file_count={len(file_names)} vector_store_path={evsp}",
            )
            return upload_documents_response(record, file_count=len(file_names), vector_store_path=evsp)
        except ValueError as e:
            if temp_paths:
                cleanup_temp_paths(temp_paths)
            audit_security_event("upload_documents", request, result="rejected", details=str(e))
            raise HTTPException(status_code=400, detail=str(e))
        except Exception as e:
            if temp_paths:
                cleanup_temp_paths(temp_paths)
            logger.exception("Document upload failed")
            raise HTTPException(status_code=500, detail=str(e))

    @router.get("/api/documents/stats")
    async def get_document_stats(request: Request, path: Optional[str] = None):
        from doc_pipeline import DocPipeline
        require_remote_viewer(request)
        pipeline = DocPipeline(vector_store_path=effective_vector_store_path(path))
        try:
            pipeline.load_store()
            stats = pipeline.get_stats()
            stats.setdefault("store_path", pipeline.vector_store_path)
            audit_security_event("get_document_stats", request, details=f"path={pipeline.vector_store_path}")
            return stats
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    # ── 异步任务 ──────────────────────────────────

    @router.post("/api/tasks")
    async def create_task(request: create_task_request_model):
        payload = await enqueue_task(
            tasks, tasks_lock,
            task_type=request.task_type,
            params=request.params,
            session_id=request.session_id,
            prune_in_memory=prune_task_records_locked,
            persist_record=persist_task_record,
            prune_persisted=prune_persisted_tasks,
            run_task=run_task,
            spawn_background_task=asyncio.create_task,
            logger=logger,
            on_record_created=(
                persist_web_research_task_placeholder
                if request.task_type == "web_research"
                else None
            ),
        )
        return payload

    @router.get("/api/tasks/{task_id}")
    async def get_task(task_id: str):
        async with tasks_lock:
            prune_task_records_locked()
            record = tasks.get(task_id)
        if record is None:
            prune_persisted_tasks()
            record = get_task_store().get(task_id)
        if record is None:
            raise HTTPException(status_code=404, detail="未找到任务")
        return task_record_payload(record)

    @router.get("/api/tasks")
    async def list_tasks(limit: int = 20):
        async with tasks_lock:
            prune_task_records_locked()
            in_memory_tasks = list(tasks.values())
        prune_persisted_tasks()
        return list_tasks_payload(
            in_memory_tasks=in_memory_tasks,
            persisted_tasks=get_task_store().list_recent(limit=max(limit, task_history_limit)),
            limit=limit,
        )

    # ── 演示稿 ──────────────────────────────────

    @router.post("/api/decks")
    async def create_deck(request: create_deck_request_model):
        from chat_store import SQLiteChatMessageHistory
        history = SQLiteChatMessageHistory(session_id=request.session_id)
        try:
            messages = resolve_report_messages_fn(
                history, answer_group_id=request.answer_group_id, panel_id=request.panel_id,
            )
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested deck scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="该会话没有对话记录")
        try:
            deck = await build_deck(
                messages=messages,
                **build_create_deck_kwargs(
                    request,
                    resolve_active_prompt_runtime=resolve_active_prompt_runtime,
                    normalize_deck_theme=normalize_deck_theme,
                ),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        deck_store.save(deck)
        art = create_deck_artifact(deck)
        payload = deck.model_dump(mode="json")
        payload["artifact_id"] = art.artifact_id
        return payload

    @router.get("/api/decks/{deck_id}")
    async def get_deck(deck_id: str):
        try:
            deck = deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc
        return deck.model_dump(mode="json")

    @router.patch("/api/decks/{deck_id}")
    async def update_deck(deck_id: str, request: update_deck_request_model):
        try:
            deck = deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc
        if request.slides is not None and not request.slides:
            raise HTTPException(status_code=400, detail="演示稿至少需要保留一页")
        apply_deck_update(deck, request, normalize_deck_theme=normalize_deck_theme)
        deck_store.save(deck)
        sync_deck_artifacts(deck)
        return deck.model_dump(mode="json")

    @router.post("/api/decks/{deck_id}/slides/{slide_id}/regenerate")
    async def regenerate_saved_deck_slide(
        deck_id: str, slide_id: str, request: regenerate_deck_slide_request_model,
    ):
        from chat_store import SQLiteChatMessageHistory
        try:
            deck = deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc
        history = SQLiteChatMessageHistory(session_id=deck.meta.session_id)
        try:
            messages = resolve_report_messages_fn(
                history,
                answer_group_id=getattr(deck.meta, "source_answer_group_id", None),
                panel_id=getattr(deck.meta, "source_panel_id", None),
            )
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested deck scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="该会话没有对话记录")
        regenerate_kwargs = build_regenerate_deck_kwargs(
            deck, request,
            normalize_model_config=normalize_model_config,
            resolve_active_prompt_runtime=resolve_active_prompt_runtime,
        )
        try:
            regenerated_slide = await regenerate_deck_slide(deck=deck, slide_id=slide_id, messages=messages, **regenerate_kwargs)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到页面") from exc
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        replace_deck_slide(deck, regenerated_slide)
        deck_store.save(deck)
        sync_deck_artifacts(deck)
        return deck.model_dump(mode="json")

    @router.get("/api/decks/{deck_id}/export")
    async def export_deck(deck_id: str, format: str = "pptx"):
        if format != "pptx":
            raise HTTPException(status_code=400, detail="当前仅支持导出 PPTX")
        try:
            deck = deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到该演示稿草稿") from exc
        try:
            ep = export_deck_payload(deck, export_deck_to_pptx=export_deck_to_pptx, build_export_filename=build_export_filename)
        except RuntimeError as exc:
            raise HTTPException(status_code=500, detail=str(exc)) from exc
        return Response(
            content=ep["content"],
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": build_download_content_disposition(ep["filename"])},
        )

    @router.post("/api/decks/{deck_id}/share", response_model=share_link_response_model)
    async def create_deck_share_link(deck_id: str, request: Request):
        require_remote_share_secret(request)
        share_secret = current_share_link_secret()
        try:
            deck_store.get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到演示稿") from exc
        payload = create_share_link_payload_fn(
            "deck", deck_id, request,
            secret=share_secret,
            encode_share_token=encode_share_token,
            build_share_url=build_share_url,
        )
        record = share_link_store.upsert(
            share_token=payload["share_token"],
            resource_type="deck",
            resource_id=deck_id,
            expires_at=time.time() + share_link_ttl_seconds,
            created_by_ip=request_client_ip(request),
            created_user_agent=request_user_agent(request),
        )
        audit_security_event("create_deck_share_link", request, details=f"deck_id={deck_id}")
        return share_link_response_model(**payload, expires_at=record.expires_at)

    # ── 报告 ──────────────────────────────────

    @router.post("/api/reports/generate")
    async def generate_report(request: generate_report_request_model):
        from chat_store import SQLiteChatMessageHistory
        history = SQLiteChatMessageHistory(session_id=request.session_id)
        try:
            msgs = resolve_report_messages_fn(history, answer_group_id=request.answer_group_id, panel_id=request.panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到对应的报告消息范围。") from exc
        if not msgs:
            raise HTTPException(status_code=400, detail="该会话没有对话记录")
        try:
            artifact, title, markdown = create_report_artifact(
                session_id=request.session_id,
                messages=msgs,
                answer_group_id=str(request.answer_group_id or "").strip(),
                panel_id=str(request.panel_id or "").strip(),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        return {"markdown": markdown, "title": title, "artifact_id": artifact.artifact_id}

    @router.get("/api/reports/download/{session_id}")
    async def download_report_pptx(
        session_id: str,
        answer_group_id: Optional[str] = None,
        panel_id: Optional[str] = None,
    ):
        from chat_store import SQLiteChatMessageHistory
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="python-pptx 未安装，请在 requirements.txt 中添加 python-pptx 并重新安装依赖",
            )
        history = SQLiteChatMessageHistory(session_id=session_id)
        try:
            msgs = resolve_report_messages_fn(history, answer_group_id=answer_group_id, panel_id=panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="未找到对应的报告消息范围。") from exc
        try:
            ensure_deckable_chat(msgs)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        if not msgs:
            raise HTTPException(status_code=400, detail="该会话没有对话记录")
        rp = report_download_payload(
            msgs,
            ensure_deckable_chat=ensure_deckable_chat,
            build_chat_report_title=build_chat_report_title,
            presentation_factory=Presentation,
            body_font_size=Pt(12),
            populate_chat_report_presentation=populate_chat_report_presentation,
            safe_report_filename=safe_report_filename,
        )
        buf = io.BytesIO()
        rp["presentation"].save(buf)
        buf.seek(0)
        return Response(
            content=buf.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": build_download_content_disposition(rp["filename"])},
        )

    # ── Artifacts ──────────────────────────────────

    @router.get("/api/artifacts/{artifact_id}")
    async def get_artifact(artifact_id: str):
        try:
            artifact = artifact_store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        return artifact_payload(artifact)

    @router.patch("/api/artifacts/{artifact_id}")
    async def update_artifact(artifact_id: str, request: update_artifact_request_model):
        try:
            artifact = artifact_store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        next_title = str(request.title or "").strip()
        if artifact.artifact_type == "report":
            if next_title:
                artifact.title = next_title
            if request.markdown is not None:
                artifact.content["markdown"] = str(request.markdown or "").strip()
            artifact_store.save(artifact)
            return artifact_payload(artifact)
        if artifact.artifact_type == "deck":
            if request.markdown is not None:
                raise HTTPException(status_code=400, detail="Deck artifact does not support markdown patching.")
            deck_id = str(artifact.linked_resource_id or artifact.content.get("deck_id") or "").strip()
            if not deck_id:
                raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
            try:
                deck = deck_store.get(deck_id)
            except KeyError as exc:
                raise HTTPException(status_code=404, detail="Deck was not found.") from exc
            if next_title:
                deck.meta.title = next_title
                if deck.slides and deck.slides[0].type == "cover":
                    deck.slides[0].title = next_title
                deck_store.save(deck)
            sync_deck_artifacts(deck)
            return artifact_payload(artifact_store.get(artifact_id))
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    @router.get("/api/artifacts/{artifact_id}/export")
    async def export_artifact(artifact_id: str, format: str = ""):
        try:
            artifact = artifact_store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        if artifact.artifact_type == "report":
            export_format = str(format or "md").strip().lower() or "md"
            if export_format == "md":
                filename = f"{safe_report_filename(artifact.title)}.md"
                return Response(
                    content=str(artifact.content.get("markdown") or ""),
                    media_type="text/markdown; charset=utf-8",
                    headers={"Content-Disposition": build_download_content_disposition(filename)},
                )
            if export_format != "pptx":
                raise HTTPException(status_code=400, detail="Report artifact only supports md / pptx.")
            try:
                from pptx import Presentation
                from pptx.util import Pt
            except ImportError as exc:
                raise HTTPException(status_code=500, detail="python-pptx is not installed.") from exc
            raw_pairs = (
                artifact.content.get("qa_pairs")
                if isinstance(artifact.content.get("qa_pairs"), list)
                else []
            )
            qa_pairs = [
                (str(item.get("question") or "").strip(), str(item.get("answer") or "").strip())
                for item in raw_pairs
                if isinstance(item, dict)
                and (str(item.get("question") or "").strip() or str(item.get("answer") or "").strip())
            ]
            if not qa_pairs:
                raise HTTPException(status_code=400, detail="Report artifact has no exportable content.")
            presentation = Presentation()
            populate_chat_report_presentation(presentation, title=artifact.title, qa_pairs=qa_pairs, body_font_size=Pt(12))
            buffer = io.BytesIO()
            presentation.save(buffer)
            buffer.seek(0)
            filename = f"{safe_report_filename(artifact.title)}.pptx"
            return Response(
                content=buffer.read(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": build_download_content_disposition(filename)},
            )
        if artifact.artifact_type == "deck":
            if str(format or "pptx").strip().lower() != "pptx":
                raise HTTPException(status_code=400, detail="Deck artifact only supports pptx.")
            deck_id = str(artifact.linked_resource_id or artifact.content.get("deck_id") or "").strip()
            if not deck_id:
                raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
            try:
                deck = deck_store.get(deck_id)
            except KeyError as exc:
                raise HTTPException(status_code=404, detail="Deck was not found.") from exc
            ep = export_deck_payload(deck, export_deck_to_pptx=export_deck_to_pptx, build_export_filename=build_export_filename)
            return Response(
                content=ep["content"],
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": build_download_content_disposition(ep["filename"])},
            )
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    @router.post("/api/artifacts/generate")
    async def generate_artifact(request: generate_artifact_request_model):
        from chat_store import SQLiteChatMessageHistory
        history = SQLiteChatMessageHistory(session_id=request.session_id)
        try:
            messages = resolve_report_messages_fn(history, answer_group_id=request.answer_group_id, panel_id=request.panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested artifact scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="No usable messages were found for artifact generation.")
        if request.artifact_type == "report":
            try:
                artifact, _, _ = create_report_artifact(
                    session_id=request.session_id,
                    messages=messages,
                    answer_group_id=str(request.answer_group_id or "").strip(),
                    panel_id=str(request.panel_id or "").strip(),
                )
            except ValueError as exc:
                raise HTTPException(status_code=400, detail=str(exc)) from exc
            return artifact_payload(artifact)
        if request.artifact_type == "deck":
            if request.panel_config is None:
                raise HTTPException(status_code=400, detail="Deck artifact requires panel_config.")
            try:
                deck = await build_deck(
                    messages=messages,
                    **build_create_deck_kwargs(
                        request,
                        resolve_active_prompt_runtime=resolve_active_prompt_runtime,
                        normalize_deck_theme=normalize_deck_theme,
                    ),
                )
            except ValueError as exc:
                raise HTTPException(status_code=400, detail=str(exc)) from exc
            deck_store.save(deck)
            artifact = create_deck_artifact(deck)
            return artifact_payload(artifact)
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    # ── 分享链接 ──────────────────────────────────

    @router.get("/api/share-links", response_model=share_link_audit_list_response_model)
    async def list_share_links(
        request: Request,
        resource_type: str = "",
        active_only: bool = False,
        limit: int = 100,
        offset: int = 0,
    ):
        require_remote_admin(request)
        records = share_link_store.list_links(resource_type=resource_type, active_only=active_only, limit=limit, offset=offset)
        payload_records = [share_link_audit_payload(record) for record in records]
        payload = {
            "share_links": payload_records,
            "total": len(payload_records),
            "active_count": sum(1 for item in payload_records if item["is_active"]),
        }
        audit_security_event(
            "list_share_links", request,
            details=f"resource_type={resource_type or '<all>'} active_only={active_only} total={payload['total']}",
        )
        return payload

    @router.delete("/api/share-links/{share_token}", response_model=revoke_share_link_response_model)
    async def revoke_share_link(share_token: str, request: Request):
        require_remote_admin(request)
        if not share_link_store.revoke(share_token):
            raise HTTPException(status_code=404, detail="未找到分享链接")
        audit_security_event("revoke_share_link", request, details=f"share_token_fp={token_fingerprint(share_token)}")
        return revoke_share_link_response_model(ok=True)

    @router.get("/shared/{share_token}")
    async def open_shared_resource(share_token: str, request: Request):
        require_remote_share_secret(request)
        share_secret = current_share_link_secret()
        try:
            link_record = share_link_store.get_active(share_token)
            if link_record is None:
                raise ValueError("分享链接不存在、已过期或已撤销")
            decoded_type, decoded_id = decode_share_token(share_token, share_secret)
            if link_record.resource_type != decoded_type or link_record.resource_id != decoded_id:
                raise ValueError("分享链接无效")
            shared_payload = open_shared_resource_payload(
                share_token, request,
                secret=share_secret,
                decode_share_token=decode_share_token,
                build_share_url=build_share_url,
                build_session_messages_payload=build_session_messages_payload,
                render_shared_session_html=render_shared_session_html,
                get_deck=deck_store.get,
                render_shared_deck_html=render_shared_deck_html,
            )
            share_link_store.record_access(
                share_token,
                accessed_ip=request_client_ip(request),
                accessed_user_agent=request_user_agent(request),
            )
            audit_security_event(
                "open_shared_resource", request,
                details=f"resource_type={decoded_type} share_token_fp={token_fingerprint(share_token)}",
            )
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc
        except KeyError as exc:
            detail = str(exc.args[0]) if exc.args else "Not found"
            raise HTTPException(status_code=404, detail=detail) from exc
        return Response(content=shared_payload["content"], media_type=shared_payload["media_type"])

    return router
