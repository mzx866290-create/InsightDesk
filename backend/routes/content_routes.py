"""Content route utilities."""

import io
import json
import logging
import base64
import time
from typing import Any, Awaitable, Callable, Optional
from urllib.parse import unquote_to_bytes

from fastapi import APIRouter, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import Response

from backend.routes.resource_access_helpers import (
    filter_visible_resources,
    grant_resource_owner,
    inherit_resource_grants,
    require_resource_access,
)
from backend.helpers.deck_report_helpers import (
    DeckExportGateError,
    attach_deck_delivery_audit,
    build_deck_delivery_response,
    update_deck_block_refs,
)
from backend.schemas.api_models import ApprovalTaskBatchDecisionRequest
from backend.tasks.backends import dispatch_task_record


def build_content_router(
    *,
    artifact_store: Any | Callable[[], Any],
    deck_store: Any | Callable[[], Any],
    share_link_store: Any,
    access_store: Any | Callable[[], Any],
    identity_store: Any | Callable[[], Any],
    tasks: dict[str, Any] | Callable[[], dict[str, Any]],
    tasks_lock: Any,
    suppressed_task_ids: set[str],
    prune_task_records_locked: Callable[..., None],
    persist_task_record: Callable[..., None],
    prune_persisted_tasks: Callable[[], None],
    get_app_config_store: Callable[[], Any],
    get_task_store: Callable[[], Any],
    run_task: Callable[..., Awaitable[None]],
    enqueue_task: Callable[..., Any],
    task_record_payload: Callable[..., dict[str, Any]],
    list_tasks_payload: Callable[..., dict[str, Any]],
    task_history_limit: int,
    artifact_payload: Callable[[Any], dict[str, Any]],
    artifact_export_formats: Callable[[Any], list[str]],
    build_deck_artifact: Callable[..., Any],
    build_report_artifact: Callable[..., Any],
    sync_deck_artifact: Callable[..., None],
    require_remote_viewer: Callable[[Request], dict[str, Any]],
    require_remote_editor: Callable[[Request], dict[str, Any]],
    require_remote_admin: Callable[[Request], dict[str, Any]],
    require_remote_share_secret: Callable[[Request], None],
    current_share_link_secret: Callable[[], str],
    audit_security_event: Callable[..., Any],
    token_fingerprint: Callable[[str], str],
    encode_share_token: Callable[..., str],
    decode_share_token: Callable[..., tuple[str, str]],
    build_share_url: Callable[..., str],
    create_share_link_payload_fn: Callable[..., dict[str, Any]],
    share_link_ttl_seconds: int | Callable[[], int],
    request_client_ip: Callable[[Request], str],
    request_user_agent: Callable[[Request], str],
    share_link_audit_payload: Callable[[Any], dict[str, Any]],
    share_link_response_model: type,
    revoke_share_link_response_model: type,
    share_link_audit_list_response_model: type,
    open_shared_resource_payload: Callable[..., dict[str, Any]],
    build_session_messages_payload: Callable[..., dict[str, Any]],
    render_shared_session_html: Callable[..., Any],
    render_shared_deck_html: Callable[..., Any],
    build_download_content_disposition: Callable[[str], str],
    build_chat_report_title: Callable[..., str],
    build_report_markdown: Callable[..., str],
    ensure_deckable_chat: Callable[..., Any],
    populate_chat_report_presentation: Callable[..., None],
    safe_report_filename: Callable[[str], str],
    stage_upload_files: Callable[..., Awaitable[Any]],
    build_upload_documents_task_record: Callable[..., Any],
    cleanup_temp_paths: Callable[..., None],
    upload_documents_response: Callable[..., dict[str, Any]],
    effective_vector_store_path: Callable[[Optional[str]], str],
    resolve_report_messages: Callable[..., list[Any]],
    resolve_active_prompt_runtime: Callable[..., Any],
    normalize_model_config: Callable[..., Any],
    build_deck: Callable[..., Awaitable[Any]],
    build_create_deck_kwargs: Callable[..., dict[str, Any]],
    build_regenerate_deck_kwargs: Callable[..., dict[str, Any]],
    apply_deck_update: Callable[..., None],
    replace_deck_slide: Callable[..., None],
    export_deck_payload: Callable[..., dict[str, Any]],
    export_deck_to_pptx: Callable[..., Any],
    build_export_filename: Callable[..., str],
    normalize_deck_theme: Callable[..., str],
    regenerate_deck_slide: Callable[..., Awaitable[Any]],
    sync_deck_artifacts: Callable[..., None],
    report_download_payload: Callable[..., dict[str, Any]],
    resolve_report_messages_fn: Callable[..., list[Any]],
    persist_web_research_task_placeholder: Callable[..., None],
    persist_multi_agent_workflow_task_placeholder: Callable[..., None],
    document_upload_max_count: int,
    document_upload_max_file_bytes: int,
    document_upload_max_total_bytes: int,
    create_task_request_model: type,
    create_multi_agent_workflow_request_model: type,
    approval_policy_request_model: type,
    approval_task_decision_request_model: type,
    create_deck_request_model: type,
    update_deck_request_model: type,
    regenerate_deck_slide_request_model: type,
    generate_report_request_model: type,
    update_artifact_request_model: type,
    generate_artifact_request_model: type,
    logger: logging.Logger,
    task_backend: str | Callable[[], str] = "memory",
    enqueue_external_task: Callable[[Any], Awaitable[Any]] | None = None,
    arq_queue_health_payload: Callable[[], Awaitable[dict[str, Any]]] | None = None,
) -> APIRouter:
    import asyncio

    router = APIRouter()
    workflow_data_row_sample_limit = 500
    approval_policy_config_key = "task_approval_policy"
    default_approval_policy = {
        "enabled": False,
        "required_task_types": [],
        "high_risk_requires_approval": True,
        "default_reviewer_role": "admin",
        "updated_at": None,
    }

    def resolve_artifact_store() -> Any:
        if callable(artifact_store):
            return artifact_store()
        return artifact_store

    def resolve_deck_store() -> Any:
        if callable(deck_store):
            return deck_store()
        return deck_store

    def resolve_tasks() -> dict[str, Any]:
        if callable(tasks):
            return tasks()
        return tasks

    def resolve_task_backend() -> str:
        value = task_backend() if callable(task_backend) else task_backend
        return str(value or "memory").strip().lower() or "memory"

    def resolve_build_deck() -> Callable[..., Awaitable[Any]]:
        return build_deck

    def resolve_share_link_ttl_seconds() -> int:
        if callable(share_link_ttl_seconds):
            return int(share_link_ttl_seconds())
        return int(share_link_ttl_seconds)

    def normalized_approval_policy(raw_policy: Any) -> dict[str, Any]:
        if hasattr(raw_policy, "model_dump"):
            data = raw_policy.model_dump()
        elif hasattr(raw_policy, "dict"):
            data = raw_policy.dict()
        elif isinstance(raw_policy, dict):
            data = dict(raw_policy)
        else:
            data = {}

        task_types: list[str] = []
        seen_task_types: set[str] = set()
        raw_task_types = data.get("required_task_types")
        if isinstance(raw_task_types, list):
            for raw_task_type in raw_task_types:
                task_type = str(raw_task_type or "").strip()
                if not task_type or task_type in seen_task_types:
                    continue
                task_types.append(task_type)
                seen_task_types.add(task_type)
        if len(task_types) > 20:
            raise HTTPException(
                status_code=400,
                detail="required_task_types must contain at most 20 items.",
            )

        reviewer_role = str(data.get("default_reviewer_role") or "").strip() or "admin"
        return {
            "enabled": bool(data.get("enabled", False)),
            "required_task_types": task_types,
            "high_risk_requires_approval": bool(
                data.get("high_risk_requires_approval", True)
            ),
            "default_reviewer_role": reviewer_role,
            "updated_at": data.get("updated_at"),
        }

    def approval_policy_payload() -> dict[str, Any]:
        record = get_app_config_store().get(approval_policy_config_key)
        if record is None:
            return dict(default_approval_policy)
        try:
            stored_policy = json.loads(record.value or "{}")
        except json.JSONDecodeError:
            logger.warning("Stored task approval policy is not valid JSON")
            return dict(default_approval_policy)
        payload = normalized_approval_policy(stored_policy)
        payload["updated_at"] = float(record.updated_at or 0.0) or None
        return payload

    def save_approval_policy_payload(raw_policy: Any) -> dict[str, Any]:
        payload = normalized_approval_policy(raw_policy)
        payload.pop("updated_at", None)
        encoded = json.dumps(payload, ensure_ascii=False, separators=(",", ":"))
        record = get_app_config_store().set(approval_policy_config_key, encoded)
        payload["updated_at"] = float(record.updated_at or 0.0) or None
        return payload

    def require_session_access(
        request: Request, session_id: str, minimum_role: str = "viewer"
    ) -> dict[str, Any]:
        role_guard = require_remote_viewer
        if minimum_role == "editor":
            role_guard = require_remote_editor
        elif minimum_role in {"admin", "owner"}:
            role_guard = require_remote_admin
        return require_resource_access(
            request,
            resource_type="session",
            resource_id=session_id,
            minimum_role=minimum_role,
            require_remote_role=role_guard,
            access_store=access_store,
            identity_store=identity_store,
            audit_security_event=audit_security_event,
        )

    def require_deck_access(
        request: Request, deck_id: str, minimum_role: str = "viewer"
    ) -> dict[str, Any]:
        role_guard = require_remote_viewer
        if minimum_role == "editor":
            role_guard = require_remote_editor
        elif minimum_role in {"admin", "owner"}:
            role_guard = require_remote_admin
        return require_resource_access(
            request,
            resource_type="deck",
            resource_id=deck_id,
            minimum_role=minimum_role,
            require_remote_role=role_guard,
            access_store=access_store,
            identity_store=identity_store,
            audit_security_event=audit_security_event,
        )

    def require_artifact_access(
        request: Request, artifact_id: str, minimum_role: str = "viewer"
    ) -> dict[str, Any]:
        role_guard = require_remote_viewer
        if minimum_role == "editor":
            role_guard = require_remote_editor
        elif minimum_role in {"admin", "owner"}:
            role_guard = require_remote_admin
        return require_resource_access(
            request,
            resource_type="artifact",
            resource_id=artifact_id,
            minimum_role=minimum_role,
            require_remote_role=role_guard,
            access_store=access_store,
            identity_store=identity_store,
            audit_security_event=audit_security_event,
        )

    def decode_workflow_data_url(data_url: str) -> bytes:
        if not data_url.startswith("data:") or "," not in data_url:
            return b""
        header, encoded = data_url.split(",", 1)
        try:
            if ";base64" in header:
                return base64.b64decode(encoded, validate=True)
            return unquote_to_bytes(encoded)
        except Exception:
            return b""

    def rows_from_workflow_excel_payload(payload: bytes) -> list[dict[str, Any]]:
        if not payload:
            return []
        try:
            import pandas as pd
        except ImportError:
            return []
        try:
            frame = pd.read_excel(io.BytesIO(payload), nrows=500)
        except Exception:
            return []
        frame = frame.where(pd.notnull(frame), None)
        rows = frame.to_dict(orient="records")
        return [{str(key): value for key, value in row.items()} for row in rows]

    def rows_from_workflow_data_file(raw_file: Any) -> list[dict[str, Any]]:
        if hasattr(raw_file, "model_dump"):
            file_payload = raw_file.model_dump(mode="json")
        elif hasattr(raw_file, "dict"):
            file_payload = raw_file.dict()
        elif isinstance(raw_file, dict):
            file_payload = dict(raw_file)
        else:
            return []

        from backend.agent.agents.data_analysis import DataAnalysisAgent

        analyzer = DataAnalysisAgent()
        extracted_text = str(file_payload.get("extracted_text") or "").strip()
        if extracted_text:
            rows = analyzer._coerce_rows(extracted_text)
            if rows:
                return rows

        data_url = str(file_payload.get("data_url") or "").strip()
        payload = decode_workflow_data_url(data_url)
        if not payload:
            return []
        file_name = str(file_payload.get("name") or "").strip().lower()
        media_type = str(file_payload.get("media_type") or "").strip().lower()
        if (
            file_name.endswith((".xlsx", ".xls"))
            or media_type
            in {
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            }
        ):
            return rows_from_workflow_excel_payload(payload)
        text = payload.decode("utf-8-sig", errors="replace")
        return analyzer._coerce_rows(text)

    def enrich_workflow_data_context(
        context: dict[str, Any],
        data_files: list[Any],
    ) -> tuple[dict[str, Any], list[dict[str, Any]]]:
        if not data_files:
            return context, []

        next_context = dict(context)
        summaries: list[dict[str, Any]] = []
        for raw_file in data_files:
            if hasattr(raw_file, "model_dump"):
                file_payload = raw_file.model_dump(mode="json")
            elif hasattr(raw_file, "dict"):
                file_payload = raw_file.dict()
            elif isinstance(raw_file, dict):
                file_payload = dict(raw_file)
            else:
                continue
            rows = rows_from_workflow_data_file(file_payload)
            if not rows:
                continue
            name = str(file_payload.get("name") or "workflow-data-file").strip()
            sampled_rows = rows[:workflow_data_row_sample_limit]
            summary = {"name": name, "row_count": len(rows)}
            if len(rows) > len(sampled_rows):
                summary.update(
                    {
                        "sampled": True,
                        "sampled_row_count": len(sampled_rows),
                        "sample_limit": workflow_data_row_sample_limit,
                    }
                )
            summaries.append(summary)
            if "rows" not in next_context and "data" not in next_context:
                next_context["rows"] = sampled_rows
                next_context["data_source"] = name
                if summary.get("sampled"):
                    next_context["data_sampling"] = summary
        if summaries:
            next_context["data_files"] = summaries
        return next_context, summaries

    def ensure_data_analysis_plan_step(
        plan: list[dict[str, Any]],
        *,
        user_request: str,
        has_data_rows: bool,
    ) -> list[dict[str, Any]]:
        if not has_data_rows:
            return plan
        if any(str(step.get("agent") or step.get("task_type") or "") == "data_analysis" for step in plan):
            return plan

        step = {
            "id": "step-data-analysis",
            "agent": "data_analysis",
            "task_type": "data_analysis",
            "description": user_request,
            "input": user_request,
            "status": "pending",
            "requires_approval": False,
            "metadata": {"planner": "workflow_data_files"},
        }
        next_plan = [dict(item) for item in plan]
        insert_at = len(next_plan)
        for index, current in enumerate(next_plan):
            if str(current.get("agent") or "") in {"writing", "review"}:
                insert_at = index
                break
        next_plan.insert(insert_at, step)
        for index, current in enumerate(next_plan, start=1):
            current["id"] = str(current.get("id") or f"step-{index}")
        return next_plan

    def create_report_artifact_for_messages(
        *,
        session_id: str,
        messages: list[Any],
        answer_group_id: str = "",
        panel_id: str = "",
    ) -> tuple[Any, str, str]:
        qa_pairs = ensure_deckable_chat(messages)
        title = build_chat_report_title(messages)
        markdown = build_report_markdown(messages, title)
        artifact = build_report_artifact(
            session_id=session_id,
            title=title,
            markdown=markdown,
            qa_pairs=qa_pairs,
            answer_group_id=answer_group_id,
            panel_id=panel_id,
        )
        resolve_artifact_store().save(artifact)
        return artifact, title, markdown

    def create_deck_artifact_for_deck(deck: Any) -> Any:
        artifact = build_deck_artifact(deck)
        resolve_artifact_store().save(artifact)
        return artifact

    def artifact_content(artifact: Any) -> dict[str, Any]:
        content = getattr(artifact, "content", {})
        return dict(content) if isinstance(content, dict) else {}

    def research_report_content(artifact: Any) -> dict[str, Any]:
        content = artifact_content(artifact)
        report = content.get("research_report")
        if isinstance(report, dict):
            return dict(report)
        return {}

    def is_research_archive_artifact(artifact: Any) -> bool:
        content = artifact_content(artifact)
        report = research_report_content(artifact)
        artifact_type = str(getattr(artifact, "artifact_type", "") or "").strip()
        if artifact_type == "research_archive":
            return True
        if bool(content.get("research_archive")):
            return True
        return (
            str(report.get("type") or "").strip() == "research_report"
            and str(report.get("version") or "").strip().lower() == "v2"
        )

    def coerce_items(value: Any) -> list[dict[str, Any]]:
        if not isinstance(value, list):
            return []
        return [dict(item) for item in value if isinstance(item, dict)]

    def research_archive_sources(
        content: dict[str, Any],
        report: dict[str, Any],
    ) -> list[dict[str, Any]]:
        sources = coerce_items(report.get("sources"))
        if sources:
            return sources
        return coerce_items(content.get("sources"))

    def research_archive_task_id(artifact: Any, content: dict[str, Any]) -> str:
        task_id = str(content.get("task_id") or "").strip()
        if task_id:
            return task_id
        linked_type = str(getattr(artifact, "linked_resource_type", "") or "").strip()
        if linked_type == "task":
            return str(getattr(artifact, "linked_resource_id", "") or "").strip()
        return ""

    def compact_text(value: Any, max_length: int = 280) -> str:
        text = " ".join(str(value or "").strip().split())
        if len(text) <= max_length:
            return text
        return f"{text[: max(0, max_length - 3)].rstrip()}..."

    def safe_int(value: Any) -> int:
        try:
            return int(value or 0)
        except (TypeError, ValueError):
            return 0

    def string_items(value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, (str, int, float)):
            text = str(value).strip()
            return [text] if text else []
        if isinstance(value, dict):
            for key in ("source_id", "claim_id", "id", "anchor_id"):
                text = str(value.get(key) or "").strip()
                if text:
                    return [text]
            return []
        if not isinstance(value, list):
            return []
        items: list[str] = []
        seen: set[str] = set()
        for item in value:
            for text in string_items(item):
                if text and text not in seen:
                    items.append(text)
                    seen.add(text)
        return items

    def collect_id_fields(item: dict[str, Any], keys: tuple[str, ...]) -> list[str]:
        values: list[str] = []
        seen: set[str] = set()
        for key in keys:
            for text in string_items(item.get(key)):
                if text and text not in seen:
                    values.append(text)
                    seen.add(text)
        return values

    def source_identifier(source: dict[str, Any]) -> str:
        return str(source.get("source_id") or source.get("id") or source.get("url") or "").strip()

    def claim_identifier(claim: dict[str, Any], index: int) -> str:
        claim_id = str(claim.get("claim_id") or claim.get("id") or "").strip()
        return claim_id or f"claim-{index + 1}"

    def coerce_citation_sections(value: Any) -> list[dict[str, Any]]:
        if isinstance(value, dict):
            sections: list[dict[str, Any]] = []
            for key, raw_item in value.items():
                if isinstance(raw_item, dict):
                    item = dict(raw_item)
                    item.setdefault("anchor_id", key)
                    item.setdefault("paragraph_id", key)
                    sections.append(item)
                elif isinstance(raw_item, list):
                    for nested in raw_item:
                        if isinstance(nested, dict):
                            item = dict(nested)
                            item.setdefault("anchor_id", key)
                            item.setdefault("paragraph_id", key)
                            sections.append(item)
                elif str(raw_item or "").strip():
                    sections.append(
                        {
                            "anchor_id": key,
                            "paragraph_id": key,
                            "text": str(raw_item).strip(),
                        }
                    )
            return sections
        return coerce_items(value)

    def research_paragraph_citations(
        content: dict[str, Any],
        report: dict[str, Any],
    ) -> list[dict[str, Any]]:
        raw_sections: list[dict[str, Any]] = []
        for container in (content, report):
            for key in ("paragraph_citations", "citation_map", "paragraphs", "sections"):
                raw_sections.extend(coerce_citation_sections(container.get(key)))

        citations: list[dict[str, Any]] = []
        seen: set[tuple[str, str, str]] = set()
        for index, item in enumerate(raw_sections):
            paragraph_id = str(
                item.get("paragraph_id")
                or item.get("paragraphId")
                or item.get("paragraph")
                or ""
            ).strip()
            section_id = str(
                item.get("section_id")
                or item.get("sectionId")
                or item.get("section")
                or ""
            ).strip()
            anchor_id = str(item.get("anchor_id") or item.get("anchorId") or item.get("id") or "").strip()
            text = compact_text(
                item.get("text")
                or item.get("content")
                or item.get("body")
                or item.get("summary")
                or item.get("title"),
                600,
            )
            claim_ids = collect_id_fields(
                item,
                ("claim_ids", "claimIds", "claims", "claim_id", "claimId"),
            )
            source_ids = collect_id_fields(
                item,
                (
                    "source_ids",
                    "sourceIds",
                    "sources",
                    "citations",
                    "evidence_source_ids",
                    "evidence",
                    "source_id",
                    "sourceId",
                ),
            )
            if not any((paragraph_id, section_id, anchor_id, text, claim_ids, source_ids)):
                continue
            if not paragraph_id and not section_id and not anchor_id:
                paragraph_id = f"paragraph-{index + 1}"
            dedupe_key = (paragraph_id or section_id, anchor_id, text)
            if dedupe_key in seen:
                continue
            seen.add(dedupe_key)
            citations.append(
                {
                    "paragraph_id": paragraph_id,
                    "section_id": section_id,
                    "text": text,
                    "claim_ids": claim_ids,
                    "source_ids": source_ids,
                    "anchor_id": anchor_id or paragraph_id or section_id,
                }
            )
        return citations

    def claim_source_ids(chain: dict[str, Any]) -> list[str]:
        return collect_id_fields(
            chain,
            (
                "source_ids",
                "sourceIds",
                "supporting_source_ids",
                "supportingSourceIds",
                "supporting_sources",
                "supportingSources",
                "sources",
                "evidence_sources",
                "evidenceSources",
                "evidence",
                "citations",
            ),
        )

    def research_navigation_index(
        paragraph_citations: list[dict[str, Any]],
    ) -> dict[str, Any]:
        paragraph_to_claims: dict[str, list[str]] = {}
        paragraph_to_sources: dict[str, list[str]] = {}
        claim_to_paragraphs: dict[str, list[str]] = {}
        source_to_paragraphs: dict[str, list[str]] = {}
        links: list[dict[str, str]] = []
        seen_links: set[tuple[str, str, str]] = set()

        def add_unique(mapping: dict[str, list[str]], key: str, value: str) -> None:
            if not key or not value:
                return
            items = mapping.setdefault(key, [])
            if value not in items:
                items.append(value)

        for paragraph in paragraph_citations:
            paragraph_id = str(
                paragraph.get("paragraph_id")
                or paragraph.get("anchor_id")
                or paragraph.get("section_id")
                or ""
            ).strip()
            if not paragraph_id:
                continue
            anchor_id = str(paragraph.get("anchor_id") or paragraph_id).strip()
            for claim_id in string_items(paragraph.get("claim_ids")):
                add_unique(paragraph_to_claims, paragraph_id, claim_id)
                add_unique(claim_to_paragraphs, claim_id, paragraph_id)
                link_key = (paragraph_id, claim_id, "claim")
                if link_key not in seen_links:
                    links.append(
                        {
                            "paragraph_id": paragraph_id,
                            "anchor_id": anchor_id,
                            "claim_id": claim_id,
                            "link_type": "claim",
                        }
                    )
                    seen_links.add(link_key)
            for source_id in string_items(paragraph.get("source_ids")):
                add_unique(paragraph_to_sources, paragraph_id, source_id)
                add_unique(source_to_paragraphs, source_id, paragraph_id)
                link_key = (paragraph_id, source_id, "source")
                if link_key not in seen_links:
                    links.append(
                        {
                            "paragraph_id": paragraph_id,
                            "anchor_id": anchor_id,
                            "source_id": source_id,
                            "link_type": "source",
                        }
                    )
                    seen_links.add(link_key)

        return {
            "paragraph_to_claims": paragraph_to_claims,
            "paragraph_to_sources": paragraph_to_sources,
            "claim_to_paragraphs": claim_to_paragraphs,
            "source_to_paragraphs": source_to_paragraphs,
            "links": links,
        }

    def research_citation_graph(
        chains: list[dict[str, Any]],
        sources: list[dict[str, Any]],
        paragraph_citations: list[dict[str, Any]] | None = None,
    ) -> dict[str, list[dict[str, Any]]]:
        nodes: list[dict[str, Any]] = []
        edges: list[dict[str, Any]] = []
        node_ids: set[str] = set()
        edge_ids: set[tuple[str, str, str]] = set()

        for paragraph in paragraph_citations or []:
            paragraph_id = str(
                paragraph.get("paragraph_id")
                or paragraph.get("anchor_id")
                or paragraph.get("section_id")
                or ""
            ).strip()
            if not paragraph_id:
                continue
            node_id = f"paragraph:{paragraph_id}"
            if node_id not in node_ids:
                nodes.append(
                    {
                        "id": node_id,
                        "type": "paragraph",
                        "paragraph_id": paragraph_id,
                        "anchor_id": str(paragraph.get("anchor_id") or paragraph_id).strip(),
                        "text": compact_text(paragraph.get("text"), 240),
                    }
                )
                node_ids.add(node_id)
            for claim_id in string_items(paragraph.get("claim_ids")):
                claim_node_id = f"claim:{claim_id}"
                if claim_node_id not in node_ids:
                    nodes.append(
                        {
                            "id": claim_node_id,
                            "type": "claim",
                            "claim_id": claim_id,
                        }
                    )
                    node_ids.add(claim_node_id)
                edge_key = (node_id, claim_node_id, "mentions")
                if edge_key not in edge_ids:
                    edges.append(
                        {
                            "source": node_id,
                            "target": claim_node_id,
                            "type": "mentions",
                        }
                    )
                    edge_ids.add(edge_key)
            for source_id in string_items(paragraph.get("source_ids")):
                source_node_id = f"source:{source_id}"
                if source_node_id not in node_ids:
                    nodes.append(
                        {
                            "id": source_node_id,
                            "type": "source",
                            "source_id": source_id,
                        }
                    )
                    node_ids.add(source_node_id)
                edge_key = (node_id, source_node_id, "cites")
                if edge_key not in edge_ids:
                    edges.append(
                        {
                            "source": node_id,
                            "target": source_node_id,
                            "type": "cites",
                        }
                    )
                    edge_ids.add(edge_key)

        for index, chain in enumerate(chains):
            claim_id = claim_identifier(chain, index)
            node_id = f"claim:{claim_id}"
            claim_text = compact_text(
                chain.get("claim_text")
                or chain.get("claim")
                or chain.get("text")
                or chain.get("statement"),
                360,
            )
            if node_id not in node_ids:
                nodes.append(
                    {
                        "id": node_id,
                        "type": "claim",
                        "claim_id": claim_id,
                        "text": claim_text,
                    }
                )
                node_ids.add(node_id)
            else:
                for node in nodes:
                    if node.get("id") == node_id and claim_text:
                        node.setdefault("text", claim_text)
                        break
            for source_id in claim_source_ids(chain):
                source_node_id = f"source:{source_id}"
                if source_node_id not in node_ids:
                    nodes.append(
                        {
                            "id": source_node_id,
                            "type": "source",
                            "source_id": source_id,
                        }
                    )
                    node_ids.add(source_node_id)
                edge_key = (node_id, source_node_id, "supports")
                if edge_key not in edge_ids:
                    edges.append(
                        {
                            "source": source_node_id,
                            "target": node_id,
                            "type": "supports",
                        }
                    )
                    edge_ids.add(edge_key)

        for source in sources:
            source_id = source_identifier(source)
            if not source_id:
                continue
            node_id = f"source:{source_id}"
            if node_id in node_ids:
                for node in nodes:
                    if node.get("id") == node_id:
                        node.update(
                            {
                                "title": compact_text(
                                    source.get("title") or source.get("name") or source.get("url"),
                                    180,
                                ),
                                "url": str(source.get("url") or source.get("href") or "").strip(),
                            }
                        )
                        break
                continue
            nodes.append(
                {
                    "id": node_id,
                    "type": "source",
                    "source_id": source_id,
                    "title": compact_text(source.get("title") or source.get("name") or source.get("url"), 180),
                    "url": str(source.get("url") or source.get("href") or "").strip(),
                }
            )
            node_ids.add(node_id)
        return {"nodes": nodes, "edges": edges}

    def research_conflict_review_records(content: dict[str, Any]) -> list[dict[str, Any]]:
        records = coerce_items(content.get("conflict_review_resolutions"))
        normalized: list[dict[str, Any]] = []
        for record in records:
            conflict_id = str(record.get("conflict_id") or "").strip()
            claim_id = str(record.get("claim_id") or "").strip()
            if not conflict_id and not claim_id:
                continue
            normalized.append(
                {
                    "conflict_id": conflict_id or claim_id,
                    "claim_id": claim_id,
                    "status": str(record.get("status") or "reviewed").strip() or "reviewed",
                    "resolution": compact_text(record.get("resolution"), 500),
                    "note": compact_text(record.get("note"), 500),
                    "reviewer": compact_text(record.get("reviewer"), 120),
                    "updated_at": float(record.get("updated_at") or 0),
                }
            )
        return normalized

    def research_conflict_summary(
        chains: list[dict[str, Any]],
        review_records: list[dict[str, Any]] | None = None,
        archive_conflict_review: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        conflict_statuses = ("contradiction", "contradicted", "conflict", "conflicting", "needs_attention")
        items: list[dict[str, Any]] = []
        conflicting_claims: list[str] = []
        reviews_by_claim = {
            str(record.get("claim_id") or "").strip(): record
            for record in (review_records or [])
            if str(record.get("claim_id") or "").strip()
        }
        reviews_by_conflict = {
            str(record.get("conflict_id") or "").strip(): record
            for record in (review_records or [])
            if str(record.get("conflict_id") or "").strip()
        }
        for index, chain in enumerate(chains):
            status = str(chain.get("status") or "").strip().lower()
            conflict_text = compact_text(
                chain.get("conflict_text")
                or chain.get("contradiction_text")
                or chain.get("needs_attention_reason")
                or chain.get("attention_reason")
                or chain.get("conflict")
                or chain.get("contradiction"),
                500,
            )
            has_signal = any(token in status for token in conflict_statuses) or any(
                bool(chain.get(key))
                for key in (
                    "conflict",
                    "conflicting",
                    "contradiction",
                    "contradicted",
                    "has_conflict",
                    "needs_attention",
                )
            )
            if not has_signal and not conflict_text:
                continue
            claim_id = claim_identifier(chain, index)
            conflict_id = str(chain.get("conflict_id") or chain.get("id") or claim_id).strip()
            review = reviews_by_conflict.get(conflict_id) or reviews_by_claim.get(claim_id)
            claim_text = compact_text(
                chain.get("claim_text")
                or chain.get("claim")
                or chain.get("text")
                or chain.get("statement"),
                500,
            )
            conflicting_claims.append(claim_id)
            items.append(
                {
                    "conflict_id": conflict_id,
                    "claim_id": claim_id,
                    "claim_text": claim_text,
                    "status": status,
                    "text": compact_text(
                        conflict_text
                        or claim_text,
                        500,
                    ),
                    "source_ids": claim_source_ids(chain),
                    "review_status": str(review.get("status") or "unreviewed").strip()
                    if review
                    else "unreviewed",
                    "review": review or None,
                }
            )
        archive_conflicts = archive_conflict_review.get("conflicts") if isinstance(archive_conflict_review, dict) else []
        if isinstance(archive_conflicts, list):
            for conflict in archive_conflicts:
                if not isinstance(conflict, dict):
                    continue
                claim_id = str(conflict.get("claim_id") or "").strip()
                archive_id = str(conflict.get("archive_id") or "").strip()
                archive_claim_id = str(conflict.get("archive_claim_id") or "").strip()
                if not claim_id or not archive_id or not archive_claim_id:
                    continue
                conflict_id = str(
                    conflict.get("conflict_id")
                    or f"archive:{archive_id}:{archive_claim_id}:{claim_id}"
                ).strip()
                review = reviews_by_conflict.get(conflict_id) or reviews_by_claim.get(claim_id)
                claim_text = compact_text(
                    conflict.get("claim_text")
                    or conflict.get("current_claim_text")
                    or conflict.get("text"),
                    500,
                )
                archive_claim_text = compact_text(conflict.get("archive_claim_text"), 500)
                text = compact_text(
                    conflict.get("conflict_text")
                    or " ".join(part for part in [claim_text, archive_claim_text] if part),
                    500,
                )
                conflicting_claims.append(claim_id)
                items.append(
                    {
                        "conflict_id": conflict_id,
                        "claim_id": claim_id,
                        "claim_text": claim_text,
                        "archive_id": archive_id,
                        "archive_claim_id": archive_claim_id,
                        "archive_claim_text": archive_claim_text,
                        "status": str(conflict.get("severity") or "needs_review").strip().lower(),
                        "text": text,
                        "source_ids": [
                            str(source_id).strip()
                            for source_id in (
                                conflict.get("source_ids")
                                if isinstance(conflict.get("source_ids"), list)
                                else [conflict.get("source_id")]
                            )
                            if str(source_id or "").strip()
                        ],
                        "review_status": str(review.get("status") or "unreviewed").strip()
                        if review
                        else "unreviewed",
                        "review": review or None,
                    }
                )
        return {
            "total": len(items),
            "conflicting_claims": conflicting_claims,
            "items": items,
            "reviewed": sum(1 for item in items if item.get("review_status") != "unreviewed"),
            "unreviewed": sum(1 for item in items if item.get("review_status") == "unreviewed"),
        }

    def research_claim_preview(
        chains: list[dict[str, Any]],
        report: dict[str, Any],
        limit: int = 3,
    ) -> list[dict[str, Any]]:
        previews: list[dict[str, Any]] = []
        for chain in chains:
            claim_text = (
                chain.get("claim_text")
                or chain.get("claim")
                or chain.get("text")
                or chain.get("statement")
            )
            if not str(claim_text or "").strip():
                continue
            previews.append(
                {
                    "claim_id": str(chain.get("claim_id") or chain.get("id") or "").strip(),
                    "claim_text": compact_text(claim_text),
                    "status": str(chain.get("status") or "").strip(),
                    "supporting_source_count": safe_int(chain.get("supporting_source_count")),
                    "has_primary_source": bool(chain.get("has_primary_source")),
                }
            )
            if len(previews) >= limit:
                return previews

        for claim in coerce_items(report.get("atomic_claims")):
            claim_text = claim.get("claim_text") or claim.get("claim") or claim.get("text")
            if not str(claim_text or "").strip():
                continue
            previews.append(
                {
                    "claim_id": str(claim.get("claim_id") or claim.get("id") or "").strip(),
                    "claim_text": compact_text(claim_text),
                    "status": str(claim.get("status") or "").strip(),
                    "supporting_source_count": 0,
                    "has_primary_source": False,
                }
            )
            if len(previews) >= limit:
                break
        return previews

    def research_source_preview(
        sources: list[dict[str, Any]],
        limit: int = 3,
    ) -> list[dict[str, Any]]:
        previews: list[dict[str, Any]] = []
        for source in sources:
            title = source.get("title") or source.get("name") or source.get("url")
            url = source.get("url") or source.get("href")
            if not str(title or url or "").strip():
                continue
            previews.append(
                {
                    "source_id": str(source.get("source_id") or source.get("id") or "").strip(),
                    "title": compact_text(title, 180),
                    "url": str(url or "").strip(),
                    "source_tier": str(source.get("source_tier") or "").strip(),
                    "freshness_band": str(source.get("freshness_band") or "").strip(),
                    "snippet": compact_text(
                        source.get("snippet") or source.get("summary") or source.get("content"),
                        220,
                    ),
                }
            )
            if len(previews) >= limit:
                break
        return previews

    def source_capability_items(value: Any) -> list[str]:
        if isinstance(value, dict):
            return [
                str(key).strip()
                for key, enabled in value.items()
                if str(key).strip() and enabled is not False and enabled is not None
            ]
        return string_items(value)

    def research_provider_capability_coverage(
        sources: list[dict[str, Any]],
    ) -> dict[str, Any]:
        items: list[dict[str, Any]] = []
        providers: list[str] = []
        for source in sources:
            source_id = source_identifier(source)
            provider = str(
                source.get("provider")
                or source.get("provider_name")
                or source.get("search_provider")
                or ""
            ).strip()
            capabilities = source_capability_items(
                source.get("capabilities")
                or source.get("capability_scopes")
                or source.get("declared_capabilities")
            )
            if not provider and not capabilities:
                continue
            if provider and provider not in providers:
                providers.append(provider)
            items.append(
                {
                    "source_id": source_id,
                    "provider": provider,
                    "capabilities": capabilities,
                    "declared": bool(provider and capabilities),
                }
            )
        return {
            "total_sources": len(sources),
            "declared_sources": sum(1 for item in items if item.get("declared")),
            "providers": providers,
            "items": items,
        }

    def research_citation_panel_payload(
        chains: list[dict[str, Any]],
        sources: list[dict[str, Any]],
        verification_summary: dict[str, Any],
    ) -> dict[str, Any]:
        source_index: dict[str, dict[str, Any]] = {}
        claim_source_links: list[dict[str, str]] = []
        seen_links: set[tuple[str, str]] = set()

        for index, source in enumerate(sources, start=1):
            source_id = source_identifier(source) or str(source.get("doc_id") or f"source-{index}").strip()
            if not source_id:
                continue
            source_index[source_id] = {
                "source_id": source_id,
                "source_index": index,
                "title": compact_text(source.get("title") or source.get("name") or source.get("url"), 180),
                "url": str(source.get("url") or source.get("href") or "").strip(),
                "domain": str(source.get("domain") or "").strip(),
                "provider": str(source.get("provider") or source.get("provider_name") or "").strip(),
                "source_tier": str(source.get("source_tier") or "").strip(),
                "source_family": str(source.get("source_family") or "").strip(),
                "freshness_band": str(source.get("freshness_band") or "").strip(),
                "published_at": str(source.get("published_at") or "").strip(),
                "snippet": compact_text(source.get("snippet") or source.get("summary") or source.get("content"), 500),
            }

        for index, chain in enumerate(chains):
            claim_id = claim_identifier(chain, index)
            for source_id in claim_source_ids(chain):
                if not claim_id or not source_id:
                    continue
                link_key = (claim_id, source_id)
                if link_key in seen_links:
                    continue
                claim_source_links.append(
                    {
                        "claim_id": claim_id,
                        "source_id": source_id,
                        "link_type": "supports",
                    }
                )
                seen_links.add(link_key)

        return {
            "version": "v2",
            "claim_evidence_chains": chains,
            "claim_verification_summary": verification_summary,
            "source_index": source_index,
            "claim_source_links": claim_source_links,
        }

    def research_archive_search_text(
        artifact: Any,
        content: dict[str, Any],
        report: dict[str, Any],
        sources: list[dict[str, Any]],
        chains: list[dict[str, Any]],
    ) -> str:
        paragraph_citations = research_paragraph_citations(content, report)
        review_records = research_conflict_review_records(content)
        conflict_summary = research_conflict_summary(
            chains,
            review_records,
            archive_conflict_review=report.get("archive_conflict_review")
            if isinstance(report.get("archive_conflict_review"), dict)
            else None,
        )
        parts: list[str] = [
            getattr(artifact, "artifact_id", ""),
            getattr(artifact, "title", ""),
            getattr(artifact, "session_id", ""),
            content.get("markdown", ""),
            report.get("query", ""),
            report.get("summary", ""),
        ]
        for claim in chains:
            parts.extend(
                [
                    claim.get("claim_id", ""),
                    claim.get("claim_text", ""),
                    claim.get("claim", ""),
                    claim.get("text", ""),
                    claim.get("conflict_text", ""),
                    claim.get("contradiction_text", ""),
                    claim.get("needs_attention_reason", ""),
                    claim.get("attention_reason", ""),
                    claim.get("conflict", ""),
                    claim.get("contradiction", ""),
                ]
            )
        for source in sources:
            parts.extend(
                [
                    source.get("source_id", ""),
                    source.get("id", ""),
                    source.get("title", ""),
                    source.get("url", ""),
                    source.get("snippet", ""),
                    source.get("summary", ""),
                ]
            )
        for paragraph in paragraph_citations:
            parts.extend(
                [
                    paragraph.get("paragraph_id", ""),
                    paragraph.get("section_id", ""),
                    paragraph.get("anchor_id", ""),
                    paragraph.get("text", ""),
                    " ".join(paragraph.get("claim_ids", [])),
                    " ".join(paragraph.get("source_ids", [])),
                ]
            )
        for item in conflict_summary["items"]:
            parts.extend(
                [
                    item.get("conflict_id", ""),
                    item.get("claim_id", ""),
                    item.get("status", ""),
                    item.get("text", ""),
                    item.get("review_status", ""),
                    " ".join(item.get("source_ids", [])),
                ]
            )
        for review in review_records:
            parts.extend(
                [
                    review.get("conflict_id", ""),
                    review.get("claim_id", ""),
                    review.get("status", ""),
                    review.get("resolution", ""),
                    review.get("note", ""),
                    review.get("reviewer", ""),
                ]
            )
        return " ".join(str(part or "") for part in parts).lower()

    def research_archive_payload(artifact: Any) -> dict[str, Any]:
        content = artifact_content(artifact)
        report = research_report_content(artifact)
        chains = coerce_items(content.get("claim_evidence_chains")) or coerce_items(
            report.get("claim_evidence_chains")
        )
        sources = research_archive_sources(content, report)
        verification_summary = (
            dict(content.get("claim_verification_summary"))
            if isinstance(content.get("claim_verification_summary"), dict)
            else {}
        ) or (
            dict(report.get("claim_verification_summary"))
            if isinstance(report.get("claim_verification_summary"), dict)
            else {}
        )
        delivery_quality = (
            dict(report.get("delivery_quality"))
            if isinstance(report.get("delivery_quality"), dict)
            else {}
        )
        paragraph_citations = research_paragraph_citations(content, report)
        navigation_index = research_navigation_index(paragraph_citations)
        citation_graph = research_citation_graph(chains, sources, paragraph_citations)
        review_records = research_conflict_review_records(content)
        conflict_summary = research_conflict_summary(
            chains,
            review_records,
            archive_conflict_review=report.get("archive_conflict_review")
            if isinstance(report.get("archive_conflict_review"), dict)
            else None,
        )
        return {
            "archive_id": str(getattr(artifact, "artifact_id", "") or ""),
            "artifact_id": str(getattr(artifact, "artifact_id", "") or ""),
            "title": str(getattr(artifact, "title", "") or ""),
            "session_id": str(getattr(artifact, "session_id", "") or ""),
            "task_id": research_archive_task_id(artifact, content),
            "created_at": float(getattr(artifact, "created_at", 0) or 0),
            "updated_at": float(getattr(artifact, "updated_at", 0) or 0),
            "claim_count": len(chains),
            "source_count": len(sources),
            "verification_summary": verification_summary,
            "delivery_quality": delivery_quality,
            "preview_claims": research_claim_preview(chains, report),
            "preview_sources": research_source_preview(sources),
            "provider_capabilities": research_provider_capability_coverage(sources),
            "paragraph_citations": paragraph_citations,
            "paragraph_claim_links": navigation_index["links"],
            "navigation_index": navigation_index,
            "citation_graph": citation_graph,
            "conflict_summary": conflict_summary,
            "conflict_review_resolutions": review_records,
            "citation_panel": research_citation_panel_payload(
                chains,
                sources,
                verification_summary,
            ),
        }

    def normalize_conflict_key_text(value: Any) -> str:
        text = " ".join(str(value or "").strip().lower().split())
        return "".join(ch for ch in text if ch.isalnum() or ch.isspace()).strip()

    def research_conflict_groups(archives: list[Any]) -> list[dict[str, Any]]:
        groups: dict[tuple[str, str, str], dict[str, Any]] = {}
        for artifact in archives:
            payload = research_archive_payload(artifact)
            artifact_id = str(payload.get("artifact_id") or "")
            title = str(payload.get("title") or "")
            for item in payload.get("conflict_summary", {}).get("items", []):
                if not isinstance(item, dict):
                    continue
                source_ids = string_items(item.get("source_ids"))
                normalized_claim = normalize_conflict_key_text(
                    item.get("claim_text") or item.get("claim_id")
                )
                normalized_source = normalize_conflict_key_text(" ".join(sorted(source_ids)))
                normalized_conflict = normalize_conflict_key_text(item.get("text"))
                key = (normalized_claim, normalized_source, normalized_conflict)
                group = groups.setdefault(
                    key,
                    {
                        "group_id": f"conflict-group-{len(groups) + 1}",
                        "normalized_claim": normalized_claim,
                        "normalized_source": normalized_source,
                        "normalized_conflict_text": normalized_conflict,
                        "conflict_text": str(item.get("text") or ""),
                        "claim_ids": [],
                        "source_ids": [],
                        "archives": [],
                        "review_statuses": [],
                        "total": 0,
                    },
                )
                claim_id = str(item.get("claim_id") or "").strip()
                if claim_id and claim_id not in group["claim_ids"]:
                    group["claim_ids"].append(claim_id)
                for source_id in source_ids:
                    if source_id not in group["source_ids"]:
                        group["source_ids"].append(source_id)
                review_status = str(item.get("review_status") or "unreviewed").strip()
                if review_status not in group["review_statuses"]:
                    group["review_statuses"].append(review_status)
                group["archives"].append(
                    {
                        "artifact_id": artifact_id,
                        "archive_id": str(payload.get("archive_id") or artifact_id),
                        "title": title,
                        "claim_id": claim_id,
                        "conflict_id": str(item.get("conflict_id") or claim_id),
                        "review_status": review_status,
                    }
                )
                group["total"] += 1
        return sorted(
            groups.values(),
            key=lambda group: (-int(group.get("total") or 0), str(group.get("group_id") or "")),
        )

    def matches_research_archive_filters(
        artifact: Any,
        *,
        q: str,
        session_id: str,
        task_id: str,
    ) -> bool:
        content = artifact_content(artifact)
        report = research_report_content(artifact)
        chains = coerce_items(content.get("claim_evidence_chains")) or coerce_items(
            report.get("claim_evidence_chains")
        )
        sources = research_archive_sources(content, report)
        if session_id and str(getattr(artifact, "session_id", "") or "").strip() != session_id:
            return False
        if task_id and research_archive_task_id(artifact, content) != task_id:
            return False
        if q and q.lower() not in research_archive_search_text(
            artifact,
            content,
            report,
            sources,
            chains,
        ):
            return False
        return True

    async def create_background_task_payload(
        *,
        http_request: Request,
        task_type: str,
        params: dict[str, Any],
        session_id: str | None = None,
        on_record_created: Callable[..., None] | None = None,
    ) -> dict[str, Any]:
        if session_id:
            require_session_access(http_request, session_id, "editor")
        else:
            require_remote_editor(http_request)
        task_state = resolve_tasks()
        payload = await enqueue_task(
            task_state,
            tasks_lock,
            task_type=task_type,
            params=params,
            session_id=session_id,
            prune_in_memory=prune_task_records_locked,
            persist_record=persist_task_record,
            prune_persisted=prune_persisted_tasks,
            run_task=run_task,
            spawn_background_task=asyncio.create_task,
            logger=logger,
            task_backend=resolve_task_backend(),
            enqueue_external_task=enqueue_external_task,
            on_record_created=on_record_created,
        )
        if session_id and payload.get("task_id"):
            inherit_resource_grants(
                source_resource_type="session",
                source_resource_id=session_id,
                target_resource_type="task",
                target_resource_id=str(payload.get("task_id") or ""),
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
                request=http_request,
            )
            grant_resource_owner(
                http_request,
                resource_type="task",
                resource_id=str(payload.get("task_id") or ""),
                require_remote_role=require_remote_editor,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
            )
        return payload

    async def dispatch_existing_task_record(record: Any) -> str:
        backend = await dispatch_task_record(
            record,
            task_backend=resolve_task_backend(),
            run_task=run_task,
            spawn_background_task=asyncio.create_task,
            enqueue_external_task=enqueue_external_task,
        )
        logger.info(
            "task_id=%s task_type=%s dispatched backend=%s",
            getattr(record, "task_id", ""),
            getattr(record, "task_type", ""),
            backend,
        )
        return backend

    # 鈹€鈹€ 鏂囨。绠＄悊 鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.post("/api/documents/upload")
    async def upload_documents(
        request: Request,
        files: list[UploadFile] = File(...),
        vector_store_path: Optional[str] = Form(default=None),
    ):
        require_remote_editor(request)
        temp_paths: list[str] = []
        try:
            evsp = effective_vector_store_path(vector_store_path)
            temp_paths, file_names = await stage_upload_files(
                files,
                max_file_count=document_upload_max_count,
                max_file_bytes=document_upload_max_file_bytes,
                max_total_bytes=document_upload_max_total_bytes,
            )
            record = build_upload_documents_task_record(
                temp_paths=temp_paths,
                file_names=file_names,
                vector_store_path=evsp,
            )
            task_state = resolve_tasks()
            async with tasks_lock:
                task_state[record.task_id] = record
                prune_task_records_locked(record.created_at)
            persist_task_record(record)
            prune_persisted_tasks()
            await dispatch_existing_task_record(record)
            logger.info("task_id=%s task_type=upload_documents created", record.task_id)
            audit_security_event(
                "upload_documents", request,
                details=f"file_count={len(file_names)} vector_store_path={evsp}",
            )
            return upload_documents_response(record, file_count=len(file_names), vector_store_path=evsp)
        except ValueError as e:
            if temp_paths:
                cleanup_temp_paths(temp_paths)
            audit_security_event("upload_documents", request, result="rejected", details=str(e))
            raise HTTPException(status_code=400, detail=str(e))
        except Exception as e:
            if temp_paths:
                cleanup_temp_paths(temp_paths)
            logger.exception("Document upload failed")
            raise HTTPException(status_code=500, detail=str(e))

    @router.get("/api/documents/stats")
    async def get_document_stats(request: Request, path: Optional[str] = None):
        from backend.services.doc_pipeline import DocPipeline
        require_remote_viewer(request)
        pipeline = DocPipeline(vector_store_path=effective_vector_store_path(path))
        try:
            pipeline.load_store()
            stats = pipeline.get_stats()
            stats.setdefault("store_path", pipeline.vector_store_path)
            audit_security_event("get_document_stats", request, details=f"path={pipeline.vector_store_path}")
            return stats
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    # 鈹€鈹€ 寮傛浠诲姟 鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.post("/api/tasks")
    async def create_task(http_request: Request, request: create_task_request_model):
        on_record_created = None
        if request.task_type == "web_research":
            on_record_created = persist_web_research_task_placeholder
        elif request.task_type == "multi_agent_workflow":
            on_record_created = persist_multi_agent_workflow_task_placeholder
        return await create_background_task_payload(
            http_request=http_request,
            task_type=request.task_type,
            params=request.params,
            session_id=request.session_id,
            on_record_created=on_record_created,
        )

    @router.post("/api/tasks/multi-agent-workflow")
    async def create_multi_agent_workflow_task(
        http_request: Request,
        request: create_multi_agent_workflow_request_model,
    ):
        user_request = str(request.user_request or "").strip()
        if not user_request:
            raise HTTPException(status_code=400, detail="Workflow user_request cannot be empty.")

        context = dict(request.context or {})
        context, data_file_summaries = enrich_workflow_data_context(
            context,
            list(request.data_files or []),
        )
        if request.session_id and not str(context.get("session_id") or "").strip():
            context["session_id"] = request.session_id
        task_approval_policy = approval_policy_payload()
        if task_approval_policy.get("enabled"):
            context["task_approval_policy"] = task_approval_policy
        plan = ensure_data_analysis_plan_step(
            [dict(step) for step in request.plan if isinstance(step, dict)],
            user_request=user_request,
            has_data_rows=bool(context.get("rows")),
        )

        params: dict[str, Any] = {
            "user_request": user_request,
            "panel_id": str(request.panel_id or "").strip(),
            "answer_group_id": str(request.answer_group_id or "").strip(),
            "model_id": str(request.model_id or "multi_agent_workflow").strip()
            or "multi_agent_workflow",
            "context": context,
            "plan": plan,
            "research_mode": str(request.research_mode or "deep").strip().lower() or "deep",
            "providers": [str(item).strip() for item in request.providers if str(item).strip()],
            "max_rounds": max(1, int(request.max_rounds or 2)),
            "max_results_per_query": max(1, int(request.max_results_per_query or 4)),
            "max_fetch_pages": max(1, int(request.max_fetch_pages or 3)),
            "time_range": str(request.time_range or "").strip(),
            "use_kb_context": bool(request.use_kb_context),
            "vector_store_path": str(request.vector_store_path or "").strip(),
            "allow_quick_fallback": bool(request.allow_quick_fallback),
        }
        if data_file_summaries:
            params["data_files"] = data_file_summaries
        if request.panel_config is not None:
            params["panel_config"] = (
                request.panel_config.model_dump(mode="json")
                if hasattr(request.panel_config, "model_dump")
                else request.panel_config.dict()
                if hasattr(request.panel_config, "dict")
                else request.panel_config
            )

        return await create_background_task_payload(
            http_request=http_request,
            task_type="multi_agent_workflow",
            params=params,
            session_id=request.session_id,
            on_record_created=persist_multi_agent_workflow_task_placeholder,
        )

    @router.get("/api/tasks")
    async def list_tasks(request: Request, limit: int = 20, status: str = ""):
        task_state = resolve_tasks()
        async with tasks_lock:
            prune_task_records_locked()
            in_memory_tasks = list(task_state.values())
        prune_persisted_tasks()
        persisted_tasks = get_task_store().list_recent(limit=max(limit, task_history_limit))
        filtered_in_memory_tasks = []
        for record in in_memory_tasks:
            if getattr(record, "session_id", None):
                try:
                    require_session_access(request, str(record.session_id), "viewer")
                except HTTPException as exc:
                    if exc.status_code == 403:
                        continue
                    raise
            else:
                require_remote_viewer(request)
            filtered_in_memory_tasks.append(record)
        filtered_persisted_tasks = []
        for record in persisted_tasks:
            if getattr(record, "session_id", None):
                try:
                    require_session_access(request, str(record.session_id), "viewer")
                except HTTPException as exc:
                    if exc.status_code == 403:
                        continue
                    raise
            else:
                require_remote_viewer(request)
            filtered_persisted_tasks.append(record)
        queue_health = None
        if resolve_task_backend() in {"arq", "redis"} and arq_queue_health_payload is not None:
            queue_health = await arq_queue_health_payload()
        payload = list_tasks_payload(
            in_memory_tasks=filtered_in_memory_tasks,
            persisted_tasks=filtered_persisted_tasks,
            limit=limit,
            status_filter=status,
            queue_health=queue_health,
        )
        return payload

    @router.get("/api/tasks/approval-policy")
    async def get_task_approval_policy(request: Request):
        require_remote_admin(request)
        return approval_policy_payload()

    @router.put("/api/tasks/approval-policy")
    async def update_task_approval_policy(
        http_request: Request,
        request: approval_policy_request_model,
    ):
        require_remote_admin(http_request)
        payload = save_approval_policy_payload(request)
        audit_security_event(
            "task_approval_policy_update",
            http_request,
            details=(
                f"enabled={payload['enabled']} "
                f"required_task_types={len(payload['required_task_types'])} "
                f"high_risk_requires_approval={payload['high_risk_requires_approval']} "
                f"default_reviewer_role={payload['default_reviewer_role']}"
            ),
        )
        return payload

    @router.get("/api/tasks/{task_id}")
    async def get_task(task_id: str, request: Request):
        task_state = resolve_tasks()
        async with tasks_lock:
            prune_task_records_locked()
            record = task_state.get(task_id)
        if record is None:
            prune_persisted_tasks()
            record = get_task_store().get(task_id)
        if record is None:
            raise HTTPException(status_code=404, detail="Task was not found.")
        if getattr(record, "session_id", None):
            require_session_access(request, str(record.session_id), "viewer")
        else:
            require_remote_viewer(request)
        return task_record_payload(record)

    @router.post("/api/tasks/{task_id}/approval")
    async def decide_task_approval(
        task_id: str,
        http_request: Request,
        request: approval_task_decision_request_model,
    ):
        return await apply_task_approval_decision(task_id, http_request, request)

    async def apply_task_approval_decision(
        task_id: str,
        http_request: Request,
        request: Any,
    ) -> dict[str, Any]:
        task_state = resolve_tasks()
        async with tasks_lock:
            prune_task_records_locked()
            record = task_state.get(task_id)
        if record is None:
            prune_persisted_tasks()
            record = get_task_store().get(task_id)
        if record is None:
            raise HTTPException(status_code=404, detail="Task was not found.")
        if record.task_type != "multi_agent_workflow":
            raise HTTPException(status_code=400, detail="Task does not support approval decisions.")
        if record.status != getattr(type(record.status), "WAITING_APPROVAL", record.status):
            if str(getattr(record.status, "value", record.status)) != "waiting_approval":
                raise HTTPException(status_code=400, detail="Task is not waiting for approval.")
        if getattr(record, "session_id", None):
            require_session_access(http_request, str(record.session_id), "editor")
        else:
            require_remote_admin(http_request)

        params = dict(record.params or {})
        params["approval_decision"] = str(request.decision or "").strip().lower()
        params["approval_reviewer"] = str(request.reviewer or "").strip()
        params["approval_comment"] = str(request.comment or "").strip()
        record.params = params
        record.status = getattr(type(record.status), "PENDING", record.status)
        record.progress = min(100, max(10, int(getattr(record, "progress", 0) or 0)))
        record.error = None
        record.result = ""
        record.updated_at = time.time()

        async with tasks_lock:
            task_state[task_id] = record
            prune_task_records_locked(record.updated_at)
        persist_task_record(record)
        prune_persisted_tasks()
        await dispatch_existing_task_record(record)
        audit_security_event(
            "task_approval_decision",
            http_request,
            details=(
                f"task_id={task_id} decision={params['approval_decision']} "
                f"session_id={getattr(record, 'session_id', '') or '<none>'}"
            ),
        )
        return task_record_payload(record)

    # 鈹€鈹€ 婕旂ず绋?鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.post("/api/tasks/approvals/batch")
    async def decide_task_approvals_batch(
        http_request: Request,
        request: ApprovalTaskBatchDecisionRequest,
    ):
        results: list[dict[str, Any]] = []
        succeeded = 0
        for task_id in request.task_ids:
            normalized_task_id = str(task_id or "").strip()
            if not normalized_task_id:
                results.append(
                    {"task_id": normalized_task_id, "ok": False, "error": "Task id is required."}
                )
                continue
            try:
                task_payload = await apply_task_approval_decision(
                    normalized_task_id,
                    http_request,
                    request,
                )
            except HTTPException as exc:
                results.append(
                    {
                        "task_id": normalized_task_id,
                        "ok": False,
                        "error": str(exc.detail),
                    }
                )
                continue
            succeeded += 1
            results.append(
                {
                    "task_id": normalized_task_id,
                    "ok": True,
                    "task": task_payload,
                }
            )

        failed = len(results) - succeeded
        audit_security_event(
            "task_approval_batch_decision",
            http_request,
            details=(
                f"total={len(results)} succeeded={succeeded} failed={failed} "
                f"decision={str(request.decision or '').strip().lower()}"
            ),
        )
        return {
            "total": len(results),
            "succeeded": succeeded,
            "failed": failed,
            "results": results,
        }

    @router.post("/api/decks")
    async def create_deck(http_request: Request, request: create_deck_request_model):
        from backend.stores.factory import create_chat_message_history
        require_session_access(http_request, request.session_id, "editor")
        history = create_chat_message_history(session_id=request.session_id)
        try:
            messages = resolve_report_messages_fn(
                history, answer_group_id=request.answer_group_id, panel_id=request.panel_id,
            )
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested deck scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="No messages were found in this session.")
        try:
            deck = await resolve_build_deck()(
                messages=messages,
                **build_create_deck_kwargs(
                    request,
                    resolve_active_prompt_runtime=resolve_active_prompt_runtime,
                    normalize_deck_theme=normalize_deck_theme,
                ),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        attach_deck_delivery_audit(deck)
        resolve_deck_store().save(deck)
        art = create_deck_artifact_for_deck(deck)
        payload = deck.model_dump(mode="json")
        payload["artifact_id"] = art.artifact_id
        inherit_resource_grants(
            source_resource_type="session",
            source_resource_id=request.session_id,
            target_resource_type="deck",
            target_resource_id=deck.deck_id,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
            request=http_request,
        )
        inherit_resource_grants(
            source_resource_type="session",
            source_resource_id=request.session_id,
            target_resource_type="artifact",
            target_resource_id=art.artifact_id,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
            request=http_request,
        )
        grant_resource_owner(
            http_request,
            resource_type="deck",
            resource_id=deck.deck_id,
            require_remote_role=require_remote_editor,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
        )
        grant_resource_owner(
            http_request,
            resource_type="artifact",
            resource_id=art.artifact_id,
            require_remote_role=require_remote_editor,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
        )
        return payload

    @router.get("/api/decks")
    async def list_decks(request: Request, limit: int = 100):
        safe_limit = max(1, min(500, int(limit or 100)))
        decks = resolve_deck_store().list_recent(limit=safe_limit)
        visible_decks = filter_visible_resources(
            request,
            decks,
            resource_type="deck",
            resource_id_getter=lambda deck: str(getattr(deck, "deck_id", "") or ""),
            require_remote_role=require_remote_viewer,
            access_store=access_store,
            identity_store=identity_store,
        )
        return {
            "decks": [
                attach_deck_delivery_audit(deck).model_dump(mode="json")
                for deck in visible_decks
            ],
            "total": len(visible_decks),
            "limit": safe_limit,
        }

    @router.get("/api/decks/{deck_id}")
    async def get_deck(deck_id: str, request: Request):
        require_deck_access(request, deck_id, "viewer")
        try:
            deck = resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        return attach_deck_delivery_audit(deck).model_dump(mode="json")

    @router.patch("/api/decks/{deck_id}")
    async def update_deck(deck_id: str, http_request: Request, request: update_deck_request_model):
        require_deck_access(http_request, deck_id, "editor")
        try:
            deck = resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        if request.slides is not None and not request.slides:
            raise HTTPException(status_code=400, detail="Deck must keep at least one slide.")
        apply_deck_update(deck, request, normalize_deck_theme=normalize_deck_theme)
        resolve_deck_store().save(deck)
        sync_deck_artifacts(deck)
        return deck.model_dump(mode="json")

    @router.patch("/api/decks/{deck_id}/slides/{slide_id}/blocks/{block_id}/refs")
    async def update_saved_deck_block_refs(
        deck_id: str,
        slide_id: str,
        block_id: str,
        http_request: Request,
        payload: dict[str, Any],
    ):
        require_deck_access(http_request, deck_id, "editor")
        try:
            deck = resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        try:
            result = update_deck_block_refs(deck, slide_id, block_id, payload)
        except KeyError as exc:
            missing_id = str(exc.args[0] if exc.args else "")
            raise HTTPException(status_code=404, detail=f"Slide or block was not found: {missing_id}") from exc
        resolve_deck_store().save(deck)
        sync_deck_artifacts(deck)
        block = result["block"]
        return {
            "deck": deck.model_dump(mode="json"),
            "slide_id": result["slide_id"],
            "block_id": result["block_id"],
            "block": block.model_dump(mode="json") if hasattr(block, "model_dump") else block,
            "citation_validation": result["citation_validation"],
            "evidence_review": result["evidence_review"],
            "export_gate": result["export_gate"],
            "slide_delivery": result["slide_delivery"],
        }

    @router.post("/api/decks/{deck_id}/slides/{slide_id}/regenerate")
    async def regenerate_saved_deck_slide(
        deck_id: str, slide_id: str, http_request: Request, request: regenerate_deck_slide_request_model,
    ):
        from backend.stores.factory import create_chat_message_history
        require_deck_access(http_request, deck_id, "editor")
        try:
            deck = resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        history = create_chat_message_history(session_id=deck.meta.session_id)
        try:
            messages = resolve_report_messages_fn(
                history,
                answer_group_id=getattr(deck.meta, "source_answer_group_id", None),
                panel_id=getattr(deck.meta, "source_panel_id", None),
            )
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested deck scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="No messages were found in this session.")
        regenerate_kwargs = build_regenerate_deck_kwargs(
            deck, request,
            normalize_model_config=normalize_model_config,
            resolve_active_prompt_runtime=resolve_active_prompt_runtime,
        )
        try:
            regenerated_slide = await regenerate_deck_slide(deck=deck, slide_id=slide_id, messages=messages, **regenerate_kwargs)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Slide was not found.") from exc
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        replace_deck_slide(deck, regenerated_slide)
        resolve_deck_store().save(deck)
        sync_deck_artifacts(deck)
        return build_deck_delivery_response(deck, focus_slide_id=slide_id)

    @router.get("/api/decks/{deck_id}/export")
    async def export_deck(
        deck_id: str,
        request: Request,
        format: str = "pptx",
        allow_unsafe_export: bool = False,
        override_reason: str = "",
    ):
        require_deck_access(request, deck_id, "viewer")
        if format != "pptx":
            raise HTTPException(status_code=400, detail="褰撳墠浠呮敮鎸佸鍑?PPTX")
        try:
            deck = resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Deck was not found.") from exc
        try:
            ep = export_deck_payload(
                deck,
                export_deck_to_pptx=export_deck_to_pptx,
                build_export_filename=build_export_filename,
                allow_unsafe_export=allow_unsafe_export,
                override_reason=override_reason,
            )
        except DeckExportGateError as exc:
            raise HTTPException(status_code=409, detail=exc.payload) from exc
        except RuntimeError as exc:
            raise HTTPException(status_code=500, detail=str(exc)) from exc
        return Response(
            content=ep["content"],
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": build_download_content_disposition(ep["filename"])},
        )

    @router.post("/api/decks/{deck_id}/share", response_model=share_link_response_model)
    async def create_deck_share_link(deck_id: str, request: Request):
        require_remote_share_secret(request)
        require_deck_access(request, deck_id, "viewer")
        share_secret = current_share_link_secret()
        try:
            resolve_deck_store().get(deck_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="鏈壘鍒版紨绀虹") from exc
        payload = create_share_link_payload_fn(
            "deck", deck_id, request,
            secret=share_secret,
            encode_share_token=encode_share_token,
            build_share_url=build_share_url,
        )
        record = share_link_store.upsert(
            share_token=payload["share_token"],
            resource_type="deck",
            resource_id=deck_id,
            expires_at=time.time() + resolve_share_link_ttl_seconds(),
            created_by_ip=request_client_ip(request),
            created_user_agent=request_user_agent(request),
        )
        audit_security_event("create_deck_share_link", request, details=f"deck_id={deck_id}")
        return share_link_response_model(**payload, expires_at=record.expires_at)

    # 鈹€鈹€ 鎶ュ憡 鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.post("/api/reports/generate")
    async def generate_report(http_request: Request, request: generate_report_request_model):
        from backend.stores.factory import create_chat_message_history
        require_session_access(http_request, request.session_id, "editor")
        history = create_chat_message_history(session_id=request.session_id)
        try:
            msgs = resolve_report_messages_fn(history, answer_group_id=request.answer_group_id, panel_id=request.panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested report scope was not found.") from exc
        if not msgs:
            raise HTTPException(status_code=400, detail="No messages were found in this session.")
        try:
            artifact, title, markdown = create_report_artifact_for_messages(
                session_id=request.session_id,
                messages=msgs,
                answer_group_id=str(request.answer_group_id or "").strip(),
                panel_id=str(request.panel_id or "").strip(),
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        inherit_resource_grants(
            source_resource_type="session",
            source_resource_id=request.session_id,
            target_resource_type="artifact",
            target_resource_id=artifact.artifact_id,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
            request=http_request,
        )
        grant_resource_owner(
            http_request,
            resource_type="artifact",
            resource_id=artifact.artifact_id,
            require_remote_role=require_remote_editor,
            access_store=access_store,
            now=time.time,
            audit_security_event=audit_security_event,
        )
        return {"markdown": markdown, "title": title, "artifact_id": artifact.artifact_id}

    @router.get("/api/reports/download/{session_id}")
    async def download_report_pptx(
        session_id: str,
        request: Request,
        answer_group_id: Optional[str] = None,
        panel_id: Optional[str] = None,
    ):
        from backend.stores.factory import create_chat_message_history
        require_session_access(request, session_id, "viewer")
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="python-pptx is not installed. Please install it and try again.",
            )
        history = create_chat_message_history(session_id=session_id)
        try:
            msgs = resolve_report_messages_fn(history, answer_group_id=answer_group_id, panel_id=panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested report scope was not found.") from exc
        try:
            ensure_deckable_chat(msgs)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        if not msgs:
            raise HTTPException(status_code=400, detail="No messages were found in this session.")
        rp = report_download_payload(
            msgs,
            ensure_deckable_chat=ensure_deckable_chat,
            build_chat_report_title=build_chat_report_title,
            presentation_factory=Presentation,
            body_font_size=Pt(12),
            populate_chat_report_presentation=populate_chat_report_presentation,
            safe_report_filename=safe_report_filename,
        )
        buf = io.BytesIO()
        rp["presentation"].save(buf)
        buf.seek(0)
        return Response(
            content=buf.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": build_download_content_disposition(rp["filename"])},
        )

    # 鈹€鈹€ Artifacts 鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.get("/api/artifacts")
    async def list_artifacts(request: Request, limit: int = 100, artifact_type: str = ""):
        safe_limit = max(1, min(500, int(limit or 100)))
        artifacts = resolve_artifact_store().list_recent(
            limit=safe_limit,
            artifact_type=artifact_type,
        )
        visible_artifacts = filter_visible_resources(
            request,
            artifacts,
            resource_type="artifact",
            resource_id_getter=lambda artifact: str(getattr(artifact, "artifact_id", "") or ""),
            require_remote_role=require_remote_viewer,
            access_store=access_store,
            identity_store=identity_store,
        )
        return {
            "artifacts": [artifact_payload(artifact) for artifact in visible_artifacts],
            "total": len(visible_artifacts),
            "limit": safe_limit,
        }

    @router.get("/api/research/archives")
    async def list_research_archives(
        request: Request,
        q: str = "",
        session_id: str = "",
        task_id: str = "",
        limit: int = 100,
    ):
        safe_limit = max(1, min(500, int(limit or 100)))
        query_text = str(q or "").strip()
        session_filter = str(session_id or "").strip()
        task_filter = str(task_id or "").strip()
        artifacts = [
            artifact
            for artifact in resolve_artifact_store().list_recent(limit=500)
            if is_research_archive_artifact(artifact)
            and matches_research_archive_filters(
                artifact,
                q=query_text,
                session_id=session_filter,
                task_id=task_filter,
            )
        ]
        visible_artifacts = filter_visible_resources(
            request,
            artifacts,
            resource_type="artifact",
            resource_id_getter=lambda artifact: str(getattr(artifact, "artifact_id", "") or ""),
            require_remote_role=require_remote_viewer,
            access_store=access_store,
            identity_store=identity_store,
        )
        limited_artifacts = visible_artifacts[:safe_limit]
        return {
            "archives": [
                research_archive_payload(artifact)
                for artifact in limited_artifacts
            ],
            "conflict_groups": research_conflict_groups(visible_artifacts),
            "total": len(visible_artifacts),
            "limit": safe_limit,
        }

    @router.post("/api/research/archives/{artifact_id}/conflict-resolutions")
    async def upsert_research_conflict_resolution(
        artifact_id: str,
        request: Request,
    ):
        require_artifact_access(request, artifact_id, "editor")
        store = resolve_artifact_store()
        try:
            artifact = store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        if not is_research_archive_artifact(artifact):
            raise HTTPException(status_code=400, detail="Artifact is not a research archive.")

        body = await request.json()
        if not isinstance(body, dict):
            raise HTTPException(status_code=400, detail="Resolution payload must be an object.")
        conflict_id = str(body.get("conflict_id") or "").strip()
        claim_id = str(body.get("claim_id") or "").strip()
        if not conflict_id and not claim_id:
            raise HTTPException(status_code=400, detail="conflict_id or claim_id is required.")
        status = str(body.get("status") or "resolved").strip() or "resolved"
        if status not in {"resolved", "dismissed", "needs_followup", "reviewed"}:
            raise HTTPException(status_code=400, detail="Unsupported resolution status.")

        content = artifact_content(artifact)
        records = research_conflict_review_records(content)
        record = {
            "conflict_id": conflict_id or claim_id,
            "claim_id": claim_id,
            "status": status,
            "resolution": compact_text(body.get("resolution"), 1000),
            "note": compact_text(body.get("note"), 1000),
            "reviewer": compact_text(body.get("reviewer"), 120),
            "updated_at": time.time(),
        }
        replaced = False
        for index, existing in enumerate(records):
            if (
                str(existing.get("conflict_id") or "") == record["conflict_id"]
                or (
                    record["claim_id"]
                    and str(existing.get("claim_id") or "") == record["claim_id"]
                )
            ):
                records[index] = record
                replaced = True
                break
        if not replaced:
            records.append(record)
        content["conflict_review_resolutions"] = records
        artifact.content = content
        store.save(artifact)
        return {
            "resolution": record,
            "archive": research_archive_payload(artifact),
        }

    @router.get("/api/artifacts/{artifact_id}")
    async def get_artifact(artifact_id: str, request: Request):
        require_artifact_access(request, artifact_id, "viewer")
        store = resolve_artifact_store()
        try:
            artifact = store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        return artifact_payload(artifact)

    @router.patch("/api/artifacts/{artifact_id}")
    async def update_artifact(artifact_id: str, http_request: Request, request: update_artifact_request_model):
        require_artifact_access(http_request, artifact_id, "editor")
        store = resolve_artifact_store()
        try:
            artifact = store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        next_title = str(request.title or "").strip()
        if artifact.artifact_type == "report":
            if next_title:
                artifact.title = next_title
            if request.markdown is not None:
                artifact.content["markdown"] = str(request.markdown or "").strip()
            store.save(artifact)
            return artifact_payload(artifact)
        if artifact.artifact_type == "deck":
            if request.markdown is not None:
                raise HTTPException(status_code=400, detail="Deck artifact does not support markdown patching.")
            deck_id = str(artifact.linked_resource_id or artifact.content.get("deck_id") or "").strip()
            if not deck_id:
                raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
            try:
                deck = resolve_deck_store().get(deck_id)
            except KeyError as exc:
                raise HTTPException(status_code=404, detail="Deck was not found.") from exc
            if next_title:
                deck.meta.title = next_title
                if deck.slides and deck.slides[0].type == "cover":
                    deck.slides[0].title = next_title
                resolve_deck_store().save(deck)
            sync_deck_artifacts(deck)
            return artifact_payload(store.get(artifact_id))
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    @router.get("/api/artifacts/{artifact_id}/export")
    async def export_artifact(
        artifact_id: str,
        request: Request,
        format: str = "",
        allow_unsafe_export: bool = False,
        override_reason: str = "",
    ):
        require_artifact_access(request, artifact_id, "viewer")
        store = resolve_artifact_store()
        try:
            artifact = store.get(artifact_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Artifact was not found.") from exc
        if artifact.artifact_type == "report":
            export_format = str(format or "md").strip().lower() or "md"
            if export_format == "md":
                filename = f"{safe_report_filename(artifact.title)}.md"
                return Response(
                    content=str(artifact.content.get("markdown") or ""),
                    media_type="text/markdown; charset=utf-8",
                    headers={"Content-Disposition": build_download_content_disposition(filename)},
                )
            if export_format == "docx":
                try:
                    from backend.artifact_service import export_report_to_docx
                except ImportError as exc:
                    raise HTTPException(status_code=500, detail="python-docx is not installed.") from exc
                filename = f"{safe_report_filename(artifact.title)}.docx"
                try:
                    content = export_report_to_docx(artifact)
                except RuntimeError as exc:
                    raise HTTPException(status_code=500, detail=str(exc)) from exc
                return Response(
                    content=content,
                    media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    headers={"Content-Disposition": build_download_content_disposition(filename)},
                )
            if export_format == "xlsx":
                try:
                    from backend.artifact_service import report_has_tables, export_report_to_xlsx
                except ImportError as exc:
                    raise HTTPException(status_code=500, detail="openpyxl is not installed.") from exc
                if not report_has_tables(artifact):
                    raise HTTPException(status_code=400, detail="Report artifact has no table content to export as xlsx.")
                filename = f"{safe_report_filename(artifact.title)}.xlsx"
                try:
                    content = export_report_to_xlsx(artifact)
                except RuntimeError as exc:
                    raise HTTPException(status_code=500, detail=str(exc)) from exc
                return Response(
                    content=content,
                    media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    headers={"Content-Disposition": build_download_content_disposition(filename)},
                )
            if export_format != "pptx":
                raise HTTPException(status_code=400, detail="Report artifact only supports md / docx / xlsx / pptx.")
            try:
                from pptx import Presentation
                from pptx.util import Pt
            except ImportError as exc:
                raise HTTPException(status_code=500, detail="python-pptx is not installed.") from exc
            raw_pairs = (
                artifact.content.get("qa_pairs")
                if isinstance(artifact.content.get("qa_pairs"), list)
                else []
            )
            qa_pairs = [
                (str(item.get("question") or "").strip(), str(item.get("answer") or "").strip())
                for item in raw_pairs
                if isinstance(item, dict)
                and (str(item.get("question") or "").strip() or str(item.get("answer") or "").strip())
            ]
            if not qa_pairs:
                raise HTTPException(status_code=400, detail="Report artifact has no exportable content.")
            presentation = Presentation()
            populate_chat_report_presentation(presentation, title=artifact.title, qa_pairs=qa_pairs, body_font_size=Pt(12))
            buffer = io.BytesIO()
            presentation.save(buffer)
            buffer.seek(0)
            filename = f"{safe_report_filename(artifact.title)}.pptx"
            return Response(
                content=buffer.read(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": build_download_content_disposition(filename)},
            )
        if artifact.artifact_type == "deck":
            if str(format or "pptx").strip().lower() != "pptx":
                raise HTTPException(status_code=400, detail="Deck artifact only supports pptx.")
            deck_id = str(artifact.linked_resource_id or artifact.content.get("deck_id") or "").strip()
            if not deck_id:
                raise HTTPException(status_code=400, detail="Deck artifact is missing deck_id.")
            try:
                deck = resolve_deck_store().get(deck_id)
            except KeyError as exc:
                raise HTTPException(status_code=404, detail="Deck was not found.") from exc
            try:
                ep = export_deck_payload(
                    deck,
                    export_deck_to_pptx=export_deck_to_pptx,
                    build_export_filename=build_export_filename,
                    allow_unsafe_export=allow_unsafe_export,
                    override_reason=override_reason,
                )
            except DeckExportGateError as exc:
                raise HTTPException(status_code=409, detail=exc.payload) from exc
            return Response(
                content=ep["content"],
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": build_download_content_disposition(ep["filename"])},
            )
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    @router.post("/api/artifacts/generate")
    async def generate_artifact(http_request: Request, request: generate_artifact_request_model):
        from backend.stores.factory import create_chat_message_history
        require_session_access(http_request, request.session_id, "editor")
        history = create_chat_message_history(session_id=request.session_id)
        try:
            messages = resolve_report_messages_fn(history, answer_group_id=request.answer_group_id, panel_id=request.panel_id)
        except KeyError as exc:
            raise HTTPException(status_code=404, detail="Requested artifact scope was not found.") from exc
        if not messages:
            raise HTTPException(status_code=400, detail="No usable messages were found for artifact generation.")
        if request.artifact_type == "report":
            try:
                artifact, _, _ = create_report_artifact_for_messages(
                    session_id=request.session_id,
                    messages=messages,
                    answer_group_id=str(request.answer_group_id or "").strip(),
                    panel_id=str(request.panel_id or "").strip(),
                )
            except ValueError as exc:
                raise HTTPException(status_code=400, detail=str(exc)) from exc
            inherit_resource_grants(
                source_resource_type="session",
                source_resource_id=request.session_id,
                target_resource_type="artifact",
                target_resource_id=artifact.artifact_id,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
                request=http_request,
            )
            grant_resource_owner(
                http_request,
                resource_type="artifact",
                resource_id=artifact.artifact_id,
                require_remote_role=require_remote_editor,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
            )
            return artifact_payload(artifact)
        if request.artifact_type == "deck":
            if request.panel_config is None:
                raise HTTPException(status_code=400, detail="Deck artifact requires panel_config.")
            try:
                deck = await resolve_build_deck()(
                    messages=messages,
                    **build_create_deck_kwargs(
                        request,
                        resolve_active_prompt_runtime=resolve_active_prompt_runtime,
                        normalize_deck_theme=normalize_deck_theme,
                    ),
                )
            except ValueError as exc:
                raise HTTPException(status_code=400, detail=str(exc)) from exc
            attach_deck_delivery_audit(deck)
            resolve_deck_store().save(deck)
            artifact = create_deck_artifact_for_deck(deck)
            inherit_resource_grants(
                source_resource_type="session",
                source_resource_id=request.session_id,
                target_resource_type="deck",
                target_resource_id=deck.deck_id,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
                request=http_request,
            )
            inherit_resource_grants(
                source_resource_type="session",
                source_resource_id=request.session_id,
                target_resource_type="artifact",
                target_resource_id=artifact.artifact_id,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
                request=http_request,
            )
            grant_resource_owner(
                http_request,
                resource_type="deck",
                resource_id=deck.deck_id,
                require_remote_role=require_remote_editor,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
            )
            grant_resource_owner(
                http_request,
                resource_type="artifact",
                resource_id=artifact.artifact_id,
                require_remote_role=require_remote_editor,
                access_store=access_store,
                now=time.time,
                audit_security_event=audit_security_event,
            )
            return artifact_payload(artifact)
        raise HTTPException(status_code=400, detail="Unsupported artifact type.")

    # 鈹€鈹€ 鍒嗕韩閾炬帴 鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€鈹€

    @router.get("/api/share-links", response_model=share_link_audit_list_response_model)
    async def list_share_links(
        request: Request,
        resource_type: str = "",
        active_only: bool = False,
        limit: int = 100,
        offset: int = 0,
    ):
        require_remote_admin(request)
        records = share_link_store.list_links(resource_type=resource_type, active_only=active_only, limit=limit, offset=offset)
        payload_records = [share_link_audit_payload(record) for record in records]
        payload = {
            "share_links": payload_records,
            "total": len(payload_records),
            "active_count": sum(1 for item in payload_records if item["is_active"]),
        }
        audit_security_event(
            "list_share_links", request,
            details=f"resource_type={resource_type or '<all>'} active_only={active_only} total={payload['total']}",
        )
        return payload

    @router.delete("/api/share-links/{share_token}", response_model=revoke_share_link_response_model)
    async def revoke_share_link(share_token: str, request: Request):
        require_remote_admin(request)
        if not share_link_store.revoke(share_token):
            raise HTTPException(status_code=404, detail="Share link was not found.")
        audit_security_event("revoke_share_link", request, details=f"share_token_fp={token_fingerprint(share_token)}")
        return revoke_share_link_response_model(ok=True)

    @router.get("/shared/{share_token}")
    async def open_shared_resource(share_token: str, request: Request):
        require_remote_share_secret(request)
        share_secret = current_share_link_secret()
        try:
            link_record = share_link_store.get_active(share_token)
            if link_record is None:
                raise ValueError("鍒嗕韩閾炬帴涓嶅瓨鍦ㄣ€佸凡杩囨湡鎴栧凡鎾ら攢")
            decoded_type, decoded_id = decode_share_token(share_token, share_secret)
            if link_record.resource_type != decoded_type or link_record.resource_id != decoded_id:
                raise ValueError("鍒嗕韩閾炬帴鏃犳晥")
            shared_payload = open_shared_resource_payload(
                share_token, request,
                secret=share_secret,
                decode_share_token=decode_share_token,
                build_share_url=build_share_url,
                build_session_messages_payload=build_session_messages_payload,
                render_shared_session_html=render_shared_session_html,
                get_deck=resolve_deck_store().get,
                render_shared_deck_html=render_shared_deck_html,
            )
            share_link_store.record_access(
                share_token,
                accessed_ip=request_client_ip(request),
                accessed_user_agent=request_user_agent(request),
            )
            audit_security_event(
                "open_shared_resource", request,
                details=f"resource_type={decoded_type} share_token_fp={token_fingerprint(share_token)}",
            )
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc
        except KeyError as exc:
            detail = str(exc.args[0]) if exc.args else "Not found"
            raise HTTPException(status_code=404, detail=detail) from exc
        return Response(content=shared_payload["content"], media_type=shared_payload["media_type"])

    return router

